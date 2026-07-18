"""Generate and offline-sign the v1.1 first-party Office template assets.

The signing key is supplied as an external PEM file and is never written into
the repository.  The generated OOXML ZIP metadata and document properties are
normalized so identical inputs produce identical asset bytes.
"""

from __future__ import annotations

import argparse
import base64
import hashlib
import json
import zipfile
from copy import copy
from datetime import UTC, datetime
from io import BytesIO
from pathlib import Path
from typing import Any, Callable

from cryptography.hazmat.primitives import serialization
from cryptography.hazmat.primitives.asymmetric.ed25519 import Ed25519PrivateKey
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.formatting.rule import FormulaRule
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.worksheet.datavalidation import DataValidation
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt

from app.office_templates.validation import inspect_ooxml_package


FIXED_TIMESTAMP = datetime(2026, 1, 1, 0, 0, 0, tzinfo=UTC)
ZIP_TIMESTAMP = (1980, 1, 1, 0, 0, 0)
KEY_ID = "suxiaoyou-office-templates-2026-01"
CJK_FONT = "Arial Unicode MS"


def _canonical_json(value: object) -> bytes:
    return (
        json.dumps(
            value,
            ensure_ascii=False,
            allow_nan=False,
            separators=(",", ":"),
            sort_keys=True,
        ).encode("utf-8")
        + b"\n"
    )


def _normalize_ooxml(payload: bytes) -> bytes:
    output = BytesIO()
    with zipfile.ZipFile(BytesIO(payload), "r") as source:
        with zipfile.ZipFile(output, "w", allowZip64=False) as target:
            for original in sorted(source.infolist(), key=lambda item: item.filename):
                if original.is_dir():
                    continue
                info = zipfile.ZipInfo(original.filename, date_time=ZIP_TIMESTAMP)
                info.compress_type = zipfile.ZIP_DEFLATED
                info.create_system = 0
                info.external_attr = 0
                info.flag_bits = 0
                target.writestr(
                    info,
                    source.read(original),
                    compress_type=zipfile.ZIP_DEFLATED,
                    compresslevel=9,
                )
    return output.getvalue()


def _set_word_font(style: Any, name: str, size: float, color: str) -> None:
    style.font.name = name
    style.font.size = Pt(size)
    style.font.color.rgb = RGBColor.from_string(color)
    style.element.rPr.rFonts.set(qn("w:eastAsia"), name)


def _set_cell_shading(cell: Any, fill: str) -> None:
    properties = cell._tc.get_or_add_tcPr()
    shading = properties.find(qn("w:shd"))
    if shading is None:
        shading = OxmlElement("w:shd")
        properties.append(shading)
    shading.set(qn("w:fill"), fill)


def _set_cell_margins(cell: Any, *, top: int, start: int, bottom: int, end: int) -> None:
    properties = cell._tc.get_or_add_tcPr()
    margins = properties.first_child_found_in("w:tcMar")
    if margins is None:
        margins = OxmlElement("w:tcMar")
        properties.append(margins)
    for tag, value in (("top", top), ("start", start), ("bottom", bottom), ("end", end)):
        node = margins.find(qn(f"w:{tag}"))
        if node is None:
            node = OxmlElement(f"w:{tag}")
            margins.append(node)
        node.set(qn("w:w"), str(value))
        node.set(qn("w:type"), "dxa")


def _make_docx() -> bytes:
    document = Document()
    properties = document.core_properties
    properties.author = "苏小有"
    properties.last_modified_by = "苏小有"
    properties.created = FIXED_TIMESTAMP
    properties.modified = FIXED_TIMESTAMP
    properties.title = "商务简报模板"

    section = document.sections[0]
    section.start_type = WD_SECTION.NEW_PAGE
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(1)
    section.right_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.header_distance = Inches(0.492)
    section.footer_distance = Inches(0.492)

    styles = document.styles
    normal = styles["Normal"]
    _set_word_font(normal, CJK_FONT, 11, "172033")
    normal.paragraph_format.space_after = Pt(6)
    normal.paragraph_format.line_spacing = 1.1
    title_style = styles["Title"]
    _set_word_font(title_style, CJK_FONT, 26, "0F766E")
    title_style.font.bold = True
    title_style.paragraph_format.space_after = Pt(8)
    for name, size, color, before, after in (
        ("Heading 1", 16, "0F766E", 16, 8),
        ("Heading 2", 13, "115E59", 12, 6),
        ("Heading 3", 12, "134E4A", 8, 4),
    ):
        style = styles[name]
        _set_word_font(style, CJK_FONT, size, color)
        style.font.bold = True
        style.paragraph_format.space_before = Pt(before)
        style.paragraph_format.space_after = Pt(after)

    title = document.add_paragraph(style="Title")
    title.add_run("{{title}}")
    metadata = document.add_paragraph()
    metadata.alignment = WD_ALIGN_PARAGRAPH.LEFT
    metadata.paragraph_format.space_after = Pt(14)
    metadata_run = metadata.add_run("面向：{{recipient}}    日期：{{report_date}}")
    metadata_run.font.name = CJK_FONT
    metadata_run.font.size = Pt(10)
    metadata_run.font.color.rgb = RGBColor(0x5B, 0x64, 0x74)

    document.add_heading("执行摘要", level=1)
    summary = document.add_paragraph("{{summary}}")
    summary.paragraph_format.space_after = Pt(12)

    document.add_heading("行动信息", level=1)
    table = document.add_table(rows=2, cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.LEFT
    table.autofit = False
    table.style = "Table Grid"
    widths = (Inches(1.5), Inches(5.0))
    values = (("负责人", "{{owner}}"), ("下一步", "{{next_step}}"))
    for row_index, row_values in enumerate(values):
        for column_index, value in enumerate(row_values):
            cell = table.cell(row_index, column_index)
            cell.width = widths[column_index]
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            _set_cell_margins(cell, top=100, start=120, bottom=100, end=120)
            paragraph = cell.paragraphs[0]
            paragraph.paragraph_format.space_after = Pt(0)
            run = paragraph.add_run(value)
            run.font.name = CJK_FONT
            run.font.size = Pt(10.5)
            if column_index == 0:
                _set_cell_shading(cell, "E6FFFB")
                run.bold = True
                run.font.color.rgb = RGBColor(0x11, 0x5E, 0x59)

    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    footer_run = footer.add_run("{{classification}}")
    footer_run.font.name = CJK_FONT
    footer_run.font.size = Pt(9)
    footer_run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    buffer = BytesIO()
    document.save(buffer)
    return _normalize_ooxml(buffer.getvalue())


def _make_xlsx() -> bytes:
    workbook = Workbook()
    workbook.properties.creator = "苏小有"
    workbook.properties.lastModifiedBy = "苏小有"
    workbook.properties.created = FIXED_TIMESTAMP.replace(tzinfo=None)
    workbook.properties.modified = FIXED_TIMESTAMP.replace(tzinfo=None)
    sheet = workbook.active
    sheet.title = "项目跟踪"
    sheet.sheet_view.showGridLines = False
    sheet.freeze_panes = "A5"
    sheet.merge_cells("A1:F1")
    sheet["A1"] = "{{project_name}} 项目跟踪"
    sheet["A1"].font = Font(name=CJK_FONT, size=20, bold=True, color="FFFFFF")
    sheet["A1"].fill = PatternFill("solid", fgColor="0F766E")
    sheet["A1"].alignment = Alignment(horizontal="left", vertical="center")
    sheet.row_dimensions[1].height = 34

    sheet["A2"] = "报告日期"
    sheet["B2"] = "{{report_date}}"
    sheet["D2"] = "负责人"
    sheet["E2"] = "{{owner}}"
    for address in ("A2", "D2"):
        sheet[address].font = Font(name=CJK_FONT, size=10, bold=True, color="115E59")
        sheet[address].fill = PatternFill("solid", fgColor="E6FFFB")
    for address in ("B2", "E2"):
        sheet[address].font = Font(name=CJK_FONT, size=10, color="172033")

    headers = ("任务", "负责角色", "状态", "开始", "截止", "进度")
    for column, header in enumerate(headers, start=1):
        cell = sheet.cell(row=4, column=column, value=header)
        cell.font = Font(name=CJK_FONT, size=10, bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="115E59")
        cell.alignment = Alignment(horizontal="center", vertical="center")
    rows = (
        ("需求冻结", "产品", "已完成", "2026-01-05", "2026-01-09", 1.0),
        ("交付开发", "研发", "进行中", "2026-01-10", "2026-01-24", 0.65),
        ("上线验收", "交付", "待开始", "2026-01-25", "2026-01-31", 0.0),
    )
    thin = Side(style="thin", color="D7DEE8")
    for row_index, row in enumerate(rows, start=5):
        for column_index, value in enumerate(row, start=1):
            cell = sheet.cell(row=row_index, column=column_index, value=value)
            cell.font = Font(name=CJK_FONT, size=10, color="172033")
            cell.border = Border(bottom=thin)
            cell.alignment = Alignment(
                horizontal="left" if column_index <= 3 else "center",
                vertical="center",
            )
        sheet.cell(row=row_index, column=6).number_format = "0%"
    sheet.auto_filter.ref = "A4:F7"

    validation = DataValidation(
        type="list",
        formula1='"待开始,进行中,已完成"',
        allow_blank=False,
    )
    sheet.add_data_validation(validation)
    validation.add("C5:C50")
    sheet.conditional_formatting.add(
        "C5:C50",
        FormulaRule(
            formula=['C5="已完成"'],
            fill=PatternFill("solid", fgColor="DCFCE7"),
            font=Font(color="166534"),
        ),
    )
    sheet.conditional_formatting.add(
        "C5:C50",
        FormulaRule(
            formula=['C5="进行中"'],
            fill=PatternFill("solid", fgColor="FEF3C7"),
            font=Font(color="92400E"),
        ),
    )

    sheet["H2"] = "平均完成率"
    sheet["H3"] = "=AVERAGE(F5:F7)"
    sheet["H2"].font = Font(name=CJK_FONT, size=10, bold=True, color="115E59")
    sheet["H3"].font = Font(name=CJK_FONT, size=18, bold=True, color="0F766E")
    sheet["H3"].number_format = "0%"

    chart = BarChart()
    chart.type = "col"
    chart.style = 10
    chart.title = "任务进度"
    chart.x_axis.title = "任务"
    chart.y_axis.title = "完成率"
    chart.y_axis.scaling.min = 0
    chart.y_axis.scaling.max = 1
    chart.y_axis.numFmt = "0%"
    chart.height = 4.4
    chart.width = 8.4
    data = Reference(sheet, min_col=6, min_row=4, max_row=7)
    categories = Reference(sheet, min_col=1, min_row=5, max_row=7)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(categories)
    chart.legend = None
    sheet.add_chart(chart, "A10")

    for column, width in {
        "A": 22,
        "B": 14,
        "C": 13,
        "D": 13,
        "E": 13,
        "F": 10,
        "G": 3,
        "H": 15,
        "I": 12,
        "J": 12,
        "K": 12,
        "L": 12,
        "M": 12,
    }.items():
        sheet.column_dimensions[column].width = width
    for row in range(4, 8):
        sheet.row_dimensions[row].height = 24
    sheet.print_area = "A1:I26"
    sheet.sheet_properties.pageSetUpPr.fitToPage = True
    sheet.page_setup.orientation = "landscape"
    sheet.page_setup.fitToWidth = 1
    sheet.page_setup.fitToHeight = 1
    sheet.page_margins.left = 0.25
    sheet.page_margins.right = 0.25
    sheet.page_margins.top = 0.35
    sheet.page_margins.bottom = 0.35

    buffer = BytesIO()
    workbook.save(buffer)
    return _normalize_ooxml(buffer.getvalue())


def _set_textbox(
    shape: Any,
    text: str,
    *,
    size: float,
    color: tuple[int, int, int],
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = PptxInches(0.12)
    frame.margin_right = PptxInches(0.12)
    frame.margin_top = PptxInches(0.06)
    frame.margin_bottom = PptxInches(0.06)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = alignment
    run = paragraph.add_run()
    run.text = text
    run.font.name = CJK_FONT
    run.font.size = PptxPt(size)
    run.font.bold = bold
    run.font.color.rgb = PptxRGBColor(*color)


def _make_pptx() -> bytes:
    presentation = Presentation()
    presentation.slide_width = PptxInches(13.333333)
    presentation.slide_height = PptxInches(7.5)
    properties = presentation.core_properties
    properties.author = "苏小有"
    properties.last_modified_by = "苏小有"
    properties.created = FIXED_TIMESTAMP
    properties.modified = FIXED_TIMESTAMP
    properties.title = "项目状态更新模板"

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = PptxRGBColor(0xF8, 0xFA, 0xFC)

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        PptxInches(0.18),
        presentation.slide_height,
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = PptxRGBColor(0x0F, 0x76, 0x6E)
    accent.line.fill.background()

    title = slide.shapes.add_textbox(
        PptxInches(0.78),
        PptxInches(0.58),
        PptxInches(11.7),
        PptxInches(0.8),
    )
    _set_textbox(
        title,
        "{{project_name}}",
        size=50,
        color=(0x0F, 0x17, 0x2A),
        bold=True,
    )
    subtitle = slide.shapes.add_textbox(
        PptxInches(0.8),
        PptxInches(1.42),
        PptxInches(11.4),
        PptxInches(0.4),
    )
    _set_textbox(
        subtitle,
        "{{period}}  ·  {{owner}}",
        size=20,
        color=(0x64, 0x74, 0x8B),
    )

    summary_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        PptxInches(0.8),
        PptxInches(2.12),
        PptxInches(6.1),
        PptxInches(3.95),
    )
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = PptxRGBColor(0xFF, 0xFF, 0xFF)
    summary_box.line.color.rgb = PptxRGBColor(0xD7, 0xDE, 0xE8)
    summary_box.line.width = PptxPt(1)
    summary_box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    _set_textbox(
        summary_box,
        "{{summary}}",
        size=24,
        color=(0x17, 0x20, 0x33),
    )

    table_shape = slide.shapes.add_table(
        3,
        2,
        PptxInches(7.22),
        PptxInches(2.12),
        PptxInches(5.28),
        PptxInches(3.95),
    )
    table = table_shape.table
    table.columns[0].width = PptxInches(1.35)
    table.columns[1].width = PptxInches(3.93)
    rows = (("状态", "{{status}}"), ("下一步", "{{next_step}}"), ("负责人", "{{owner}}"))
    for row_index, values in enumerate(rows):
        table.rows[row_index].height = PptxInches(1.31)
        for column_index, value in enumerate(values):
            cell = table.cell(row_index, column_index)
            cell.margin_left = PptxInches(0.16)
            cell.margin_right = PptxInches(0.16)
            cell.margin_top = PptxInches(0.1)
            cell.margin_bottom = PptxInches(0.1)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = PptxRGBColor(
                0xE6 if column_index == 0 else 0xFF,
                0xFF,
                0xFB if column_index == 0 else 0xFF,
            )
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            run = paragraph.add_run()
            run.text = value
            run.font.name = CJK_FONT
            run.font.size = PptxPt(18)
            run.font.bold = column_index == 0
            run.font.color.rgb = PptxRGBColor(
                *( (0x11, 0x5E, 0x59) if column_index == 0 else (0x17, 0x20, 0x33) )
            )

    footer = slide.shapes.add_textbox(
        PptxInches(0.8),
        PptxInches(6.72),
        PptxInches(11.7),
        PptxInches(0.28),
    )
    _set_textbox(
        footer,
        "苏小有 · 项目状态更新",
        size=16,
        color=(0x94, 0xA3, 0xB8),
        alignment=PP_ALIGN.RIGHT,
    )

    buffer = BytesIO()
    presentation.save(buffer)
    return _normalize_ooxml(buffer.getvalue())


def _placeholder(
    name: str,
    description: str,
    *,
    min_chars: int = 1,
    max_chars: int,
) -> dict[str, Any]:
    return {
        "name": name,
        "type": "text",
        "required": True,
        "min_chars": min_chars,
        "max_chars": max_chars,
        "description": description,
    }


def _entry(
    *,
    template_id: str,
    format_name: str,
    title: str,
    description: str,
    filename: str,
    content: bytes,
    placeholders: list[dict[str, Any]],
    baseline_id: str,
    unit_kind: str,
) -> dict[str, Any]:
    return {
        "template_id": template_id,
        "template_version": "1.0.0",
        "format": format_name,
        "title": title,
        "description": description,
        "asset_path": f"templates/{filename}",
        "source_sha256": hashlib.sha256(content).hexdigest(),
        "license": "Apache-2.0",
        "provenance": (
            "Generated deterministically by backend/scripts/"
            "generate_office_template_assets.py from first-party source"
        ),
        "placeholders": placeholders,
        "allowed_operations": ["instantiate_text"],
        "allowed_output_rules": {
            "extensions": [f".{format_name}"],
            "max_output_bytes": 20 * 1024 * 1024,
            "allow_overwrite": False,
        },
        "expected_render_baseline": {
            "baseline_id": baseline_id,
            "unit_kind": unit_kind,
            "min_units": 1,
            "max_units": 1,
        },
    }


def _deterministic(factory: Callable[[], bytes]) -> bytes:
    first = factory()
    second = factory()
    if first != second:
        raise RuntimeError(f"{factory.__name__} did not produce deterministic bytes")
    return first


def generate(output_root: Path, signing_key_path: Path) -> str:
    private = serialization.load_pem_private_key(
        signing_key_path.read_bytes(),
        password=None,
    )
    if not isinstance(private, Ed25519PrivateKey):
        raise TypeError("template signing key must be Ed25519")

    docx = _deterministic(_make_docx)
    xlsx = _deterministic(_make_xlsx)
    pptx = _deterministic(_make_pptx)
    inspect_ooxml_package(
        docx,
        "docx",
        expected_placeholders=(
            "classification",
            "next_step",
            "owner",
            "recipient",
            "report_date",
            "summary",
            "title",
        ),
    )
    inspect_ooxml_package(
        xlsx,
        "xlsx",
        expected_placeholders=("owner", "project_name", "report_date"),
    )
    inspect_ooxml_package(
        pptx,
        "pptx",
        expected_placeholders=(
            "next_step",
            "owner",
            "period",
            "project_name",
            "status",
            "summary",
        ),
    )

    entries = [
        _entry(
            template_id="business-brief",
            format_name="docx",
            title="商务简报",
            description="单页执行摘要与行动信息模板。",
            filename="business-brief.docx",
            content=docx,
            placeholders=[
                _placeholder("classification", "页脚分类标记。", max_chars=40),
                _placeholder("next_step", "明确的下一步行动。", max_chars=600),
                _placeholder("owner", "主要负责人或团队。", max_chars=80),
                _placeholder("recipient", "简报面向的受众。", max_chars=120),
                _placeholder("report_date", "报告日期文本。", max_chars=40),
                _placeholder("summary", "执行摘要正文。", max_chars=2400),
                _placeholder("title", "简报标题。", max_chars=160),
            ],
            baseline_id="business-brief-1.0.0-default",
            unit_kind="pages",
        ),
        _entry(
            template_id="project-tracker",
            format_name="xlsx",
            title="项目跟踪表",
            description="含公式、状态验证、条件格式和原生进度图的单表模板。",
            filename="project-tracker.xlsx",
            content=xlsx,
            placeholders=[
                _placeholder("owner", "项目负责人。", max_chars=80),
                _placeholder("project_name", "项目名称。", max_chars=120),
                _placeholder("report_date", "报告日期文本。", max_chars=40),
            ],
            baseline_id="project-tracker-1.0.0-default",
            unit_kind="worksheets",
        ),
        _entry(
            template_id="status-update",
            format_name="pptx",
            title="项目状态更新",
            description="单页项目状态、摘要和行动信息演示模板。",
            filename="status-update.pptx",
            content=pptx,
            placeholders=[
                _placeholder("next_step", "下一步行动。", max_chars=300),
                _placeholder("owner", "项目负责人。", max_chars=80),
                _placeholder("period", "状态更新周期。", max_chars=80),
                _placeholder("project_name", "项目名称。", max_chars=120),
                _placeholder("status", "当前状态。", max_chars=120),
                _placeholder("summary", "可直接向受众展示的状态摘要。", max_chars=700),
            ],
            baseline_id="status-update-1.0.0-default",
            unit_kind="slides",
        ),
    ]
    catalog = {
        "schema_version": 1,
        "catalog_id": "suxiaoyou-office-templates",
        "catalog_version": "1.0.0",
        "templates": entries,
    }
    catalog_bytes = _canonical_json(catalog)
    signature = private.sign(catalog_bytes)
    signature_envelope = {
        "schema_version": 1,
        "algorithm": "Ed25519",
        "key_id": KEY_ID,
        "catalog_sha256": hashlib.sha256(catalog_bytes).hexdigest(),
        "signature": base64.b64encode(signature).decode("ascii"),
    }

    template_root = output_root / "templates"
    template_root.mkdir(parents=True, exist_ok=True)
    (template_root / "business-brief.docx").write_bytes(docx)
    (template_root / "project-tracker.xlsx").write_bytes(xlsx)
    (template_root / "status-update.pptx").write_bytes(pptx)
    (output_root / "catalog.json").write_bytes(catalog_bytes)
    (output_root / "catalog.sig.json").write_bytes(
        _canonical_json(signature_envelope)
    )
    public = private.public_key().public_bytes(
        encoding=serialization.Encoding.Raw,
        format=serialization.PublicFormat.Raw,
    )
    return base64.b64encode(public).decode("ascii")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output-root", type=Path, required=True)
    parser.add_argument("--signing-key", type=Path, required=True)
    arguments = parser.parse_args()
    public_key = generate(
        arguments.output_root.expanduser().resolve(),
        arguments.signing_key.expanduser().resolve(strict=True),
    )
    print(f"trust-root {KEY_ID} {public_key}")


if __name__ == "__main__":
    main()
