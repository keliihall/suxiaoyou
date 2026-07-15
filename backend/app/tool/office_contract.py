"""Offline create/edit/reopen contract for the restricted Office tool.

The release pipeline executes this module through the frozen PyInstaller
backend on every native target.  It deliberately exercises the public tool
boundary rather than the format helpers so missing bundled dependencies,
platform-specific atomic replacement problems, and file-version regressions
all fail the installer job before an artifact can be published.
"""

from __future__ import annotations

import asyncio
import hashlib
import os
import platform
import re
import subprocess
import sys
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Awaitable, Callable

from app.schemas.agent import AgentInfo
from app.storage.file_versions import FileVersionStore
from app.tool.builtin.office import OfficeTool
from app.tool.context import ToolContext


OFFICE_CONTRACT_SCHEMA_VERSION = 1
OFFICE_CONTRACT_VERSION = "v1.0-restricted-office-1"
SUPPORTED_PLATFORMS = frozenset(
    {
        "windows-x64",
        "macos-arm64",
        "macos-x64",
        "linux-arm64",
        "linux-x64",
    }
)
_COMMIT_PATTERN = re.compile(r"^(?!0{40}$)[0-9a-f]{40}$")


class OfficeContractError(RuntimeError):
    """The release Office contract did not satisfy an asserted postcondition."""


def native_platform_id(
    *, system_name: str | None = None, machine_name: str | None = None
) -> str:
    """Return the scorecard platform ID for the executing native runtime."""

    system = (system_name or platform.system()).strip().casefold()
    machine = (machine_name or platform.machine()).strip().casefold()
    operating_system = {
        "darwin": "macos",
        "linux": "linux",
        "windows": "windows",
    }.get(system)
    architecture = {
        "amd64": "x64",
        "x86_64": "x64",
        "arm64": "arm64",
        "aarch64": "arm64",
    }.get(machine)
    result = f"{operating_system}-{architecture}"
    if result not in SUPPORTED_PLATFORMS:
        raise OfficeContractError(
            f"Unsupported native Office contract platform: system={system!r}, "
            f"machine={machine!r}"
        )
    return result


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _source_commit() -> str:
    candidate = os.environ.get("SUXIAOYOU_RELEASE_COMMIT", "").strip().casefold()
    if _COMMIT_PATTERN.fullmatch(candidate):
        return candidate
    try:
        candidate = subprocess.run(
            ["git", "rev-parse", "HEAD^{commit}"],
            check=True,
            capture_output=True,
            text=True,
            timeout=10,
        ).stdout.strip().casefold()
    except (OSError, subprocess.SubprocessError):
        return "unknown"
    return candidate if _COMMIT_PATTERN.fullmatch(candidate) else "unknown"


def _context(workspace: Path, call_id: str) -> ToolContext:
    return ToolContext(
        session_id="office-release-contract",
        message_id="office-release-contract",
        agent=AgentInfo(
            name="office-release-contract",
            description="Offline release verification",
            mode="primary",
        ),
        call_id=call_id,
        language="en",
        workspace=str(workspace),
    )


async def _execute(
    tool: OfficeTool,
    workspace: Path,
    call_id: str,
    args: dict[str, Any],
) -> dict[str, Any]:
    result = await tool(args, _context(workspace, call_id))
    if not result.success:
        raise OfficeContractError(result.error or "Office tool failed without an error")
    if result.metadata.get("reopened_and_validated") is not True:
        raise OfficeContractError("Office tool did not report reopen validation")
    if result.metadata.get("atomic_install") is not True:
        raise OfficeContractError("Office tool did not report atomic installation")
    return result.metadata


def _verify_version_snapshot(
    workspace: Path,
    target: Path,
    initial_sha256: str,
    edit_metadata: dict[str, Any],
) -> None:
    if edit_metadata.get("previous_sha256") != initial_sha256:
        raise OfficeContractError("Office edit did not version the pre-edit bytes")
    version_id = edit_metadata.get("previous_version_id")
    if not isinstance(version_id, str) or not version_id:
        raise OfficeContractError("Office edit did not return a durable version ID")
    versions = FileVersionStore(workspace).list_versions(file_path=target)
    if not any(
        item.id == version_id
        and item.sha256 == initial_sha256
        and item.operation == "office.edit"
        for item in versions
    ):
        raise OfficeContractError("Office edit version is absent from durable history")


def _format_result(
    *,
    target: Path,
    initial_sha256: str,
    edit_metadata: dict[str, Any],
) -> dict[str, Any]:
    final_sha256 = _sha256(target)
    if final_sha256 == initial_sha256:
        raise OfficeContractError("Office edit left the file byte-for-byte unchanged")
    return {
        "created": True,
        "edited": True,
        "reopened_and_validated": True,
        "independent_reopen_validated": True,
        "atomic_install": True,
        "version_snapshot_verified": True,
        "initial_sha256": initial_sha256,
        "final_sha256": final_sha256,
        "final_size": target.stat().st_size,
        "previous_version_id": edit_metadata["previous_version_id"],
    }


async def _exercise_docx(workspace: Path) -> dict[str, Any]:
    from docx import Document

    tool = OfficeTool()
    created = await _execute(
        tool,
        workspace,
        "office-contract-docx-create",
        {
            "file_path": "office-contract.docx",
            "operation": "create",
            "document": {
                "title": "suyo Office contract",
                "paragraphs": [
                    {"text": "DOCX-CREATE", "style": "heading1"},
                    {"text": "Cross-platform declarative document"},
                ],
                "tables": [
                    {
                        "headers": ["stage", "value"],
                        "rows": [["created", 1]],
                    }
                ],
            },
        },
    )
    target = Path(created["file_path"])
    initial_sha256 = _sha256(target)
    edited = await _execute(
        tool,
        workspace,
        "office-contract-docx-edit",
        {
            "file_path": str(target),
            "operation": "edit",
            "document": {"paragraphs": [{"text": "DOCX-APPEND"}]},
            "replacements": [
                {"old_text": "DOCX-CREATE", "new_text": "DOCX-EDITED"}
            ],
        },
    )
    _verify_version_snapshot(workspace, target, initial_sha256, edited)
    reopened = Document(str(target))
    paragraphs = [paragraph.text for paragraph in reopened.paragraphs]
    if "DOCX-EDITED" not in paragraphs or "DOCX-APPEND" not in paragraphs:
        raise OfficeContractError("Independent DOCX reopen lost edited content")
    if reopened.tables[0].cell(1, 0).text != "created":
        raise OfficeContractError("Independent DOCX reopen lost the original table")
    return _format_result(
        target=target,
        initial_sha256=initial_sha256,
        edit_metadata=edited,
    )


async def _exercise_xlsx(workspace: Path) -> dict[str, Any]:
    from openpyxl import load_workbook

    tool = OfficeTool()
    created = await _execute(
        tool,
        workspace,
        "office-contract-xlsx-create",
        {
            "file_path": "office-contract.xlsx",
            "operation": "create",
            "workbook": {
                "sheets": [
                    {
                        "name": "Data",
                        "rows": [
                            ["stage", "value"],
                            ["XLSX-CREATE", 2],
                            ["formula", "=B2*2"],
                        ],
                    },
                    {"name": "Remove", "rows": [["temporary"]]},
                ],
                "cells": [
                    {
                        "sheet": "Data",
                        "cell": "B2",
                        "style": {
                            "number_format": "0.00",
                            "font": {"bold": True},
                        },
                    }
                ],
            },
        },
    )
    target = Path(created["file_path"])
    initial_sha256 = _sha256(target)
    edited = await _execute(
        tool,
        workspace,
        "office-contract-xlsx-edit",
        {
            "file_path": str(target),
            "operation": "edit",
            "workbook": {
                "delete_sheets": ["Remove"],
                "sheets": [
                    {
                        "name": "Data",
                        "action": "append",
                        "rows": [["XLSX-APPEND", 3]],
                    },
                    {
                        "name": "Notes",
                        "action": "create",
                        "rows": [["XLSX-EDITED"]],
                    },
                ],
                "cells": [
                    {"sheet": "Data", "cell": "A2", "value": "XLSX-EDITED"}
                ],
            },
        },
    )
    _verify_version_snapshot(workspace, target, initial_sha256, edited)
    reopened = load_workbook(
        str(target),
        read_only=False,
        data_only=False,
        keep_vba=False,
        keep_links=False,
    )
    try:
        if reopened.sheetnames != ["Data", "Notes"]:
            raise OfficeContractError("Independent XLSX reopen lost sheet edits")
        if reopened["Data"]["A2"].value != "XLSX-EDITED":
            raise OfficeContractError("Independent XLSX reopen lost a cell edit")
        if reopened["Data"]["A4"].value != "XLSX-APPEND":
            raise OfficeContractError("Independent XLSX reopen lost an appended row")
        if reopened["Data"]["B3"].value != "=B2*2":
            raise OfficeContractError("Independent XLSX reopen lost the formula")
        if reopened["Data"]["B2"].font.bold is not True:
            raise OfficeContractError("Independent XLSX reopen lost the basic style")
        if reopened["Notes"]["A1"].value != "XLSX-EDITED":
            raise OfficeContractError("Independent XLSX reopen lost the new sheet")
    finally:
        reopened.close()
    return _format_result(
        target=target,
        initial_sha256=initial_sha256,
        edit_metadata=edited,
    )


async def _exercise_pptx(workspace: Path) -> dict[str, Any]:
    from pptx import Presentation

    tool = OfficeTool()
    created = await _execute(
        tool,
        workspace,
        "office-contract-pptx-create",
        {
            "file_path": "office-contract.pptx",
            "operation": "create",
            "presentation": {
                "slides": [
                    {
                        "title": "suyo Office contract",
                        "bullets": ["PPTX-CREATE", "Declarative presentation"],
                        "tables": [
                            {
                                "left_inches": 1,
                                "top_inches": 4,
                                "width_inches": 5,
                                "height_inches": 1.5,
                                "headers": ["stage", "value"],
                                "rows": [["created", 1]],
                            }
                        ],
                    }
                ]
            },
        },
    )
    target = Path(created["file_path"])
    initial_sha256 = _sha256(target)
    edited = await _execute(
        tool,
        workspace,
        "office-contract-pptx-edit",
        {
            "file_path": str(target),
            "operation": "edit",
            "presentation": {
                "slides": [{"title": "PPTX-APPEND", "subtitle": "Validated"}]
            },
            "replacements": [
                {"old_text": "PPTX-CREATE", "new_text": "PPTX-EDITED"}
            ],
        },
    )
    _verify_version_snapshot(workspace, target, initial_sha256, edited)
    reopened = Presentation(str(target))
    if len(reopened.slides) != 2:
        raise OfficeContractError("Independent PPTX reopen lost an appended slide")
    first_slide_text = "\n".join(
        shape.text
        for shape in reopened.slides[0].shapes
        if hasattr(shape, "text")
    )
    if "PPTX-EDITED" not in first_slide_text:
        raise OfficeContractError("Independent PPTX reopen lost the text edit")
    if reopened.slides[1].shapes.title.text != "PPTX-APPEND":
        raise OfficeContractError("Independent PPTX reopen lost the appended title")
    tables = [
        shape.table
        for shape in reopened.slides[0].shapes
        if getattr(shape, "has_table", False)
    ]
    if len(tables) != 1 or tables[0].cell(1, 0).text != "created":
        raise OfficeContractError("Independent PPTX reopen lost the original table")
    return _format_result(
        target=target,
        initial_sha256=initial_sha256,
        edit_metadata=edited,
    )


async def run_office_contract(
    *,
    expected_platform: str | None = None,
    source_commit: str | None = None,
    release_ref: str | None = None,
) -> dict[str, Any]:
    """Run every format offline and return machine-readable release evidence."""

    started_at = datetime.now(timezone.utc)
    native_platform = native_platform_id()
    if expected_platform is not None:
        if expected_platform not in SUPPORTED_PLATFORMS:
            raise OfficeContractError(
                f"Unknown expected Office contract platform: {expected_platform}"
            )
        if expected_platform != native_platform:
            raise OfficeContractError(
                f"Office contract ran on {native_platform}, expected {expected_platform}"
            )

    exercises: tuple[
        tuple[str, Callable[[Path], Awaitable[dict[str, Any]]]], ...
    ] = (
        ("docx", _exercise_docx),
        ("xlsx", _exercise_xlsx),
        ("pptx", _exercise_pptx),
    )
    results: dict[str, dict[str, Any]] = {}
    previous_private_dir = os.environ.get("SUXIAOYOU_PRIVATE_DATA_DIR")
    try:
        with tempfile.TemporaryDirectory(prefix="suyo-office-contract-") as temporary:
            root = Path(temporary)
            workspace = root / "workspace"
            private = root / "private"
            workspace.mkdir()
            # Mirror the managed-workspace invariant before the first Office
            # operation. Windows targeted transactions intentionally reject a
            # missing output parent instead of creating it after validation.
            (workspace / "suxiaoyou_written").mkdir()
            private.mkdir()
            os.environ["SUXIAOYOU_PRIVATE_DATA_DIR"] = str(private)
            for format_name, exercise in exercises:
                try:
                    results[format_name] = await exercise(workspace)
                except Exception as exc:
                    results[format_name] = {
                        "created": False,
                        "edited": False,
                        "reopened_and_validated": False,
                        "independent_reopen_validated": False,
                        "atomic_install": False,
                        "version_snapshot_verified": False,
                        "error_type": type(exc).__name__,
                        "error": str(exc),
                    }
    finally:
        if previous_private_dir is None:
            os.environ.pop("SUXIAOYOU_PRIVATE_DATA_DIR", None)
        else:
            os.environ["SUXIAOYOU_PRIVATE_DATA_DIR"] = previous_private_dir

    all_passed = all(
        all(
            result.get(field) is True
            for field in (
                "created",
                "edited",
                "reopened_and_validated",
                "independent_reopen_validated",
                "atomic_install",
                "version_snapshot_verified",
            )
        )
        for result in results.values()
    ) and set(results) == {"docx", "xlsx", "pptx"}
    completed_at = datetime.now(timezone.utc)
    return {
        "schema_version": OFFICE_CONTRACT_SCHEMA_VERSION,
        "contract_version": OFFICE_CONTRACT_VERSION,
        "status": "ok" if all_passed else "failed",
        "all_passed": all_passed,
        "platform": native_platform,
        "source_commit": (source_commit or _source_commit()).strip().casefold(),
        "release_ref": release_ref or os.environ.get("GITHUB_REF_NAME") or "local",
        "started_at": started_at.isoformat(),
        "completed_at": completed_at.isoformat(),
        "runner": {
            "system": platform.system(),
            "machine": platform.machine(),
            "python_version": platform.python_version(),
            "frozen_backend": bool(getattr(sys, "frozen", False)),
        },
        "formats": results,
    }


def run_office_contract_sync(
    *,
    expected_platform: str | None = None,
    source_commit: str | None = None,
    release_ref: str | None = None,
) -> dict[str, Any]:
    """Synchronous CLI/test wrapper."""

    return asyncio.run(
        run_office_contract(
            expected_platform=expected_platform,
            source_commit=source_commit,
            release_ref=release_ref,
        )
    )
