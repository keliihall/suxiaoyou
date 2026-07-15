"""Read tool — read file contents with optional offset/limit paging."""

from __future__ import annotations

import csv
import errno
import heapq
import json
import logging
import os
import re
import stat
from pathlib import Path
from typing import Any

from app.tool.base import ToolDefinition, ToolResult
from app.tool.context import ToolContext
from app.tool.extractors import extract_document, is_supported_binary
from app.tool.workspace import WorkspaceViolation, resolve_and_validate

logger = logging.getLogger(__name__)

_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp", ".svg"}
_DATA_EXTENSIONS = {".csv", ".tsv", ".xlsx"}
_DATA_SAMPLE_ROWS = 5
_MAX_DATA_ROWS = 2000
_MAX_SELECTED_PAGES = 200
_MAX_IMAGE_BYTES = 20 * 1024 * 1024
_MAX_DIRECTORY_ENTRIES = 1000
_READ_CHUNK_BYTES = 1024 * 1024
_WINDOWS_REPARSE_POINT = 0x400
_PAGE_TOKEN_RE = re.compile(r"^(\d+)(?:\s*-\s*(\d+))?$")


class ReadTool(ToolDefinition):

    @property
    def id(self) -> str:
        return "read"

    @property
    def is_concurrency_safe(self) -> bool:
        return True

    @property
    def description(self) -> str:
        return (
            "Read a file from the filesystem. Supports offset/limit for paging "
            "through large files. Can also list directory contents. "
            "Handles text, PDF, DOCX, XLSX, PPTX, CSV/TSV, and common image formats. "
            "Use 'pages' to read specific PDF pages/PPTX slides or one named XLSX sheet. "
            "Use 'format=json' for bounded structured XLSX/CSV/TSV output."
        )

    def parameters_schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Absolute or relative path to the file to read",
                },
                "offset": {
                    "type": "integer",
                    "description": "Line number to start reading from (1-based)",
                    "default": 1,
                    "minimum": 1,
                },
                "limit": {
                    "type": "integer",
                    "description": "Maximum lines or spreadsheet rows to read (1-2000)",
                    "default": 2000,
                    "minimum": 1,
                    "maximum": _MAX_DATA_ROWS,
                },
                "pages": {
                    "type": "string",
                    "description": "Page range for PDFs/PPTX (e.g. '1-3' or '5'), or sheet name for XLSX (e.g. 'Revenue')",
                },
                "format": {
                    "type": "string",
                    "enum": ["json"],
                    "description": "Set to 'json' for structured spreadsheet output (XLSX/CSV/TSV)",
                },
            },
            "required": ["file_path"],
        }

    async def execute(self, args: dict[str, Any], ctx: ToolContext) -> ToolResult:
        file_path = args.get("file_path", "")
        if not isinstance(file_path, str) or not file_path.strip():
            return ToolResult(error="file_path must be a non-empty string")

        try:
            self._paging(args)
        except ValueError as exc:
            return ToolResult(error=str(exc))

        ext = os.path.splitext(file_path)[1].lower()
        pages_requested = "pages" in args and args.get("pages") is not None
        if pages_requested and (
            not isinstance(args.get("pages"), str) or not args["pages"].strip()
        ):
            return ToolResult(error="pages must be a non-empty string")

        output_format = args.get("format")
        if output_format is not None and output_format != "json":
            return ToolResult(error="format must be 'json'")
        if output_format and ext not in _DATA_EXTENSIONS:
            return ToolResult(error="format=json is supported only for XLSX, CSV, and TSV files")
        if pages_requested and ext not in {".pdf", ".pptx", ".xlsx"}:
            return ToolResult(
                error="pages is supported only for PDF pages, PPTX slides, or an XLSX sheet name"
            )

        if ext == ".xls":
            return ToolResult(
                error="Legacy .xls files are not supported; save the workbook as .xlsx or CSV first"
            )

        # Images: return as base64 for multimodal LLM
        if ext in _IMAGE_EXTENSIONS:
            return self._read_image(file_path, ctx)

        # Data files: return a bounded summary/page or structured JSON.
        if ext in _DATA_EXTENSIONS:
            return self._read_data_file(file_path, ext, args, ctx)

        return await self._filesystem_execute(args, ctx)

    # ------------------------------------------------------------------
    # Data files: schema + sample (CSV, XLSX)
    # ------------------------------------------------------------------

    def _read_data_file(
        self,
        file_path: str,
        ext: str,
        args: dict[str, Any],
        ctx: ToolContext,
    ) -> ToolResult:
        """Return a bounded schema/sample, row page, or structured JSON."""
        try:
            resolved, source = self._resolve_read_path(file_path, ctx)
        except WorkspaceViolation as e:
            return ToolResult(error=str(e))

        if not os.path.exists(resolved):
            return ToolResult(error=f"File not found: {file_path}")

        try:
            if ext in (".csv", ".tsv"):
                return self._summarise_csv(resolved, file_path, args, source)
            else:
                return self._summarise_xlsx(resolved, file_path, args, source)
        except Exception as e:
            return ToolResult(error=f"Cannot read {os.path.basename(file_path)}: {e}")

    def _summarise_csv(
        self,
        resolved: str,
        file_path: str,
        args: dict[str, Any],
        source: str,
    ) -> ToolResult:
        offset, requested_limit = self._paging(args)
        structured = args.get("format") == "json"
        explicit_paging = "offset" in args or "limit" in args
        row_limit = requested_limit if structured or explicit_paging else _DATA_SAMPLE_ROWS

        with open(resolved, "r", encoding="utf-8", errors="replace") as f:
            # Sniff delimiter
            sample_text = f.read(8192)
            f.seek(0)
            try:
                dialect = csv.Sniffer().sniff(sample_text)
                delimiter = dialect.delimiter
            except csv.Error:
                delimiter = "," if file_path.lower().endswith(".csv") else "\t"

            reader = csv.reader(f, delimiter=delimiter)
            headers = next(reader, None)
            if headers is None:
                empty_output = (
                    json.dumps(
                        {"file": os.path.basename(file_path), "columns": [], "rows": []},
                        ensure_ascii=False,
                        indent=2,
                    )
                    if structured
                    else "(Empty file)"
                )
                return ToolResult(
                    output=empty_output,
                    title=os.path.basename(file_path),
                    metadata={
                        "source": source,
                        "format": "csv" if Path(file_path).suffix.lower() == ".csv" else "tsv",
                        "total_rows": 0,
                        "columns": [],
                        "shown": 0,
                    },
                )

            data_rows: list[list[str]] = []
            total_rows = 0
            for row_number, row in enumerate(reader, 1):
                total_rows = row_number
                if offset <= row_number < offset + row_limit:
                    data_rows.append(row)

        metadata = {
            "source": source,
            "format": "csv" if Path(file_path).suffix.lower() == ".csv" else "tsv",
            "total_rows": total_rows,
            "columns": headers,
            "offset": offset,
            "shown": len(data_rows),
            "has_more": offset - 1 + len(data_rows) < total_rows,
        }

        if structured:
            return ToolResult(
                output=json.dumps(
                    {
                        "file": os.path.basename(file_path),
                        "columns": headers,
                        "offset": offset,
                        "total_rows": total_rows,
                        "rows": data_rows,
                    },
                    ensure_ascii=False,
                    indent=2,
                ),
                title=os.path.basename(file_path),
                metadata=metadata,
            )

        # Build output
        parts = [
            f"File: {os.path.basename(file_path)}",
            f"Rows: {total_rows:,}  |  Columns: {len(headers)}",
            "",
            "Columns: " + ", ".join(headers),
            "",
            (
                f"Rows {offset}-{offset + len(data_rows) - 1}:"
                if explicit_paging and data_rows
                else "Sample rows:"
            ),
        ]
        display_headers = [self._cell_text(header) for header in headers]
        # Header
        parts.append(" | ".join(display_headers))
        parts.append(" | ".join(["---"] * len(headers)))
        for row in data_rows:
            # Pad/truncate to match header count
            display_row = [self._cell_text(cell) for cell in row]
            padded = display_row[:len(headers)] + [""] * max(0, len(headers) - len(row))
            parts.append(" | ".join(padded))

        parts.append("")
        if metadata["has_more"]:
            parts.append(f"More rows are available after row {offset + len(data_rows) - 1}.")

        return ToolResult(
            output="\n".join(parts),
            title=os.path.basename(file_path),
            metadata=metadata,
        )

    def _summarise_xlsx(
        self,
        resolved: str,
        file_path: str,
        args: dict[str, Any],
        source: str,
    ) -> ToolResult:
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise ImportError("openpyxl is not installed")

        wb = load_workbook(resolved, read_only=True, data_only=True)
        try:
            selected_sheet = args.get("pages")
            if selected_sheet and selected_sheet not in wb.sheetnames:
                return ToolResult(
                    error=(
                        f"Sheet not found: {selected_sheet}. Available sheets: "
                        + ", ".join(wb.sheetnames)
                    )
                )
            sheet_names = [selected_sheet] if selected_sheet else list(wb.sheetnames)
            structured = args.get("format") == "json"
            explicit_paging = "offset" in args or "limit" in args
            offset, requested_limit = self._paging(args)
            remaining_budget = requested_limit

            parts = [
                f"File: {os.path.basename(file_path)}",
                f"Sheets: {', '.join(wb.sheetnames)}",
                "",
            ]
            structured_sheets: list[dict[str, Any]] = []
            all_metadata: dict[str, Any] = {
                "source": source,
                "format": "xlsx",
                "selected_sheet": selected_sheet,
                "offset": offset,
                "row_limit": requested_limit,
                "row_limit_scope": "workbook" if structured or explicit_paging else "sample_per_sheet",
                "sheets": {},
            }

            for sheet_name in sheet_names:
                ws = wb[sheet_name]
                row_iter = ws.iter_rows(values_only=True)
                header_values = next(row_iter, None)
                if header_values is None or (
                    all(value is None for value in header_values) and (ws.max_row or 0) <= 1
                ):
                    headers: list[str] = []
                    data_rows: list[list[Any]] = []
                    total_rows = 0
                else:
                    headers = [self._cell_text(value) for value in header_values]
                    data_rows = []
                    total_rows = 0
                    if structured or explicit_paging:
                        sheet_limit = remaining_budget
                    else:
                        sheet_limit = _DATA_SAMPLE_ROWS
                    for row_number, row in enumerate(row_iter, 1):
                        total_rows = row_number
                        if offset <= row_number < offset + sheet_limit:
                            data_rows.append(list(row))
                    if structured or explicit_paging:
                        remaining_budget = max(0, remaining_budget - len(data_rows))

                sheet_metadata = {
                    "total_rows": total_rows,
                    "columns": headers,
                    "shown": len(data_rows),
                    "has_more": offset - 1 + len(data_rows) < total_rows,
                }
                all_metadata["sheets"][sheet_name] = sheet_metadata
                structured_sheets.append(
                    {
                        "name": sheet_name,
                        "columns": headers,
                        "offset": offset,
                        "total_rows": total_rows,
                        "rows": data_rows,
                        "has_more": sheet_metadata["has_more"],
                    }
                )

                if not structured:
                    if not headers:
                        parts.append(f"=== {sheet_name}: (empty) ===")
                        continue
                    parts.append(
                        f"=== {sheet_name} ({total_rows:,} rows, {len(headers)} cols) ==="
                    )
                    parts.append("Columns: " + ", ".join(headers))
                    parts.append("")
                    parts.append(" | ".join(headers))
                    parts.append(" | ".join(["---"] * len(headers)))
                    for row in data_rows:
                        cells = [self._cell_text(value) for value in row]
                        padded = cells[:len(headers)] + [""] * max(0, len(headers) - len(cells))
                        parts.append(" | ".join(padded))
                    if sheet_metadata["has_more"]:
                        if data_rows:
                            parts.append(
                                f"More rows are available after row {offset + len(data_rows) - 1}."
                            )
                        else:
                            parts.append("Rows omitted because the workbook row limit was reached.")
                    parts.append("")

            if structured:
                output = json.dumps(
                    {
                        "file": os.path.basename(file_path),
                        "available_sheets": wb.sheetnames,
                        "sheets": structured_sheets,
                    },
                    ensure_ascii=False,
                    indent=2,
                    default=str,
                )
            else:
                output = "\n".join(parts)

            return ToolResult(
                output=output,
                title=os.path.basename(file_path),
                metadata=all_metadata,
            )
        finally:
            wb.close()

    # ------------------------------------------------------------------
    # Image path (base64 for multimodal LLM)
    # ------------------------------------------------------------------

    def _read_image(self, file_path: str, ctx: ToolContext) -> ToolResult:
        """Return an image as base64 for the LLM to see visually.

        Stores the data URL in metadata so the message builder can convert
        the tool result into multimodal content (text + image_url).
        """
        import base64
        import mimetypes

        source = "filesystem"
        try:
            resolved, source = self._resolve_read_path(file_path, ctx)
        except WorkspaceViolation as exc:
            return ToolResult(error=str(exc))

        try:
            raw = self._read_bounded_image(resolved, file_path)
            b64 = base64.b64encode(raw).decode("utf-8")
            mime_type, _ = mimetypes.guess_type(resolved)
            if not mime_type or not mime_type.startswith("image/"):
                ext = os.path.splitext(resolved)[1].lstrip(".")
                mime_type = f"image/{ext}"

            data_url = f"data:{mime_type};base64,{b64}"

            return ToolResult(
                output=f"[Image: {os.path.basename(file_path)}]",
                title=os.path.basename(file_path),
                metadata={
                    "source": source,
                    "format": os.path.splitext(file_path)[1].lower(),
                    "image_data_url": data_url,
                },
            )
        except FileNotFoundError:
            return ToolResult(error=f"Image not found: {file_path}")
        except PermissionError:
            return ToolResult(error=f"Permission denied: {file_path}")
        except ValueError as exc:
            return ToolResult(error=str(exc))
        except OSError as e:
            return ToolResult(error=f"Cannot read image {file_path}: {e}")

    # ------------------------------------------------------------------
    # Filesystem fallback path
    # ------------------------------------------------------------------

    async def _filesystem_execute(self, args: dict[str, Any], ctx: ToolContext) -> ToolResult:
        file_path = args["file_path"]

        # Workspace restriction check
        try:
            file_path, source = self._resolve_read_path(file_path, ctx)
        except WorkspaceViolation as e:
            return ToolResult(error=str(e))

        offset, limit = self._paging(args)

        try:
            path_info = os.lstat(file_path)
        except FileNotFoundError:
            return ToolResult(error=f"File not found: {file_path}")
        except PermissionError:
            return ToolResult(error=f"Permission denied: {file_path}")

        if self._is_link_or_reparse(path_info):
            return ToolResult(
                error=f"Path changed to a symbolic link during validation: {file_path}"
            )

        # Directory listing
        if stat.S_ISDIR(path_info.st_mode):
            if source == "session_attachment":
                return ToolResult(error="Session attachment access does not allow directory listing")
            try:
                entries, total_entries = self._bounded_directory_listing(file_path)
                listing = "\n".join(entries)
                directory_name = os.path.basename(file_path)
                truncated = total_entries > len(entries)
                return ToolResult(
                    output=listing,
                    title=ctx.tr(
                        f"已列出 {directory_name} 中的 {len(entries)} 个条目",
                        f"Listed {len(entries)} entries in {directory_name}",
                    ),
                    metadata={
                        "source": source,
                        "total_entries": total_entries,
                        "shown": len(entries),
                        "entry_limit": _MAX_DIRECTORY_ENTRIES,
                        "directory_truncated": truncated,
                        "truncated": truncated,
                        "has_more": truncated,
                    },
                )
            except PermissionError:
                return ToolResult(error=f"Permission denied: {file_path}")
            except FileNotFoundError:
                return ToolResult(error=f"Directory not found: {file_path}")
            except NotADirectoryError:
                return ToolResult(error=f"Directory changed while being listed: {file_path}")
            except ValueError as exc:
                return ToolResult(error=str(exc))
            except OSError as exc:
                return ToolResult(error=f"Cannot list directory {file_path}: {exc}")

        # Binary document extraction (PDF, DOCX, XLSX, PPTX)
        if is_supported_binary(file_path):
            try:
                selected_pages: list[int] | None = None
                total_pages: int | None = None
                ext = Path(file_path).suffix.lower()
                if args.get("pages") and ext == ".pdf":
                    text, selected_pages, total_pages = self._extract_pdf_pages(
                        file_path, args["pages"]
                    )
                elif args.get("pages") and ext == ".pptx":
                    text, selected_pages, total_pages = self._extract_pptx_slides(
                        file_path, args["pages"]
                    )
                else:
                    text = extract_document(file_path)
            except ImportError as e:
                return ToolResult(error=str(e))
            except ValueError as e:
                return ToolResult(error=str(e))
            except Exception as e:
                return ToolResult(
                    error=f"Cannot read {os.path.basename(file_path)}: {e}"
                )

            result = self._format_lines(text, file_path, offset, limit)
            if result.metadata is None:
                result.metadata = {}
            result.metadata["source"] = source
            if selected_pages is not None:
                selection_key = "selected_pages" if ext == ".pdf" else "selected_slides"
                total_key = "total_pages" if ext == ".pdf" else "total_slides"
                result.metadata[selection_key] = selected_pages
                result.metadata[total_key] = total_pages
            return result

        # Text file reading
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                selected: list[tuple[int, str]] = []
                total_lines = 0
                for line_number, line in enumerate(f, 1):
                    total_lines = line_number
                    if offset <= line_number < offset + limit:
                        selected.append((line_number, line))

            # Format with line numbers (cat -n style)
            output_lines = []
            for line_number, line in selected:
                # Truncate very long lines
                line_content = line.rstrip("\n\r")
                if len(line_content) > 2000:
                    line_content = line_content[:2000] + "..."
                output_lines.append(f"{line_number:>6}\t{line_content}")

            output = "\n".join(output_lines)

            last_shown = offset - 1 + len(selected)
            if last_shown < total_lines:
                output += f"\n\n... ({total_lines - last_shown} more lines)"

            return ToolResult(
                output=output,
                title=os.path.basename(file_path),
                metadata={"total_lines": total_lines, "shown": len(selected), "source": source},
            )

        except UnicodeDecodeError:
            return ToolResult(error=f"Cannot read binary file: {file_path}")
        except PermissionError:
            return ToolResult(error=f"Permission denied: {file_path}")

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _is_link_or_reparse(info: os.stat_result) -> bool:
        return stat.S_ISLNK(info.st_mode) or bool(
            getattr(info, "st_file_attributes", 0) & _WINDOWS_REPARSE_POINT
        )

    @classmethod
    def _read_bounded_image(cls, resolved: str, display_path: str) -> bytes:
        """Read one stable regular image without following a replaced link."""

        flags = (
            os.O_RDONLY
            | getattr(os, "O_BINARY", 0)
            | getattr(os, "O_NOFOLLOW", 0)
            | getattr(os, "O_NONBLOCK", 0)
        )
        try:
            fd = os.open(resolved, flags)
        except OSError as exc:
            if exc.errno == errno.ELOOP:
                raise ValueError(
                    f"Image changed to a symbolic link during validation: {display_path}"
                ) from exc
            raise
        try:
            before = os.fstat(fd)
            if not stat.S_ISREG(before.st_mode):
                raise ValueError(f"Image path is not a regular file: {display_path}")
            if before.st_size > _MAX_IMAGE_BYTES:
                raise ValueError(
                    f"Image exceeds the {_MAX_IMAGE_BYTES // (1024 * 1024)} MiB "
                    f"read limit: {display_path}"
                )

            chunks: list[bytes] = []
            total = 0
            while True:
                remaining = _MAX_IMAGE_BYTES + 1 - total
                if remaining <= 0:
                    raise ValueError(
                        f"Image exceeds the {_MAX_IMAGE_BYTES // (1024 * 1024)} MiB "
                        f"read limit: {display_path}"
                    )
                chunk = os.read(fd, min(_READ_CHUNK_BYTES, remaining))
                if not chunk:
                    break
                total += len(chunk)
                chunks.append(chunk)
                if total > _MAX_IMAGE_BYTES:
                    raise ValueError(
                        f"Image exceeds the {_MAX_IMAGE_BYTES // (1024 * 1024)} MiB "
                        f"read limit: {display_path}"
                    )

            after = os.fstat(fd)
            try:
                current = os.lstat(resolved)
            except FileNotFoundError as exc:
                raise ValueError(
                    f"Image changed while being read; try again: {display_path}"
                ) from exc
            if cls._is_link_or_reparse(current) or not cls._same_file_identity(
                after,
                current,
            ):
                raise ValueError(
                    f"Image changed while being read; try again: {display_path}"
                )
            if cls._mutable_file_identity(before) != cls._mutable_file_identity(after):
                raise ValueError(
                    f"Image changed while being read; try again: {display_path}"
                )
            return b"".join(chunks)
        finally:
            os.close(fd)

    @classmethod
    def _bounded_directory_listing(cls, directory: str) -> tuple[list[str], int]:
        """Return the globally sorted first N names with O(N) memory."""

        flags = (
            os.O_RDONLY
            | getattr(os, "O_DIRECTORY", 0)
            | getattr(os, "O_NOFOLLOW", 0)
            | getattr(os, "O_NONBLOCK", 0)
        )
        directory_fd: int | None = None
        before: os.stat_result
        if os.name == "nt":  # pragma: no cover - exercised by Windows CI
            before = os.lstat(directory)
            if cls._is_link_or_reparse(before) or not stat.S_ISDIR(before.st_mode):
                raise NotADirectoryError(directory)
            scan_target: str | int = directory
        else:
            try:
                directory_fd = os.open(directory, flags)
            except OSError as exc:
                if exc.errno == errno.ELOOP:
                    raise ValueError(
                        f"Directory changed to a symbolic link during validation: {directory}"
                    ) from exc
                raise
            before = os.fstat(directory_fd)
            if not stat.S_ISDIR(before.st_mode):
                os.close(directory_fd)
                raise NotADirectoryError(directory)
            scan_target = directory_fd

        total_entries = 0
        try:
            with os.scandir(scan_target) as iterator:

                def names():
                    nonlocal total_entries
                    for entry in iterator:
                        total_entries += 1
                        yield entry.name

                selected = heapq.nsmallest(
                    _MAX_DIRECTORY_ENTRIES + 1,
                    names(),
                )

            after = (
                os.fstat(directory_fd)
                if directory_fd is not None
                else os.lstat(directory)
            )
            try:
                current = os.lstat(directory)
            except FileNotFoundError as exc:
                raise ValueError(
                    f"Directory changed while being listed; try again: {directory}"
                ) from exc
            if (
                cls._is_link_or_reparse(current)
                or not cls._same_file_identity(after, current)
                or cls._mutable_directory_identity(before)
                != cls._mutable_directory_identity(after)
            ):
                raise ValueError(
                    f"Directory changed while being listed; try again: {directory}"
                )
            return selected[:_MAX_DIRECTORY_ENTRIES], total_entries
        finally:
            if directory_fd is not None:
                os.close(directory_fd)

    @staticmethod
    def _same_file_identity(left: os.stat_result, right: os.stat_result) -> bool:
        return (left.st_dev, left.st_ino) == (right.st_dev, right.st_ino)

    @staticmethod
    def _mutable_file_identity(
        info: os.stat_result,
    ) -> tuple[int, int, int, int, int, int]:
        return (
            info.st_dev,
            info.st_ino,
            info.st_size,
            info.st_mtime_ns,
            info.st_ctime_ns,
            info.st_mode,
        )

    @staticmethod
    def _mutable_directory_identity(info: os.stat_result) -> tuple[int, int, int, int]:
        return (
            info.st_dev,
            info.st_ino,
            info.st_mtime_ns,
            info.st_ctime_ns,
        )

    @staticmethod
    def _paging(args: dict[str, Any]) -> tuple[int, int]:
        offset = args.get("offset", 1)
        limit = args.get("limit", _MAX_DATA_ROWS)
        if isinstance(offset, bool) or not isinstance(offset, int) or offset < 1:
            raise ValueError("offset must be an integer greater than or equal to 1")
        if (
            isinstance(limit, bool)
            or not isinstance(limit, int)
            or limit < 1
            or limit > _MAX_DATA_ROWS
        ):
            raise ValueError(f"limit must be an integer between 1 and {_MAX_DATA_ROWS}")
        return offset, limit

    @staticmethod
    def _cell_text(value: Any) -> str:
        if value is None:
            return ""
        text = str(value).replace("\r", " ").replace("\n", " ")
        return text if len(text) <= 2000 else text[:2000] + "..."

    @staticmethod
    def _resolve_read_path(file_path: str, ctx: ToolContext) -> tuple[str, str]:
        """Resolve a workspace path or an exact user-selected attachment.

        A registered attachment authorizes only that canonical file path. It
        never grants access to a sibling path or to the contents of a selected
        directory.
        """

        try:
            return resolve_and_validate(file_path, ctx.workspace), "filesystem"
        except WorkspaceViolation:
            candidate = str(Path(file_path).expanduser().resolve())
            if candidate not in ctx.attachment_paths:
                raise
            return candidate, "session_attachment"

    @staticmethod
    def _parse_page_selection(selection: str, total: int, label: str) -> list[int]:
        if total < 1:
            raise ValueError(f"Cannot select {label}: the document has no {label}")

        selected: list[int] = []
        seen: set[int] = set()
        for raw_token in selection.split(","):
            token = raw_token.strip()
            match = _PAGE_TOKEN_RE.fullmatch(token)
            if not match:
                raise ValueError(
                    f"Invalid {label} selection '{selection}'; use values such as '1-3,5'"
                )
            start = int(match.group(1))
            end = int(match.group(2) or start)
            if start < 1 or end < start or end > total:
                raise ValueError(
                    f"Invalid {label} selection '{selection}'; valid range is 1-{total}"
                )
            for page_number in range(start, end + 1):
                if page_number not in seen:
                    selected.append(page_number)
                    seen.add(page_number)
                if len(selected) > _MAX_SELECTED_PAGES:
                    raise ValueError(
                        f"At most {_MAX_SELECTED_PAGES} {label} may be selected per read"
                    )
        return selected

    @classmethod
    def _extract_pdf_pages(
        cls, file_path: str, selection: str
    ) -> tuple[str, list[int], int]:
        try:
            from pypdf import PdfReader
        except ImportError:
            raise ImportError("Cannot read PDF files: pypdf is not installed")

        reader = PdfReader(file_path)
        selected = cls._parse_page_selection(selection, len(reader.pages), "PDF pages")
        parts: list[str] = []
        for page_number in selected:
            text = reader.pages[page_number - 1].extract_text() or ""
            parts.append(
                f"--- Page {page_number} ---\n"
                + (text if text.strip() else "(No text content found on this page)")
            )
        return "\n\n".join(parts), selected, len(reader.pages)

    @classmethod
    def _extract_pptx_slides(
        cls, file_path: str, selection: str
    ) -> tuple[str, list[int], int]:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("Cannot read PPTX files: python-pptx is not installed")

        presentation = Presentation(file_path)
        selected = cls._parse_page_selection(selection, len(presentation.slides), "PPTX slides")
        parts: list[str] = []
        for slide_number in selected:
            slide = presentation.slides[slide_number - 1]
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.extend(
                        paragraph.text.strip()
                        for paragraph in shape.text_frame.paragraphs
                        if paragraph.text.strip()
                    )
                if shape.has_table:
                    for row in shape.table.rows:
                        texts.append(" | ".join(cell.text.strip() for cell in row.cells))
            parts.append(
                f"--- Slide {slide_number} ---\n"
                + ("\n".join(texts) if texts else "(No text content found on this slide)")
            )
        return "\n\n".join(parts), selected, len(presentation.slides)

    @staticmethod
    def _format_lines(
        text: str, file_path: str, offset: int, limit: int
    ) -> ToolResult:
        """Format extracted text with line numbers, applying offset/limit."""
        lines = text.split("\n")
        total_lines = len(lines)

        start = offset - 1
        end = start + limit
        selected = lines[start:end]

        output_lines = []
        for i, line in enumerate(selected, start=offset):
            line_content = line.rstrip("\n\r")
            if len(line_content) > 2000:
                line_content = line_content[:2000] + "..."
            output_lines.append(f"{i:>6}\t{line_content}")

        output = "\n".join(output_lines)

        if end < total_lines:
            output += f"\n\n... ({total_lines - end} more lines)"

        ext = os.path.splitext(file_path)[1].lower()
        return ToolResult(
            output=output,
            title=os.path.basename(file_path),
            metadata={
                "total_lines": total_lines,
                "shown": len(selected),
                "format": ext,
            },
        )
