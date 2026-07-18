"""Restricted declarative creation and editing for OOXML Office files.

This tool intentionally exposes a small data model instead of arbitrary Python
or shell execution.  Every output is written beside the destination, reopened
with the corresponding Office library, inspected for unsafe OOXML features,
and only then atomically installed.
"""

from __future__ import annotations

import asyncio
import copy
from dataclasses import dataclass, field, replace
import hashlib
import io
import json
import logging
import math
import os
import re
import secrets
import stat
import tempfile
import zipfile
from collections.abc import Iterable, Iterator, Mapping, Sequence
from pathlib import Path, PurePosixPath
from typing import Any, Final
from xml.etree import ElementTree

from sqlalchemy import select, text
from sqlalchemy.ext.asyncio import AsyncSession, async_sessionmaker

from app import release_features
from app.models.office_user_template import OfficeUserTemplate
from app.models.session import Session
from app.models.workspace_instance import WorkspaceInstance
from app.office_templates.errors import OfficeTemplateError, TemplateContractError
from app.office_templates.user import (
    UserOfficeTemplateService,
    UserTemplatePlaceholder,
    decode_user_template_placeholder_schema,
    get_user_office_template_service,
    validate_user_template_ref,
    validate_user_template_values,
)
from app.storage.checkpoints import inspect_workspace_identity
from app.storage.file_versions import FileVersionStore, default_file_version_storage_root
from app.office_validation.draft import OfficeDraftSeal, OfficeDraftValidationResult
from app.office_validation.errors import OfficeValidationError
from app.office_validation.models import OfficeValidationReport
from app.office_validation.precommit import (
    OfficeCreateValidationPlan,
    OfficeEditMutationIntent,
    OfficePrecommitCoordinator,
    OfficePrecommitRequest,
    OfficePrecommitValidationSession,
    get_office_precommit_coordinator,
)
from app.office_validation.precommit_repair import (
    OfficePrecommitRepairError,
    OfficePrecommitRepairRequest,
    OfficePrecommitRepairer,
    build_precommit_repair_request,
    copy_replacement_args,
)
from app.tool.base import ToolDefinition, ToolResult
from app.tool.context import ToolContext
from app.tool.file_metadata import (
    UnsupportedFileMetadataError,
    ensure_mutation_metadata_supported,
)
from app.tool.file_versioning import version_metadata
from app.tool.workspace import WorkspaceViolation, resolve_and_validate, resolve_for_write
from app.tool.workspace_transaction import (
    WorkspaceMutationError,
    WorkspaceMutationTransaction,
)


logger = logging.getLogger(__name__)


_OFFICE_REPAIR_TIMEOUT_SECONDS: Final[float] = 120.0
_OFFICE_REPAIR_SETTLEMENT_GRACE_SECONDS: Final[float] = 1.0
_OFFICE_REPAIR_TOKEN_PREFIX: Final[str] = "sxy-office-repair:v1:"
_OFFICE_REPAIR_PATH_SCAN_MAX_DEPTH: Final[int] = 48
_OFFICE_REPAIR_PATH_SCAN_MAX_NODES: Final[int] = 600_000


class _OfficeRepairTimeoutError(TimeoutError):
    """A bounded repair worker exceeded its server-owned deadline."""


_FORMATS = {
    ".docx": "document",
    ".xlsx": "workbook",
    ".pptx": "presentation",
}
_MIME_TYPES = {
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def _office_precommit_coordinator(
    ctx: ToolContext,
) -> OfficePrecommitCoordinator | None:
    """Read only a server-injected coordinator, never a tool argument."""

    app_state = getattr(ctx, "_app_state", None) or {}
    if "office_precommit_coordinator" in app_state:
        coordinator = app_state.get("office_precommit_coordinator")
    else:
        coordinator = get_office_precommit_coordinator()
    if coordinator is None:
        return None
    if not isinstance(coordinator, OfficePrecommitCoordinator):
        raise TypeError("Office precommit coordinator is invalid")
    return coordinator


def _office_precommit_repairer(ctx: ToolContext) -> OfficePrecommitRepairer | None:
    """Read an optional repair boundary only from trusted server state."""

    app_state = getattr(ctx, "_app_state", None) or {}
    repairer = app_state.get("office_precommit_repairer")
    if repairer is None:
        return None
    if not isinstance(repairer, OfficePrecommitRepairer):
        raise TypeError("Office precommit repairer is invalid")
    return repairer


def _office_precommit_request(
    args: Mapping[str, Any],
    ctx: ToolContext,
    file_path: str,
    *,
    user_template: _UserTemplateLease | None = None,
) -> OfficePrecommitRequest:
    """Bind a v1.1 write to checkpoint and transaction-owned identities."""

    identities = (
        ctx.root_turn_id,
        ctx.turn_run_id,
        ctx.checkpoint_id,
        ctx.workspace_instance_id,
    )
    if any(not isinstance(value, str) or not value.strip() for value in identities):
        raise OfficeInputError(
            "Office v1.1 写入缺少受信任的 checkpoint 身份，原文件未更改。",
            (
                "Office v1.1 writing is missing its trusted checkpoint identity; "
                "the original file was not changed."
            ),
        )
    workspace = Path(ctx.workspace or "").resolve()
    relative = Path(file_path).relative_to(workspace).as_posix()
    raw_template = args.get("first_party_template")
    template_id: str | None = None
    template_version: str | None = None
    if isinstance(raw_template, Mapping):
        raw_id = raw_template.get("template_id")
        raw_version = raw_template.get("template_version")
        if isinstance(raw_id, str) and isinstance(raw_version, str):
            template_id = raw_id
            template_version = raw_version
    trusted_create_plan: OfficeCreateValidationPlan | None = None
    if user_template is not None:
        template_id = user_template.template_ref
        template_version = str(user_template.revision)
        trusted_create_plan = user_template.precommit_plan
        if trusted_create_plan is None:
            raise TypeError("User template lease has no trusted create plan")
    suffix = Path(file_path).suffix.casefold()
    return OfficePrecommitRequest(
        operation=str(args["operation"]),  # type: ignore[arg-type]
        document_format=suffix[1:],  # type: ignore[arg-type]
        relative_path=relative,
        session_id=ctx.session_id,
        message_id=ctx.message_id,
        call_id=ctx.call_id,
        root_turn_id=identities[0],
        turn_run_id=identities[1],
        checkpoint_id=identities[2],
        workspace_instance_id=identities[3],
        template_id=template_id,
        template_version=template_version,
        trusted_create_plan=trusted_create_plan,
    )


def _office_edit_mutation_intent(
    args: Mapping[str, Any],
    summary: Mapping[str, Any],
    document_format: str,
) -> OfficeEditMutationIntent:
    """Derive a narrow visual/structural envelope from the completed writer."""

    if args.get("operation") != "edit" or document_format not in {
        "docx",
        "xlsx",
        "pptx",
    }:
        raise OfficeInputError(
            "Office 编辑意图无效。",
            "The Office edit intent is invalid.",
        )

    def count(name: str) -> int:
        value = summary.get(name, 0)
        if (
            not isinstance(value, int)
            or isinstance(value, bool)
            or value < 0
        ):
            raise TypeError(f"Office writer summary field {name} is invalid")
        return value

    replacements = _parse_replacements(args.get("replacements"))
    removed_characters = sum(
        max(0, len(item["old_text"]) - len(item["new_text"]))
        for item in replacements
    )

    if document_format == "pptx":
        slides_added = count("slides_added")
        changed_ratio = 0.45 if replacements else 0.0
        return OfficeEditMutationIntent(
            document_format="pptx",
            max_added_pages=slides_added,
            max_removed_pages=0,
            required_page_delta=slides_added,
            expected_logical_unit_delta=slides_added,
            max_outside_changed_ratio=changed_ratio,
            max_total_changed_ratio=changed_ratio,
            max_blank_fraction_increase=0.10,
        )

    if document_format == "xlsx":
        sheets_created = count("sheets_created")
        sheets_deleted = count("sheets_deleted")
        rows_appended = count("rows_appended")
        cells_written = count("cells_written")
        charts_added = count("charts_added")
        layout_changes = sum(
            count(name)
            for name in (
                "merged_ranges_added",
                "dimensions_changed",
                "freeze_panes_changed",
                "auto_filters_changed",
                "conditional_formats_added",
                "data_validations_added",
                "named_ranges_added",
            )
        )
        max_added = min(
            50,
            sheets_created * 8
            + charts_added * 2
            + math.ceil(rows_appended / 50)
            + math.ceil(cells_written / 500)
            + int(layout_changes > 0),
        )
        max_removed = min(
            50,
            sheets_deleted * 8
            + min(2, math.ceil(cells_written / 500))
            + int(layout_changes > 0),
        )
        broad = bool(sheets_created or sheets_deleted or charts_added or layout_changes)
        changed_ratio = 0.55 if broad else 0.30
        return OfficeEditMutationIntent(
            document_format="xlsx",
            max_added_pages=max_added,
            max_removed_pages=max_removed,
            expected_logical_unit_delta=sheets_created - sheets_deleted,
            max_outside_changed_ratio=changed_ratio,
            max_total_changed_ratio=changed_ratio,
            max_blank_fraction_increase=0.10,
        )

    paragraphs_added = count("paragraphs_added")
    page_breaks_added = count("page_breaks_added")
    tables_added = count("tables_added")
    images_added = count("images_added")
    charts_added = count("static_charts_added")
    sections_changed = count("sections_changed")
    appended_units = paragraphs_added + tables_added + images_added + charts_added
    max_added = min(
        50,
        page_breaks_added
        + math.ceil(appended_units / 8)
        + sections_changed * 4,
    )
    max_removed = min(
        50,
        math.ceil(removed_characters / 2_000) + sections_changed * 4,
    )
    changed_ratio = (
        0.75
        if sections_changed
        else 0.60
        if replacements
        else 0.45
    )
    return OfficeEditMutationIntent(
        document_format="docx",
        max_added_pages=max_added,
        max_removed_pages=max_removed,
        max_outside_changed_ratio=changed_ratio,
        max_total_changed_ratio=changed_ratio,
        max_blank_fraction_increase=0.10,
    )


async def _await_transaction_thread(
    function: Any,
    *args: Any,
) -> Any:
    """Settle a transaction worker before propagating coroutine cancellation."""

    worker = asyncio.create_task(asyncio.to_thread(function, *args))
    cancellation: asyncio.CancelledError | None = None
    while True:
        try:
            result = await asyncio.shield(worker)
            break
        except asyncio.CancelledError as exc:
            cancellation = cancellation or exc
            if worker.done():
                break
    try:
        result = worker.result()
    except BaseException as exc:
        if cancellation is not None:
            raise cancellation from exc
        raise
    if cancellation is not None:
        raise cancellation
    return result


async def _await_repair_worker(
    repairer: OfficePrecommitRepairer,
    request: OfficePrecommitRepairRequest,
    *,
    timeout_seconds: float | None = None,
    settlement_grace_seconds: float | None = None,
) -> Any:
    """Bound one capability-free repair call without trusting cancellation.

    Production uses code-owned limits.  Optional values exist solely for
    deterministic unit tests and are never read from tool/model arguments or
    application configuration.  A worker that suppresses cancellation may be
    detached after the grace period because its only input is an immutable,
    JSON-only request with no staging path, handle, policy, or seal authority.
    """

    _assert_detachable_repair_request(request)
    timeout, grace = _repair_worker_limits(
        timeout_seconds=timeout_seconds,
        settlement_grace_seconds=settlement_grace_seconds,
    )
    worker = asyncio.ensure_future(repairer.repair(request))
    try:
        done, _pending = await asyncio.wait({worker}, timeout=timeout)
    except asyncio.CancelledError as cancellation:
        await _cancel_and_settle_repair_worker(worker, grace=grace)
        raise cancellation
    if done:
        return worker.result()

    settlement_cancellation = await _cancel_and_settle_repair_worker(
        worker,
        grace=grace,
    )
    if settlement_cancellation is not None:
        raise settlement_cancellation
    raise _OfficeRepairTimeoutError("Office repair worker timed out")


def _repair_worker_limits(
    *,
    timeout_seconds: float | None,
    settlement_grace_seconds: float | None,
) -> tuple[float, float]:
    production_timeout = timeout_seconds is None
    timeout = (
        _OFFICE_REPAIR_TIMEOUT_SECONDS
        if production_timeout
        else timeout_seconds
    )
    grace = (
        _OFFICE_REPAIR_SETTLEMENT_GRACE_SECONDS
        if settlement_grace_seconds is None
        else settlement_grace_seconds
    )
    if (
        not isinstance(timeout, (int, float))
        or isinstance(timeout, bool)
        or not math.isfinite(float(timeout))
        or float(timeout) <= 0
        or float(timeout) > 300
        or (production_timeout and float(timeout) < 1)
    ):
        raise RuntimeError("Office repair timeout is invalid")
    if (
        not isinstance(grace, (int, float))
        or isinstance(grace, bool)
        or not math.isfinite(float(grace))
        or float(grace) <= 0
        or float(grace) > 10
    ):
        raise RuntimeError("Office repair settlement grace is invalid")
    return float(timeout), float(grace)


def _assert_detachable_repair_request(
    request: OfficePrecommitRepairRequest,
) -> None:
    """Prove the detached worker owns immutable data, never staging authority."""

    from types import MappingProxyType

    if not isinstance(request, OfficePrecommitRepairRequest):
        raise TypeError("Office repair request is invalid")

    def immutable_json(value: object) -> bool:
        if value is None or isinstance(value, (str, bool, int)):
            return True
        if isinstance(value, float):
            return math.isfinite(value)
        if isinstance(value, MappingProxyType):
            return all(
                isinstance(key, str) and immutable_json(item)
                for key, item in value.items()
            )
        if isinstance(value, tuple):
            return all(immutable_json(item) for item in value)
        return False

    if not immutable_json(request.tokenized_args):
        raise TypeError("Office repair request contains mutable authority")


async def _cancel_and_settle_repair_worker(
    worker: asyncio.Future[Any],
    *,
    grace: float,
) -> asyncio.CancelledError | None:
    """Cancel once, wait to one fixed deadline, then safely detach."""

    worker.cancel()
    deadline = asyncio.get_running_loop().time() + grace
    cancellation: asyncio.CancelledError | None = None
    while not worker.done():
        remaining = deadline - asyncio.get_running_loop().time()
        if remaining <= 0:
            break
        try:
            await asyncio.wait({worker}, timeout=remaining)
        except asyncio.CancelledError as exc:
            cancellation = cancellation or exc
    if worker.done():
        _consume_repair_worker_result(worker)
    else:
        worker.add_done_callback(_consume_repair_worker_result)
    return cancellation


def _consume_repair_worker_result(worker: asyncio.Future[Any]) -> None:
    """Retrieve a detached worker exception without logging private text."""

    try:
        worker.exception()
    except BaseException:
        pass


async def _await_transaction_commit(
    transaction: WorkspaceMutationTransaction,
    *,
    seal: object | None,
) -> tuple[Any, asyncio.CancelledError | None]:
    """Return the known commit outcome even when cancellation arrived mid-commit."""

    function = (
        transaction.commit
        if seal is None
        else lambda: transaction.commit_with_precommit_office_seal(seal)
    )
    worker = asyncio.create_task(asyncio.to_thread(function))
    cancellation: asyncio.CancelledError | None = None
    while True:
        try:
            result = await asyncio.shield(worker)
            break
        except asyncio.CancelledError as exc:
            cancellation = cancellation or exc
            if worker.done():
                break
    try:
        result = worker.result()
    except BaseException as exc:
        if cancellation is not None:
            transaction.abort()
            raise cancellation from exc
        raise
    return result, cancellation


@dataclass(frozen=True, slots=True)
class _OfficeRepairPathField:
    """One schema-owned filesystem field in an Office argument object."""

    location: tuple[str | int, ...]
    domain: str


@dataclass(frozen=True, slots=True)
class _OfficeRepairPathTokenTable:
    """Per-call path authority retained only by ``OfficeTool.execute``."""

    nonce: str = field(repr=False)
    target_token: str
    target_raw_path: str = field(repr=False)
    read_token_raw_paths: tuple[tuple[str, str], ...] = field(repr=False)


def _is_sequence_container(value: object) -> bool:
    return isinstance(value, Sequence) and not isinstance(
        value,
        (str, bytes, bytearray, memoryview),
    )


def _is_supported_office_path_location(
    location: tuple[str | int, ...],
) -> str | None:
    if location == ("file_path",):
        return "target"
    if (
        len(location) == 4
        and location[0] == "document"
        and location[1] in {"images", "charts"}
        and isinstance(location[2], int)
        and location[3] == "path"
    ):
        return "read"
    if (
        len(location) == 6
        and location[0] == "presentation"
        and location[1] == "slides"
        and isinstance(location[2], int)
        and location[3] == "images"
        and isinstance(location[4], int)
        and location[5] == "path"
    ):
        return "read"
    return None


def _looks_like_office_path_key(key: str) -> bool:
    return key in {"path", "paths"} or key.endswith("_path") or key.endswith(
        "_paths"
    )


def _enumerate_office_path_fields(
    args: Mapping[str, Any],
    *,
    reject_unsupported: bool = True,
) -> tuple[_OfficeRepairPathField, ...]:
    """Enumerate every supported path and fail closed on new structured ones.

    The same enumeration feeds sparse staging, request tokenization, and
    replacement unmasking.  Consequently a future Office schema path cannot
    cross the repair boundary until it is explicitly assigned a token domain.
    """

    fields: list[_OfficeRepairPathField] = []
    active: set[int] = set()
    nodes = 0

    def visit(value: object, location: tuple[str | int, ...], depth: int) -> None:
        nonlocal nodes
        is_mapping = isinstance(value, Mapping)
        is_sequence = _is_sequence_container(value)
        if not is_mapping and not is_sequence:
            return
        nodes += 1
        if (
            depth > _OFFICE_REPAIR_PATH_SCAN_MAX_DEPTH
            or nodes > _OFFICE_REPAIR_PATH_SCAN_MAX_NODES
        ):
            raise OfficePrecommitRepairError(
                "Office repair path structure exceeds bounds"
            )
        if isinstance(value, Mapping):
            identity = id(value)
            if identity in active:
                raise OfficePrecommitRepairError(
                    "Office repair path structure contains a reference cycle"
                )
            active.add(identity)
            try:
                for key, child in value.items():
                    if not isinstance(key, str):
                        raise OfficePrecommitRepairError(
                            "Office repair argument keys must be strings"
                        )
                    child_location = location + (key,)
                    domain = _is_supported_office_path_location(child_location)
                    if domain is not None:
                        fields.append(
                            _OfficeRepairPathField(
                                location=child_location,
                                domain=domain,
                            )
                        )
                    elif (
                        reject_unsupported
                        and _looks_like_office_path_key(key)
                        and (
                            len(child_location) == 1
                            or child_location[0]
                            in {"document", "workbook", "presentation"}
                        )
                    ):
                        raise OfficePrecommitRepairError(
                            "Office repair encountered an unsupported path field"
                        )
                    visit(child, child_location, depth + 1)
            finally:
                active.remove(identity)
            return
        if _is_sequence_container(value):
            identity = id(value)
            if identity in active:
                raise OfficePrecommitRepairError(
                    "Office repair path structure contains a reference cycle"
                )
            active.add(identity)
            try:
                for index, child in enumerate(value):
                    visit(child, location + (index,), depth + 1)
            finally:
                active.remove(identity)

    visit(args, (), 0)
    target_fields = [item for item in fields if item.domain == "target"]
    if len(target_fields) != 1:
        raise OfficePrecommitRepairError(
            "Office repair requires exactly one target path field"
        )
    return tuple(fields)


def _office_path_value(
    args: Mapping[str, Any],
    location: tuple[str | int, ...],
) -> object:
    value: object = args
    for component in location:
        if isinstance(component, int):
            if not _is_sequence_container(value):
                raise OfficePrecommitRepairError(
                    "Office repair path structure changed unexpectedly"
                )
            value = value[component]  # type: ignore[index]
        else:
            if not isinstance(value, Mapping) or component not in value:
                raise OfficePrecommitRepairError(
                    "Office repair path structure changed unexpectedly"
                )
            value = value[component]
    return value


def _set_office_path_value(
    args: dict[str, Any],
    location: tuple[str | int, ...],
    replacement: str,
) -> None:
    value: Any = args
    for component in location[:-1]:
        value = value[component]
    value[location[-1]] = replacement


def _declared_local_image_paths(
    args: Mapping[str, Any],
    workspace: str,
) -> tuple[str, ...]:
    """Collect existing local image inputs for sparse transactional staging."""

    resolved: set[str] = set()
    for path_field in _enumerate_office_path_fields(
        args,
        reject_unsupported=False,
    ):
        if path_field.domain != "read":
            continue
        raw_path = _office_path_value(args, path_field.location)
        if not isinstance(raw_path, str) or not raw_path.strip() or "://" in raw_path:
            continue
        try:
            path = resolve_and_validate(raw_path, workspace)
        except WorkspaceViolation:
            continue
        if Path(path).is_file():
            resolved.add(path)
    return tuple(sorted(resolved))


def _tokenize_office_repair_args(
    args: Mapping[str, Any],
) -> tuple[dict[str, Any], _OfficeRepairPathTokenTable]:
    """Replace target/read paths with per-call, domain-separated tokens."""

    copied = copy_replacement_args(args)
    fields = _enumerate_office_path_fields(copied)
    nonce = secrets.token_urlsafe(24)
    target_token = f"{_OFFICE_REPAIR_TOKEN_PREFIX}target:{nonce}"
    target_raw_path: str | None = None
    read_token_by_raw_path: dict[str, str] = {}
    read_token_raw_paths: list[tuple[str, str]] = []
    for path_field in fields:
        raw_path = _office_path_value(copied, path_field.location)
        if not isinstance(raw_path, str) or not raw_path.strip():
            raise OfficePrecommitRepairError(
                "Office repair path fields must be non-empty strings"
            )
        if path_field.domain == "target":
            target_raw_path = raw_path
            token = target_token
        else:
            token = read_token_by_raw_path.get(raw_path, "")
            if not token:
                token = (
                    f"{_OFFICE_REPAIR_TOKEN_PREFIX}read:{nonce}:"
                    f"{len(read_token_raw_paths)}"
                )
                read_token_by_raw_path[raw_path] = token
                read_token_raw_paths.append((token, raw_path))
        _set_office_path_value(copied, path_field.location, token)
    if target_raw_path is None:
        raise OfficePrecommitRepairError("Office repair target path is missing")
    table = _OfficeRepairPathTokenTable(
        nonce=nonce,
        target_token=target_token,
        target_raw_path=target_raw_path,
        read_token_raw_paths=tuple(read_token_raw_paths),
    )
    return copied, table


def _reject_misplaced_office_repair_tokens(
    value: object,
    *,
    path_locations: frozenset[tuple[str | int, ...]],
    location: tuple[str | int, ...] = (),
) -> None:
    if isinstance(value, str):
        if (
            value.startswith(_OFFICE_REPAIR_TOKEN_PREFIX)
            and location not in path_locations
        ):
            raise OfficePrecommitRepairError(
                "Office repair token appeared outside a path field"
            )
        return
    if isinstance(value, Mapping):
        for key, child in value.items():
            if isinstance(key, str):
                _reject_misplaced_office_repair_tokens(
                    child,
                    path_locations=path_locations,
                    location=location + (key,),
                )
        return
    if _is_sequence_container(value):
        for index, child in enumerate(value):
            _reject_misplaced_office_repair_tokens(
                child,
                path_locations=path_locations,
                location=location + (index,),
            )


def _unmask_office_replacement_args(
    replacement: dict[str, Any],
    token_table: _OfficeRepairPathTokenTable,
) -> dict[str, Any]:
    """Accept only this call's exact target/read tokens, then restore paths."""

    fields = _enumerate_office_path_fields(replacement)
    path_locations = frozenset(item.location for item in fields)
    _reject_misplaced_office_repair_tokens(
        replacement,
        path_locations=path_locations,
    )
    raw_path_by_read_token = dict(token_table.read_token_raw_paths)
    for path_field in fields:
        token = _office_path_value(replacement, path_field.location)
        if not isinstance(token, str):
            raise OfficePrecommitRepairError(
                "Office repair returned a non-token path value"
            )
        if path_field.domain == "target":
            if token != token_table.target_token:
                raise OfficePrecommitRepairError(
                    "Office repair target token is invalid"
                )
            raw_path = token_table.target_raw_path
        else:
            raw_path = raw_path_by_read_token.get(token)
            if raw_path is None:
                raise OfficePrecommitRepairError(
                    "Office repair read token is invalid"
                )
        _set_office_path_value(replacement, path_field.location, raw_path)
    return replacement


def _office_template_identity(args: Mapping[str, Any]) -> tuple[str, str] | None:
    if "first_party_template" not in args:
        return None
    template = args.get("first_party_template")
    if not isinstance(template, Mapping):
        raise OfficePrecommitRepairError(
            "Office repair template identity is invalid"
        )
    template_id = template.get("template_id")
    template_version = template.get("template_version")
    if not isinstance(template_id, str) or not isinstance(template_version, str):
        raise OfficePrecommitRepairError(
            "Office repair template identity is invalid"
        )
    return template_id, template_version


def _office_repair_semantic_projection(
    args: Mapping[str, Any],
) -> dict[str, Any]:
    """Remove only reviewed presentation/layout fields from an argument copy.

    Everything else is immutable across automatic repair.  This deny-by-
    default projection means a future Office schema field is semantic until a
    reviewed source change explicitly classifies it as layout-only.
    """

    projected = copy_replacement_args(args)

    def mapping(value: object) -> dict[str, Any] | None:
        return value if isinstance(value, dict) else None

    def mappings(value: object) -> Iterator[dict[str, Any]]:
        if isinstance(value, list):
            for item in value:
                if isinstance(item, dict):
                    yield item

    document = mapping(projected.get("document"))
    if document is not None:
        for paragraph in mappings(document.get("paragraphs")):
            for key in ("style", "page_break_after", "format", "list"):
                paragraph.pop(key, None)
            for run in mappings(paragraph.get("runs")):
                for key in (
                    "bold",
                    "italic",
                    "underline",
                    "color",
                    "size",
                    "font",
                ):
                    run.pop(key, None)
        for table in mappings(document.get("tables")):
            table.pop("format", None)
        for image in mappings(document.get("images")):
            image.pop("width_inches", None)
        for section in mappings(document.get("sections")):
            for key in ("start", "orientation", "paper_size", "margins"):
                section.pop(key, None)
        for chart in mappings(document.get("charts")):
            chart.pop("width_inches", None)

    workbook = mapping(projected.get("workbook"))
    if workbook is not None:
        for cell in mappings(workbook.get("cells")):
            cell.pop("style", None)
        for key in ("row_heights", "column_widths", "freeze_panes"):
            workbook.pop(key, None)
        for conditional_format in mappings(workbook.get("conditional_formats")):
            conditional_format.pop("fill_color", None)
        for chart in mappings(workbook.get("charts")):
            chart.pop("anchor", None)

    presentation = mapping(projected.get("presentation"))
    if presentation is not None:
        for slide in mappings(presentation.get("slides")):
            for key in ("layout_index", "layout_name", "title_style"):
                slide.pop(key, None)
            for text_box in mappings(slide.get("text_boxes")):
                for key in (
                    "left_inches",
                    "top_inches",
                    "width_inches",
                    "height_inches",
                    "font_size",
                    "style",
                ):
                    text_box.pop(key, None)
            for table in mappings(slide.get("tables")):
                for key in (
                    "left_inches",
                    "top_inches",
                    "width_inches",
                    "height_inches",
                    "style",
                ):
                    table.pop(key, None)
            for image in mappings(slide.get("images")):
                for key in (
                    "left_inches",
                    "top_inches",
                    "width_inches",
                    "height_inches",
                    "crop_left",
                    "crop_top",
                    "crop_right",
                    "crop_bottom",
                    "align",
                ):
                    image.pop(key, None)
            for shape in mappings(slide.get("shapes")):
                for key in (
                    "left_inches",
                    "top_inches",
                    "width_inches",
                    "height_inches",
                    "fill_color",
                    "line_color",
                    "text_style",
                ):
                    shape.pop(key, None)
            for chart in mappings(slide.get("charts")):
                for key in (
                    "left_inches",
                    "top_inches",
                    "width_inches",
                    "height_inches",
                ):
                    chart.pop(key, None)

    return projected


def _validate_office_replacement_args(
    replacement: dict[str, Any],
    *,
    original: Mapping[str, Any],
    workspace: str,
    original_read_paths: frozenset[str],
) -> dict[str, Any]:
    """Keep target/authority fixed and prevent expansion of the staged read set."""

    unknown = set(replacement) - _ALLOWED_TOP_LEVEL_ARGS
    if unknown:
        raise OfficePrecommitRepairError(
            "Office repair replacement contains unsupported fields"
        )
    if "file_path" not in replacement or "operation" not in replacement:
        raise OfficePrecommitRepairError(
            "Office repair replacement is not a complete argument object"
        )
    if (
        replacement["file_path"] != original.get("file_path")
        or replacement["operation"] != original.get("operation")
        or replacement.get("overwrite", False) != original.get("overwrite", False)
        or _office_template_identity(replacement)
        != _office_template_identity(original)
    ):
        raise OfficePrecommitRepairError(
            "Office repair replacement changed a fixed request identity"
        )
    repaired_read_paths = frozenset(
        _declared_local_image_paths(replacement, workspace)
    )
    if not repaired_read_paths.issubset(original_read_paths):
        raise OfficePrecommitRepairError(
            "Office repair replacement expanded the local image read set"
        )
    if _office_repair_semantic_projection(
        replacement
    ) != _office_repair_semantic_projection(original):
        raise OfficePrecommitRepairError(
            "Office repair replacement changed semantic content"
        )
    return replacement


def _office_validation_failure_result(
    ctx: ToolContext,
    failures: Sequence[tuple[int, OfficeValidationReport]],
    *,
    repair_attempts: int,
    repair_status: str,
) -> ToolResult:
    """Return location-only diagnostics with no private validation authority."""

    checks: list[dict[str, Any]] = []
    for round_number, report in failures:
        for check in report.checks:
            if check.outcome == "pass":
                continue
            box = check.box
            checks.append(
                {
                    "round": round_number,
                    "code": check.code,
                    "outcome": check.outcome,
                    "page_number": box.page_number if box is not None else None,
                    "box": (
                        {
                            "x": box.x,
                            "y": box.y,
                            "width": box.width,
                            "height": box.height,
                        }
                        if box is not None
                        else None
                    ),
                }
            )
    latest = failures[-1][1]
    return ToolResult(
        error=ctx.tr(
            "Office 权威预提交验证未通过，原文件未更改。",
            (
                "Authoritative Office precommit validation did not pass; "
                "the original file was not changed."
            ),
        ),
        metadata={
            "office_visual_validation": "failed",
            "office_validation_verdict": latest.verdict,
            "office_validation_failure": {
                "rounds": len(failures),
                "repair_attempts": repair_attempts,
                "repair_status": repair_status,
                "checks": checks,
            },
        },
    )
_REQUIRED_PARTS = {
    ".docx": "word/document.xml",
    ".xlsx": "xl/workbook.xml",
    ".pptx": "ppt/presentation.xml",
}
_MACRO_OR_TEMPLATE_EXTENSIONS = {
    ".docm",
    ".dotm",
    ".dotx",
    ".xls",
    ".xlsb",
    ".xlsm",
    ".xltm",
    ".xltx",
    ".potm",
    ".potx",
    ".ppam",
    ".pps",
    ".ppsm",
    ".ppsx",
    ".ppt",
    ".pptm",
}
_UNSUPPORTED_RELATIONSHIP_KINDS = frozenset(
    {
        "activexcontrol",
        "activexcontrolbinary",
        "attachedtemplate",
        "audio",
        "control",
        "controlprop",
        "ctrlprop",
        "embeddedobject",
        "embeddedpackage",
        "externallink",
        "media",
        "oleobject",
        "package",
        "vbaproject",
        "vbaprojectsignature",
        "vbaprojectsignatureagile",
        "vbaprojectsignaturev3",
        "video",
    }
)
_UNSUPPORTED_EMBEDDED_PATH_SEGMENTS = frozenset(
    {"activex", "controls", "ctrlprops", "embeddings"}
)
_UNSUPPORTED_EMBEDDED_CONTENT_TYPE_MARKERS = (
    b"controlproperties",
    b"ms-office.activex",
    b"officedocument.oleobject",
)
_CUSTOM_XML_RELATIONSHIP_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/customXml"
)
_COMMON_EDIT_RELATIONSHIP_TYPES = frozenset(
    {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument",
        "http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties",
        "http://schemas.openxmlformats.org/package/2006/relationships/metadata/thumbnail",
    }
)
_ALLOWED_EDIT_RELATIONSHIP_TYPES = {
    ".docx": _COMMON_EDIT_RELATIONSHIP_TYPES
    | frozenset(
        {
            "http://schemas.microsoft.com/office/2007/relationships/stylesWithEffects",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/fontTable",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/footer",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/header",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/numbering",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/settings",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/webSettings",
        }
    ),
    ".xlsx": _COMMON_EDIT_RELATIONSHIP_TYPES
    | frozenset(
        {
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet",
        }
    ),
    ".pptx": _COMMON_EDIT_RELATIONSHIP_TYPES
    | frozenset(
        {
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/presProps",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/printerSettings",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/tableStyles",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/viewProps",
        }
    ),
}
_COMMON_EDIT_PART_PATTERNS = (
    re.compile(r"\[Content_Types\]\.xml"),
    re.compile(r"_rels/\.rels"),
    re.compile(r"docProps/(?:app|core)\.xml"),
    re.compile(r"docProps/thumbnail\.(?:jpe?g|png)"),
)
_ALLOWED_EDIT_PART_PATTERNS = {
    ".docx": _COMMON_EDIT_PART_PATTERNS
    + (
        re.compile(r"word/document\.xml"),
        re.compile(r"word/_rels/document\.xml\.rels"),
        re.compile(
            r"word/(?:fontTable|numbering|settings|styles|stylesWithEffects|webSettings)\.xml"
        ),
        re.compile(r"word/theme/theme\d+\.xml"),
        re.compile(r"word/(?:header|footer)\d+\.xml"),
        re.compile(r"word/_rels/(?:header|footer)\d+\.xml\.rels"),
        re.compile(
            r"word/media/[^/]+\.(?:bmp|emf|gif|jpe?g|png|tiff?|wmf)",
            re.IGNORECASE,
        ),
    ),
    ".xlsx": _COMMON_EDIT_PART_PATTERNS
    + (
        re.compile(r"xl/workbook\.xml"),
        re.compile(r"xl/_rels/workbook\.xml\.rels"),
        re.compile(r"xl/styles\.xml"),
        re.compile(r"xl/theme/theme\d+\.xml"),
        re.compile(r"xl/worksheets/sheet\d+\.xml"),
        re.compile(r"xl/worksheets/_rels/sheet\d+\.xml\.rels"),
    ),
    ".pptx": _COMMON_EDIT_PART_PATTERNS
    + (
        re.compile(r"ppt/presentation\.xml"),
        re.compile(r"ppt/_rels/presentation\.xml\.rels"),
        re.compile(r"ppt/(?:presProps|tableStyles|viewProps)\.xml"),
        re.compile(r"ppt/printerSettings/printerSettings\d+\.bin"),
        re.compile(r"ppt/slideMasters/slideMaster\d+\.xml"),
        re.compile(r"ppt/slideMasters/_rels/slideMaster\d+\.xml\.rels"),
        re.compile(r"ppt/slideLayouts/slideLayout\d+\.xml"),
        re.compile(r"ppt/slideLayouts/_rels/slideLayout\d+\.xml\.rels"),
        re.compile(r"ppt/slides/slide\d+\.xml"),
        re.compile(r"ppt/slides/_rels/slide\d+\.xml\.rels"),
        re.compile(r"ppt/theme/theme\d+\.xml"),
        re.compile(
            r"ppt/media/[^/]+\.(?:bmp|emf|gif|jpe?g|png|tiff?|wmf)",
            re.IGNORECASE,
        ),
    ),
}
_V2_ALLOWED_EDIT_PART_PATTERNS = {
    ".docx": (),
    ".xlsx": (
        re.compile(r"xl/charts/chart\d+\.xml"),
        re.compile(r"xl/drawings/drawing\d+\.xml"),
        re.compile(r"xl/drawings/_rels/drawing\d+\.xml\.rels"),
    ),
    ".pptx": (
        re.compile(r"ppt/charts/chart\d+\.xml"),
        re.compile(r"ppt/charts/_rels/chart\d+\.xml\.rels"),
        re.compile(r"ppt/embeddings/Microsoft_Excel_(?:Work)?Sheet\d+\.xlsx"),
        re.compile(r"ppt/notesMasters/notesMaster\d+\.xml"),
        re.compile(r"ppt/notesMasters/_rels/notesMaster\d+\.xml\.rels"),
        re.compile(r"ppt/notesSlides/notesSlide\d+\.xml"),
        re.compile(r"ppt/notesSlides/_rels/notesSlide\d+\.xml\.rels"),
    ),
}
_V2_ALLOWED_RELATIONSHIP_TYPES = {
    ".docx": frozenset(),
    ".xlsx": frozenset(
        {
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/drawing",
        }
    ),
    ".pptx": frozenset(
        {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/package",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme",
        }
    ),
}
_V2_ALLOWED_CHART_ELEMENT_NAMES = frozenset(
    {
        # Exact union emitted by the pinned openpyxl/python-pptx release for the
        # four v1.1 chart kinds.  A namespace allowlist alone is insufficient:
        # both libraries silently discard unfamiliar elements in the standard
        # chart namespace during a round trip.
        "auto",
        "autoTitleDeleted",
        "autoUpdate",
        "axId",
        "axPos",
        "barChart",
        "barDir",
        "bodyPr",
        "cat",
        "catAx",
        "chart",
        "chartSpace",
        "crossAx",
        "crossBetween",
        "crosses",
        "date1904",
        "defRPr",
        "delete",
        "dispBlanksAs",
        "endParaRPr",
        "externalData",
        "f",
        "firstSliceAng",
        "formatCode",
        "gapWidth",
        "grouping",
        "idx",
        "layout",
        "lblAlgn",
        "lblOffset",
        "legend",
        "legendPos",
        "lineChart",
        "ln",
        "lstStyle",
        "majorGridlines",
        "majorTickMark",
        "marker",
        "max",
        "min",
        "minorTickMark",
        "noFill",
        "noMultiLvlLbl",
        "numCache",
        "numFmt",
        "numRef",
        "order",
        "orientation",
        "overlay",
        "p",
        "pPr",
        "pieChart",
        "plotArea",
        "plotVisOnly",
        "prstDash",
        "pt",
        "ptCount",
        "r",
        "rich",
        "scaling",
        "scatterChart",
        "scatterStyle",
        "ser",
        "showDLblsOverMax",
        "smooth",
        "spPr",
        "strCache",
        "strRef",
        "style",
        "symbol",
        "t",
        "tickLblPos",
        "title",
        "tx",
        "txPr",
        "v",
        "val",
        "valAx",
        "varyColors",
        "xVal",
        "yVal",
    }
)
_V2_EMBEDDED_WORKBOOK_PART_PATTERNS = (
    re.compile(r"\[Content_Types\]\.xml"),
    re.compile(r"_rels/\.rels"),
    re.compile(r"docProps/(?:app|core)\.xml"),
    re.compile(r"xl/workbook\.xml"),
    re.compile(r"xl/_rels/workbook\.xml\.rels"),
    re.compile(r"xl/worksheets/sheet\d+\.xml"),
    re.compile(r"xl/sharedStrings\.xml"),
    re.compile(r"xl/styles\.xml"),
    re.compile(r"xl/theme/theme\d+\.xml"),
)
_V2_EMBEDDED_WORKBOOK_RELATIONSHIP_TYPES = frozenset(
    {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/sharedStrings",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet",
        "http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties",
    }
)
_V2_EMBEDDED_WORKBOOK_CONTENT_TYPES = frozenset(
    {
        "application/vnd.openxmlformats-officedocument.extended-properties+xml",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sharedstrings+xml",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.styles+xml",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml",
        "application/vnd.openxmlformats-officedocument.theme+xml",
        "application/vnd.openxmlformats-package.core-properties+xml",
        "application/vnd.openxmlformats-package.relationships+xml",
        "application/xml",
    }
)
_XLSX_WORKSHEET_PART = re.compile(
    r"xl/worksheets/(?:_rels/)?sheet\d+\.xml(?:\.rels)?"
)
_ALLOWED_TOP_LEVEL_ARGS = {
    "file_path",
    "operation",
    "overwrite",
    "document",
    "workbook",
    "presentation",
    "replacements",
    "first_party_template",
    "user_template",
}
_DOCX_STYLE_NAMES = {
    "normal": "Normal",
    "title": "Title",
    "subtitle": "Subtitle",
    "heading1": "Heading 1",
    "heading2": "Heading 2",
    "heading3": "Heading 3",
    "bullet": "List Bullet",
    "numbered": "List Number",
}
_INVALID_SHEET_TITLE = re.compile(r"[\\/*?:\[\]]")
_EXTERNAL_FORMULA = re.compile(
    r"(?:https?://|file://|\\\\|\[[^\]]+\.(?:csv|xls|xlsb|xlsm|xlsx)\]"
    r"|\|[^!\r\n]{0,8192}!)",
    re.IGNORECASE,
)
_EXCEL4_MACRO_FORMULA = re.compile(
    r"(?:^|[=+\-*/^&,;@(<>\s])"
    r"(?:(?:_xlfn|_xlws)\.)*(?:CALL|EXEC|REGISTER(?:\.ID)?|RUN)\s*\(",
    re.IGNORECASE,
)

MAX_INPUT_FILE_BYTES = 50 * 1024 * 1024
MAX_OUTPUT_FILE_BYTES = 75 * 1024 * 1024
MAX_ARCHIVE_ENTRIES = 5_000
MAX_ARCHIVE_MEMBER_BYTES = 100 * 1024 * 1024
MAX_ARCHIVE_TOTAL_BYTES = 250 * 1024 * 1024
MAX_RELATIONSHIP_BYTES = 2 * 1024 * 1024
MAX_TEXT_CHARS = 1_000_000
MAX_TOTAL_TEXT_CHARS = 5_000_000
MAX_DECLARATIVE_ITEMS = 500_000
MAX_PARAGRAPHS = 2_000
MAX_TABLES = 200
MAX_TABLE_CELLS = 100_000
MAX_TABLE_COLUMNS = 256
MAX_SHEETS = 100
MAX_WORKBOOK_CELLS = 200_000
MAX_SLIDES = 300
MAX_BULLETS_PER_SLIDE = 500
MAX_PPTX_TABLES_PER_SLIDE = 20
MAX_PPTX_TABLE_CELLS_PER_SLIDE = 10_000
MAX_PPTX_TABLE_COLUMNS = 50
MAX_TEXT_BOXES_PER_SLIDE = 100
MAX_REPLACEMENTS = 200
MAX_IMAGES_PER_FILE = 100
MAX_IMAGE_BYTES = 20 * 1024 * 1024
MAX_TOTAL_IMAGE_BYTES = 50 * 1024 * 1024
MAX_IMAGE_PIXELS = 40_000_000
MAX_V2_FORMAT_ITEMS = 1_000
MAX_CHARTS_PER_FILE = 50
MAX_CHART_SERIES = 20
MAX_CHART_POINTS = 10_000
MAX_MERGED_CELLS_PER_RANGE = 100_000
MAX_TOTAL_MERGED_CELLS = 200_000
MAX_EMBEDDED_WORKBOOK_ENTRIES = 500
MAX_EMBEDDED_WORKBOOK_MEMBER_BYTES = 20 * 1024 * 1024
MAX_EMBEDDED_WORKBOOK_TOTAL_BYTES = 50 * 1024 * 1024
MAX_EMBEDDED_WORKBOOK_COMPRESSION_RATIO = 500
_IMAGE_EXTENSIONS = {".bmp", ".gif", ".jpeg", ".jpg", ".png", ".tif", ".tiff"}
_HEX_COLOR = re.compile(r"^(?:[0-9a-fA-F]{6}|[0-9a-fA-F]{8})$")
_CELL_RANGE = re.compile(r"^[A-Z]{1,3}[1-9][0-9]{0,6}:[A-Z]{1,3}[1-9][0-9]{0,6}$")
_CELL_COORDINATE = re.compile(r"^[A-Z]{1,3}[1-9][0-9]{0,6}$")
_COLUMN_COORDINATE = re.compile(r"^[A-Z]{1,3}$")
_DEFINED_NAME = re.compile(r"^[A-Za-z_][A-Za-z0-9_.]{0,254}$")


def _office_v2_enabled() -> bool:
    """Require the complete source-owned Office authoring dependency graph."""

    from app.release_readiness import v11_capability_released

    return v11_capability_released("office_authoring")


def _user_office_templates_enabled() -> bool:
    """Keep user imports independently gated from first-party authoring."""

    from app.release_readiness import v11_capability_released

    return v11_capability_released("user_office_templates")


class OfficeInputError(ValueError):
    """A safe, request-localizable Office tool error."""

    def __init__(self, zh: str, en: str):
        self.zh = zh
        self.en = en
        super().__init__(en)


_USER_TEMPLATE_SHA256 = re.compile(r"^[0-9a-f]{64}$")


@dataclass(frozen=True, slots=True)
class _UserTemplateLease:
    """Server-owned, path-free snapshot used for one create transaction."""

    record_id: str
    template_ref: str
    revision: int
    expected_state_version: int
    workspace_instance_id: str
    format: str
    source_sha256: str
    manifest_sha256: str
    render_manifest_sha256: str
    render_cache_key: str
    renderer_id: str
    renderer_version: str
    font_digest: str
    render_parameters_version: str
    render_parameters_sha256: str
    placeholder_schema: tuple[UserTemplatePlaceholder, ...]
    placeholder_parts: tuple[str, ...]
    values: tuple[tuple[str, str], ...]
    record_fingerprint: str
    precommit_plan: OfficeCreateValidationPlan | None = field(repr=False)
    service: UserOfficeTemplateService = field(repr=False, compare=False)
    session_factory: async_sessionmaker[AsyncSession] = field(
        repr=False,
        compare=False,
    )

    def values_dict(self) -> dict[str, str]:
        return dict(self.values)


def _user_template_error() -> OfficeInputError:
    return OfficeInputError(
        "用户 Office 模板未获批准、已变更或不属于当前工作区，原文件未更改。",
        (
            "The user Office template is not approved, changed, or does not "
            "belong to the current workspace; the original file was not changed."
        ),
    )


def _parse_user_template_request(
    args: Mapping[str, Any],
) -> tuple[str, int, int, Mapping[str, object]]:
    if not _user_office_templates_enabled():
        raise OfficeInputError(
            "用户 Office 模板 Beta 尚未发布。",
            "User Office template Beta is not released.",
        )
    if args.get("operation") != "create":
        raise OfficeInputError(
            "用户 Office 模板仅支持 create 操作。",
            "User Office templates support only the create operation.",
        )
    if "overwrite" in args:
        raise OfficeInputError(
            "用户 Office 模板不能与 overwrite 参数同时使用。",
            "A user Office template cannot be combined with overwrite.",
        )
    conflicting = [
        field_name
        for field_name in (
            "document",
            "workbook",
            "presentation",
            "replacements",
            "first_party_template",
        )
        if args.get(field_name) is not None
    ]
    if conflicting:
        raise OfficeInputError(
            "用户 Office 模板不能与普通内容或第一方模板同时使用。",
            (
                "A user Office template cannot be combined with ordinary content "
                "or a first-party template."
            ),
        )
    raw = args.get("user_template")
    if not isinstance(raw, Mapping) or set(raw) != {
        "template_ref",
        "revision",
        "expected_state_version",
        "values",
    }:
        raise _user_template_error()
    try:
        template_ref = validate_user_template_ref(raw.get("template_ref"))
    except TemplateContractError as exc:
        raise _user_template_error() from exc
    revision = raw.get("revision")
    state_version = raw.get("expected_state_version")
    values = raw.get("values")
    if (
        not isinstance(revision, int)
        or isinstance(revision, bool)
        or revision < 1
        or not isinstance(state_version, int)
        or isinstance(state_version, bool)
        or state_version < 1
        or not isinstance(values, Mapping)
    ):
        raise _user_template_error()
    return template_ref, revision, state_version, values


def _user_template_service(ctx: ToolContext) -> UserOfficeTemplateService:
    app_state = getattr(ctx, "_app_state", None) or {}
    service = app_state.get("office_user_template_service")
    if service is None:
        service = get_user_office_template_service()
    if not isinstance(service, UserOfficeTemplateService):
        raise OfficeInputError(
            "用户 Office 模板运行时不可用，原文件未更改。",
            (
                "The user Office template runtime is unavailable; "
                "the original file was not changed."
            ),
        )
    return service


def _user_template_record_fingerprint(record: OfficeUserTemplate) -> str:
    def timestamp(value: object) -> str | None:
        return value.isoformat() if hasattr(value, "isoformat") else None

    payload = {
        "id": record.id,
        "template_ref": record.template_ref,
        "revision": record.revision,
        "state_version": record.state_version,
        "workspace_instance_id": record.workspace_instance_id,
        "created_by_session_id": record.created_by_session_id,
        "import_idempotency_key": record.import_idempotency_key,
        "import_request_sha256": record.import_request_sha256,
        "display_name": record.display_name,
        "format": record.format,
        "source_sha256": record.source_sha256,
        "source_size_bytes": record.source_size_bytes,
        "manifest_sha256": record.manifest_sha256,
        "placeholder_schema": record.placeholder_schema,
        "placeholder_parts": record.placeholder_parts,
        "allowed_operations": record.allowed_operations,
        "status": record.status,
        "render_quality": record.render_quality,
        "renderer_id": record.renderer_id,
        "renderer_version": record.renderer_version,
        "font_digest": record.font_digest,
        "render_parameters_version": record.render_parameters_version,
        "render_parameters_sha256": record.render_parameters_sha256,
        "render_cache_key": record.render_cache_key,
        "render_manifest_sha256": record.render_manifest_sha256,
        "render_page_count": record.render_page_count,
        "validation_report": record.validation_report,
        "time_approved": timestamp(record.time_approved),
        "time_tombstoned": timestamp(record.time_tombstoned),
        "time_created": timestamp(record.time_created),
        "time_updated": timestamp(record.time_updated),
    }
    encoded = json.dumps(
        payload,
        ensure_ascii=False,
        allow_nan=False,
        separators=(",", ":"),
        sort_keys=True,
    ).encode("utf-8")
    return hashlib.sha256(encoded).hexdigest()


def _validate_approved_user_template_record(
    record: OfficeUserTemplate,
    *,
    template_ref: str,
    revision: int,
    expected_state_version: int,
    workspace_instance_id: str,
) -> tuple[tuple[UserTemplatePlaceholder, ...], str]:
    hashes = (
        record.source_sha256,
        record.manifest_sha256,
        record.render_manifest_sha256,
        record.render_cache_key,
        record.font_digest,
        record.render_parameters_sha256,
    )
    report = record.validation_report
    if (
        record.template_ref != template_ref
        or record.revision != revision
        or record.state_version != expected_state_version
        or record.workspace_instance_id != workspace_instance_id
        or record.status != "approved"
        or record.render_quality != "authoritative"
        or record.time_approved is None
        or record.time_tombstoned is not None
        or record.format not in {"docx", "xlsx", "pptx"}
        or record.allowed_operations != ["instantiate_text"]
        or record.render_page_count < 1
        or any(
            not isinstance(value, str)
            or _USER_TEMPLATE_SHA256.fullmatch(value) is None
            for value in hashes
        )
        or not isinstance(report, dict)
        or report.get("ooxml_safety") != "pass"
        or report.get("placeholder_contract") != "pass"
        or report.get("independent_reopen") != "pass"
        or report.get("render_quality") != "authoritative"
        or report.get("approval_eligible") is not True
    ):
        raise _user_template_error()
    try:
        schema = decode_user_template_placeholder_schema(record.placeholder_schema)
    except TemplateContractError as exc:
        raise _user_template_error() from exc
    if not isinstance(record.placeholder_parts, list) or not record.placeholder_parts:
        raise _user_template_error()
    return schema, _user_template_record_fingerprint(record)


async def _load_user_template_record(
    db: AsyncSession,
    ctx: ToolContext,
    *,
    template_ref: str,
    revision: int,
    for_update: bool,
) -> OfficeUserTemplate:
    workspace_instance_id = ctx.workspace_instance_id
    if (
        not isinstance(workspace_instance_id, str)
        or not workspace_instance_id
        or not isinstance(ctx.workspace, str)
        or not ctx.workspace
    ):
        raise _user_template_error()
    if for_update:
        session = (
            await db.execute(
                select(Session)
                .where(Session.id == ctx.session_id)
                .with_for_update()
            )
        ).scalar_one_or_none()
        instance = (
            await db.execute(
                select(WorkspaceInstance)
                .where(WorkspaceInstance.id == workspace_instance_id)
                .with_for_update()
            )
        ).scalar_one_or_none()
    else:
        session = await db.get(Session, ctx.session_id)
        instance = await db.get(WorkspaceInstance, workspace_instance_id)
    if (
        session is None
        or instance is None
        or session.time_archived is not None
        or instance.status != "active"
        or dict(instance.details or {}).get("release_intent") is not None
        or session.project_id != instance.project_id
    ):
        raise _user_template_error()
    try:
        canonical, identity = await asyncio.to_thread(
            inspect_workspace_identity,
            ctx.workspace,
        )
        session_canonical, session_identity = await asyncio.to_thread(
            inspect_workspace_identity,
            session.directory,
        )
    except Exception as exc:
        raise _user_template_error() from exc
    if (
        canonical != session_canonical
        or identity != session_identity
        or canonical != instance.root_path
        or identity != instance.identity_token
    ):
        raise _user_template_error()
    statement = select(OfficeUserTemplate).where(
        OfficeUserTemplate.workspace_instance_id == workspace_instance_id,
        OfficeUserTemplate.template_ref == template_ref,
        OfficeUserTemplate.revision == revision,
    )
    if for_update:
        statement = statement.with_for_update()
    record = (await db.execute(statement)).scalar_one_or_none()
    if record is None:
        raise _user_template_error()
    return record


async def _verify_user_template_evidence(
    lease: _UserTemplateLease,
) -> None:
    try:
        await lease.service.verify_approval_evidence(
            template_ref=lease.template_ref,
            revision=lease.revision,
            source_sha256=lease.source_sha256,
            manifest_sha256=lease.manifest_sha256,
            render_manifest_sha256=lease.render_manifest_sha256,
            format_name=lease.format,  # type: ignore[arg-type]
            placeholder_schema=lease.placeholder_schema,
            placeholder_parts=lease.placeholder_parts,
            render_cache_key=lease.render_cache_key,
            renderer_id=lease.renderer_id,
            renderer_version=lease.renderer_version,
            font_digest=lease.font_digest,
            render_parameters_version=lease.render_parameters_version,
            render_parameters_sha256=lease.render_parameters_sha256,
        )
    except (OfficeTemplateError, OfficeValidationError) as exc:
        raise _user_template_error() from exc


async def _resolve_user_template_lease(
    args: Mapping[str, Any],
    ctx: ToolContext,
) -> _UserTemplateLease:
    template_ref, revision, state_version, raw_values = (
        _parse_user_template_request(args)
    )
    app_state = getattr(ctx, "_app_state", None) or {}
    session_factory = app_state.get("session_factory")
    if not isinstance(session_factory, async_sessionmaker):
        raise OfficeInputError(
            "用户 Office 模板数据库运行时不可用，原文件未更改。",
            (
                "The user Office template database runtime is unavailable; "
                "the original file was not changed."
            ),
        )
    async with session_factory() as db:
        record = await _load_user_template_record(
            db,
            ctx,
            template_ref=template_ref,
            revision=revision,
            for_update=False,
        )
        schema, fingerprint = _validate_approved_user_template_record(
            record,
            template_ref=template_ref,
            revision=revision,
            expected_state_version=state_version,
            workspace_instance_id=ctx.workspace_instance_id or "",
        )
        try:
            values = validate_user_template_values(schema, raw_values)
        except TemplateContractError as exc:
            raise OfficeInputError(
                "用户 Office 模板字段不满足已批准的文本长度约束。",
                (
                    "User Office template values do not satisfy the approved "
                    "text bounds."
                ),
            ) from exc
        record_id = record.id
        format_name = record.format
        source_sha256 = record.source_sha256
        manifest_sha256 = record.manifest_sha256
        render_manifest_sha256 = record.render_manifest_sha256
        render_cache_key = record.render_cache_key
        renderer_id = record.renderer_id
        renderer_version = record.renderer_version
        font_digest = record.font_digest
        parameters_version = record.render_parameters_version
        parameters_sha256 = record.render_parameters_sha256
        placeholder_parts = tuple(record.placeholder_parts)

    service = _user_template_service(ctx)
    provisional = _UserTemplateLease(
        record_id=record_id,
        template_ref=template_ref,
        revision=revision,
        expected_state_version=state_version,
        workspace_instance_id=ctx.workspace_instance_id or "",
        format=format_name,
        source_sha256=source_sha256,
        manifest_sha256=manifest_sha256,
        render_manifest_sha256=render_manifest_sha256,
        render_cache_key=render_cache_key,
        renderer_id=renderer_id,
        renderer_version=renderer_version,
        font_digest=font_digest,
        render_parameters_version=parameters_version,
        render_parameters_sha256=parameters_sha256,
        placeholder_schema=schema,
        placeholder_parts=placeholder_parts,
        values=tuple(sorted(values.items())),
        record_fingerprint=fingerprint,
        precommit_plan=None,
        service=service,
        session_factory=session_factory,
    )
    # Rendering can be comparatively expensive.  It runs only from the scalar,
    # immutable lease after the ORM session is closed, so no detached/lazy ORM
    # access leaks across the trust boundary.
    await _verify_user_template_evidence(provisional)
    try:
        plan = await asyncio.to_thread(
            service.build_precommit_plan,
            template_ref=template_ref,
            revision=revision,
            source_sha256=source_sha256,
            manifest_sha256=manifest_sha256,
            format_name=format_name,
            placeholder_schema=schema,
            placeholder_parts=placeholder_parts,
            renderer_id=renderer_id,
            renderer_version=renderer_version,
            font_digest=font_digest,
            render_parameters_version=parameters_version,
            render_parameters_sha256=parameters_sha256,
        )
    except (OfficeTemplateError, OfficeValidationError) as exc:
        raise _user_template_error() from exc
    if not isinstance(plan, OfficeCreateValidationPlan):
        raise _user_template_error()
    return _UserTemplateLease(
        record_id=record_id,
        template_ref=template_ref,
        revision=revision,
        expected_state_version=state_version,
        workspace_instance_id=ctx.workspace_instance_id or "",
        format=format_name,
        source_sha256=source_sha256,
        manifest_sha256=manifest_sha256,
        render_manifest_sha256=render_manifest_sha256,
        render_cache_key=render_cache_key,
        renderer_id=renderer_id,
        renderer_version=renderer_version,
        font_digest=font_digest,
        render_parameters_version=parameters_version,
        render_parameters_sha256=parameters_sha256,
        placeholder_schema=schema,
        placeholder_parts=placeholder_parts,
        values=tuple(sorted(values.items())),
        record_fingerprint=fingerprint,
        precommit_plan=plan,
        service=service,
        session_factory=session_factory,
    )


async def _commit_user_template_with_revalidation(
    transaction: WorkspaceMutationTransaction,
    *,
    seal: object,
    lease: _UserTemplateLease,
    ctx: ToolContext,
) -> tuple[Any, asyncio.CancelledError | None]:
    """Linearize template state and the sealed filesystem commit."""

    # Recapture the approved render before acquiring a database write/row lock.
    # The locked phase below repeats the exact DB fingerprint and immutable
    # registry/source contract, but never performs a renderer invocation.
    await _verify_user_template_evidence(lease)

    filesystem_committed = False
    async with lease.session_factory() as db:
        try:
            dialect = db.get_bind().dialect.name
            if dialect == "sqlite":
                await db.execute(text("BEGIN IMMEDIATE"))
            else:
                await db.begin()
            record = await _load_user_template_record(
                db,
                ctx,
                template_ref=lease.template_ref,
                revision=lease.revision,
                for_update=dialect != "sqlite",
            )
            schema, fingerprint = _validate_approved_user_template_record(
                record,
                template_ref=lease.template_ref,
                revision=lease.revision,
                expected_state_version=lease.expected_state_version,
                workspace_instance_id=lease.workspace_instance_id,
            )
            if schema != lease.placeholder_schema or fingerprint != lease.record_fingerprint:
                raise _user_template_error()
            try:
                await asyncio.to_thread(
                    lease.service.verify_registry_contract,
                    template_ref=lease.template_ref,
                    revision=lease.revision,
                    source_sha256=lease.source_sha256,
                    manifest_sha256=lease.manifest_sha256,
                    format_name=lease.format,  # type: ignore[arg-type]
                    placeholder_schema=lease.placeholder_schema,
                    placeholder_parts=lease.placeholder_parts,
                )
            except (OfficeTemplateError, OfficeValidationError) as exc:
                raise _user_template_error() from exc
            commit, deferred_cancellation = await _await_transaction_commit(
                transaction,
                seal=seal,
            )
            filesystem_committed = True
            try:
                await db.commit()
            except Exception:
                # The read lock has served its linearization purpose and the
                # filesystem commit is already durable.  Closing the session
                # releases the lock; do not report a false rollback.
                logger.exception(
                    "Could not close user-template state lock after sealed commit"
                )
            return commit, deferred_cancellation
        except BaseException:
            if not filesystem_committed:
                try:
                    await db.rollback()
                except Exception:
                    logger.exception(
                        "Could not roll back user-template state lock"
                    )
            raise


def _validate_user_template_precommit_evidence(
    lease: _UserTemplateLease,
    result: OfficeDraftValidationResult,
    seal: object,
) -> OfficeDraftSeal:
    """Bind the consumed seal to the approval's exact golden/provider identity."""

    if not isinstance(seal, OfficeDraftSeal) or seal is not result.candidate:
        raise _user_template_error()
    report = result.report
    if (
        report.verdict != "pass"
        or report.document_format != lease.format
        or report.baseline_sha256 != lease.source_sha256
        or report.renderer_id != lease.renderer_id
        or report.renderer_version != lease.renderer_version
        or report.font_digest != lease.font_digest
        or seal.renderer_id != lease.renderer_id
        or seal.renderer_version != lease.renderer_version
        or seal.font_digest != lease.font_digest
        or seal.parameters_version != lease.render_parameters_version
        or seal.parameters_sha256 != lease.render_parameters_sha256
        or seal.quality != "authoritative"
    ):
        raise _user_template_error()
    return seal


def _path_free_user_template_commit_metadata(
    raw: Mapping[str, object],
    workspace: Path,
) -> dict[str, object]:
    """Remove host paths from the user-template result envelope."""

    metadata = dict(raw)
    for key in (
        "written_files",
        "deleted_files",
        "recovery_sidecars",
        "recovery_files",
    ):
        values = metadata.get(key)
        relative_values: list[str] = []
        if isinstance(values, list):
            for value in values:
                if not isinstance(value, str):
                    continue
                try:
                    candidate = Path(value)
                    relative = (
                        candidate.resolve().relative_to(workspace).as_posix()
                        if candidate.is_absolute()
                        else PurePosixPath(value).as_posix()
                    )
                    parsed = PurePosixPath(relative)
                    if (
                        parsed.is_absolute()
                        or any(part in {"", ".", ".."} for part in parsed.parts)
                        or "\\" in relative
                    ):
                        continue
                except (OSError, ValueError):
                    continue
                relative_values.append(relative)
        metadata[key] = relative_values
    return metadata


class OfficeTool(ToolDefinition):
    """Create or make bounded edits to macro-free OOXML Office files."""

    @property
    def id(self) -> str:
        return "office"

    @property
    def description(self) -> str:
        base = (
            "Safely create or make limited declarative edits to .docx, .xlsx, and "
            ".pptx files inside the selected workspace. This tool does not run "
            "Python or shell commands and does not accept external templates or "
            "macro-enabled formats. DOCX supports paragraphs, tables, page breaks, "
            "local images, append, and exact text replacement; XLSX supports sheets, "
            "rows, cell updates, basic styles, and sheet deletion; PPTX supports "
            "title/bullet slides, text boxes, tables, local images, append, and exact "
            "text replacement. "
            "XLSX formulas beginning with '=' are stored but are never recalculated "
            "by this tool. Outputs are reopened and validated before atomic install."
        )
        if not _office_v2_enabled():
            return base
        description = (
            base
            + " Office v1.1 additionally supports bounded page/section formatting, "
            "explicit text and table styles, worksheet layout/validation rules, named "
            "ranges, native bar/line/pie/scatter charts, slide-layout selection, "
            "speaker notes, static DOCX chart images with alt/source metadata, and "
            "signed first-party DOCX/XLSX/PPTX templates."
        )
        if _user_office_templates_enabled():
            description += (
                " Approved workspace-scoped user templates can be instantiated "
                "through an opaque immutable revision."
            )
        return description

    def parameters_schema(self) -> dict[str, Any]:
        scalar_schema: dict[str, Any] = {
            "oneOf": [
                {"type": "string"},
                {"type": "number"},
                {"type": "boolean"},
                {"type": "null"},
            ]
        }
        schema = {
            "type": "object",
            "additionalProperties": False,
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": (
                        "Destination path. Relative paths are written under the "
                        "workspace's suxiaoyou_written directory. Only .docx, .xlsx, "
                        "and .pptx are accepted."
                    ),
                },
                "operation": {
                    "type": "string",
                    "enum": ["create", "edit"],
                    "description": (
                        "Create a new file, or edit an existing macro-free OOXML file."
                    ),
                },
                "overwrite": {
                    "type": "boolean",
                    "default": False,
                    "description": (
                        "For create only: explicitly allow replacing an existing file."
                    ),
                },
                "document": {
                    "type": "object",
                    "description": "DOCX content. Valid only for a .docx path.",
                    "properties": {
                        "title": {"type": "string"},
                        "paragraphs": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "text": {"type": "string"},
                                    "style": {
                                        "type": "string",
                                        "enum": list(_DOCX_STYLE_NAMES),
                                        "default": "normal",
                                    },
                                    "page_break_after": {
                                        "type": "boolean",
                                        "default": False,
                                    },
                                },
                                "required": ["text"],
                            },
                        },
                        "tables": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "headers": {
                                        "type": "array",
                                        "items": scalar_schema,
                                    },
                                    "rows": {
                                        "type": "array",
                                        "items": {
                                            "type": "array",
                                            "items": scalar_schema,
                                        },
                                    },
                                },
                                "required": ["rows"],
                            },
                        },
                        "images": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "path": {"type": "string"},
                                    "width_inches": {"type": "number"},
                                    "caption": {"type": "string"},
                                },
                                "required": ["path"],
                            },
                        },
                    },
                },
                "workbook": {
                    "type": "object",
                    "description": "XLSX content. Valid only for a .xlsx path.",
                    "properties": {
                        "sheets": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "name": {"type": "string"},
                                    "action": {
                                        "type": "string",
                                        "enum": ["create", "append"],
                                    },
                                    "rows": {
                                        "type": "array",
                                        "items": {
                                            "type": "array",
                                            "items": scalar_schema,
                                        },
                                    },
                                },
                                "required": ["name", "rows"],
                            },
                        },
                        "cells": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "sheet": {"type": "string"},
                                    "cell": {"type": "string"},
                                    "value": scalar_schema,
                                    "style": {
                                        "type": "object",
                                        "properties": {
                                            "number_format": {"type": "string"},
                                            "font": {
                                                "type": "object",
                                                "properties": {
                                                    "bold": {"type": "boolean"},
                                                    "italic": {"type": "boolean"},
                                                    "color": {"type": "string"},
                                                    "size": {"type": "number"},
                                                },
                                            },
                                            "fill": {
                                                "type": "object",
                                                "properties": {
                                                    "color": {"type": "string"},
                                                },
                                                "required": ["color"],
                                            },
                                        },
                                    },
                                },
                                "required": ["sheet", "cell"],
                            },
                        },
                        "delete_sheets": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "For edit only: exact worksheet names to delete.",
                        },
                    },
                },
                "presentation": {
                    "type": "object",
                    "description": "PPTX content. Valid only for a .pptx path.",
                    "properties": {
                        "slides": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "title": {"type": "string"},
                                    "subtitle": {"type": "string"},
                                    "bullets": {
                                        "type": "array",
                                        "items": {
                                            "oneOf": [
                                                {"type": "string"},
                                                {
                                                    "type": "object",
                                                    "properties": {
                                                        "text": {"type": "string"},
                                                        "level": {
                                                            "type": "integer",
                                                            "minimum": 0,
                                                            "maximum": 4,
                                                        },
                                                    },
                                                    "required": ["text"],
                                                },
                                            ]
                                        },
                                    },
                                    "text_boxes": {
                                        "type": "array",
                                        "items": {
                                            "type": "object",
                                            "properties": {
                                                "text": {"type": "string"},
                                                "left_inches": {"type": "number"},
                                                "top_inches": {"type": "number"},
                                                "width_inches": {"type": "number"},
                                                "height_inches": {"type": "number"},
                                                "font_size": {"type": "number"},
                                            },
                                            "required": [
                                                "text",
                                                "left_inches",
                                                "top_inches",
                                                "width_inches",
                                                "height_inches",
                                            ],
                                        },
                                    },
                                    "tables": {
                                        "type": "array",
                                        "items": {
                                            "type": "object",
                                            "properties": {
                                                "left_inches": {"type": "number"},
                                                "top_inches": {"type": "number"},
                                                "width_inches": {"type": "number"},
                                                "height_inches": {"type": "number"},
                                                "headers": {
                                                    "type": "array",
                                                    "items": scalar_schema,
                                                },
                                                "rows": {
                                                    "type": "array",
                                                    "items": {
                                                        "type": "array",
                                                        "items": scalar_schema,
                                                    },
                                                },
                                            },
                                            "required": [
                                                "left_inches",
                                                "top_inches",
                                                "width_inches",
                                                "height_inches",
                                                "rows",
                                            ],
                                        },
                                    },
                                    "images": {
                                        "type": "array",
                                        "items": {
                                            "type": "object",
                                            "properties": {
                                                "path": {"type": "string"},
                                                "left_inches": {"type": "number"},
                                                "top_inches": {"type": "number"},
                                                "width_inches": {"type": "number"},
                                                "height_inches": {"type": "number"},
                                            },
                                            "required": [
                                                "path",
                                                "left_inches",
                                                "top_inches",
                                            ],
                                        },
                                    },
                                },
                                "required": ["title"],
                            },
                        },
                    },
                },
                "replacements": {
                    "type": "array",
                    "description": (
                        "For DOCX/PPTX edit: exact body text replacements. Matches "
                        "cannot span paragraphs. By default a match must be unique."
                    ),
                    "items": {
                        "type": "object",
                        "properties": {
                            "old_text": {"type": "string"},
                            "new_text": {"type": "string"},
                            "replace_all": {"type": "boolean", "default": False},
                        },
                        "required": ["old_text", "new_text"],
                    },
                },
            },
            "required": ["file_path", "operation"],
        }
        if _office_v2_enabled():
            _extend_office_v2_schema(schema, scalar_schema)
        return schema

    async def execute(self, args: dict[str, Any], ctx: ToolContext) -> ToolResult:
        if not ctx.workspace:
            return ToolResult(
                error=ctx.tr(
                    "Office 工具需要先选择工作区。",
                    "The Office tool requires a selected workspace.",
                )
            )
        if ctx.is_aborted:
            return ToolResult(error=ctx.tr("操作已取消。", "Operation cancelled."))

        try:
            file_path = resolve_for_write(str(args.get("file_path", "")), ctx.workspace)
        except WorkspaceViolation:
            return ToolResult(
                error=ctx.tr(
                    "拒绝访问：Office 文件必须位于当前工作区内。",
                    "Access denied: Office files must stay inside the current workspace.",
                )
            )

        transaction: WorkspaceMutationTransaction | None = None
        precommit_session: OfficePrecommitValidationSession | None = None
        precommit_result: OfficeDraftValidationResult | None = None
        commit_completed = False
        try:
            original_args = copy.deepcopy(args)
            current_args = copy.deepcopy(args)
            user_template_lease = (
                await _resolve_user_template_lease(original_args, ctx)
                if "user_template" in original_args
                else None
            )
            ensure_mutation_metadata_supported(file_path)
            transaction = WorkspaceMutationTransaction(
                ctx.workspace,
                ctx,
                operation=f"office.{args.get('operation', 'unknown')}",
            )
            image_paths = _declared_local_image_paths(args, ctx.workspace)
            original_read_paths = frozenset(image_paths)
            # The closure covers the keyword-only read set and lets
            # cancellation wait for sparse staging to settle before abort.
            staged_workspace = await _await_transaction_thread(
                lambda: transaction.prepare_paths(
                    [file_path],
                    read_paths=image_paths,
                )
            )
            staged_target = transaction.staged_path(file_path)
            coordinator: OfficePrecommitCoordinator | None = None
            request: OfficePrecommitRequest | None = None
            view = None
            repairer: OfficePrecommitRepairer | None = None
            tokenized_repair_args: dict[str, Any] | None = None
            repair_path_tokens: _OfficeRepairPathTokenTable | None = None
            if (
                _office_v2_enabled()
                and args.get("operation") in {"create", "edit"}
                and Path(file_path).suffix.casefold() in _FORMATS
            ):
                coordinator = _office_precommit_coordinator(ctx)
                if coordinator is None:
                    raise OfficeInputError(
                        "Office v1.1 权威预提交验证运行时不可用，原文件未更改。",
                        (
                            "The authoritative Office v1.1 precommit validator is "
                            "unavailable; the original file was not changed."
                        ),
                    )
                request = _office_precommit_request(
                    args,
                    ctx,
                    file_path,
                    user_template=user_template_lease,
                )
                view = transaction.arm_office_precommit_validation(file_path)
                # User-template authority and values are immutable for this
                # transaction.  Until repair has a separately reviewed fixed-
                # identity contract, a validation failure is terminal.
                repairer = (
                    None
                    if (
                        user_template_lease is not None
                        or "first_party_template" in original_args
                    )
                    else _office_precommit_repairer(ctx)
                )

            failures: list[tuple[int, OfficeValidationReport]] = []
            repair_attempts = 0
            used_precommit_sessions: list[OfficePrecommitValidationSession] = []
            used_precommit_results: list[OfficeDraftValidationResult] = []
            used_candidate_seals: list[object] = []
            for candidate_round in range(1, 4):
                try:
                    if user_template_lease is None:
                        summary = await _await_transaction_thread(
                            _run_office_operation,
                            staged_target,
                            current_args,
                            ctx,
                            staged_workspace,
                        )
                    else:
                        summary = await _await_transaction_thread(
                            lambda: _run_office_operation(
                                staged_target,
                                current_args,
                                ctx,
                                staged_workspace,
                                user_template=user_template_lease,
                            )
                        )
                except OfficeInputError:
                    if candidate_round == 1 or not failures:
                        raise
                    transaction.abort()
                    return _office_validation_failure_result(
                        ctx,
                        failures,
                        repair_attempts=repair_attempts,
                        repair_status="rejected",
                    )
                if ctx.is_aborted:
                    if precommit_session is not None:
                        precommit_session.abort()
                    transaction.abort()
                    return ToolResult(
                        error=ctx.tr("操作已取消。", "Operation cancelled.")
                    )

                seal: object | None = None
                if coordinator is not None:
                    assert request is not None and view is not None
                    round_request = request
                    if request.operation == "edit":
                        round_request = replace(
                            request,
                            trusted_edit_intent=_office_edit_mutation_intent(
                                current_args,
                                summary,
                                request.document_format,
                            ),
                        )
                    precommit_session = await coordinator.begin(
                        request=round_request,
                        view=view,
                    )
                    if not isinstance(
                        precommit_session,
                        OfficePrecommitValidationSession,
                    ):
                        raise TypeError(
                            "Office precommit coordinator returned an invalid session"
                        )
                    if any(
                        precommit_session is previous
                        for previous in used_precommit_sessions
                    ):
                        raise TypeError(
                            "Office precommit coordinator reused a validation session"
                        )
                    used_precommit_sessions.append(precommit_session)
                    precommit_result = await precommit_session.validate_candidate()
                    if not isinstance(precommit_result, OfficeDraftValidationResult):
                        raise TypeError(
                            "Office precommit session returned an invalid result"
                        )
                    if any(
                        precommit_result is previous
                        for previous in used_precommit_results
                    ) or any(
                        precommit_result.candidate is previous
                        for previous in used_candidate_seals
                    ):
                        raise TypeError(
                            "Office precommit session reused candidate evidence"
                        )
                    used_precommit_results.append(precommit_result)
                    used_candidate_seals.append(precommit_result.candidate)
                    report = precommit_result.report
                    if report.verdict != "pass":
                        failures.append((candidate_round, report))
                        precommit_session.abort()
                        precommit_session = None
                        precommit_result = None
                        if ctx.is_aborted:
                            transaction.abort()
                            return ToolResult(
                                error=ctx.tr("操作已取消。", "Operation cancelled.")
                            )
                        if repairer is None:
                            transaction.abort()
                            return _office_validation_failure_result(
                                ctx,
                                failures,
                                repair_attempts=repair_attempts,
                                repair_status="unavailable",
                            )
                        if repair_attempts >= 2:
                            transaction.abort()
                            return _office_validation_failure_result(
                                ctx,
                                failures,
                                repair_attempts=repair_attempts,
                                repair_status="limit_reached",
                            )

                        repair_attempts += 1
                        try:
                            if tokenized_repair_args is None:
                                tokenized_repair_args, repair_path_tokens = (
                                    _tokenize_office_repair_args(original_args)
                                )
                            assert tokenized_repair_args is not None
                            assert repair_path_tokens is not None
                            repair_request: OfficePrecommitRepairRequest = (
                                build_precommit_repair_request(
                                    tokenized_args=tokenized_repair_args,
                                    report=report,
                                    attempt=repair_attempts,  # type: ignore[arg-type]
                                )
                            )
                            raw_replacement = await _await_repair_worker(
                                repairer,
                                repair_request,
                            )
                            replacement = _unmask_office_replacement_args(
                                copy_replacement_args(raw_replacement),
                                repair_path_tokens,
                            )
                            replacement = _validate_office_replacement_args(
                                replacement,
                                original=original_args,
                                workspace=ctx.workspace,
                                original_read_paths=original_read_paths,
                            )
                        except asyncio.CancelledError:
                            raise
                        except _OfficeRepairTimeoutError:
                            transaction.abort()
                            return _office_validation_failure_result(
                                ctx,
                                failures,
                                repair_attempts=repair_attempts,
                                repair_status="timeout",
                            )
                        except Exception as exc:
                            logger.warning(
                                "Office precommit repair attempt was rejected (%s)",
                                type(exc).__name__,
                            )
                            transaction.abort()
                            return _office_validation_failure_result(
                                ctx,
                                failures,
                                repair_attempts=repair_attempts,
                                repair_status="rejected",
                            )
                        if ctx.is_aborted:
                            transaction.abort()
                            return ToolResult(
                                error=ctx.tr("操作已取消。", "Operation cancelled.")
                            )
                        await _await_transaction_thread(
                            transaction.reset_office_precommit_target,
                            file_path,
                        )
                        current_args = replacement
                        continue
                    seal = precommit_session.consume_commit_seal(precommit_result)
                    if seal is not precommit_result.candidate:
                        raise TypeError(
                            "Office precommit session returned unrelated commit evidence"
                        )

                ensure_mutation_metadata_supported(file_path)
                if user_template_lease is not None:
                    if seal is None or precommit_result is None:
                        raise TypeError(
                            "User template commit has no authoritative seal"
                        )
                    seal = _validate_user_template_precommit_evidence(
                        user_template_lease,
                        precommit_result,
                        seal,
                    )
                    commit, deferred_cancellation = (
                        await _commit_user_template_with_revalidation(
                            transaction,
                            seal=seal,
                            lease=user_template_lease,
                            ctx=ctx,
                        )
                    )
                else:
                    commit, deferred_cancellation = await _await_transaction_commit(
                        transaction,
                        seal=seal,
                    )
                commit_completed = True
                if precommit_session is not None and precommit_result is not None:
                    precommit_session.mark_committed(precommit_result)
                    report = precommit_result.report
                    summary.update(
                        {
                            "office_visual_validation": "authoritative",
                            "office_validation_verdict": report.verdict,
                            "office_validation_schema_version": report.schema_version,
                            "office_validation_renderer_id": report.renderer_id,
                            "office_validation_renderer_version": report.renderer_version,
                            "office_validation_font_digest": report.font_digest,
                            "office_validation_checkpoint_id": report.checkpoint_id,
                            "office_validation_root_turn_id": report.root_turn_id,
                            "office_validation_repair_attempts": repair_attempts,
                            # SessionProcessor consumes this private, JSON-safe
                            # envelope before ToolResult metadata is persisted or
                            # emitted.  It is bound to the exact checkpoint change
                            # so rewind can invalidate stale validation evidence.
                            "_office_validation_report": report.to_dict(),
                        }
                    )
                if deferred_cancellation is not None:
                    raise deferred_cancellation
                break
            else:  # pragma: no cover - the third failed candidate returns above.
                raise RuntimeError("Office repair loop exceeded its fixed bound")
        except asyncio.CancelledError:
            if precommit_session is not None:
                precommit_session.abort()
            if transaction is not None and not commit_completed:
                transaction.abort()
            raise
        except OfficeInputError as exc:
            if precommit_session is not None:
                precommit_session.abort()
            if transaction is not None:
                transaction.abort()
            return ToolResult(error=ctx.tr(exc.zh, exc.en))
        except OfficePrecommitRepairError as exc:
            if precommit_session is not None:
                precommit_session.abort()
            if transaction is not None:
                transaction.abort()
            logger.warning(
                "Office repair boundary rejected declarative input (%s)",
                type(exc).__name__,
            )
            return ToolResult(
                error=ctx.tr(
                    "Office 修复边界拒绝了不安全的路径结构，原文件未更改。",
                    (
                        "The Office repair boundary rejected an unsafe path "
                        "structure; the original file was not changed."
                    ),
                )
            )
        except OfficeValidationError:
            if precommit_session is not None:
                precommit_session.abort()
            if transaction is not None:
                transaction.abort()
            return ToolResult(
                error=ctx.tr(
                    "Office 权威预提交验证未通过，原文件未更改。",
                    (
                        "Authoritative Office precommit validation did not pass; "
                        "the original file was not changed."
                    ),
                )
            )
        except (UnsupportedFileMetadataError, WorkspaceMutationError) as exc:
            if precommit_session is not None:
                precommit_session.abort()
            if transaction is not None:
                transaction.abort()
            return ToolResult(error=str(exc))
        except PermissionError:
            if precommit_session is not None:
                precommit_session.abort()
            if transaction is not None:
                transaction.abort()
            return ToolResult(
                error=ctx.tr(
                    f"没有权限写入：{file_path}",
                    f"Permission denied writing: {file_path}",
                )
            )
        except OSError:
            if precommit_session is not None:
                precommit_session.abort()
            if transaction is not None:
                transaction.abort()
            logger.exception("Office file write failed for %s", file_path)
            return ToolResult(
                error=ctx.tr(
                    "Office 文件写入失败，原文件未更改。",
                    "Office file write failed; the original file was not changed.",
                )
            )
        except Exception:
            if precommit_session is not None:
                precommit_session.abort()
            if transaction is not None:
                transaction.abort()
            logger.exception("Office processing failed for %s", file_path)
            return ToolResult(
                error=ctx.tr(
                    "Office 文件处理或重新打开校验失败，原文件未更改。",
                    (
                        "Office processing or reopen validation failed; "
                        "the original file was not changed."
                    ),
                )
            )

        previous_version = None
        if commit.previous_version_ids:
            try:
                version_ids = set(commit.previous_version_ids)
                previous_version = next(
                    (
                        version
                        for version in FileVersionStore(Path(ctx.workspace)).list_versions()
                        if version.id in version_ids
                    ),
                    None,
                )
            except Exception:
                # The workspace commit is already durable.  A read-only metadata
                # lookup failure must not turn a completed write into a false
                # failure response or attempt an impossible post-commit rollback.
                logger.exception("Could not load Office version metadata for %s", file_path)
        summary = {**summary, **version_metadata(previous_version)}

        operation = str(args["operation"])
        action_zh = "已创建" if operation == "create" else "已编辑"
        action_en = "Created" if operation == "create" else "Edited"
        name = Path(file_path).name
        output = ctx.tr(
            f"{action_zh}并校验 {file_path}。文件已重新打开校验并原子替换。",
            (
                f"{action_en} and validated {file_path}. "
                "The file was reopened, verified, and atomically installed."
            ),
        )
        if summary["format"] == "xlsx":
            output += ctx.tr(
                " 公式会保存，但本工具不会重算公式结果。",
                " Formulas are stored, but this tool does not recalculate results.",
            )

        metadata_file_path = file_path
        commit_metadata = commit.metadata
        if summary.get("user_template") is True:
            metadata_file_path = Path(file_path).relative_to(
                Path(ctx.workspace).resolve()
            ).as_posix()
            commit_metadata = _path_free_user_template_commit_metadata(
                commit_metadata,
                Path(ctx.workspace).resolve(),
            )
        metadata = {
            "file_path": metadata_file_path,
            "mime_type": _MIME_TYPES[Path(file_path).suffix.lower()],
            "operation": operation,
            "format": summary["format"],
            "reopened_and_validated": True,
            "atomic_install": True,
            "macros_allowed": False,
            "external_templates_allowed": False,
            **commit_metadata,
            **summary,
        }
        if summary["format"] == "xlsx":
            metadata["formulas_recalculated"] = False

        return ToolResult(
            output=output,
            title=ctx.tr(f"{action_zh} {name}", f"{action_en} {name}"),
            metadata=metadata,
        )


def _first_party_template_schema(
    template_id: str,
    fields: Mapping[str, tuple[int, int]],
) -> dict[str, Any]:
    """Build the exact model-facing contract for one signed catalog entry."""

    value_properties = {
        name: {"type": "string", "minLength": minimum, "maxLength": maximum}
        for name, (minimum, maximum) in fields.items()
    }
    return {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "template_id": {"const": template_id},
            "template_version": {"const": "1.0.0"},
            "values": {
                "type": "object",
                "additionalProperties": False,
                "properties": value_properties,
                "required": sorted(value_properties),
            },
        },
        "required": ["template_id", "template_version", "values"],
    }


def _extend_office_v2_schema(
    schema: dict[str, Any],
    scalar_schema: dict[str, Any],
) -> None:
    """Expose only the reviewed v1.1 declarative subset while its gate is on."""

    properties = schema["properties"]
    properties["first_party_template"] = {
        "description": (
            "Create from a signed bundled first-party template. Valid only with "
            "operation=create, overwrite=false, and no ordinary Office payload."
        ),
        "oneOf": [
            _first_party_template_schema(
                "business-brief",
                {
                    "classification": (1, 40),
                    "next_step": (1, 600),
                    "owner": (1, 80),
                    "recipient": (1, 120),
                    "report_date": (1, 40),
                    "summary": (1, 2400),
                    "title": (1, 160),
                },
            ),
            _first_party_template_schema(
                "project-tracker",
                {
                    "owner": (1, 80),
                    "project_name": (1, 120),
                    "report_date": (1, 40),
                },
            ),
            _first_party_template_schema(
                "status-update",
                {
                    "next_step": (1, 300),
                    "owner": (1, 80),
                    "period": (1, 80),
                    "project_name": (1, 120),
                    "status": (1, 120),
                    "summary": (1, 700),
                },
            ),
        ],
    }
    if _user_office_templates_enabled():
        properties["user_template"] = {
            "type": "object",
            "additionalProperties": False,
            "description": (
                "Create from one approved user template in the current workspace. "
                "Do not combine with overwrite, ordinary content, or a first-party "
                "template."
            ),
            "properties": {
                "template_ref": {
                    "type": "string",
                    "pattern": r"^utpl-[0-9a-z]{26}$",
                },
                "revision": {"type": "integer", "minimum": 1},
                "expected_state_version": {
                    "type": "integer",
                    "minimum": 1,
                },
                "values": {
                    "type": "object",
                    "additionalProperties": {"type": "string"},
                    "minProperties": 1,
                    "maxProperties": 256,
                },
            },
            "required": [
                "template_ref",
                "revision",
                "expected_state_version",
                "values",
            ],
        }
    properties["document"]["additionalProperties"] = False
    document = properties["document"]["properties"]
    paragraph = document["paragraphs"]["items"]
    paragraph["additionalProperties"] = False
    paragraph.pop("required", None)
    paragraph["oneOf"] = [{"required": ["text"]}, {"required": ["runs"]}]
    paragraph["properties"].update(
        {
            "runs": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "text": {"type": "string"},
                        "bold": {"type": "boolean"},
                        "italic": {"type": "boolean"},
                        "underline": {"type": "boolean"},
                        "color": {"type": "string"},
                        "size": {"type": "number"},
                        "font": {"type": "string"},
                    },
                    "required": ["text"],
                },
            },
            "format": {
                "type": "object",
                "additionalProperties": False,
                "properties": {
                    "alignment": {
                        "type": "string",
                        "enum": ["left", "center", "right", "justify"],
                    },
                    "keep_with_next": {"type": "boolean"},
                    "keep_together": {"type": "boolean"},
                    "page_break_before": {"type": "boolean"},
                    "widow_control": {"type": "boolean"},
                    "space_before_pt": {"type": "number"},
                    "space_after_pt": {"type": "number"},
                    "line_spacing": {"type": "number"},
                },
            },
            "list": {
                "type": "object",
                "additionalProperties": False,
                "properties": {
                    "level": {"type": "integer", "minimum": 0, "maximum": 8},
                    "ordered": {"type": "boolean"},
                },
                "required": ["level", "ordered"],
            },
        }
    )
    table = document["tables"]["items"]
    table["additionalProperties"] = False
    document["images"]["items"]["additionalProperties"] = False
    table["properties"].update(
        {
            "merges": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "start": {
                            "type": "array",
                            "items": {"type": "integer", "minimum": 0},
                            "minItems": 2,
                            "maxItems": 2,
                        },
                        "end": {
                            "type": "array",
                            "items": {"type": "integer", "minimum": 0},
                            "minItems": 2,
                            "maxItems": 2,
                        },
                    },
                    "required": ["start", "end"],
                },
            },
            "format": {
                "type": "object",
                "additionalProperties": False,
                "properties": {
                    "border_color": {"type": "string"},
                    "border_size": {"type": "integer"},
                    "header_fill_color": {"type": "string"},
                    "body_fill_color": {"type": "string"},
                },
            },
        }
    )
    document["sections"] = {
        "type": "array",
        "items": {
            "type": "object",
            "additionalProperties": False,
            "properties": {
                "action": {"type": "string", "enum": ["configure", "add"]},
                "index": {"type": "integer", "minimum": 0},
                "start": {
                    "type": "string",
                    "enum": ["new_page", "continuous", "even_page", "odd_page"],
                },
                "orientation": {"type": "string", "enum": ["portrait", "landscape"]},
                "paper_size": {"type": "string", "enum": ["a4", "letter", "legal"]},
                "margins": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "top_inches": {"type": "number"},
                        "bottom_inches": {"type": "number"},
                        "left_inches": {"type": "number"},
                        "right_inches": {"type": "number"},
                    },
                },
                "header": {"type": "string"},
                "footer": {"type": "string"},
            },
            "required": ["action"],
        },
    }
    document["charts"] = {
        "type": "array",
        "items": {
            "type": "object",
            "additionalProperties": False,
            "properties": {
                "path": {"type": "string"},
                "width_inches": {"type": "number"},
                "alt_text": {"type": "string"},
                "source": {"type": "string"},
            },
            "required": ["path", "alt_text", "source"],
        },
    }

    workbook = properties["workbook"]["properties"]
    properties["workbook"]["additionalProperties"] = False
    workbook["sheets"]["items"]["additionalProperties"] = False
    workbook["cells"]["items"]["additionalProperties"] = False
    sheet_range = {
        "type": "object",
        "additionalProperties": False,
        "properties": {"sheet": {"type": "string"}, "range": {"type": "string"}},
        "required": ["sheet", "range"],
    }
    workbook.update(
        {
            "merged_cells": {"type": "array", "items": copy.deepcopy(sheet_range)},
            "row_heights": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "sheet": {"type": "string"},
                        "row": {"type": "integer"},
                        "height": {"type": "number"},
                    },
                    "required": ["sheet", "row", "height"],
                },
            },
            "column_widths": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "sheet": {"type": "string"},
                        "column": {"type": "string"},
                        "width": {"type": "number"},
                    },
                    "required": ["sheet", "column", "width"],
                },
            },
            "freeze_panes": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "sheet": {"type": "string"},
                        "cell": {
                            "oneOf": [{"type": "string"}, {"type": "null"}]
                        },
                    },
                    "required": ["sheet", "cell"],
                },
            },
            "auto_filters": {"type": "array", "items": copy.deepcopy(sheet_range)},
            "conditional_formats": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "sheet": {"type": "string"},
                        "range": {"type": "string"},
                        "operator": {
                            "type": "string",
                            "enum": [
                                "equal", "notEqual", "greaterThan", "lessThan",
                                "greaterThanOrEqual", "lessThanOrEqual", "between", "notBetween",
                            ],
                        },
                        "formula": {"type": "array", "items": {"type": "string"}},
                        "fill_color": {"type": "string"},
                    },
                    "required": ["sheet", "range", "operator", "formula", "fill_color"],
                },
            },
            "data_validations": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "sheet": {"type": "string"},
                        "range": {"type": "string"},
                        "type": {
                            "type": "string",
                            "enum": ["list", "whole", "decimal", "date", "textLength", "custom"],
                        },
                        "operator": {"type": "string"},
                        "formula1": {"type": "string"},
                        "formula2": {"type": "string"},
                        "allow_blank": {"type": "boolean"},
                    },
                    "required": ["sheet", "range", "type", "formula1"],
                },
            },
            "named_ranges": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "name": {"type": "string"},
                        "sheet": {"type": "string"},
                        "range": {"type": "string"},
                    },
                    "required": ["name", "sheet", "range"],
                },
            },
            "charts": {"type": "array", "items": _xlsx_chart_schema()},
        }
    )

    slide = properties["presentation"]["properties"]["slides"]["items"]
    properties["presentation"]["additionalProperties"] = False
    slide["additionalProperties"] = False
    slide["properties"]["text_boxes"]["items"]["additionalProperties"] = False
    slide["properties"]["tables"]["items"]["additionalProperties"] = False
    slide["properties"]["images"]["items"]["additionalProperties"] = False
    properties["replacements"]["items"]["additionalProperties"] = False
    slide["properties"].update(
        {
            "layout_index": {"type": "integer", "minimum": 0},
            "layout_name": {"type": "string"},
            "title_style": _pptx_text_style_schema(),
            "shapes": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "type": {"type": "string", "enum": ["rectangle", "rounded_rectangle", "ellipse", "line"]},
                        "left_inches": {"type": "number"},
                        "top_inches": {"type": "number"},
                        "width_inches": {"type": "number"},
                        "height_inches": {"type": "number"},
                        "text": {"type": "string"},
                        "fill_color": {"type": "string"},
                        "line_color": {"type": "string"},
                        "text_style": _pptx_text_style_schema(),
                    },
                    "required": ["type", "left_inches", "top_inches", "width_inches", "height_inches"],
                },
            },
            "charts": {"type": "array", "items": _pptx_chart_schema()},
            "speaker_notes": {"type": "string"},
        }
    )
    slide["properties"]["text_boxes"]["items"]["properties"]["style"] = _pptx_text_style_schema()
    slide["properties"]["tables"]["items"]["properties"]["style"] = {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "header_fill_color": {"type": "string"},
            "body_fill_color": {"type": "string"},
            "font_size": {"type": "number"},
        },
    }
    slide["properties"]["images"]["items"]["properties"].update(
        {
            "crop_left": {"type": "number"},
            "crop_top": {"type": "number"},
            "crop_right": {"type": "number"},
            "crop_bottom": {"type": "number"},
            "align": {"type": "string", "enum": ["none", "center", "right"]},
        }
    )


def _pptx_text_style_schema() -> dict[str, Any]:
    return {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "font_size": {"type": "number"},
            "bold": {"type": "boolean"},
            "italic": {"type": "boolean"},
            "color": {"type": "string"},
            "font": {"type": "string"},
            "alignment": {"type": "string", "enum": ["left", "center", "right"]},
        },
    }


def _xlsx_chart_schema() -> dict[str, Any]:
    return {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "sheet": {"type": "string"},
            "type": {"type": "string", "enum": ["bar", "line", "pie", "scatter"]},
            "data_range": {"type": "string"},
            "categories_range": {"type": "string"},
            "x_range": {"type": "string"},
            "y_ranges": {"type": "array", "items": {"type": "string"}},
            "titles_from_data": {"type": "boolean"},
            "series_from": {"type": "string", "enum": ["columns", "rows"]},
            "title": {"type": "string"},
            "anchor": {"type": "string"},
        },
        "required": ["sheet", "type", "anchor"],
    }


def _pptx_chart_schema() -> dict[str, Any]:
    return {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "type": {"type": "string", "enum": ["bar", "line", "pie", "scatter"]},
            "categories": {"type": "array", "items": scalar_schema_for_chart()},
            "series": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "name": {"type": "string"},
                        "values": {"type": "array", "items": {"type": "number"}},
                        "x_values": {"type": "array", "items": {"type": "number"}},
                        "y_values": {"type": "array", "items": {"type": "number"}},
                    },
                    "required": ["name"],
                },
            },
            "title": {"type": "string"},
            "left_inches": {"type": "number"},
            "top_inches": {"type": "number"},
            "width_inches": {"type": "number"},
            "height_inches": {"type": "number"},
        },
        "required": ["type", "series", "left_inches", "top_inches", "width_inches", "height_inches"],
    }


def scalar_schema_for_chart() -> dict[str, Any]:
    return {"oneOf": [{"type": "string"}, {"type": "number"}, {"type": "boolean"}]}


def _run_office_operation(
    target: Path,
    args: Mapping[str, Any],
    ctx: ToolContext,
    staged_workspace: Path,
    *,
    user_template: _UserTemplateLease | None = None,
) -> dict[str, Any]:
    """Build and validate one Office file entirely in private staging."""

    unknown = sorted(set(args) - _ALLOWED_TOP_LEVEL_ARGS)
    if "template_path" in args or "template" in args:
        raise OfficeInputError(
            "Office 工具不接受外部模板。",
            "External templates are not accepted by the Office tool.",
        )
    if unknown:
        raise OfficeInputError(
            f"不支持的 Office 参数：{', '.join(unknown)}",
            f"Unsupported Office parameters: {', '.join(unknown)}",
        )
    _validate_office_v2_request(args)
    _validate_request_budget(args)

    operation = args.get("operation")
    if operation not in {"create", "edit"}:
        raise OfficeInputError(
            "operation 必须是 create 或 edit。",
            "operation must be create or edit.",
        )
    overwrite = args.get("overwrite", False)
    if not isinstance(overwrite, bool):
        raise OfficeInputError("overwrite 必须是布尔值。", "overwrite must be a boolean.")
    if operation == "edit" and overwrite:
        raise OfficeInputError(
            "edit 操作不使用 overwrite。",
            "overwrite is not used with the edit operation.",
        )

    suffix = target.suffix.lower()
    if suffix in _MACRO_OR_TEMPLATE_EXTENSIONS:
        raise OfficeInputError(
            "不支持宏、模板或旧版 Office 格式；仅允许 .docx、.xlsx 和 .pptx。",
            (
                "Macro-enabled, template, and legacy Office formats are not supported; "
                "only .docx, .xlsx, and .pptx are allowed."
            ),
        )
    if suffix not in _FORMATS:
        raise OfficeInputError(
            "Office 文件扩展名必须是 .docx、.xlsx 或 .pptx。",
            "Office file extension must be .docx, .xlsx, or .pptx.",
        )

    if "user_template" in args:
        if user_template is None:
            raise _user_template_error()
        return _run_user_template_operation(
            target,
            args,
            staged_workspace=staged_workspace,
            operation=operation,
            suffix=suffix,
            lease=user_template,
        )
    if user_template is not None:
        raise _user_template_error()

    if "first_party_template" in args:
        return _run_first_party_template_operation(
            target,
            args,
            staged_workspace=staged_workspace,
            operation=operation,
            overwrite=overwrite,
            suffix=suffix,
        )

    expected_payload = _FORMATS[suffix]
    present_payloads = [name for name in _FORMATS.values() if args.get(name) is not None]
    replacements_only_edit = (
        operation == "edit"
        and suffix in {".docx", ".pptx"}
        and args.get("replacements") is not None
        and not present_payloads
    )
    if present_payloads != [expected_payload] and not replacements_only_edit:
        raise OfficeInputError(
            f"{suffix} 必须且只能提供 {expected_payload} 内容。",
            f"{suffix} requires exactly one {expected_payload} payload.",
        )
    if suffix == ".xlsx" and args.get("replacements") is not None:
        raise OfficeInputError(
            "XLSX 请使用 workbook.cells 或 workbook.sheets 进行编辑。",
            "Edit XLSX through workbook.cells or workbook.sheets.",
        )
    if operation == "create" and args.get("replacements") is not None:
        raise OfficeInputError(
            "replacements 仅用于编辑现有 DOCX/PPTX。",
            "replacements is only for editing an existing DOCX/PPTX.",
        )

    target.parent.mkdir(parents=True, exist_ok=True)
    target_exists = target.exists()
    target_mode = (
        stat.S_IMODE(target.stat().st_mode)
        if target_exists and target.is_file()
        else None
    )
    if operation == "edit":
        if not target_exists or not target.is_file():
            raise OfficeInputError(
                f"找不到要编辑的文件：{target}",
                f"File to edit was not found: {target}",
            )
        if target.stat().st_size > MAX_INPUT_FILE_BYTES:
            raise OfficeInputError(
                "Office 输入文件超过 50 MiB 限制。",
                "Office input exceeds the 50 MiB limit.",
            )
        source_parts = _inspect_ooxml_archive(target, suffix, audit_for_edit=True)
        source_digests = (
            _archive_part_digests(target) if _office_v2_enabled() else None
        )
    else:
        source_parts = None
        source_digests = None
    if operation != "edit" and target_exists and not overwrite:
        raise OfficeInputError(
            f"文件已存在：{target}。如需替换，请显式设置 overwrite=true。",
            f"File already exists: {target}. Set overwrite=true to replace it explicitly.",
        )
    elif operation != "edit" and target.exists() and not target.is_file():
        raise OfficeInputError(
            f"目标不是普通文件：{target}",
            f"Target is not a regular file: {target}",
        )

    fd, temporary_name = tempfile.mkstemp(
        prefix=f".{target.name}.",
        suffix=suffix,
        dir=target.parent,
    )
    os.close(fd)
    temporary = Path(temporary_name)
    try:
        logical_workspace = Path(ctx.workspace or "").resolve()
        if suffix == ".docx":
            summary, expected = _write_docx(
                temporary,
                target,
                args,
                operation,
                logical_workspace,
                staged_workspace,
            )
        elif suffix == ".xlsx":
            summary, expected = _write_xlsx(temporary, target, args, operation)
        else:
            summary, expected = _write_pptx(
                temporary,
                target,
                args,
                operation,
                logical_workspace,
                staged_workspace,
            )

        _flush_file(temporary)
        if temporary.stat().st_size > MAX_OUTPUT_FILE_BYTES:
            raise OfficeInputError(
                "Office 输出文件超过 75 MiB 限制。",
                "Office output exceeds the 75 MiB limit.",
            )
        output_parts = _inspect_ooxml_archive(temporary, suffix)
        _reopen_and_verify(temporary, suffix, expected)
        if source_parts is not None:
            _verify_edit_part_preservation(
                source_parts,
                output_parts,
                suffix,
                args,
            )
        if source_digests is not None:
            output_digests = _archive_part_digests(temporary)
            preservation = _verify_untouched_part_digests(
                target,
                suffix,
                args,
                source_digests,
                output_digests,
            )
            summary.update(preservation)

        if target_mode is not None:
            try:
                os.chmod(temporary, target_mode)
            except OSError:
                pass
        _atomic_replace(temporary, target)
        _fsync_directory(target.parent)
        return summary
    finally:
        try:
            temporary.unlink()
        except FileNotFoundError:
            pass


def _run_user_template_operation(
    target: Path,
    args: Mapping[str, Any],
    *,
    staged_workspace: Path,
    operation: str,
    suffix: str,
    lease: _UserTemplateLease,
) -> dict[str, Any]:
    """Instantiate one fixed approved revision inside private WMT staging."""

    template_ref, revision, state_version, raw_values = (
        _parse_user_template_request(args)
    )
    try:
        values = validate_user_template_values(
            lease.placeholder_schema,
            raw_values,
        )
    except TemplateContractError as exc:
        raise OfficeInputError(
            "用户 Office 模板字段不满足已批准的文本长度约束。",
            (
                "User Office template values do not satisfy the approved "
                "text bounds."
            ),
        ) from exc
    if (
        operation != "create"
        or template_ref != lease.template_ref
        or revision != lease.revision
        or state_version != lease.expected_state_version
        or tuple(sorted(values.items())) != lease.values
        or suffix != f".{lease.format}"
    ):
        raise _user_template_error()
    target.parent.mkdir(parents=True, exist_ok=True)
    if target.exists() or target.is_symlink():
        raise OfficeInputError(
            "目标文件已存在；用户 Office 模板不会覆盖它。",
            (
                "The destination already exists; a user Office template "
                "will not overwrite it."
            ),
        )
    try:
        result = lease.service.instantiate_approved(
            template_ref=lease.template_ref,
            revision=lease.revision,
            placeholder_schema=lease.placeholder_schema,
            values=values,
            staging_root=staged_workspace,
            output_path=target,
        )
    except OfficeTemplateError as exc:
        raise _user_template_error() from exc
    if (
        result.template_id != lease.template_ref
        or result.template_version != str(lease.revision)
        or result.source_sha256 != lease.source_sha256
        or result.template_sha256 != lease.manifest_sha256
    ):
        try:
            target.unlink()
        except OSError:
            pass
        raise _user_template_error()
    if target.stat().st_size > MAX_OUTPUT_FILE_BYTES:
        try:
            target.unlink()
        except OSError:
            pass
        raise OfficeInputError(
            "Office 输出文件超过 75 MiB 限制。",
            "Office output exceeds the 75 MiB limit.",
        )
    try:
        _inspect_ooxml_archive(target, suffix)
        _reopen_user_template(target, suffix)
        output_sha256 = hashlib.sha256(target.read_bytes()).hexdigest()
    except Exception:
        try:
            target.unlink()
        except OSError:
            pass
        raise
    if output_sha256 != result.output_sha256:
        try:
            target.unlink()
        except OSError:
            pass
        raise _user_template_error()
    return {
        "format": lease.format,
        "user_template": True,
        "user_template_ref": lease.template_ref,
        "user_template_revision": lease.revision,
        "user_template_state_version": lease.expected_state_version,
        "user_template_source_sha256": lease.source_sha256,
        "user_template_manifest_sha256": lease.manifest_sha256,
        "user_template_render_manifest_sha256": lease.render_manifest_sha256,
        "user_template_output_sha256": output_sha256,
        "user_template_placeholders_replaced": sum(
            change.occurrences for change in result.changes
        ),
    }


def _reopen_user_template(path: Path, suffix: str) -> None:
    try:
        if suffix == ".docx":
            from docx import Document

            Document(str(path))
        elif suffix == ".xlsx":
            from openpyxl import load_workbook

            workbook = load_workbook(
                str(path),
                read_only=False,
                data_only=False,
                keep_vba=False,
                keep_links=False,
            )
            workbook.close()
        else:
            from pptx import Presentation

            Presentation(str(path))
    except Exception as exc:
        raise OfficeInputError(
            "用户 Office 模板输出无法重新打开校验。",
            "The user Office template output failed independent reopen validation.",
        ) from exc


def _run_first_party_template_operation(
    target: Path,
    args: Mapping[str, Any],
    *,
    staged_workspace: Path,
    operation: str,
    overwrite: bool,
    suffix: str,
) -> dict[str, Any]:
    """Instantiate one signed bundled template inside private transaction staging."""

    from app.office_templates import (
        BundledOfficeTemplateService,
        OfficeTemplateError,
        TemplateContractError,
        TemplateFeatureDisabledError,
        TemplateIntegrityError,
        TemplateNotFoundError,
        TemplateSecurityError,
    )

    if operation != "create":
        raise OfficeInputError(
            "第一方 Office 模板仅支持 create 操作。",
            "First-party Office templates support only the create operation.",
        )
    if overwrite:
        raise OfficeInputError(
            "第一方 Office 模板不允许覆盖现有文件。",
            "First-party Office templates cannot overwrite an existing file.",
        )
    conflicting = [
        field
        for field in ("document", "workbook", "presentation", "replacements")
        if args.get(field) is not None
    ]
    if conflicting:
        raise OfficeInputError(
            "第一方 Office 模板不能与普通内容负载同时使用。",
            "A first-party Office template cannot be combined with an ordinary content payload.",
        )
    raw = args.get("first_party_template")
    template = _mapping(raw, "first_party_template")
    _strict_fields(
        template,
        {"template_id", "template_version", "values"},
        "first_party_template",
    )
    template_id = _required_text(
        template.get("template_id"),
        "first_party_template.template_id",
    )
    template_version = _required_text(
        template.get("template_version"),
        "first_party_template.template_version",
    )
    values = template.get("values")
    if not isinstance(values, Mapping):
        raise OfficeInputError(
            "first_party_template.values 必须是字段对象。",
            "first_party_template.values must be a field object.",
        )

    target.parent.mkdir(parents=True, exist_ok=True)
    if target.exists():
        raise OfficeInputError(
            f"文件已存在：{target}。第一方模板不会覆盖它。",
            f"File already exists: {target}. A first-party template will not overwrite it.",
        )

    registry_root = (
        default_file_version_storage_root().parent
        / "v1.1"
        / "office-template-registry"
    ).absolute()
    try:
        service = BundledOfficeTemplateService(registry_root)
        descriptor = next(
            (
                item
                for item in service.list_templates()
                if item.immutable_key == (template_id, template_version)
            ),
            None,
        )
        if descriptor is None:
            raise TemplateNotFoundError("first-party Office template was not found")
        expected_suffix = f".{descriptor.manifest.format}"
        if suffix != expected_suffix:
            raise OfficeInputError(
                "目标扩展名与第一方模板格式不匹配。",
                "The destination extension does not match the first-party template format.",
            )
        result = service.instantiate(
            template_id,
            template_version,
            values,
            staging_root=staged_workspace,
            output_path=target,
        )
    except OfficeInputError:
        raise
    except TemplateFeatureDisabledError as exc:
        raise OfficeInputError(
            "第一方 Office 模板尚未发布。",
            "First-party Office templates are not released.",
        ) from exc
    except (TemplateContractError, TemplateNotFoundError) as exc:
        raise OfficeInputError(
            "第一方 Office 模板参数与签名目录不匹配。",
            "The first-party Office template request does not match the signed catalog.",
        ) from exc
    except (TemplateIntegrityError, TemplateSecurityError) as exc:
        raise OfficeInputError(
            "第一方 Office 模板完整性或安全校验失败。",
            "First-party Office template integrity or safety validation failed.",
        ) from exc
    except OfficeTemplateError as exc:
        raise OfficeInputError(
            "第一方 Office 模板无法安全生成。",
            "The first-party Office template could not be generated safely.",
        ) from exc

    if target.stat().st_size > MAX_OUTPUT_FILE_BYTES:
        try:
            target.unlink()
        except OSError:
            pass
        raise OfficeInputError(
            "Office 输出文件超过 75 MiB 限制。",
            "Office output exceeds the 75 MiB limit.",
        )
    try:
        _inspect_ooxml_archive(target, suffix)
        _reopen_first_party_template(target, suffix)
    except Exception:
        try:
            target.unlink()
        except OSError:
            pass
        raise
    if hashlib.sha256(target.read_bytes()).hexdigest() != result.output_sha256:
        try:
            target.unlink()
        except OSError:
            pass
        raise OfficeInputError(
            "第一方 Office 模板输出在校验期间发生变化。",
            "The first-party Office template output changed during validation.",
        )
    return {
        "format": descriptor.manifest.format,
        "first_party_template": True,
        "template_id": result.template_id,
        "template_version": result.template_version,
        "template_source_sha256": result.source_sha256,
        "template_manifest_sha256": result.template_sha256,
        "template_output_sha256": result.output_sha256,
        "template_placeholders_replaced": sum(
            change.occurrences for change in result.changes
        ),
        "template_render_baseline_id": (
            descriptor.expected_render_baseline.baseline_id
        ),
    }


def _reopen_first_party_template(path: Path, suffix: str) -> None:
    """Prove the instantiated copy opens in the same libraries used by edits."""

    try:
        if suffix == ".docx":
            from docx import Document

            Document(str(path))
        elif suffix == ".xlsx":
            from openpyxl import load_workbook

            workbook = load_workbook(
                str(path),
                read_only=False,
                data_only=False,
                keep_vba=False,
                keep_links=False,
            )
            workbook.close()
        else:
            from pptx import Presentation

            Presentation(str(path))
    except Exception as exc:
        raise OfficeInputError(
            "第一方 Office 模板输出无法重新打开校验。",
            "The first-party Office template output failed independent reopen validation.",
        ) from exc


def _validate_office_v2_request(args: Mapping[str, Any]) -> None:
    """Make v1.1 inputs impossible to smuggle through while the gate is off."""

    if "user_template" in args and not _user_office_templates_enabled():
        raise OfficeInputError(
            "用户 Office 模板 Beta 尚未发布。",
            "User Office template Beta is not released.",
        )
    found: list[str] = []
    if "first_party_template" in args:
        found.append("first_party_template")
    document = args.get("document")
    if isinstance(document, Mapping):
        found.extend(f"document.{key}" for key in {"sections", "charts"} & set(document))
        for index, paragraph in enumerate(_safe_mapping_sequence(document.get("paragraphs"))):
            found.extend(
                f"document.paragraphs[{index}].{key}"
                for key in {"runs", "format", "list"} & set(paragraph)
            )
        for index, table in enumerate(_safe_mapping_sequence(document.get("tables"))):
            found.extend(
                f"document.tables[{index}].{key}"
                for key in {"merges", "format"} & set(table)
            )
    workbook = args.get("workbook")
    if isinstance(workbook, Mapping):
        found.extend(
            f"workbook.{key}"
            for key in {
                "merged_cells", "row_heights", "column_widths", "freeze_panes",
                "auto_filters", "conditional_formats", "data_validations",
                "named_ranges", "charts",
            }
            & set(workbook)
        )
    presentation = args.get("presentation")
    if isinstance(presentation, Mapping):
        for index, slide in enumerate(_safe_mapping_sequence(presentation.get("slides"))):
            found.extend(
                f"presentation.slides[{index}].{key}"
                for key in {"layout_index", "layout_name", "title_style", "shapes", "charts", "speaker_notes"}
                & set(slide)
            )
            for field in ("text_boxes", "tables", "images"):
                for item_index, item in enumerate(_safe_mapping_sequence(slide.get(field))):
                    extra = {
                        "text_boxes": {"style"},
                        "tables": {"style"},
                        "images": {"crop_left", "crop_top", "crop_right", "crop_bottom", "align"},
                    }[field] & set(item)
                    found.extend(
                        f"presentation.slides[{index}].{field}[{item_index}].{key}"
                        for key in extra
                    )
    if found and not _office_v2_enabled():
        preview = ", ".join(sorted(found)[:3])
        raise OfficeInputError(
            f"Office v1.1 功能尚未发布：{preview}",
            f"Office v1.1 fields are not released: {preview}",
        )


def _safe_mapping_sequence(value: Any) -> Iterator[Mapping[str, Any]]:
    if isinstance(value, Sequence) and not isinstance(value, (str, bytes)):
        for item in value:
            if isinstance(item, Mapping):
                yield item


def _write_docx(
    temporary: Path,
    target: Path,
    args: Mapping[str, Any],
    operation: str,
    workspace: Path,
    staged_workspace: Path,
) -> tuple[dict[str, Any], dict[str, Any]]:
    from docx import Document

    document = Document(str(target)) if operation == "edit" else Document()
    if operation == "create":
        _drop_default_docx_custom_xml(document)
    raw_payload = args.get("document")
    payload = (
        {}
        if operation == "edit" and raw_payload is None
        else _mapping(raw_payload, "document")
    )
    if _office_v2_enabled():
        _strict_fields(
            payload,
            {"title", "paragraphs", "tables", "images", "sections", "charts"},
            "document",
        )
    replacements = _parse_replacements(args.get("replacements"))
    paragraphs = _parse_docx_paragraphs(payload.get("paragraphs"))
    tables = _parse_tables(payload.get("tables"))
    images = _parse_docx_images(
        payload.get("images"),
        workspace,
        staged_workspace,
    )
    sections = (
        _parse_docx_sections(payload.get("sections"))
        if _office_v2_enabled()
        else []
    )
    charts = (
        _parse_docx_charts(
            payload.get("charts"),
            workspace,
            staged_workspace,
        )
        if _office_v2_enabled()
        else []
    )
    _validate_total_image_bytes(
        sum(len(image["data"]) for image in [*images, *charts])
    )
    title = _optional_text(payload.get("title"), "document.title")

    if operation == "create" and not (
        title or paragraphs or tables or images or sections or charts
    ):
        raise OfficeInputError(
            "DOCX 创建至少需要 title、paragraphs 或 tables 之一。",
            "DOCX creation requires title, paragraphs, or tables.",
        )
    if operation == "edit" and not (
        replacements or title or paragraphs or tables or images or sections or charts
    ):
        raise OfficeInputError(
            "DOCX 编辑至少需要一项变更。",
            "DOCX editing requires at least one change.",
        )

    replaced = _apply_replacements(
        list(_iter_docx_paragraphs(document)), replacements, "DOCX"
    )
    if title:
        document.core_properties.title = title
        document.add_heading(title, level=0)
    for item in paragraphs:
        paragraph = document.add_paragraph()
        paragraph.style = _DOCX_STYLE_NAMES[item["style"]]
        if item.get("runs") is not None:
            for run_data in item["runs"]:
                _add_docx_run(paragraph, run_data)
        else:
            paragraph.add_run(item["text"])
        if item.get("format"):
            _apply_docx_paragraph_format(paragraph, item["format"])
        if item.get("list"):
            _apply_docx_multilevel_list(document, paragraph, item["list"])
        if item["page_break_after"]:
            document.add_page_break()
    for table_data in tables:
        headers = table_data["headers"]
        rows = table_data["rows"]
        column_count = max([len(headers), *(len(row) for row in rows)])
        table = document.add_table(rows=1 if headers else 0, cols=column_count)
        table.style = "Table Grid"
        if headers:
            for index, value in enumerate(headers):
                table.rows[0].cells[index].text = _cell_text(value)
        for row in rows:
            cells = table.add_row().cells
            for index, value in enumerate(row):
                cells[index].text = _cell_text(value)
        _apply_docx_table_v2(table, table_data)

    from docx.shared import Inches

    for image in images:
        width = Inches(image["width_inches"]) if image["width_inches"] else None
        document.add_picture(io.BytesIO(image["data"]), width=width)
        if image["caption"]:
            caption = document.add_paragraph(image["caption"])
            try:
                caption.style = "Caption"
            except KeyError:
                # Default python-docx templates include Caption.  A safe edit
                # of an unusual file should still retain the requested text.
                pass

    for chart in charts:
        width = Inches(chart["width_inches"]) if chart["width_inches"] else None
        inline_shape = document.add_picture(io.BytesIO(chart["data"]), width=width)
        inline_shape._inline.docPr.set("descr", chart["alt_text"])
        inline_shape._inline.docPr.set("title", chart["source"])

    if sections:
        _apply_docx_sections(document, sections)

    expected = {
        "semantic_digest": _semantic_digest(_docx_semantic(document)),
        "inline_shapes": len(document.inline_shapes),
        "page_breaks": _docx_page_break_count(document),
    }
    if _office_v2_enabled():
        expected["v2_snapshot"] = _docx_v2_snapshot(document)
    document.save(str(temporary))
    return (
        {
            "format": "docx",
            "paragraphs_added": len(paragraphs) + (1 if title else 0),
            "tables_added": len(tables),
            "images_added": len(images),
            **(
                {
                    "static_charts_added": len(charts),
                    "sections_changed": len(sections),
                }
                if _office_v2_enabled()
                else {}
            ),
            "page_breaks_added": sum(
                int(item["page_break_after"]) for item in paragraphs
            ),
            "replacements": replaced,
        },
        expected,
    )


def _write_xlsx(
    temporary: Path,
    target: Path,
    args: Mapping[str, Any],
    operation: str,
) -> tuple[dict[str, Any], dict[str, Any]]:
    from openpyxl import Workbook, load_workbook

    if operation == "edit":
        workbook = load_workbook(
            str(target),
            read_only=False,
            data_only=False,
            keep_vba=False,
            keep_links=True,
        )
    else:
        workbook = Workbook()
        workbook.remove(workbook.active)

    try:
        return _write_xlsx_workbook(temporary, workbook, args, operation)
    finally:
        workbook.close()


def _write_xlsx_workbook(
    temporary: Path,
    workbook: Any,
    args: Mapping[str, Any],
    operation: str,
) -> tuple[dict[str, Any], dict[str, Any]]:
    from openpyxl.utils import get_column_letter
    from openpyxl.utils.cell import coordinate_to_tuple

    payload = _mapping(args.get("workbook"), "workbook")
    if _office_v2_enabled():
        _strict_fields(
            payload,
            {
                "sheets", "cells", "delete_sheets", "merged_cells", "row_heights",
                "column_widths", "freeze_panes", "auto_filters", "conditional_formats",
                "data_validations", "named_ranges", "charts",
            },
            "workbook",
        )
    sheets = _sequence(payload.get("sheets", []), "workbook.sheets")
    cells = _sequence(payload.get("cells", []), "workbook.cells")
    delete_sheets = _sequence(
        payload.get("delete_sheets", []), "workbook.delete_sheets"
    )
    v2_keys = (
        "merged_cells",
        "row_heights",
        "column_widths",
        "freeze_panes",
        "auto_filters",
        "conditional_formats",
        "data_validations",
        "named_ranges",
        "charts",
    )
    v2_items = {
        key: _sequence(payload.get(key, []), f"workbook.{key}")
        for key in v2_keys
    }
    if operation == "create" and not sheets:
        raise OfficeInputError(
            "XLSX 创建至少需要一个 workbook.sheets 项。",
            "XLSX creation requires at least one workbook.sheets item.",
        )
    if operation == "create" and delete_sheets:
        raise OfficeInputError(
            "workbook.delete_sheets 仅用于编辑现有 XLSX。",
            "workbook.delete_sheets is only for editing an existing XLSX file.",
        )
    if operation == "edit" and not (
        sheets or cells or delete_sheets or any(v2_items.values())
    ):
        raise OfficeInputError(
            "XLSX 编辑至少需要一项变更。",
            "XLSX editing requires at least one change.",
        )
    if len(sheets) > MAX_SHEETS:
        raise OfficeInputError(
            f"单次最多处理 {MAX_SHEETS} 个工作表。",
            f"At most {MAX_SHEETS} sheets may be processed in one call.",
        )

    written_cells: dict[tuple[str, str], Any] = {}
    written_styles: dict[tuple[str, str], dict[str, Any]] = {}
    formula_count = 0
    cell_count = 0
    created_sheets = 0
    deleted_sheets = 0
    appended_rows = 0
    styles_applied = 0
    requested_names: set[str] = set()

    delete_names: list[str] = []
    for index, raw_name in enumerate(delete_sheets):
        name = _required_text(raw_name, f"workbook.delete_sheets[{index}]")
        if name in delete_names:
            raise OfficeInputError(
                f"重复的工作表删除项：{name}",
                f"Duplicate sheet deletion: {name}",
            )
        if name not in workbook.sheetnames:
            raise OfficeInputError(
                f"找不到要删除的工作表：{name}",
                f"Sheet to delete was not found: {name}",
            )
        delete_names.append(name)
    if delete_names and len(delete_names) >= len(workbook.sheetnames):
        raise OfficeInputError(
            "XLSX 必须至少保留一个工作表。",
            "An XLSX file must retain at least one sheet.",
        )
    for name in delete_names:
        workbook.remove(workbook[name])
        deleted_sheets += 1

    for index, raw_sheet in enumerate(sheets):
        sheet = _mapping(raw_sheet, f"workbook.sheets[{index}]")
        if _office_v2_enabled():
            _strict_fields(sheet, {"name", "action", "rows"}, f"workbook.sheets[{index}]")
        name = _required_text(sheet.get("name"), f"workbook.sheets[{index}].name")
        _validate_sheet_name(name)
        folded = name.casefold()
        if folded in requested_names:
            raise OfficeInputError(
                f"工作表名重复：{name}",
                f"Duplicate sheet name: {name}",
            )
        requested_names.add(folded)

        action = sheet.get("action", "create" if operation == "create" else "append")
        if action not in {"create", "append"}:
            raise OfficeInputError(
                "workbook.sheets[].action 必须是 create 或 append。",
                "workbook.sheets[].action must be create or append.",
            )
        if operation == "create" and action != "create":
            raise OfficeInputError(
                "创建 XLSX 时，sheet action 只能是 create。",
                "Sheet action must be create when creating an XLSX file.",
            )
        if action == "create":
            if name.casefold() in {existing.casefold() for existing in workbook.sheetnames}:
                raise OfficeInputError(
                    f"工作表已存在：{name}",
                    f"Sheet already exists: {name}",
                )
            worksheet = workbook.create_sheet(name)
            created_sheets += 1
        else:
            if name not in workbook.sheetnames:
                raise OfficeInputError(
                    f"找不到工作表：{name}",
                    f"Sheet not found: {name}",
                )
            worksheet = workbook[name]

        rows = _sequence(sheet.get("rows"), f"workbook.sheets[{index}].rows")
        for row_index, raw_row in enumerate(rows):
            row = _row_values(raw_row, f"workbook.sheets[{index}].rows[{row_index}]")
            cell_count += len(row)
            if cell_count > MAX_WORKBOOK_CELLS:
                raise OfficeInputError(
                    f"单次最多写入 {MAX_WORKBOOK_CELLS} 个单元格。",
                    f"At most {MAX_WORKBOOK_CELLS} cells may be written in one call.",
                )
            for value in row:
                formula_count += int(_is_formula(value))
            worksheet.append(row)
            appended_rows += 1
            installed_row = worksheet.max_row
            for column_index, value in enumerate(row, 1):
                coordinate = f"{get_column_letter(column_index)}{installed_row}"
                written_cells[(name, coordinate)] = value

    for index, raw_cell in enumerate(cells):
        item = _mapping(raw_cell, f"workbook.cells[{index}]")
        if _office_v2_enabled():
            _strict_fields(item, {"sheet", "cell", "value", "style"}, f"workbook.cells[{index}]")
        sheet_name = _required_text(item.get("sheet"), f"workbook.cells[{index}].sheet")
        coordinate = _required_text(item.get("cell"), f"workbook.cells[{index}].cell").upper()
        if sheet_name not in workbook.sheetnames:
            raise OfficeInputError(
                f"找不到工作表：{sheet_name}",
                f"Sheet not found: {sheet_name}",
            )
        try:
            row_number, column_number = coordinate_to_tuple(coordinate)
        except (TypeError, ValueError):
            raise OfficeInputError(
                f"无效单元格坐标：{coordinate}",
                f"Invalid cell coordinate: {coordinate}",
            ) from None
        if not 1 <= row_number <= 1_048_576 or not 1 <= column_number <= 16_384:
            raise OfficeInputError(
                f"单元格坐标超出 XLSX 范围：{coordinate}",
                f"Cell coordinate exceeds XLSX limits: {coordinate}",
            )
        has_value = "value" in item
        has_style = item.get("style") is not None
        if not has_value and not has_style:
            raise OfficeInputError(
                f"workbook.cells[{index}] 至少需要 value 或 style。",
                f"workbook.cells[{index}] requires value or style.",
            )
        cell = workbook[sheet_name][coordinate]
        if has_value:
            value = _scalar(item["value"], f"workbook.cells[{index}].value")
            formula_count += int(_is_formula(value))
            cell.value = value
            written_cells[(sheet_name, coordinate)] = value
        if has_style:
            _apply_xlsx_style(cell, item["style"], f"workbook.cells[{index}].style")
            written_styles[(sheet_name, coordinate)] = _xlsx_style_snapshot(cell)
            styles_applied += 1
        cell_count += 1
        if cell_count > MAX_WORKBOOK_CELLS:
            raise OfficeInputError(
                f"单次最多写入 {MAX_WORKBOOK_CELLS} 个单元格。",
                f"At most {MAX_WORKBOOK_CELLS} cells may be written in one call.",
            )

    v2_summary: dict[str, int] = {}
    if _office_v2_enabled():
        v2_summary = _apply_xlsx_v2(workbook, v2_items)

    if not any(sheet.sheet_state == "visible" for sheet in workbook.worksheets):
        raise OfficeInputError(
            "XLSX 必须至少保留一个可见工作表。",
            "An XLSX file must retain at least one visible sheet.",
        )
    workbook.save(str(temporary))
    expected = {
        "sheet_names": tuple(workbook.sheetnames),
        "written_cells": written_cells,
        "written_styles": written_styles,
    }
    if _office_v2_enabled():
        expected["v2_snapshot"] = _xlsx_v2_snapshot(workbook)
    return (
        {
            "format": "xlsx",
            "sheets_created": created_sheets,
            "sheets_deleted": deleted_sheets,
            "rows_appended": appended_rows,
            "cells_written": cell_count,
            "styles_applied": styles_applied,
            "formulas_written": formula_count,
            **v2_summary,
        },
        expected,
    )


def _strict_fields(value: Mapping[str, Any], allowed: set[str], name: str) -> None:
    unknown = sorted(set(value) - allowed)
    if unknown:
        raise OfficeInputError(
            f"{name} 包含不支持的字段：{', '.join(unknown)}",
            f"{name} contains unsupported fields: {', '.join(unknown)}",
        )


def _xlsx_sheet(workbook: Any, raw: Any, name: str) -> Any:
    sheet_name = _required_text(raw, name)
    if sheet_name not in workbook.sheetnames:
        raise OfficeInputError(
            f"找不到工作表：{sheet_name}",
            f"Sheet not found: {sheet_name}",
        )
    return workbook[sheet_name]


def _xlsx_range(value: Any, name: str) -> tuple[str, tuple[int, int, int, int]]:
    from openpyxl.utils.cell import range_boundaries

    raw = _required_text(value, name).upper()
    if not _CELL_RANGE.fullmatch(raw):
        raise OfficeInputError(
            f"{name} 必须是明确的 A1 区间（例如 A1:D10）。",
            f"{name} must be an explicit A1 range such as A1:D10.",
        )
    min_col, min_row, max_col, max_row = range_boundaries(raw)
    if min_col > max_col or min_row > max_row or max_col > 16_384 or max_row > 1_048_576:
        raise OfficeInputError(
            f"{name} 超出 XLSX 范围或起止顺序无效。",
            f"{name} exceeds XLSX bounds or has reversed endpoints.",
        )
    return raw, (min_col, min_row, max_col, max_row)


def _xlsx_cell(value: Any, name: str) -> str:
    from openpyxl.utils.cell import coordinate_to_tuple

    raw = _required_text(value, name).upper()
    if not _CELL_COORDINATE.fullmatch(raw):
        raise OfficeInputError(
            f"{name} 必须是单个 A1 单元格坐标。",
            f"{name} must be one A1 cell coordinate.",
        )
    row, column = coordinate_to_tuple(raw)
    if row > 1_048_576 or column > 16_384:
        raise OfficeInputError(
            f"{name} 超出 XLSX 范围。",
            f"{name} exceeds XLSX bounds.",
        )
    return raw


def _rectangles_overlap(
    left: tuple[int, int, int, int],
    right: tuple[int, int, int, int],
) -> bool:
    return not (
        left[2] < right[0]
        or right[2] < left[0]
        or left[3] < right[1]
        or right[3] < left[1]
    )


def _xlsx_merged_range_area(bounds: tuple[int, int, int, int]) -> int:
    return (bounds[2] - bounds[0] + 1) * (bounds[3] - bounds[1] + 1)


def _apply_xlsx_v2(
    workbook: Any,
    items: Mapping[str, Sequence[Any]],
) -> dict[str, int]:
    from openpyxl.formatting.rule import CellIsRule
    from openpyxl.styles import PatternFill
    from openpyxl.utils import quote_sheetname
    from openpyxl.utils.cell import absolute_coordinate
    from openpyxl.worksheet.datavalidation import DataValidation
    from openpyxl.workbook.defined_name import DefinedName

    if sum(len(value) for value in items.values()) > MAX_V2_FORMAT_ITEMS:
        raise OfficeInputError(
            f"Office v1.1 复杂格式项最多 {MAX_V2_FORMAT_ITEMS} 项。",
            f"Office v1.1 accepts at most {MAX_V2_FORMAT_ITEMS} complex-format items.",
        )
    existing_chart_count = sum(
        len(worksheet._charts) for worksheet in workbook.worksheets
    )
    if existing_chart_count + len(items["charts"]) > MAX_CHARTS_PER_FILE:
        raise OfficeInputError(
            f"单个工作簿操作后最多包含 {MAX_CHARTS_PER_FILE} 个图表。",
            (
                f"A workbook may contain at most {MAX_CHARTS_PER_FILE} charts "
                "after this operation."
            ),
        )

    summary = {
        "merged_ranges_added": 0,
        "dimensions_changed": 0,
        "freeze_panes_changed": 0,
        "auto_filters_changed": 0,
        "conditional_formats_added": 0,
        "data_validations_added": 0,
        "named_ranges_added": 0,
        "charts_added": 0,
    }

    requested_merges: dict[str, list[tuple[int, int, int, int]]] = {}
    planned_merges: list[tuple[Any, str]] = []
    total_merged_cells = sum(
        _xlsx_merged_range_area(tuple(cell_range.bounds))
        for worksheet in workbook.worksheets
        for cell_range in worksheet.merged_cells.ranges
    )
    for index, raw in enumerate(items["merged_cells"]):
        name = f"workbook.merged_cells[{index}]"
        item = _mapping(raw, name)
        _strict_fields(item, {"sheet", "range"}, name)
        worksheet = _xlsx_sheet(workbook, item.get("sheet"), f"{name}.sheet")
        cell_range, bounds = _xlsx_range(item.get("range"), f"{name}.range")
        if bounds[0] == bounds[2] and bounds[1] == bounds[3]:
            raise OfficeInputError(
                f"{name}.range 必须包含至少两个单元格。",
                f"{name}.range must contain at least two cells.",
            )
        area = _xlsx_merged_range_area(bounds)
        if area > MAX_MERGED_CELLS_PER_RANGE:
            raise OfficeInputError(
                f"{name}.range 合并单元格过多。",
                f"{name}.range contains too many merged cells.",
            )
        total_merged_cells += area
        if total_merged_cells > MAX_TOTAL_MERGED_CELLS:
            raise OfficeInputError(
                "workbook.merged_cells 合并单元格总量过大。",
                "workbook.merged_cells exceeds the aggregate merged-cell budget.",
            )
        existing = [tuple(r.bounds) for r in worksheet.merged_cells.ranges]
        prior = requested_merges.setdefault(worksheet.title, [])
        if any(_rectangles_overlap(bounds, other) for other in [*existing, *prior]):
            raise OfficeInputError(
                f"{name}.range 与现有或本次合并区域重叠。",
                f"{name}.range overlaps an existing or requested merged range.",
            )
        prior.append(bounds)
        planned_merges.append((worksheet, cell_range))

    # Validate every range and the aggregate allocation before openpyxl creates
    # any MergedCell objects.  A single full-sheet merge otherwise allocates
    # billions of objects before the output-size guard can run.
    for worksheet, cell_range in planned_merges:
        worksheet.merge_cells(cell_range)
        summary["merged_ranges_added"] += 1

    seen_rows: set[tuple[str, int]] = set()
    for index, raw in enumerate(items["row_heights"]):
        name = f"workbook.row_heights[{index}]"
        item = _mapping(raw, name)
        _strict_fields(item, {"sheet", "row", "height"}, name)
        worksheet = _xlsx_sheet(workbook, item.get("sheet"), f"{name}.sheet")
        row = item.get("row")
        if isinstance(row, bool) or not isinstance(row, int) or not 1 <= row <= 1_048_576:
            raise OfficeInputError(f"{name}.row 无效。", f"{name}.row is invalid.")
        key = (worksheet.title, row)
        if key in seen_rows:
            raise OfficeInputError(f"{name} 重复。", f"{name} is duplicated.")
        seen_rows.add(key)
        worksheet.row_dimensions[row].height = _bounded_number(
            item.get("height"), f"{name}.height", minimum=0.1, maximum=409
        )
        summary["dimensions_changed"] += 1

    seen_columns: set[tuple[str, str]] = set()
    for index, raw in enumerate(items["column_widths"]):
        name = f"workbook.column_widths[{index}]"
        item = _mapping(raw, name)
        _strict_fields(item, {"sheet", "column", "width"}, name)
        worksheet = _xlsx_sheet(workbook, item.get("sheet"), f"{name}.sheet")
        column = _required_text(item.get("column"), f"{name}.column").upper()
        if not _COLUMN_COORDINATE.fullmatch(column):
            raise OfficeInputError(
                f"{name}.column 必须是列字母。",
                f"{name}.column must be column letters.",
            )
        from openpyxl.utils.cell import column_index_from_string

        if column_index_from_string(column) > 16_384:
            raise OfficeInputError(f"{name}.column 超出范围。", f"{name}.column is out of bounds.")
        key = (worksheet.title, column)
        if key in seen_columns:
            raise OfficeInputError(f"{name} 重复。", f"{name} is duplicated.")
        seen_columns.add(key)
        worksheet.column_dimensions[column].width = _bounded_number(
            item.get("width"), f"{name}.width", minimum=0.1, maximum=255
        )
        summary["dimensions_changed"] += 1

    seen_freeze: set[str] = set()
    for index, raw in enumerate(items["freeze_panes"]):
        name = f"workbook.freeze_panes[{index}]"
        item = _mapping(raw, name)
        _strict_fields(item, {"sheet", "cell"}, name)
        worksheet = _xlsx_sheet(workbook, item.get("sheet"), f"{name}.sheet")
        if worksheet.title in seen_freeze:
            raise OfficeInputError(f"{name} 重复。", f"{name} is duplicated.")
        seen_freeze.add(worksheet.title)
        cell = item.get("cell")
        worksheet.freeze_panes = None if cell is None else _xlsx_cell(cell, f"{name}.cell")
        summary["freeze_panes_changed"] += 1

    seen_filters: set[str] = set()
    for index, raw in enumerate(items["auto_filters"]):
        name = f"workbook.auto_filters[{index}]"
        item = _mapping(raw, name)
        _strict_fields(item, {"sheet", "range"}, name)
        worksheet = _xlsx_sheet(workbook, item.get("sheet"), f"{name}.sheet")
        if worksheet.title in seen_filters:
            raise OfficeInputError(f"{name} 重复。", f"{name} is duplicated.")
        seen_filters.add(worksheet.title)
        cell_range, _ = _xlsx_range(item.get("range"), f"{name}.range")
        worksheet.auto_filter.ref = cell_range
        summary["auto_filters_changed"] += 1

    for index, raw in enumerate(items["conditional_formats"]):
        name = f"workbook.conditional_formats[{index}]"
        item = _mapping(raw, name)
        _strict_fields(item, {"sheet", "range", "operator", "formula", "fill_color"}, name)
        worksheet = _xlsx_sheet(workbook, item.get("sheet"), f"{name}.sheet")
        cell_range, bounds = _xlsx_range(item.get("range"), f"{name}.range")
        if (bounds[2] - bounds[0] + 1) * (bounds[3] - bounds[1] + 1) > MAX_CHART_POINTS:
            raise OfficeInputError(f"{name}.range 过大。", f"{name}.range is too large.")
        operator = item.get("operator")
        allowed_operators = {
            "equal", "notEqual", "greaterThan", "lessThan", "greaterThanOrEqual",
            "lessThanOrEqual", "between", "notBetween",
        }
        if operator not in allowed_operators:
            raise OfficeInputError(f"{name}.operator 不支持。", f"{name}.operator is unsupported.")
        formulas_raw = _sequence(item.get("formula"), f"{name}.formula")
        required_count = 2 if operator in {"between", "notBetween"} else 1
        if len(formulas_raw) != required_count:
            raise OfficeInputError(
                f"{name}.formula 必须包含 {required_count} 项。",
                f"{name}.formula must contain {required_count} item(s).",
            )
        formulas = [_xlsx_formula_text(value, f"{name}.formula[{i}]") for i, value in enumerate(formulas_raw)]
        color = _validate_hex_color(item.get("fill_color"), f"{name}.fill_color")
        fill = PatternFill(fill_type="solid", fgColor=color)
        worksheet.conditional_formatting.add(
            cell_range,
            CellIsRule(operator=operator, formula=formulas, fill=fill),
        )
        summary["conditional_formats_added"] += 1

    allowed_validation_types = {"list", "whole", "decimal", "date", "textLength", "custom"}
    allowed_validation_operators = {
        "between", "notBetween", "equal", "notEqual", "lessThan", "lessThanOrEqual",
        "greaterThan", "greaterThanOrEqual",
    }
    for index, raw in enumerate(items["data_validations"]):
        name = f"workbook.data_validations[{index}]"
        item = _mapping(raw, name)
        _strict_fields(item, {"sheet", "range", "type", "operator", "formula1", "formula2", "allow_blank"}, name)
        worksheet = _xlsx_sheet(workbook, item.get("sheet"), f"{name}.sheet")
        cell_range, bounds = _xlsx_range(item.get("range"), f"{name}.range")
        if (bounds[2] - bounds[0] + 1) * (bounds[3] - bounds[1] + 1) > MAX_CHART_POINTS:
            raise OfficeInputError(f"{name}.range 过大。", f"{name}.range is too large.")
        validation_type = item.get("type")
        if validation_type not in allowed_validation_types:
            raise OfficeInputError(f"{name}.type 不支持。", f"{name}.type is unsupported.")
        operator = item.get("operator")
        if validation_type in {"list", "custom"}:
            if operator is not None or item.get("formula2") is not None:
                raise OfficeInputError(
                    f"{name} 的 list/custom 不接受 operator 或 formula2。",
                    f"{name} list/custom does not accept operator or formula2.",
                )
        elif operator not in allowed_validation_operators:
            raise OfficeInputError(f"{name}.operator 无效。", f"{name}.operator is invalid.")
        if operator in {"between", "notBetween"} and item.get("formula2") is None:
            raise OfficeInputError(f"{name}.formula2 必填。", f"{name}.formula2 is required.")
        allow_blank = item.get("allow_blank", False)
        if not isinstance(allow_blank, bool):
            raise OfficeInputError(f"{name}.allow_blank 必须是布尔值。", f"{name}.allow_blank must be boolean.")
        validation = DataValidation(
            type=validation_type,
            operator=operator,
            formula1=_xlsx_formula_text(item.get("formula1"), f"{name}.formula1"),
            formula2=(
                _xlsx_formula_text(item.get("formula2"), f"{name}.formula2")
                if item.get("formula2") is not None
                else None
            ),
            allow_blank=allow_blank,
        )
        worksheet.add_data_validation(validation)
        validation.add(cell_range)
        summary["data_validations_added"] += 1

    requested_names: set[str] = set()
    existing_names = {defined.name.casefold() for defined in workbook.defined_names.values()}
    for index, raw in enumerate(items["named_ranges"]):
        item_name = f"workbook.named_ranges[{index}]"
        item = _mapping(raw, item_name)
        _strict_fields(item, {"name", "sheet", "range"}, item_name)
        name = _required_text(item.get("name"), f"{item_name}.name")
        if not _DEFINED_NAME.fullmatch(name) or _CELL_COORDINATE.fullmatch(name.upper()):
            raise OfficeInputError(f"{item_name}.name 无效。", f"{item_name}.name is invalid.")
        if name.casefold() in existing_names | requested_names:
            raise OfficeInputError(f"{item_name}.name 重复。", f"{item_name}.name is duplicated.")
        worksheet = _xlsx_sheet(workbook, item.get("sheet"), f"{item_name}.sheet")
        cell_range, _ = _xlsx_range(item.get("range"), f"{item_name}.range")
        destination = f"{quote_sheetname(worksheet.title)}!{absolute_coordinate(cell_range)}"
        workbook.defined_names.add(DefinedName(name, attr_text=destination))
        requested_names.add(name.casefold())
        summary["named_ranges_added"] += 1

    for index, raw in enumerate(items["charts"]):
        _add_xlsx_chart(workbook, raw, f"workbook.charts[{index}]")
        summary["charts_added"] += 1
    return summary


def _xlsx_formula_text(value: Any, name: str) -> str:
    text = _required_text(value, name)
    if _EXTERNAL_FORMULA.search(text):
        raise OfficeInputError(
            f"{name} 包含外部工作簿或网络引用。",
            f"{name} contains an external workbook or network reference.",
        )
    if len(text) > 8_192:
        raise OfficeInputError(f"{name} 过长。", f"{name} is too long.")
    return text


def _add_xlsx_chart(workbook: Any, raw: Any, name: str) -> None:
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference, ScatterChart, Series

    item = _mapping(raw, name)
    _strict_fields(
        item,
        {"sheet", "type", "data_range", "categories_range", "x_range", "y_ranges", "titles_from_data", "series_from", "title", "anchor"},
        name,
    )
    worksheet = _xlsx_sheet(workbook, item.get("sheet"), f"{name}.sheet")
    chart_type = item.get("type")
    if chart_type not in {"bar", "line", "pie", "scatter"}:
        raise OfficeInputError(f"{name}.type 不支持。", f"{name}.type is unsupported.")
    anchor = _xlsx_cell(item.get("anchor"), f"{name}.anchor")
    title = _optional_text(item.get("title"), f"{name}.title")

    if chart_type == "scatter":
        if {"data_range", "categories_range", "titles_from_data", "series_from"} & set(item):
            raise OfficeInputError(
                f"{name} 的 scatter 只接受 x_range 和 y_ranges。",
                f"{name} scatter accepts only x_range and y_ranges.",
            )
        x_range, x_bounds = _xlsx_range(item.get("x_range"), f"{name}.x_range")
        y_values = _sequence(item.get("y_ranges"), f"{name}.y_ranges")
        if not 1 <= len(y_values) <= MAX_CHART_SERIES:
            raise OfficeInputError(
                f"{name}.y_ranges 必须包含 1 到 {MAX_CHART_SERIES} 项。",
                f"{name}.y_ranges must contain 1 through {MAX_CHART_SERIES} items.",
            )
        if x_bounds[0] != x_bounds[2]:
            raise OfficeInputError(f"{name}.x_range 必须是单列。", f"{name}.x_range must be one column.")
        point_count = x_bounds[3] - x_bounds[1] + 1
        if point_count > MAX_CHART_POINTS:
            raise OfficeInputError(f"{name} 数据点过多。", f"{name} has too many points.")
        chart = ScatterChart()
        x_ref = Reference(worksheet, min_col=x_bounds[0], min_row=x_bounds[1], max_row=x_bounds[3])
        for series_index, raw_range in enumerate(y_values):
            _, bounds = _xlsx_range(raw_range, f"{name}.y_ranges[{series_index}]")
            if bounds[0] != bounds[2] or bounds[3] - bounds[1] + 1 != point_count:
                raise OfficeInputError(
                    f"{name}.y_ranges[{series_index}] 必须为与 x_range 等长的单列。",
                    f"{name}.y_ranges[{series_index}] must be one column matching x_range length.",
                )
            y_ref = Reference(worksheet, min_col=bounds[0], min_row=bounds[1], max_row=bounds[3])
            chart.series.append(Series(y_ref, x_ref, title=f"Series {series_index + 1}"))
    else:
        if {"x_range", "y_ranges"} & set(item):
            raise OfficeInputError(
                f"{name} 的 {chart_type} 不接受 x_range/y_ranges。",
                f"{name} {chart_type} does not accept x_range/y_ranges.",
            )
        _, bounds = _xlsx_range(item.get("data_range"), f"{name}.data_range")
        area = (bounds[2] - bounds[0] + 1) * (bounds[3] - bounds[1] + 1)
        if area > MAX_CHART_POINTS:
            raise OfficeInputError(f"{name}.data_range 过大。", f"{name}.data_range is too large.")
        series_from = item.get("series_from", "columns")
        if series_from not in {"columns", "rows"}:
            raise OfficeInputError(f"{name}.series_from 无效。", f"{name}.series_from is invalid.")
        series_count = (bounds[2] - bounds[0] + 1) if series_from == "columns" else (bounds[3] - bounds[1] + 1)
        if series_count > MAX_CHART_SERIES:
            raise OfficeInputError(f"{name} 数据系列过多。", f"{name} has too many data series.")
        if chart_type == "pie" and series_count != 1:
            raise OfficeInputError(
                f"{name} pie 必须恰好包含一个数据系列。",
                f"{name} pie must contain exactly one data series.",
            )
        chart = {"bar": BarChart, "line": LineChart, "pie": PieChart}[chart_type]()
        data = Reference(worksheet, min_col=bounds[0], min_row=bounds[1], max_col=bounds[2], max_row=bounds[3])
        titles = item.get("titles_from_data", True)
        if not isinstance(titles, bool):
                raise OfficeInputError(f"{name}.titles_from_data 必须是布尔值。", f"{name}.titles_from_data must be boolean.")
        point_count = (
            bounds[3] - bounds[1] + 1
            if series_from == "columns"
            else bounds[2] - bounds[0] + 1
        ) - int(titles)
        if point_count <= 0:
            raise OfficeInputError(
                f"{name}.data_range 没有可绘制的数据点。",
                f"{name}.data_range has no plottable data points.",
            )
        chart.add_data(data, titles_from_data=titles, from_rows=series_from == "rows")
        if item.get("categories_range") is not None:
            _, category_bounds = _xlsx_range(item.get("categories_range"), f"{name}.categories_range")
            if category_bounds[0] != category_bounds[2] and category_bounds[1] != category_bounds[3]:
                raise OfficeInputError(
                    f"{name}.categories_range 必须是单行或单列。",
                    f"{name}.categories_range must be one row or one column.",
                )
            category_count = max(
                category_bounds[2] - category_bounds[0] + 1,
                category_bounds[3] - category_bounds[1] + 1,
            )
            if category_count != point_count:
                raise OfficeInputError(
                    f"{name}.categories_range 必须与每个数据系列等长。",
                    f"{name}.categories_range must match each data series length.",
                )
            categories = Reference(
                worksheet,
                min_col=category_bounds[0], min_row=category_bounds[1],
                max_col=category_bounds[2], max_row=category_bounds[3],
            )
            chart.set_categories(categories)
    if title is not None:
        chart.title = title
    worksheet.add_chart(chart, anchor)


def _xlsx_chart_anchor(chart: Any) -> str:
    from openpyxl.utils import get_column_letter

    anchor = chart.anchor
    if isinstance(anchor, str):
        return anchor.upper()
    marker = getattr(anchor, "_from", None)
    if marker is None:
        return ""
    return f"{get_column_letter(marker.col + 1)}{marker.row + 1}"


def _xlsx_v2_snapshot(workbook: Any) -> dict[str, Any]:
    sheets: dict[str, Any] = {}
    for worksheet in workbook.worksheets:
        conditional: list[tuple[Any, ...]] = []
        for entry in worksheet.conditional_formatting:
            for rule in worksheet.conditional_formatting[entry]:
                conditional.append(
                    (
                        str(entry.sqref),
                        rule.type,
                        rule.operator,
                        tuple(rule.formula or ()),
                    )
                )
        validations = tuple(
            sorted(
                (
                    str(validation.sqref), validation.type, validation.operator,
                    str(validation.formula1), str(validation.formula2), bool(validation.allow_blank),
                )
                for validation in worksheet.data_validations.dataValidation
            )
        )
        freeze = worksheet.freeze_panes
        if freeze is not None and not isinstance(freeze, str):
            freeze = freeze.coordinate
        sheets[worksheet.title] = {
            "merged": tuple(sorted(str(value) for value in worksheet.merged_cells.ranges)),
            "row_heights": tuple(sorted((index, float(dimension.height)) for index, dimension in worksheet.row_dimensions.items() if dimension.height is not None)),
            "column_widths": tuple(sorted((index, float(dimension.width)) for index, dimension in worksheet.column_dimensions.items() if dimension.width is not None)),
            "freeze_panes": freeze,
            "auto_filter": worksheet.auto_filter.ref,
            "conditional_formats": tuple(sorted(conditional)),
            "data_validations": validations,
            "charts": tuple((chart.__class__.__name__, len(chart.series), _xlsx_chart_anchor(chart)) for chart in worksheet._charts),
        }
    names = tuple(sorted((defined.name, defined.attr_text) for defined in workbook.defined_names.values()))
    return {"sheets": sheets, "defined_names": names}


def _write_pptx(
    temporary: Path,
    target: Path,
    args: Mapping[str, Any],
    operation: str,
    workspace: Path,
    staged_workspace: Path,
) -> tuple[dict[str, Any], dict[str, Any]]:
    from pptx import Presentation

    presentation = Presentation(str(target)) if operation == "edit" else Presentation()
    raw_payload = args.get("presentation")
    payload = (
        {}
        if operation == "edit" and raw_payload is None
        else _mapping(raw_payload, "presentation")
    )
    if _office_v2_enabled():
        _strict_fields(payload, {"slides"}, "presentation")
    raw_slides = _sequence(payload.get("slides", []), "presentation.slides")
    replacements = _parse_replacements(args.get("replacements"))
    if len(raw_slides) > MAX_SLIDES:
        raise OfficeInputError(
            f"单次最多添加 {MAX_SLIDES} 张幻灯片。",
            f"At most {MAX_SLIDES} slides may be added in one call.",
        )
    if operation == "create" and not raw_slides:
        raise OfficeInputError(
            "PPTX 创建至少需要一张 presentation.slides。",
            "PPTX creation requires at least one presentation.slides item.",
        )
    if operation == "edit" and not (raw_slides or replacements):
        raise OfficeInputError(
            "PPTX 编辑至少需要一项变更。",
            "PPTX editing requires at least one change.",
        )
    if _office_v2_enabled():
        existing_chart_count = sum(
            int(getattr(shape, "has_chart", False))
            for slide in presentation.slides
            for shape in slide.shapes
        )
        declared_chart_count = 0
        for index, raw_slide in enumerate(raw_slides):
            slide = _mapping(raw_slide, f"presentation.slides[{index}]")
            declared_charts = _sequence(
                slide.get("charts", []),
                f"presentation.slides[{index}].charts",
            )
            declared_chart_count += len(declared_charts)
            if existing_chart_count + declared_chart_count > MAX_CHARTS_PER_FILE:
                raise OfficeInputError(
                    f"单个 PPTX 文件操作后最多包含 {MAX_CHARTS_PER_FILE} 个图表。",
                    (
                        "At most "
                        f"{MAX_CHARTS_PER_FILE} charts may exist per PPTX file "
                        "after this operation."
                    ),
                )

    replaced = _apply_replacements(
        list(_iter_pptx_paragraphs(presentation)), replacements, "PPTX"
    )
    text_boxes_added = 0
    tables_added = 0
    images_added = 0
    charts_added = 0
    shapes_added = 0
    notes_added = 0
    total_image_bytes = 0
    for index, raw_slide in enumerate(raw_slides):
        slide_data = _mapping(raw_slide, f"presentation.slides[{index}]")
        if _office_v2_enabled():
            _strict_fields(
                slide_data,
                {
                    "title", "subtitle", "bullets", "text_boxes", "tables", "images",
                    "layout_index", "layout_name", "title_style", "shapes", "charts",
                    "speaker_notes",
                },
                f"presentation.slides[{index}]",
            )
        title = _required_text(slide_data.get("title"), f"presentation.slides[{index}].title")
        subtitle = _optional_text(
            slide_data.get("subtitle"), f"presentation.slides[{index}].subtitle"
        )
        bullets = _parse_bullets(
            slide_data.get("bullets", []), f"presentation.slides[{index}].bullets"
        )
        text_boxes = _parse_pptx_text_boxes(
            slide_data.get("text_boxes", []),
            f"presentation.slides[{index}].text_boxes",
        )
        tables = _parse_pptx_tables(
            slide_data.get("tables", []), f"presentation.slides[{index}].tables"
        )
        images = _parse_pptx_images(
            slide_data.get("images", []),
            workspace,
            staged_workspace,
            f"presentation.slides[{index}].images",
        )
        images_added += len(images)
        if images_added > MAX_IMAGES_PER_FILE:
            raise OfficeInputError(
                f"单个 Office 文件最多添加 {MAX_IMAGES_PER_FILE} 张图片。",
                f"At most {MAX_IMAGES_PER_FILE} images may be added to one Office file.",
            )
        total_image_bytes += sum(len(image["data"]) for image in images)
        _validate_total_image_bytes(total_image_bytes)
        shapes = (
            _parse_pptx_shapes(
                slide_data.get("shapes", []),
                f"presentation.slides[{index}].shapes",
            )
            if _office_v2_enabled()
            else []
        )
        charts = (
            _parse_pptx_charts(
                slide_data.get("charts", []),
                f"presentation.slides[{index}].charts",
            )
            if _office_v2_enabled()
            else []
        )
        title_style = (
            _parse_pptx_text_style(
                slide_data.get("title_style"),
                f"presentation.slides[{index}].title_style",
            )
            if _office_v2_enabled() and slide_data.get("title_style") is not None
            else None
        )
        speaker_notes = (
            _optional_text(
                slide_data.get("speaker_notes"),
                f"presentation.slides[{index}].speaker_notes",
            )
            if _office_v2_enabled()
            else None
        )
        if subtitle is not None and bullets:
            raise OfficeInputError(
                "单张幻灯片不能同时提供 subtitle 和 bullets。",
                "A slide cannot provide subtitle and bullets at the same time.",
            )
        _add_slide(
            presentation,
            title,
            subtitle,
            bullets,
            text_boxes,
            tables,
            images,
            layout_index=slide_data.get("layout_index") if _office_v2_enabled() else None,
            layout_name=slide_data.get("layout_name") if _office_v2_enabled() else None,
            title_style=title_style,
            shapes=shapes,
            charts=charts,
            speaker_notes=speaker_notes,
        )
        text_boxes_added += len(text_boxes)
        tables_added += len(tables)
        shapes_added += len(shapes)
        charts_added += len(charts)
        notes_added += int(speaker_notes is not None)

    expected = {
        "semantic_digest": _semantic_digest(_pptx_semantic(presentation)),
        "shape_counts": _pptx_shape_counts(presentation),
    }
    if _office_v2_enabled():
        expected["v2_snapshot"] = _pptx_v2_snapshot(presentation)
    presentation.save(str(temporary))
    return (
        {
            "format": "pptx",
            "slides_added": len(raw_slides),
            "text_boxes_added": text_boxes_added,
            "tables_added": tables_added,
            "images_added": images_added,
            **(
                {
                    "charts_added": charts_added,
                    "shapes_added": shapes_added,
                    "speaker_notes_added": notes_added,
                }
                if _office_v2_enabled()
                else {}
            ),
            "replacements": replaced,
        },
        expected,
    )


def _reopen_and_verify(path: Path, suffix: str, expected: Mapping[str, Any]) -> None:
    """Reopen with the format library and compare a declarative semantic model."""

    try:
        if suffix == ".docx":
            from docx import Document

            reopened = Document(str(path))
            actual = _semantic_digest(_docx_semantic(reopened))
            if actual != expected["semantic_digest"]:
                raise ValueError("DOCX semantic verification mismatch")
            if len(reopened.inline_shapes) != expected["inline_shapes"]:
                raise ValueError("DOCX image verification mismatch")
            if _docx_page_break_count(reopened) != expected["page_breaks"]:
                raise ValueError("DOCX page-break verification mismatch")
            if (
                "v2_snapshot" in expected
                and _docx_v2_snapshot(reopened) != expected["v2_snapshot"]
            ):
                raise ValueError("DOCX v1.1 structure verification mismatch")
        elif suffix == ".pptx":
            from pptx import Presentation

            reopened = Presentation(str(path))
            actual = _semantic_digest(_pptx_semantic(reopened))
            if actual != expected["semantic_digest"]:
                raise ValueError("PPTX semantic verification mismatch")
            if _pptx_shape_counts(reopened) != expected["shape_counts"]:
                raise ValueError("PPTX shape verification mismatch")
            if (
                "v2_snapshot" in expected
                and _pptx_v2_snapshot(reopened) != expected["v2_snapshot"]
            ):
                raise ValueError("PPTX v1.1 structure verification mismatch")
        else:
            from openpyxl import load_workbook

            workbook = load_workbook(
                str(path),
                read_only=False,
                data_only=False,
                keep_vba=False,
                keep_links=False,
            )
            try:
                if tuple(workbook.sheetnames) != expected["sheet_names"]:
                    raise ValueError("XLSX sheet verification mismatch")
                for (sheet_name, coordinate), value in expected["written_cells"].items():
                    if workbook[sheet_name][coordinate].value != value:
                        raise ValueError("XLSX cell verification mismatch")
                for (sheet_name, coordinate), style in expected["written_styles"].items():
                    if _xlsx_style_snapshot(workbook[sheet_name][coordinate]) != style:
                        raise ValueError("XLSX style verification mismatch")
                if (
                    "v2_snapshot" in expected
                    and _xlsx_v2_snapshot(workbook) != expected["v2_snapshot"]
                ):
                    raise ValueError("XLSX v1.1 structure verification mismatch")
            finally:
                workbook.close()
    except OfficeInputError:
        raise
    except Exception as exc:
        raise OfficeInputError(
            "Office 文件重新打开校验失败，未替换原文件。",
            "Office reopen validation failed; the original file was not replaced.",
        ) from exc


def _drop_default_docx_custom_xml(document: Any) -> None:
    """Remove python-docx's bundled bibliography customXml from new files.

    Existing custom XML is rejected during the edit compatibility audit.  New
    documents must therefore not inherit the library template's unmanaged
    customXml part or they would become uneditable on their next operation.
    """

    for relationship_id, relationship in list(document.part.rels.items()):
        if relationship.reltype.rstrip("/") == _CUSTOM_XML_RELATIONSHIP_TYPE:
            document.part.drop_rel(relationship_id)


def _edit_part_is_allowed(name: str, suffix: str) -> bool:
    patterns = _ALLOWED_EDIT_PART_PATTERNS[suffix]
    if _office_v2_enabled():
        patterns += _V2_ALLOWED_EDIT_PART_PATTERNS[suffix]
    return any(pattern.fullmatch(name) for pattern in patterns)


def _v2_embedded_chart_workbook_allowed(name: str, suffix: str) -> bool:
    return bool(
        _office_v2_enabled()
        and suffix == ".pptx"
        and re.fullmatch(r"ppt/embeddings/Microsoft_Excel_(?:Work)?Sheet\d+\.xlsx", name)
    )


def _v2_chart_package_relationship_allowed(
    relationship_part: str,
    relationship: ElementTree.Element,
    suffix: str,
) -> bool:
    target = relationship.attrib.get("Target", "")
    return bool(
        _office_v2_enabled()
        and suffix == ".pptx"
        and re.fullmatch(
            r"ppt/charts/_rels/chart\d+\.xml\.rels", relationship_part
        )
        and re.fullmatch(
            r"\.\./embeddings/Microsoft_Excel_(?:Work)?Sheet\d+\.xlsx", target
        )
        and relationship.attrib.get("TargetMode", "Internal") == "Internal"
    )


def _inspect_chart_part(archive: zipfile.ZipFile, name: str) -> None:
    """Reject chart extensions, external data, combo charts, and unknown types."""

    try:
        root = ElementTree.fromstring(archive.read(name))
    except ElementTree.ParseError as exc:
        raise OfficeInputError(
            "Office 图表部件无法解析。",
            "An Office chart part could not be parsed.",
        ) from exc
    local_names = {element.tag.rsplit("}", 1)[-1] for element in root.iter()}
    allowed_namespaces = {
        "http://schemas.openxmlformats.org/drawingml/2006/chart",
        "http://schemas.openxmlformats.org/drawingml/2006/main",
        "http://schemas.openxmlformats.org/markup-compatibility/2006",
    }
    unknown_namespaces = {
        element.tag[1:].split("}", 1)[0]
        for element in root.iter()
        if element.tag.startswith("{")
        and element.tag[1:].split("}", 1)[0] not in allowed_namespaces
    }
    if unknown_namespaces:
        raise OfficeInputError(
            "Office 图表包含未知 XML 扩展命名空间。",
            "Office chart contains an unknown XML extension namespace.",
        )
    forbidden = {"ext", "extLst"}
    if not name.startswith("ppt/"):
        forbidden.add("externalData")
    if forbidden & local_names:
        raise OfficeInputError(
            "Office 图表包含未知扩展或外部数据，已拒绝。",
            "Office chart contains an unknown extension or external data and was refused.",
        )
    unsupported_elements = local_names - _V2_ALLOWED_CHART_ELEMENT_NAMES
    if unsupported_elements:
        preview = ", ".join(sorted(unsupported_elements)[:3])
        raise OfficeInputError(
            "Office 图表包含 v1.1 安全子集之外的 XML 元素。",
            (
                "Office chart contains an unsupported XML element outside the "
                f"v1.1 safe subset ({preview})."
            ),
        )
    chart_types = {
        name
        for name in local_names
        if name.endswith("Chart") and name not in {"chart", "ofPieChart"}
    }
    if not chart_types or not chart_types <= {"barChart", "lineChart", "pieChart", "scatterChart"}:
        raise OfficeInputError(
            "Office 图表类型不在 v1.1 安全子集中。",
            "Office chart type is outside the v1.1 safe subset.",
        )
    if len(chart_types) != 1:
        raise OfficeInputError(
            "Office v1.1 不支持组合图表。",
            "Office v1.1 does not support combo charts.",
        )
    for element in root.iter():
        if element.tag.rsplit("}", 1)[-1] == "f" and element.text:
            if (
                _EXTERNAL_FORMULA.search(element.text)
                or _EXCEL4_MACRO_FORMULA.search(element.text)
            ):
                raise OfficeInputError(
                    "Office 图表公式包含外部引用或 Excel 4.0 宏函数。",
                    (
                        "Office chart formula contains an external reference or "
                        "Excel 4.0 macro function."
                    ),
                )


def _inspect_xlsx_worksheet_merge_budget(
    archive: zipfile.ZipFile,
    name: str,
    *,
    aggregate_merged_cells: list[int],
) -> None:
    """Stream merge metadata before openpyxl can allocate MergedCell objects."""

    from openpyxl.utils.cell import range_boundaries

    try:
        with archive.open(name) as stream:
            for _, element in ElementTree.iterparse(stream, events=("end",)):
                if element.tag.rsplit("}", 1)[-1] == "mergeCell":
                    raw = element.attrib.get("ref", "").upper()
                    if not _CELL_RANGE.fullmatch(raw):
                        raise OfficeInputError(
                            "XLSX 包含无效的合并单元格区间。",
                            "XLSX contains an invalid merged-cell range.",
                        )
                    min_col, min_row, max_col, max_row = range_boundaries(raw)
                    if (
                        min_col > max_col
                        or min_row > max_row
                        or max_col > 16_384
                        or max_row > 1_048_576
                    ):
                        raise OfficeInputError(
                            "XLSX 合并单元格区间超出范围。",
                            "XLSX merged-cell range is outside spreadsheet bounds.",
                        )
                    area = _xlsx_merged_range_area(
                        (min_col, min_row, max_col, max_row)
                    )
                    if area > MAX_MERGED_CELLS_PER_RANGE:
                        raise OfficeInputError(
                            "XLSX 单个合并区间的单元格过多。",
                            "An XLSX merged range contains too many cells.",
                        )
                    aggregate_merged_cells[0] += area
                    if aggregate_merged_cells[0] > MAX_TOTAL_MERGED_CELLS:
                        raise OfficeInputError(
                            "XLSX 合并单元格总量过大。",
                            "XLSX exceeds the aggregate merged-cell budget.",
                        )
                element.clear()
    except OfficeInputError:
        raise
    except (OSError, ElementTree.ParseError) as exc:
        raise OfficeInputError(
            "XLSX 工作表部件无法安全解析。",
            "An XLSX worksheet part could not be safely parsed.",
        ) from exc


def _embedded_workbook_part_is_allowed(name: str) -> bool:
    return any(
        pattern.fullmatch(name)
        for pattern in _V2_EMBEDDED_WORKBOOK_PART_PATTERNS
    )


def _inspect_embedded_chart_workbook(
    archive: zipfile.ZipFile,
    name: str,
    *,
    aggregate_uncompressed: list[int] | None = None,
) -> None:
    """Audit one embedded chart workbook against one shared nested budget."""

    budget = aggregate_uncompressed if aggregate_uncompressed is not None else [0]
    try:
        with zipfile.ZipFile(io.BytesIO(archive.read(name))) as embedded:
            infos = embedded.infolist()
            if len(infos) > MAX_EMBEDDED_WORKBOOK_ENTRIES:
                raise OfficeInputError(
                    "PPTX 图表嵌入工作簿条目过多。",
                    "PPTX chart embedded workbook has too many entries.",
                )
            names: set[str] = set()
            folded_names: set[str] = set()
            for info in infos:
                pure = PurePosixPath(info.filename)
                folded = info.filename.casefold()
                if (
                    not info.filename
                    or info.filename.startswith("/")
                    or "\\" in info.filename
                    or ".." in pure.parts
                    or any(
                        marker in folded
                        for marker in (
                            "vbaproject", "externalLinks".casefold(), "connections",
                            "customxml", "embeddings", "activex", "controls",
                            "macrosheets", "intlmacrosheets", "dialogsheets",
                        )
                    )
                ):
                    raise OfficeInputError(
                        "PPTX 图表嵌入工作簿包含不安全部件。",
                        "PPTX chart embedded workbook contains an unsafe part.",
                    )
                if info.filename in names or folded in folded_names:
                    raise OfficeInputError(
                        "PPTX 图表嵌入工作簿包含重复部件。",
                        "PPTX chart embedded workbook contains duplicate parts.",
                    )
                names.add(info.filename)
                folded_names.add(folded)
                if not _embedded_workbook_part_is_allowed(info.filename):
                    raise OfficeInputError(
                        "PPTX 图表嵌入工作簿包含未知或不安全部件。",
                        (
                            "PPTX chart embedded workbook contains an unknown or "
                            "unsafe part."
                        ),
                    )
                if info.flag_bits & 0x1:
                    raise OfficeInputError(
                        "PPTX 图表嵌入工作簿不得加密。",
                        "PPTX chart embedded workbook must not be encrypted.",
                    )
                if info.compress_type not in {
                    zipfile.ZIP_STORED,
                    zipfile.ZIP_DEFLATED,
                }:
                    raise OfficeInputError(
                        "PPTX 图表嵌入工作簿压缩方式不受支持。",
                        (
                            "PPTX chart embedded workbook uses an unsupported "
                            "compression method."
                        ),
                    )
                if info.file_size > MAX_EMBEDDED_WORKBOOK_MEMBER_BYTES:
                    raise OfficeInputError(
                        "PPTX 图表嵌入工作簿部件过大。",
                        "PPTX chart embedded workbook part is too large.",
                    )
                if info.file_size and info.compress_size == 0:
                    raise OfficeInputError(
                        "PPTX 图表嵌入工作簿压缩比异常。",
                        (
                            "PPTX chart embedded workbook has a suspicious "
                            "compression ratio."
                        ),
                    )
                if (
                    info.compress_size
                    and info.file_size / info.compress_size
                    > MAX_EMBEDDED_WORKBOOK_COMPRESSION_RATIO
                ):
                    raise OfficeInputError(
                        "PPTX 图表嵌入工作簿压缩比过高。",
                        (
                            "PPTX chart embedded workbook compression ratio is "
                            "too high."
                        ),
                    )
                budget[0] += info.file_size
                if budget[0] > MAX_EMBEDDED_WORKBOOK_TOTAL_BYTES:
                    raise OfficeInputError(
                        "PPTX 图表嵌入工作簿解压后总量过大。",
                        (
                            "PPTX chart embedded workbooks exceed the aggregate "
                            "uncompressed-size budget."
                        ),
                    )

            required = {
                "[Content_Types].xml",
                "_rels/.rels",
                "xl/_rels/workbook.xml.rels",
                "xl/workbook.xml",
            }
            if not required <= names or not any(
                re.fullmatch(r"xl/worksheets/sheet\d+\.xml", item)
                for item in names
            ):
                raise OfficeInputError(
                    "PPTX 图表嵌入工作簿结构不完整。",
                    "PPTX chart embedded workbook structure is incomplete.",
                )

            for info in infos:
                if info.filename == "[Content_Types].xml":
                    content_types = embedded.read(info)
                    lowered_types = content_types.lower()
                    if any(
                        marker in lowered_types
                        for marker in (
                            b"macroenabled",
                            b"macrosheet",
                            b"dialogsheet",
                            b"vbaproject",
                            b"officedocument.oleobject",
                            b"ms-office.activex",
                        )
                    ):
                        raise OfficeInputError(
                            "PPTX 图表嵌入工作簿包含宏或嵌入对象。",
                            "PPTX chart embedded workbook contains a macro or embedded object.",
                        )
                    content_types_root = ElementTree.fromstring(content_types)
                    for element in content_types_root.iter():
                        local_name = element.tag.rsplit("}", 1)[-1]
                        if local_name not in {"Types", "Default", "Override"}:
                            raise OfficeInputError(
                                "PPTX 图表嵌入工作簿包含未知内容类型。",
                                (
                                    "PPTX chart embedded workbook contains an "
                                    "unknown content type."
                                ),
                            )
                        content_type = element.attrib.get("ContentType")
                        if (
                            content_type is not None
                            and content_type.casefold()
                            not in _V2_EMBEDDED_WORKBOOK_CONTENT_TYPES
                        ):
                            raise OfficeInputError(
                                "PPTX 图表嵌入工作簿包含未批准内容类型。",
                                (
                                    "PPTX chart embedded workbook contains an "
                                    "unapproved content type."
                                ),
                            )
                if (
                    re.fullmatch(r"xl/worksheets/sheet\d+\.xml", info.filename)
                    or info.filename == "xl/workbook.xml"
                ):
                    worksheet_root = ElementTree.fromstring(embedded.read(info))
                    if any(
                        element.tag.rsplit("}", 1)[-1] in {"f", "definedName"}
                        and element.text
                        and (
                            _EXTERNAL_FORMULA.search(element.text)
                            or _EXCEL4_MACRO_FORMULA.search(element.text)
                        )
                        for element in worksheet_root.iter()
                    ):
                        raise OfficeInputError(
                            "PPTX 图表嵌入工作簿包含外部公式或 Excel 4.0 宏公式。",
                            (
                                "PPTX chart embedded workbook contains an external "
                                "or Excel 4.0 macro formula."
                            ),
                        )
                if info.filename.endswith(".rels"):
                    rel_root = ElementTree.fromstring(embedded.read(info))
                    for relationship in rel_root.iter():
                        if not relationship.tag.endswith("Relationship"):
                            continue
                        relationship_type = relationship.attrib.get("Type", "").rstrip("/")
                        kind = relationship_type.rsplit("/", 1)[-1].casefold()
                        if (
                            relationship.attrib.get("TargetMode", "Internal")
                            != "Internal"
                            or kind in _UNSUPPORTED_RELATIONSHIP_KINDS
                            or kind
                            in {"dialogsheet", "intlmacrosheet", "xlmacrosheet"}
                            or relationship_type
                            not in _V2_EMBEDDED_WORKBOOK_RELATIONSHIP_TYPES
                        ):
                            raise OfficeInputError(
                                "PPTX 图表嵌入工作簿包含外链、宏或嵌入对象。",
                                (
                                    "PPTX chart embedded workbook contains an "
                                    "external link, macro, or embedded object."
                                ),
                            )
    except OfficeInputError:
        raise
    except (OSError, zipfile.BadZipFile, ElementTree.ParseError) as exc:
        raise OfficeInputError(
            "PPTX 图表嵌入工作簿无法安全解析。",
            "PPTX chart embedded workbook could not be safely parsed.",
        ) from exc


def _inspect_ooxml_archive(
    path: Path,
    suffix: str,
    *,
    audit_for_edit: bool = False,
) -> frozenset[str]:
    """Reject malformed, unsafe, lossy, or oversized OOXML inputs.

    Editing uses a deliberately small package-part and relationship allowlist.
    The format libraries are permitted to touch only structures whose
    round-trip behavior is covered by the Office contract; unfamiliar OOXML is
    rejected before it can be silently removed during save.
    """

    try:
        with zipfile.ZipFile(path) as archive:
            infos = archive.infolist()
            if len(infos) > MAX_ARCHIVE_ENTRIES:
                raise OfficeInputError(
                    "Office 压缩包条目过多。",
                    "Office archive contains too many entries.",
                )
            names: set[str] = set()
            folded_names: set[str] = set()
            total_bytes = 0
            for info in infos:
                name = info.filename
                pure = PurePosixPath(name)
                if (
                    not name
                    or name.startswith("/")
                    or "\\" in name
                    or ".." in pure.parts
                    or pure.is_absolute()
                ):
                    raise OfficeInputError(
                        "Office 压缩包包含不安全路径。",
                        "Office archive contains an unsafe path.",
                    )
                folded = name.casefold()
                if name in names or folded in folded_names:
                    raise OfficeInputError(
                        "Office 压缩包包含重复条目。",
                        "Office archive contains duplicate entries.",
                    )
                names.add(name)
                folded_names.add(folded)
                folded_parts = {part.casefold() for part in pure.parts}
                if audit_for_edit and "customxml" in folded_parts:
                    raise OfficeInputError(
                        "Office 文件包含本工具无法保证保真的 customXml 数据，已拒绝编辑。",
                        (
                            "Office file contains customXml data that this tool cannot "
                            "guarantee to preserve; editing was refused."
                        ),
                    )
                if (
                    folded_parts & _UNSUPPORTED_EMBEDDED_PATH_SEGMENTS
                    and not _v2_embedded_chart_workbook_allowed(name, suffix)
                ):
                    raise OfficeInputError(
                        "Office 文件包含本工具无法安全保留的嵌入对象或控件。",
                        (
                            "Office file contains an embedded object or control that "
                            "this tool cannot safely preserve."
                        ),
                    )
                if info.flag_bits & 0x1:
                    raise OfficeInputError(
                        "不支持加密的 Office 文件。",
                        "Encrypted Office files are not supported.",
                    )
                if info.compress_type not in {zipfile.ZIP_STORED, zipfile.ZIP_DEFLATED}:
                    raise OfficeInputError(
                        "Office 压缩包使用了不支持的压缩方式。",
                        "Office archive uses an unsupported compression method.",
                    )
                if info.file_size > MAX_ARCHIVE_MEMBER_BYTES:
                    raise OfficeInputError(
                        "Office 压缩包中的单个文件过大。",
                        "An Office archive member is too large.",
                    )
                total_bytes += info.file_size
                if total_bytes > MAX_ARCHIVE_TOTAL_BYTES:
                    raise OfficeInputError(
                        "Office 压缩包解压后过大。",
                        "Office archive is too large after decompression.",
                    )
                if info.file_size and info.compress_size == 0:
                    raise OfficeInputError(
                        "Office 压缩包的压缩比异常。",
                        "Office archive has a suspicious compression ratio.",
                    )
                if info.compress_size and info.file_size / info.compress_size > 500:
                    raise OfficeInputError(
                        "Office 压缩包的压缩比过高。",
                        "Office archive compression ratio is too high.",
                    )

            if "[Content_Types].xml" not in names or _REQUIRED_PARTS[suffix] not in names:
                raise OfficeInputError(
                    f"文件不是有效的 {suffix} OOXML 文档。",
                    f"File is not a valid {suffix} OOXML document.",
                )
            if any(name.casefold().endswith("vbaproject.bin") for name in names):
                raise OfficeInputError(
                    "检测到 Office 宏；本工具不处理宏。",
                    "Office macros were detected; this tool does not process macros.",
                )

            content_types_info = archive.getinfo("[Content_Types].xml")
            if content_types_info.file_size > MAX_RELATIONSHIP_BYTES:
                raise OfficeInputError(
                    "Office 内容类型清单过大。",
                    "The Office content-types manifest is too large.",
                )
            content_types = archive.read(content_types_info)
            lowered_types = content_types.lower()
            if (
                b"macroenabled" in lowered_types
                or b"vbaproject" in lowered_types
                or b"wordprocessingml.template" in lowered_types
                or b"spreadsheetml.template" in lowered_types
                or b"presentationml.template" in lowered_types
            ):
                raise OfficeInputError(
                    "检测到宏或 Office 模板内容类型。",
                    "A macro or Office template content type was detected.",
                )
            if any(
                marker in lowered_types
                for marker in _UNSUPPORTED_EMBEDDED_CONTENT_TYPE_MARKERS
            ):
                raise OfficeInputError(
                    "Office 内容类型清单中包含不支持的嵌入对象或控件。",
                    (
                        "The Office content-types manifest contains an unsupported "
                        "embedded object or control."
                    ),
                )

            if suffix == ".xlsx":
                merged_cell_budget = [0]
                for name in sorted(names):
                    if re.fullmatch(r"xl/worksheets/sheet\d+\.xml", name):
                        _inspect_xlsx_worksheet_merge_budget(
                            archive,
                            name,
                            aggregate_merged_cells=merged_cell_budget,
                        )

            if _office_v2_enabled():
                embedded_workbook_budget = [0]
                chart_parts = sorted(
                    name
                    for name in names
                    if re.fullmatch(r"(?:xl|ppt)/charts/chart\d+\.xml", name)
                )
                if len(chart_parts) > MAX_CHARTS_PER_FILE:
                    raise OfficeInputError(
                        f"单个 Office 文件最多包含 {MAX_CHARTS_PER_FILE} 个图表。",
                        (
                            f"An Office file may contain at most {MAX_CHARTS_PER_FILE} "
                            "chart parts."
                        ),
                    )
                for name in sorted(names):
                    if name in chart_parts:
                        _inspect_chart_part(archive, name)
                    if _v2_embedded_chart_workbook_allowed(name, suffix):
                        _inspect_embedded_chart_workbook(
                            archive,
                            name,
                            aggregate_uncompressed=embedded_workbook_budget,
                        )

            for info in infos:
                if not info.filename.casefold().endswith(".rels"):
                    continue
                if info.file_size > MAX_RELATIONSHIP_BYTES:
                    raise OfficeInputError(
                        "Office 关系文件过大。",
                        "An Office relationships file is too large.",
                    )
                try:
                    root = ElementTree.fromstring(archive.read(info))
                except ElementTree.ParseError as exc:
                    raise OfficeInputError(
                        "Office 关系文件无法解析。",
                        "An Office relationships file could not be parsed.",
                    ) from exc
                for relationship in root.iter():
                    if not relationship.tag.endswith("Relationship"):
                        continue
                    rel_type = relationship.attrib.get("Type", "").strip()
                    rel_kind = rel_type.rstrip("/").rsplit("/", 1)[-1].casefold()
                    if (
                        _office_v2_enabled()
                        and relationship.attrib.get("TargetMode") == "External"
                    ):
                        raise OfficeInputError(
                            "Office 文件包含外部关系。",
                            "Office file contains an external relationship.",
                        )
                    if (
                        rel_kind in _UNSUPPORTED_RELATIONSHIP_KINDS
                        and not _v2_chart_package_relationship_allowed(
                            info.filename, relationship, suffix
                        )
                    ):
                        raise OfficeInputError(
                            "Office 文件包含外部模板、外部工作簿或嵌入对象。",
                            (
                                "Office file contains an external template, external "
                                "workbook, macro, or embedded object."
                            ),
                        )
                    allowed_relationships = _ALLOWED_EDIT_RELATIONSHIP_TYPES[suffix]
                    if _office_v2_enabled():
                        allowed_relationships |= _V2_ALLOWED_RELATIONSHIP_TYPES[suffix]
                    if (
                        audit_for_edit
                        and rel_type.rstrip("/")
                        not in allowed_relationships
                    ):
                        raise OfficeInputError(
                            "Office 文件包含本工具无法保证保真的关系类型，已拒绝编辑。",
                            (
                                "Office file contains an unsupported relationship type "
                                "that this tool cannot guarantee to preserve; editing "
                                "was refused."
                            ),
                        )

            if audit_for_edit:
                unsupported_parts = sorted(
                    name for name in names if not _edit_part_is_allowed(name, suffix)
                )
                if unsupported_parts:
                    preview = ", ".join(unsupported_parts[:3])
                    raise OfficeInputError(
                        "Office 文件包含本工具无法保证保真的部件，已拒绝编辑。",
                        (
                            "Office file contains a package part that this tool cannot "
                            f"guarantee to preserve ({preview}); editing was refused."
                        ),
                    )
    except OfficeInputError:
        raise
    except (OSError, zipfile.BadZipFile) as exc:
        raise OfficeInputError(
            f"无法读取 {suffix} OOXML 文档。",
            f"Could not read the {suffix} OOXML document.",
        ) from exc
    return frozenset(names)


def _verify_edit_part_preservation(
    source_parts: frozenset[str],
    output_parts: frozenset[str],
    suffix: str,
    args: Mapping[str, Any],
) -> None:
    """Fail before installation if a supported input part disappeared."""

    missing = set(source_parts - output_parts)
    if suffix == ".xlsx":
        workbook = args.get("workbook")
        deleting_sheets = isinstance(workbook, Mapping) and bool(
            workbook.get("delete_sheets")
        )
        if deleting_sheets:
            # openpyxl may renumber worksheet part names after an explicit sheet
            # deletion.  Sheet names and requested values are verified by the
            # semantic reopen check, while every non-worksheet part remains
            # subject to exact preservation.
            missing = {
                name for name in missing if not _XLSX_WORKSHEET_PART.fullmatch(name)
            }
    if missing:
        preview = ", ".join(sorted(missing)[:3])
        raise OfficeInputError(
            "Office 编辑会删除原文件中的部件，已取消并保留原文件。",
            (
                "Office edit would remove existing package parts "
                f"({preview}); the edit was cancelled and the original was preserved."
            ),
        )


def _archive_part_digests(path: Path) -> dict[str, str]:
    with zipfile.ZipFile(path) as archive:
        return {
            info.filename: hashlib.sha256(archive.read(info)).hexdigest()
            for info in archive.infolist()
        }


def _xlsx_sheet_part_map(path: Path) -> dict[str, str]:
    relationships_namespace = "http://schemas.openxmlformats.org/package/2006/relationships"
    document_relationship_namespace = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    spreadsheet_namespace = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
    with zipfile.ZipFile(path) as archive:
        workbook = ElementTree.fromstring(archive.read("xl/workbook.xml"))
        relationships = ElementTree.fromstring(
            archive.read("xl/_rels/workbook.xml.rels")
        )
    targets = {
        relationship.attrib["Id"]: relationship.attrib.get("Target", "")
        for relationship in relationships.findall(
            f"{{{relationships_namespace}}}Relationship"
        )
    }
    result: dict[str, str] = {}
    for sheet in workbook.findall(f".//{{{spreadsheet_namespace}}}sheet"):
        rel_id = sheet.attrib.get(f"{{{document_relationship_namespace}}}id")
        target = targets.get(rel_id or "", "")
        if not target:
            continue
        normalized = str(PurePosixPath("xl") / target.lstrip("/"))
        if target.startswith("/xl/"):
            normalized = target.lstrip("/")
        result[sheet.attrib.get("name", "")] = normalized
    return result


def _requested_xlsx_sheet_mutations(args: Mapping[str, Any]) -> tuple[set[str], bool]:
    workbook = args.get("workbook")
    if not isinstance(workbook, Mapping):
        return set(), False
    if workbook.get("delete_sheets"):
        return set(), True
    names: set[str] = set()
    for raw in _safe_mapping_sequence(workbook.get("sheets")):
        if isinstance(raw.get("name"), str):
            names.add(raw["name"])
    for raw in _safe_mapping_sequence(workbook.get("cells")):
        if isinstance(raw.get("sheet"), str):
            names.add(raw["sheet"])
    for field in (
        "merged_cells", "row_heights", "column_widths", "freeze_panes", "auto_filters",
        "conditional_formats", "data_validations", "charts",
    ):
        for raw in _safe_mapping_sequence(workbook.get(field)):
            if isinstance(raw.get("sheet"), str):
                names.add(raw["sheet"])
    return names, False


def _verify_untouched_part_digests(
    source_path: Path,
    suffix: str,
    args: Mapping[str, Any],
    source: Mapping[str, str],
    output: Mapping[str, str],
) -> dict[str, Any]:
    """Prove selected non-target OOXML parts stayed byte-for-byte identical."""

    candidates = {
        name
        for name in source
        if "/media/" in name or "/printerSettings/" in name
    }
    if suffix == ".xlsx":
        # The declarative API can add charts but never edits an existing chart
        # part.  Treat every source chart as immutable so a library round trip
        # cannot silently discard an accepted chart element.
        candidates.update(
            name
            for name in source
            if re.fullmatch(r"xl/charts/chart\d+\.xml", name)
        )
        affected, renumbered = _requested_xlsx_sheet_mutations(args)
        if not renumbered:
            for sheet_name, part_name in _xlsx_sheet_part_map(source_path).items():
                if sheet_name not in affected:
                    candidates.add(part_name)
                    rels = str(
                        PurePosixPath(part_name).parent
                        / "_rels"
                        / f"{PurePosixPath(part_name).name}.rels"
                    )
                    if rels in source:
                        candidates.add(rels)
    elif suffix == ".pptx":
        candidates.update(
            name
            for name in source
            if re.fullmatch(
                r"ppt/(?:charts/(?:_rels/)?chart\d+\.xml(?:\.rels)?|"
                r"embeddings/Microsoft_Excel_(?:Work)?Sheet\d+\.xlsx)",
                name,
            )
        )
        if args.get("replacements") is None:
            candidates.update(
                name
                for name in source
                if re.fullmatch(
                    r"ppt/slides/(?:_rels/)?slide\d+\.xml(?:\.rels)?",
                    name,
                )
            )
    mismatched = sorted(
        name for name in candidates if output.get(name) != source.get(name)
    )
    if mismatched:
        preview = ", ".join(mismatched[:3])
        raise OfficeInputError(
            "Office 编辑改变了声明为未触及的 OOXML 部件，已取消。",
            (
                "Office edit changed OOXML parts declared untouched "
                f"({preview}); the edit was cancelled."
            ),
        )
    changed = sorted(
        name for name in source.keys() & output.keys() if source[name] != output[name]
    )
    added = sorted(output.keys() - source.keys())
    return {
        "package_parts_changed_count": len(changed),
        "package_parts_added_count": len(added),
        "package_parts_verified_unchanged_count": len(candidates),
        "package_parts_changed": changed[:100],
        "package_parts_added": added[:100],
        "package_parts_verified_unchanged": sorted(candidates)[:100],
        "package_part_lists_truncated": any(
            len(values) > 100 for values in (changed, added, candidates)
        ),
    }


def _mapping(value: Any, name: str) -> Mapping[str, Any]:
    if not isinstance(value, Mapping):
        raise OfficeInputError(f"{name} 必须是对象。", f"{name} must be an object.")
    return value


def _validate_request_budget(value: Mapping[str, Any]) -> None:
    """Bound aggregate declarative input before any directory or temp creation."""

    stack: list[Any] = list(value.values())
    visited: set[int] = set()
    text_chars = 0
    items = 0
    while stack:
        current = stack.pop()
        items += 1
        if items > MAX_DECLARATIVE_ITEMS:
            raise OfficeInputError(
                "Office 请求包含的声明式项目过多。",
                "The Office request contains too many declarative items.",
            )
        if isinstance(current, str):
            text_chars += len(current)
            if text_chars > MAX_TOTAL_TEXT_CHARS:
                raise OfficeInputError(
                    f"Office 请求的文本总量超过 {MAX_TOTAL_TEXT_CHARS} 个字符。",
                    (
                        "The Office request exceeds the aggregate "
                        f"{MAX_TOTAL_TEXT_CHARS}-character limit."
                    ),
                )
            continue
        if isinstance(current, Mapping):
            identity = id(current)
            if identity in visited:
                continue
            visited.add(identity)
            stack.extend(current.values())
        elif isinstance(current, Sequence) and not isinstance(current, (bytes, bytearray)):
            identity = id(current)
            if identity in visited:
                continue
            visited.add(identity)
            stack.extend(current)


def _sequence(value: Any, name: str) -> Sequence[Any]:
    if isinstance(value, (str, bytes)) or not isinstance(value, Sequence):
        raise OfficeInputError(f"{name} 必须是数组。", f"{name} must be an array.")
    return value


def _required_text(value: Any, name: str) -> str:
    if not isinstance(value, str) or not value.strip():
        raise OfficeInputError(
            f"{name} 必须是非空字符串。",
            f"{name} must be a non-empty string.",
        )
    return _bounded_text(value, name)


def _optional_text(value: Any, name: str) -> str | None:
    if value is None:
        return None
    if not isinstance(value, str):
        raise OfficeInputError(f"{name} 必须是字符串。", f"{name} must be a string.")
    return _bounded_text(value, name)


def _bounded_text(value: str, name: str) -> str:
    if len(value) > MAX_TEXT_CHARS:
        raise OfficeInputError(
            f"{name} 超过 {MAX_TEXT_CHARS} 个字符的限制。",
            f"{name} exceeds the {MAX_TEXT_CHARS}-character limit.",
        )
    if "\x00" in value:
        raise OfficeInputError(
            f"{name} 包含不允许的空字符。",
            f"{name} contains a forbidden null character.",
        )
    return value


def _scalar(value: Any, name: str) -> str | int | float | bool | None:
    if value is None or isinstance(value, (str, bool, int, float)):
        if isinstance(value, str):
            _bounded_text(value, name)
            if _is_formula(value) and _EXTERNAL_FORMULA.search(value):
                raise OfficeInputError(
                    f"{name} 公式包含外部工作簿或网络引用。",
                    f"{name} formula contains an external workbook or network reference.",
                )
        if isinstance(value, float) and not math.isfinite(value):
            raise OfficeInputError(
                f"{name} 必须是有限数值。",
                f"{name} must be a finite number.",
            )
        return value
    raise OfficeInputError(
        f"{name} 必须是字符串、数字、布尔值或 null。",
        f"{name} must be a string, number, boolean, or null.",
    )


def _bounded_number(
    value: Any,
    name: str,
    *,
    minimum: float,
    maximum: float,
) -> float:
    if isinstance(value, bool) or not isinstance(value, (int, float)):
        raise OfficeInputError(f"{name} 必须是数字。", f"{name} must be a number.")
    number = float(value)
    if not math.isfinite(number) or not minimum <= number <= maximum:
        raise OfficeInputError(
            f"{name} 必须介于 {minimum} 和 {maximum} 之间。",
            f"{name} must be between {minimum} and {maximum}.",
        )
    return number


def _read_local_image(
    value: Any,
    workspace: Path,
    staged_workspace: Path,
    name: str,
) -> dict[str, Any]:
    from PIL import Image, UnidentifiedImageError

    raw_path = _required_text(value, name)
    if "://" in raw_path:
        raise OfficeInputError(
            f"{name} 只允许工作区内的本地图片。",
            f"{name} accepts only a local image inside the workspace.",
        )
    try:
        logical_path = Path(resolve_and_validate(raw_path, str(workspace)))
    except WorkspaceViolation as exc:
        raise OfficeInputError(
            f"{name} 必须位于当前工作区内。",
            f"{name} must stay inside the current workspace.",
        ) from exc
    try:
        relative = logical_path.relative_to(workspace)
    except ValueError as exc:  # pragma: no cover - guarded by resolve_and_validate
        raise OfficeInputError(
            f"{name} 必须位于当前工作区内。",
            f"{name} must stay inside the current workspace.",
        ) from exc
    resolved = staged_workspace / relative
    if not resolved.is_file():
        raise OfficeInputError(
            f"找不到本地图片：{raw_path}",
            f"Local image was not found: {raw_path}",
        )
    if resolved.suffix.lower() not in _IMAGE_EXTENSIONS:
        raise OfficeInputError(
            f"{name} 使用了不支持的图片格式。",
            f"{name} uses an unsupported image format.",
        )
    try:
        size = resolved.stat().st_size
        if size > MAX_IMAGE_BYTES:
            raise OfficeInputError(
                f"单张 Office 图片不能超过 {MAX_IMAGE_BYTES // (1024 * 1024)} MiB。",
                (
                    "An Office image cannot exceed "
                    f"{MAX_IMAGE_BYTES // (1024 * 1024)} MiB."
                ),
            )
        data = resolved.read_bytes()
        if len(data) != size:
            raise OfficeInputError(
                "读取期间本地图片发生了变化。",
                "The local image changed while it was being read.",
            )
        with Image.open(io.BytesIO(data)) as image:
            width, height = image.size
            if width <= 0 or height <= 0 or width * height > MAX_IMAGE_PIXELS:
                raise OfficeInputError(
                    "Office 图片像素数超过安全限制。",
                    "The Office image exceeds the safe pixel limit.",
                )
            image.verify()
    except OfficeInputError:
        raise
    except (OSError, UnidentifiedImageError) as exc:
        raise OfficeInputError(
            f"无法验证本地图片：{raw_path}",
            f"Could not validate the local image: {raw_path}",
        ) from exc
    return {"path": resolved, "data": data}


def _validate_total_image_bytes(total_bytes: int) -> None:
    if total_bytes > MAX_TOTAL_IMAGE_BYTES:
        raise OfficeInputError(
            f"单个 Office 请求的图片总量不能超过 {MAX_TOTAL_IMAGE_BYTES // (1024 * 1024)} MiB。",
            (
                "Images in one Office request cannot exceed "
                f"{MAX_TOTAL_IMAGE_BYTES // (1024 * 1024)} MiB in total."
            ),
        )


def _row_values(value: Any, name: str) -> list[str | int | float | bool | None]:
    row = _sequence(value, name)
    if not row:
        raise OfficeInputError(f"{name} 不能为空。", f"{name} cannot be empty.")
    if len(row) > 16_384:
        raise OfficeInputError(
            f"{name} 超过 XLSX 最大列数。",
            f"{name} exceeds the XLSX column limit.",
        )
    return [_scalar(item, f"{name}[{index}]") for index, item in enumerate(row)]


def _parse_docx_paragraphs(value: Any) -> list[dict[str, Any]]:
    paragraphs = _sequence([] if value is None else value, "document.paragraphs")
    if len(paragraphs) > MAX_PARAGRAPHS:
        raise OfficeInputError(
            f"单次最多添加 {MAX_PARAGRAPHS} 个段落。",
            f"At most {MAX_PARAGRAPHS} paragraphs may be added in one call.",
        )
    parsed: list[dict[str, Any]] = []
    for index, raw in enumerate(paragraphs):
        item = _mapping(raw, f"document.paragraphs[{index}]")
        item_name = f"document.paragraphs[{index}]"
        if _office_v2_enabled():
            _strict_fields(
                item,
                {"text", "runs", "style", "page_break_after", "format", "list"},
                item_name,
            )
        has_text = "text" in item
        has_runs = _office_v2_enabled() and "runs" in item
        if has_text == has_runs:
            raise OfficeInputError(
                f"{item_name} 必须且只能提供 text 或 runs。",
                f"{item_name} requires exactly one of text or runs.",
            )
        text = _optional_text(item.get("text"), f"{item_name}.text") if has_text else None
        if has_text and text is None:
            raise OfficeInputError(
                f"{item_name}.text 不能缺失。",
                f"{item_name}.text is required.",
            )
        runs = _parse_docx_runs(item.get("runs"), f"{item_name}.runs") if has_runs else None
        style_name = item.get("style", "normal")
        if style_name not in _DOCX_STYLE_NAMES:
            raise OfficeInputError(
                f"不支持的 DOCX 段落样式：{style_name}",
                f"Unsupported DOCX paragraph style: {style_name}",
            )
        page_break_after = item.get("page_break_after", False)
        if not isinstance(page_break_after, bool):
            raise OfficeInputError(
                f"document.paragraphs[{index}].page_break_after 必须是布尔值。",
                f"document.paragraphs[{index}].page_break_after must be a boolean.",
            )
        paragraph_format = (
            _parse_docx_paragraph_format(item.get("format"), f"{item_name}.format")
            if _office_v2_enabled() and item.get("format") is not None
            else None
        )
        list_format = (
            _parse_docx_list(item.get("list"), f"{item_name}.list")
            if _office_v2_enabled() and item.get("list") is not None
            else None
        )
        if list_format is not None and style_name != "normal":
            raise OfficeInputError(
                f"{item_name} 的 list 不能与显式段落 style 组合。",
                f"{item_name} list cannot be combined with an explicit paragraph style.",
            )
        parsed.append(
            {
                "text": text,
                "runs": runs,
                "style": str(style_name),
                "page_break_after": page_break_after,
                "format": paragraph_format,
                "list": list_format,
            }
        )
    return parsed


def _parse_docx_runs(value: Any, name: str) -> list[dict[str, Any]]:
    runs = _sequence(value, name)
    if not runs or len(runs) > 1_000:
        raise OfficeInputError(
            f"{name} 必须包含 1 到 1000 个运行。",
            f"{name} must contain 1 through 1000 runs.",
        )
    parsed: list[dict[str, Any]] = []
    for index, raw in enumerate(runs):
        item_name = f"{name}[{index}]"
        item = _mapping(raw, item_name)
        _strict_fields(item, {"text", "bold", "italic", "underline", "color", "size", "font"}, item_name)
        run: dict[str, Any] = {"text": _optional_text(item.get("text"), f"{item_name}.text")}
        if run["text"] is None:
            raise OfficeInputError(f"{item_name}.text 必填。", f"{item_name}.text is required.")
        for field in ("bold", "italic", "underline"):
            if field in item:
                if not isinstance(item[field], bool):
                    raise OfficeInputError(f"{item_name}.{field} 必须是布尔值。", f"{item_name}.{field} must be boolean.")
                run[field] = item[field]
        if "color" in item:
            run["color"] = _validate_hex_color(item["color"], f"{item_name}.color")[-6:]
        if "size" in item:
            run["size"] = _bounded_number(item["size"], f"{item_name}.size", minimum=1, maximum=200)
        if "font" in item:
            font = _required_text(item["font"], f"{item_name}.font")
            if len(font) > 100:
                raise OfficeInputError(f"{item_name}.font 过长。", f"{item_name}.font is too long.")
            run["font"] = font
        parsed.append(run)
    return parsed


def _parse_docx_paragraph_format(value: Any, name: str) -> dict[str, Any]:
    item = _mapping(value, name)
    allowed = {
        "alignment", "keep_with_next", "keep_together", "page_break_before",
        "widow_control", "space_before_pt", "space_after_pt", "line_spacing",
    }
    _strict_fields(item, allowed, name)
    if not item:
        raise OfficeInputError(f"{name} 不能为空。", f"{name} cannot be empty.")
    parsed: dict[str, Any] = {}
    if "alignment" in item:
        if item["alignment"] not in {"left", "center", "right", "justify"}:
            raise OfficeInputError(f"{name}.alignment 无效。", f"{name}.alignment is invalid.")
        parsed["alignment"] = item["alignment"]
    for field in ("keep_with_next", "keep_together", "page_break_before", "widow_control"):
        if field in item:
            if not isinstance(item[field], bool):
                raise OfficeInputError(f"{name}.{field} 必须是布尔值。", f"{name}.{field} must be boolean.")
            parsed[field] = item[field]
    for field in ("space_before_pt", "space_after_pt"):
        if field in item:
            parsed[field] = _bounded_number(item[field], f"{name}.{field}", minimum=0, maximum=1_000)
    if "line_spacing" in item:
        parsed["line_spacing"] = _bounded_number(item["line_spacing"], f"{name}.line_spacing", minimum=0.5, maximum=10)
    return parsed


def _parse_docx_list(value: Any, name: str) -> dict[str, Any]:
    item = _mapping(value, name)
    _strict_fields(item, {"level", "ordered"}, name)
    if set(item) != {"level", "ordered"}:
        raise OfficeInputError(f"{name} 必须包含 level 和 ordered。", f"{name} requires level and ordered.")
    level = item["level"]
    if isinstance(level, bool) or not isinstance(level, int) or not 0 <= level <= 8:
        raise OfficeInputError(f"{name}.level 必须是 0 到 8 的整数。", f"{name}.level must be an integer from 0 through 8.")
    if not isinstance(item["ordered"], bool):
        raise OfficeInputError(f"{name}.ordered 必须是布尔值。", f"{name}.ordered must be boolean.")
    return {"level": level, "ordered": item["ordered"]}


def _add_docx_run(paragraph: Any, data: Mapping[str, Any]) -> None:
    from docx.shared import RGBColor
    from docx.shared import Pt

    run = paragraph.add_run(data["text"])
    for field in ("bold", "italic", "underline"):
        if field in data:
            setattr(run, field, data[field])
    if data.get("color"):
        run.font.color.rgb = RGBColor.from_string(data["color"])
    if data.get("size") is not None:
        run.font.size = Pt(data["size"])
    if data.get("font"):
        run.font.name = data["font"]


def _apply_docx_paragraph_format(paragraph: Any, data: Mapping[str, Any]) -> None:
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt

    paragraph_format = paragraph.paragraph_format
    if data.get("alignment"):
        paragraph.alignment = {
            "left": WD_ALIGN_PARAGRAPH.LEFT,
            "center": WD_ALIGN_PARAGRAPH.CENTER,
            "right": WD_ALIGN_PARAGRAPH.RIGHT,
            "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
        }[data["alignment"]]
    for field in ("keep_with_next", "keep_together", "page_break_before", "widow_control"):
        if field in data:
            setattr(paragraph_format, field, data[field])
    if "space_before_pt" in data:
        paragraph_format.space_before = Pt(data["space_before_pt"])
    if "space_after_pt" in data:
        paragraph_format.space_after = Pt(data["space_after_pt"])
    if "line_spacing" in data:
        paragraph_format.line_spacing = data["line_spacing"]


def _docx_numbering_id(document: Any, ordered: bool) -> int:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    cache = getattr(document.part, "_suxiaoyou_v2_numbering", None)
    if cache is None:
        cache = {}
        setattr(document.part, "_suxiaoyou_v2_numbering", cache)
    if ordered in cache:
        return cache[ordered]
    numbering = document.part.numbering_part.element
    abstract_ids = [int(element.get(qn("w:abstractNumId"))) for element in numbering.findall(qn("w:abstractNum"))]
    num_ids = [int(element.get(qn("w:numId"))) for element in numbering.findall(qn("w:num"))]
    abstract_id = max(abstract_ids, default=-1) + 1
    num_id = max(num_ids, default=0) + 1
    abstract = OxmlElement("w:abstractNum")
    abstract.set(qn("w:abstractNumId"), str(abstract_id))
    multi = OxmlElement("w:multiLevelType")
    multi.set(qn("w:val"), "multilevel")
    abstract.append(multi)
    bullet_chars = ("•", "◦", "▪")
    for level in range(9):
        lvl = OxmlElement("w:lvl")
        lvl.set(qn("w:ilvl"), str(level))
        start = OxmlElement("w:start")
        start.set(qn("w:val"), "1")
        lvl.append(start)
        num_fmt = OxmlElement("w:numFmt")
        num_fmt.set(qn("w:val"), "decimal" if ordered else "bullet")
        lvl.append(num_fmt)
        level_text = OxmlElement("w:lvlText")
        level_text.set(
            qn("w:val"),
            "".join(f"%{part}." for part in range(1, level + 2))
            if ordered
            else bullet_chars[level % len(bullet_chars)],
        )
        lvl.append(level_text)
        suffix = OxmlElement("w:suff")
        suffix.set(qn("w:val"), "tab")
        lvl.append(suffix)
        p_pr = OxmlElement("w:pPr")
        tabs = OxmlElement("w:tabs")
        tab = OxmlElement("w:tab")
        tab.set(qn("w:val"), "num")
        tab.set(qn("w:pos"), str(720 * (level + 1)))
        tabs.append(tab)
        p_pr.append(tabs)
        indent = OxmlElement("w:ind")
        indent.set(qn("w:left"), str(720 * (level + 1)))
        indent.set(qn("w:hanging"), "360")
        p_pr.append(indent)
        lvl.append(p_pr)
        abstract.append(lvl)
    first_num_index = next(
        (index for index, child in enumerate(numbering) if child.tag == qn("w:num")),
        len(numbering),
    )
    numbering.insert(first_num_index, abstract)
    num = OxmlElement("w:num")
    num.set(qn("w:numId"), str(num_id))
    abstract_ref = OxmlElement("w:abstractNumId")
    abstract_ref.set(qn("w:val"), str(abstract_id))
    num.append(abstract_ref)
    numbering.append(num)
    cache[ordered] = num_id
    return num_id


def _apply_docx_multilevel_list(document: Any, paragraph: Any, data: Mapping[str, Any]) -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    p_pr = paragraph._p.get_or_add_pPr()
    existing = p_pr.find(qn("w:numPr"))
    if existing is not None:
        p_pr.remove(existing)
    num_pr = OxmlElement("w:numPr")
    ilvl = OxmlElement("w:ilvl")
    ilvl.set(qn("w:val"), str(data["level"]))
    num_id = OxmlElement("w:numId")
    num_id.set(qn("w:val"), str(_docx_numbering_id(document, data["ordered"])))
    num_pr.extend((ilvl, num_id))
    p_pr.append(num_pr)


def _parse_docx_sections(value: Any) -> list[dict[str, Any]]:
    items = _sequence([] if value is None else value, "document.sections")
    if len(items) > 20:
        raise OfficeInputError("document.sections 最多 20 项。", "document.sections accepts at most 20 items.")
    parsed: list[dict[str, Any]] = []
    configured: set[int] = set()
    for index, raw in enumerate(items):
        name = f"document.sections[{index}]"
        item = _mapping(raw, name)
        _strict_fields(item, {"action", "index", "start", "orientation", "paper_size", "margins", "header", "footer"}, name)
        action = item.get("action")
        if action not in {"configure", "add"}:
            raise OfficeInputError(f"{name}.action 无效。", f"{name}.action is invalid.")
        section_index = item.get("index")
        if action == "configure":
            if isinstance(section_index, bool) or not isinstance(section_index, int) or section_index < 0:
                raise OfficeInputError(f"{name}.index 无效。", f"{name}.index is invalid.")
            if section_index in configured:
                raise OfficeInputError(f"{name}.index 重复。", f"{name}.index is duplicated.")
            configured.add(section_index)
        elif "index" in item:
            raise OfficeInputError(f"{name} add 不接受 index。", f"{name} add does not accept index.")
        start = item.get("start", "new_page")
        if start not in {"new_page", "continuous", "even_page", "odd_page"}:
            raise OfficeInputError(f"{name}.start 无效。", f"{name}.start is invalid.")
        orientation = item.get("orientation")
        if orientation is not None and orientation not in {"portrait", "landscape"}:
            raise OfficeInputError(f"{name}.orientation 无效。", f"{name}.orientation is invalid.")
        paper_size = item.get("paper_size")
        if paper_size is not None and paper_size not in {"a4", "letter", "legal"}:
            raise OfficeInputError(f"{name}.paper_size 无效。", f"{name}.paper_size is invalid.")
        margins: dict[str, float] | None = None
        if item.get("margins") is not None:
            raw_margins = _mapping(item["margins"], f"{name}.margins")
            margin_fields = {"top_inches", "bottom_inches", "left_inches", "right_inches"}
            _strict_fields(raw_margins, margin_fields, f"{name}.margins")
            if not raw_margins:
                raise OfficeInputError(f"{name}.margins 不能为空。", f"{name}.margins cannot be empty.")
            margins = {
                field: _bounded_number(raw_margins[field], f"{name}.margins.{field}", minimum=0, maximum=5)
                for field in raw_margins
            }
        if not ({"start", "orientation", "paper_size", "margins", "header", "footer"} & set(item)) and action == "configure":
            raise OfficeInputError(f"{name} 未提供变更。", f"{name} provides no change.")
        parsed.append(
            {
                "action": action,
                "index": section_index,
                "start": start,
                "has_start": "start" in item,
                "orientation": orientation,
                "paper_size": paper_size,
                "margins": margins,
                "header": _optional_text(item.get("header"), f"{name}.header") if "header" in item else None,
                "footer": _optional_text(item.get("footer"), f"{name}.footer") if "footer" in item else None,
                "has_header": "header" in item,
                "has_footer": "footer" in item,
            }
        )
    return parsed


def _apply_docx_sections(document: Any, sections: Sequence[Mapping[str, Any]]) -> None:
    from docx.enum.section import WD_ORIENT, WD_SECTION_START
    from docx.shared import Inches

    starts = {
        "new_page": WD_SECTION_START.NEW_PAGE,
        "continuous": WD_SECTION_START.CONTINUOUS,
        "even_page": WD_SECTION_START.EVEN_PAGE,
        "odd_page": WD_SECTION_START.ODD_PAGE,
    }
    papers = {
        "a4": (8.27, 11.69),
        "letter": (8.5, 11.0),
        "legal": (8.5, 14.0),
    }
    for data in sections:
        if data["action"] == "add":
            section = document.add_section(starts[data["start"]])
        else:
            if data["index"] >= len(document.sections):
                raise OfficeInputError(
                    f"document.sections index 超出现有节范围：{data['index']}",
                    f"document.sections index is outside the existing range: {data['index']}",
                )
            section = document.sections[data["index"]]
            if data["has_start"]:
                section.start_type = starts[data["start"]]
        width = float(section.page_width.inches)
        height = float(section.page_height.inches)
        if data["paper_size"]:
            width, height = papers[data["paper_size"]]
        orientation = data["orientation"]
        effective_landscape = orientation == "landscape" or (
            orientation is None and section.orientation == WD_ORIENT.LANDSCAPE
        )
        if effective_landscape:
            width, height = min(width, height), max(width, height)
            width, height = height, width
            if orientation == "landscape":
                section.orientation = WD_ORIENT.LANDSCAPE
        elif orientation == "portrait":
            width, height = min(width, height), max(width, height)
            section.orientation = WD_ORIENT.PORTRAIT
        section.page_width = Inches(width)
        section.page_height = Inches(height)
        for field, value in (data["margins"] or {}).items():
            setattr(section, field.removesuffix("_inches") + "_margin", Inches(value))
        if data["has_header"]:
            section.header.is_linked_to_previous = False
            section.header.paragraphs[0].text = data["header"] or ""
        if data["has_footer"]:
            section.footer.is_linked_to_previous = False
            section.footer.paragraphs[0].text = data["footer"] or ""


def _parse_docx_charts(
    value: Any,
    workspace: Path,
    staged_workspace: Path,
) -> list[dict[str, Any]]:
    items = _sequence([] if value is None else value, "document.charts")
    if len(items) > MAX_CHARTS_PER_FILE:
        raise OfficeInputError("document.charts 图表过多。", "document.charts contains too many charts.")
    parsed: list[dict[str, Any]] = []
    for index, raw in enumerate(items):
        name = f"document.charts[{index}]"
        item = _mapping(raw, name)
        _strict_fields(item, {"path", "width_inches", "alt_text", "source"}, name)
        image = _read_local_image(item.get("path"), workspace, staged_workspace, f"{name}.path")
        alt_text = _required_text(item.get("alt_text"), f"{name}.alt_text")
        source = _required_text(item.get("source"), f"{name}.source")
        if len(alt_text) > 1_024 or len(source) > 2_048:
            raise OfficeInputError(f"{name} 元数据过长。", f"{name} metadata is too long.")
        width = (
            _bounded_number(item["width_inches"], f"{name}.width_inches", minimum=0.1, maximum=50)
            if item.get("width_inches") is not None
            else None
        )
        parsed.append({**image, "width_inches": width, "alt_text": alt_text, "source": source})
    return parsed


def _parse_docx_images(
    value: Any,
    workspace: Path,
    staged_workspace: Path,
) -> list[dict[str, Any]]:
    images = _sequence([] if value is None else value, "document.images")
    if len(images) > MAX_IMAGES_PER_FILE:
        raise OfficeInputError(
            f"单个 Office 文件最多添加 {MAX_IMAGES_PER_FILE} 张图片。",
            f"At most {MAX_IMAGES_PER_FILE} images may be added to one Office file.",
        )
    parsed: list[dict[str, Any]] = []
    total_bytes = 0
    for index, raw in enumerate(images):
        item = _mapping(raw, f"document.images[{index}]")
        if _office_v2_enabled():
            _strict_fields(
                item,
                {"path", "width_inches", "caption"},
                f"document.images[{index}]",
            )
        image = _read_local_image(
            item.get("path"),
            workspace,
            staged_workspace,
            f"document.images[{index}].path",
        )
        total_bytes += len(image["data"])
        _validate_total_image_bytes(total_bytes)
        width = item.get("width_inches")
        if width is not None:
            width = _bounded_number(
                width,
                f"document.images[{index}].width_inches",
                minimum=0.1,
                maximum=50,
            )
        caption = _optional_text(
            item.get("caption"), f"document.images[{index}].caption"
        )
        parsed.append({**image, "width_inches": width, "caption": caption})
    return parsed


def _parse_tables(value: Any) -> list[dict[str, list[Any]]]:
    tables = _sequence([] if value is None else value, "document.tables")
    if len(tables) > MAX_TABLES:
        raise OfficeInputError(
            f"单次最多添加 {MAX_TABLES} 个表格。",
            f"At most {MAX_TABLES} tables may be added in one call.",
        )
    parsed: list[dict[str, list[Any]]] = []
    total_cells = 0
    for index, raw in enumerate(tables):
        table = _mapping(raw, f"document.tables[{index}]")
        item_name = f"document.tables[{index}]"
        if _office_v2_enabled():
            _strict_fields(table, {"headers", "rows", "merges", "format"}, item_name)
        headers = [
            _scalar(item, f"document.tables[{index}].headers[{cell_index}]")
            for cell_index, item in enumerate(
                _sequence(table.get("headers", []), f"document.tables[{index}].headers")
            )
        ]
        raw_rows = _sequence(table.get("rows"), f"document.tables[{index}].rows")
        rows = [
            _row_values(row, f"document.tables[{index}].rows[{row_index}]")
            for row_index, row in enumerate(raw_rows)
        ]
        if not headers and not rows:
            raise OfficeInputError(
                f"document.tables[{index}] 不能为空。",
                f"document.tables[{index}] cannot be empty.",
            )
        widest_row = max([len(headers), *(len(row) for row in rows)])
        if widest_row > MAX_TABLE_COLUMNS:
            raise OfficeInputError(
                f"DOCX 表格最多支持 {MAX_TABLE_COLUMNS} 列。",
                f"DOCX tables support at most {MAX_TABLE_COLUMNS} columns.",
            )
        total_cells += len(headers) + sum(len(row) for row in rows)
        if total_cells > MAX_TABLE_CELLS:
            raise OfficeInputError(
                f"单次最多添加 {MAX_TABLE_CELLS} 个 DOCX 表格单元格。",
                f"At most {MAX_TABLE_CELLS} DOCX table cells may be added in one call.",
            )
        row_count = len(rows) + (1 if headers else 0)
        merges = (
            _parse_docx_table_merges(
                table.get("merges", []),
                f"{item_name}.merges",
                row_count,
                widest_row,
            )
            if _office_v2_enabled()
            else []
        )
        table_format = (
            _parse_docx_table_format(table.get("format"), f"{item_name}.format")
            if _office_v2_enabled() and table.get("format") is not None
            else None
        )
        parsed.append(
            {"headers": headers, "rows": rows, "merges": merges, "format": table_format}
        )
    return parsed


def _parse_docx_table_merges(
    value: Any,
    name: str,
    rows: int,
    columns: int,
) -> list[dict[str, tuple[int, int]]]:
    items = _sequence(value, name)
    if len(items) > MAX_TABLE_CELLS:
        raise OfficeInputError(f"{name} 项数过多。", f"{name} contains too many items.")
    parsed: list[dict[str, tuple[int, int]]] = []
    rectangles: list[tuple[int, int, int, int]] = []
    for index, raw in enumerate(items):
        item_name = f"{name}[{index}]"
        item = _mapping(raw, item_name)
        _strict_fields(item, {"start", "end"}, item_name)
        endpoints: list[tuple[int, int]] = []
        for field in ("start", "end"):
            values = _sequence(item.get(field), f"{item_name}.{field}")
            if len(values) != 2 or any(isinstance(value, bool) or not isinstance(value, int) for value in values):
                raise OfficeInputError(
                    f"{item_name}.{field} 必须是 [row, column] 整数对。",
                    f"{item_name}.{field} must be an integer [row, column] pair.",
                )
            endpoints.append((values[0], values[1]))
        start, end = endpoints
        if (
            min(*start, *end) < 0
            or start[0] > end[0]
            or start[1] > end[1]
            or end[0] >= rows
            or end[1] >= columns
            or start == end
        ):
            raise OfficeInputError(f"{item_name} 范围无效。", f"{item_name} range is invalid.")
        rectangle = (start[1], start[0], end[1], end[0])
        if any(_rectangles_overlap(rectangle, prior) for prior in rectangles):
            raise OfficeInputError(f"{item_name} 与其他合并重叠。", f"{item_name} overlaps another merge.")
        rectangles.append(rectangle)
        parsed.append({"start": start, "end": end})
    return parsed


def _parse_docx_table_format(value: Any, name: str) -> dict[str, Any]:
    item = _mapping(value, name)
    _strict_fields(item, {"border_color", "border_size", "header_fill_color", "body_fill_color"}, name)
    if not item:
        raise OfficeInputError(f"{name} 不能为空。", f"{name} cannot be empty.")
    parsed: dict[str, Any] = {}
    for field in ("border_color", "header_fill_color", "body_fill_color"):
        if field in item:
            parsed[field] = _validate_hex_color(item[field], f"{name}.{field}")[-6:]
    if "border_size" in item:
        size = item["border_size"]
        if isinstance(size, bool) or not isinstance(size, int) or not 2 <= size <= 96:
            raise OfficeInputError(f"{name}.border_size 必须是 2 到 96 的整数。", f"{name}.border_size must be an integer from 2 through 96.")
        if "border_color" not in parsed:
            raise OfficeInputError(f"{name}.border_size 需要 border_color。", f"{name}.border_size requires border_color.")
        parsed["border_size"] = size
    if "border_color" in parsed and "border_size" not in parsed:
        parsed["border_size"] = 8
    return parsed


def _apply_docx_table_v2(table: Any, data: Mapping[str, Any]) -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    for merge in data.get("merges", []):
        table.cell(*merge["start"]).merge(table.cell(*merge["end"]))
    style = data.get("format")
    if not style:
        return
    for row_index, row in enumerate(table.rows):
        fill = style.get("header_fill_color") if data["headers"] and row_index == 0 else style.get("body_fill_color")
        for cell in row.cells:
            tc_pr = cell._tc.get_or_add_tcPr()
            if fill:
                shading = tc_pr.find(qn("w:shd"))
                if shading is None:
                    shading = OxmlElement("w:shd")
                    tc_pr.append(shading)
                shading.set(qn("w:val"), "clear")
                shading.set(qn("w:color"), "auto")
                shading.set(qn("w:fill"), fill)
            if style.get("border_color"):
                borders = tc_pr.find(qn("w:tcBorders"))
                if borders is None:
                    borders = OxmlElement("w:tcBorders")
                    tc_pr.append(borders)
                for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
                    tag = qn(f"w:{edge}")
                    border = borders.find(tag)
                    if border is None:
                        border = OxmlElement(f"w:{edge}")
                        borders.append(border)
                    border.set(qn("w:val"), "single")
                    border.set(qn("w:sz"), str(style["border_size"]))
                    border.set(qn("w:space"), "0")
                    border.set(qn("w:color"), style["border_color"])


def _parse_pptx_text_boxes(value: Any, name: str) -> list[dict[str, Any]]:
    items = _sequence(value, name)
    if len(items) > MAX_TEXT_BOXES_PER_SLIDE:
        raise OfficeInputError(
            f"单张幻灯片最多添加 {MAX_TEXT_BOXES_PER_SLIDE} 个文本框。",
            f"A slide supports at most {MAX_TEXT_BOXES_PER_SLIDE} text boxes.",
        )
    parsed: list[dict[str, Any]] = []
    for index, raw in enumerate(items):
        item_name = f"{name}[{index}]"
        item = _mapping(raw, item_name)
        if _office_v2_enabled():
            _strict_fields(
                item,
                {"text", "left_inches", "top_inches", "width_inches", "height_inches", "font_size", "style"},
                item_name,
            )
        text = _required_text(item.get("text"), f"{item_name}.text")
        parsed.append(
            {
                "text": text,
                "left_inches": _bounded_number(
                    item.get("left_inches"),
                    f"{item_name}.left_inches",
                    minimum=0,
                    maximum=100,
                ),
                "top_inches": _bounded_number(
                    item.get("top_inches"),
                    f"{item_name}.top_inches",
                    minimum=0,
                    maximum=100,
                ),
                "width_inches": _bounded_number(
                    item.get("width_inches"),
                    f"{item_name}.width_inches",
                    minimum=0.1,
                    maximum=100,
                ),
                "height_inches": _bounded_number(
                    item.get("height_inches"),
                    f"{item_name}.height_inches",
                    minimum=0.1,
                    maximum=100,
                ),
                "font_size": (
                    _bounded_number(
                        item["font_size"],
                        f"{item_name}.font_size",
                        minimum=1,
                        maximum=200,
                    )
                    if item.get("font_size") is not None
                    else None
                ),
                "style": (
                    _parse_pptx_text_style(item.get("style"), f"{item_name}.style")
                    if _office_v2_enabled() and item.get("style") is not None
                    else None
                ),
            }
        )
    return parsed


def _parse_pptx_tables(value: Any, name: str) -> list[dict[str, Any]]:
    items = _sequence(value, name)
    if len(items) > MAX_PPTX_TABLES_PER_SLIDE:
        raise OfficeInputError(
            f"单张幻灯片最多添加 {MAX_PPTX_TABLES_PER_SLIDE} 个表格。",
            f"A slide supports at most {MAX_PPTX_TABLES_PER_SLIDE} tables.",
        )
    parsed: list[dict[str, Any]] = []
    total_cells = 0
    for index, raw in enumerate(items):
        item_name = f"{name}[{index}]"
        item = _mapping(raw, item_name)
        if _office_v2_enabled():
            _strict_fields(
                item,
                {"left_inches", "top_inches", "width_inches", "height_inches", "headers", "rows", "style"},
                item_name,
            )
        headers = [
            _scalar(value, f"{item_name}.headers[{column}]")
            for column, value in enumerate(
                _sequence(item.get("headers", []), f"{item_name}.headers")
            )
        ]
        rows = [
            _row_values(row, f"{item_name}.rows[{row_index}]")
            for row_index, row in enumerate(
                _sequence(item.get("rows"), f"{item_name}.rows")
            )
        ]
        if not headers and not rows:
            raise OfficeInputError(f"{item_name} 不能为空。", f"{item_name} cannot be empty.")
        columns = max([len(headers), *(len(row) for row in rows)])
        if columns > MAX_PPTX_TABLE_COLUMNS:
            raise OfficeInputError(
                f"PPTX 表格最多支持 {MAX_PPTX_TABLE_COLUMNS} 列。",
                f"PPTX tables support at most {MAX_PPTX_TABLE_COLUMNS} columns.",
            )
        total_cells += columns * (len(rows) + (1 if headers else 0))
        if total_cells > MAX_PPTX_TABLE_CELLS_PER_SLIDE:
            raise OfficeInputError(
                f"单张幻灯片的表格最多包含 {MAX_PPTX_TABLE_CELLS_PER_SLIDE} 个单元格。",
                (
                    "Tables on one slide support at most "
                    f"{MAX_PPTX_TABLE_CELLS_PER_SLIDE} cells."
                ),
            )
        parsed.append(
            {
                "headers": headers,
                "rows": rows,
                "left_inches": _bounded_number(
                    item.get("left_inches"),
                    f"{item_name}.left_inches",
                    minimum=0,
                    maximum=100,
                ),
                "top_inches": _bounded_number(
                    item.get("top_inches"),
                    f"{item_name}.top_inches",
                    minimum=0,
                    maximum=100,
                ),
                "width_inches": _bounded_number(
                    item.get("width_inches"),
                    f"{item_name}.width_inches",
                    minimum=0.1,
                    maximum=100,
                ),
                "height_inches": _bounded_number(
                    item.get("height_inches"),
                    f"{item_name}.height_inches",
                    minimum=0.1,
                    maximum=100,
                ),
                "style": (
                    _parse_pptx_table_style(item.get("style"), f"{item_name}.style")
                    if _office_v2_enabled() and item.get("style") is not None
                    else None
                ),
            }
        )
    return parsed


def _parse_pptx_images(
    value: Any,
    workspace: Path,
    staged_workspace: Path,
    name: str,
) -> list[dict[str, Any]]:
    items = _sequence(value, name)
    if len(items) > MAX_IMAGES_PER_FILE:
        raise OfficeInputError(
            f"单个 Office 文件最多添加 {MAX_IMAGES_PER_FILE} 张图片。",
            f"At most {MAX_IMAGES_PER_FILE} images may be added to one Office file.",
        )
    parsed: list[dict[str, Any]] = []
    for index, raw in enumerate(items):
        item_name = f"{name}[{index}]"
        item = _mapping(raw, item_name)
        if _office_v2_enabled():
            _strict_fields(
                item,
                {
                    "path", "left_inches", "top_inches", "width_inches", "height_inches",
                    "crop_left", "crop_top", "crop_right", "crop_bottom", "align",
                },
                item_name,
            )
        image = _read_local_image(
            item.get("path"),
            workspace,
            staged_workspace,
            f"{item_name}.path",
        )
        width = (
            _bounded_number(
                item["width_inches"],
                f"{item_name}.width_inches",
                minimum=0.1,
                maximum=100,
            )
            if item.get("width_inches") is not None
            else None
        )
        height = (
            _bounded_number(
                item["height_inches"],
                f"{item_name}.height_inches",
                minimum=0.1,
                maximum=100,
            )
            if item.get("height_inches") is not None
            else None
        )
        parsed.append(
            {
                **image,
                "left_inches": _bounded_number(
                    item.get("left_inches"),
                    f"{item_name}.left_inches",
                    minimum=0,
                    maximum=100,
                ),
                "top_inches": _bounded_number(
                    item.get("top_inches"),
                    f"{item_name}.top_inches",
                    minimum=0,
                    maximum=100,
                ),
                "width_inches": width,
                "height_inches": height,
                **(
                    _parse_pptx_crop_and_alignment(item, item_name)
                    if _office_v2_enabled()
                    else {}
                ),
            }
        )
    return parsed


def _parse_pptx_text_style(value: Any, name: str) -> dict[str, Any]:
    style = _mapping(value, name)
    _strict_fields(style, {"font_size", "bold", "italic", "color", "font", "alignment"}, name)
    if not style:
        raise OfficeInputError(f"{name} 不能为空。", f"{name} cannot be empty.")
    parsed: dict[str, Any] = {}
    if "font_size" in style:
        parsed["font_size"] = _bounded_number(style["font_size"], f"{name}.font_size", minimum=1, maximum=200)
    for field in ("bold", "italic"):
        if field in style:
            if not isinstance(style[field], bool):
                raise OfficeInputError(f"{name}.{field} 必须是布尔值。", f"{name}.{field} must be boolean.")
            parsed[field] = style[field]
    if "color" in style:
        parsed["color"] = _validate_hex_color(style["color"], f"{name}.color")[-6:]
    if "font" in style:
        font = _required_text(style["font"], f"{name}.font")
        if len(font) > 100:
            raise OfficeInputError(f"{name}.font 过长。", f"{name}.font is too long.")
        parsed["font"] = font
    if "alignment" in style:
        if style["alignment"] not in {"left", "center", "right"}:
            raise OfficeInputError(f"{name}.alignment 无效。", f"{name}.alignment is invalid.")
        parsed["alignment"] = style["alignment"]
    return parsed


def _parse_pptx_table_style(value: Any, name: str) -> dict[str, Any]:
    style = _mapping(value, name)
    _strict_fields(style, {"header_fill_color", "body_fill_color", "font_size"}, name)
    if not style:
        raise OfficeInputError(f"{name} 不能为空。", f"{name} cannot be empty.")
    parsed: dict[str, Any] = {}
    for field in ("header_fill_color", "body_fill_color"):
        if field in style:
            parsed[field] = _validate_hex_color(style[field], f"{name}.{field}")[-6:]
    if "font_size" in style:
        parsed["font_size"] = _bounded_number(style["font_size"], f"{name}.font_size", minimum=1, maximum=200)
    return parsed


def _parse_pptx_crop_and_alignment(item: Mapping[str, Any], name: str) -> dict[str, Any]:
    result: dict[str, Any] = {}
    total_horizontal = 0.0
    total_vertical = 0.0
    for field in ("crop_left", "crop_top", "crop_right", "crop_bottom"):
        value = item.get(field, 0.0)
        parsed = _bounded_number(value, f"{name}.{field}", minimum=0, maximum=0.95)
        result[field] = parsed
        if field in {"crop_left", "crop_right"}:
            total_horizontal += parsed
        else:
            total_vertical += parsed
    if total_horizontal >= 1 or total_vertical >= 1:
        raise OfficeInputError(
            f"{name} 裁剪比例不能完全移除图片。",
            f"{name} crop ratios cannot remove the entire picture.",
        )
    align = item.get("align", "none")
    if align not in {"none", "center", "right"}:
        raise OfficeInputError(f"{name}.align 无效。", f"{name}.align is invalid.")
    result["align"] = align
    return result


def _parse_pptx_shapes(value: Any, name: str) -> list[dict[str, Any]]:
    items = _sequence(value, name)
    if len(items) > MAX_TEXT_BOXES_PER_SLIDE:
        raise OfficeInputError(f"{name} 项数过多。", f"{name} contains too many items.")
    parsed: list[dict[str, Any]] = []
    allowed_types = {"rectangle", "rounded_rectangle", "ellipse", "line"}
    for index, raw in enumerate(items):
        item_name = f"{name}[{index}]"
        item = _mapping(raw, item_name)
        _strict_fields(
            item,
            {"type", "left_inches", "top_inches", "width_inches", "height_inches", "text", "fill_color", "line_color", "text_style"},
            item_name,
        )
        shape_type = item.get("type")
        if shape_type not in allowed_types:
            raise OfficeInputError(f"{item_name}.type 不支持。", f"{item_name}.type is unsupported.")
        if shape_type == "line" and (item.get("text") is not None or item.get("fill_color") is not None):
            raise OfficeInputError(
                f"{item_name} 的 line 不接受 text/fill_color。",
                f"{item_name} line does not accept text/fill_color.",
            )
        parsed.append(
            {
                "type": shape_type,
                "left_inches": _bounded_number(item.get("left_inches"), f"{item_name}.left_inches", minimum=0, maximum=100),
                "top_inches": _bounded_number(item.get("top_inches"), f"{item_name}.top_inches", minimum=0, maximum=100),
                "width_inches": _bounded_number(item.get("width_inches"), f"{item_name}.width_inches", minimum=0.01, maximum=100),
                "height_inches": _bounded_number(item.get("height_inches"), f"{item_name}.height_inches", minimum=0.01, maximum=100),
                "text": _optional_text(item.get("text"), f"{item_name}.text"),
                "fill_color": (_validate_hex_color(item["fill_color"], f"{item_name}.fill_color")[-6:] if item.get("fill_color") is not None else None),
                "line_color": (_validate_hex_color(item["line_color"], f"{item_name}.line_color")[-6:] if item.get("line_color") is not None else None),
                "text_style": (
                    _parse_pptx_text_style(item.get("text_style"), f"{item_name}.text_style")
                    if item.get("text_style") is not None
                    else None
                ),
            }
        )
    return parsed


def _finite_chart_number(value: Any, name: str) -> float:
    if isinstance(value, bool) or not isinstance(value, (int, float)):
        raise OfficeInputError(f"{name} 必须是数字。", f"{name} must be numeric.")
    result = float(value)
    if not math.isfinite(result):
        raise OfficeInputError(f"{name} 必须是有限数。", f"{name} must be finite.")
    return result


def _parse_pptx_charts(value: Any, name: str) -> list[dict[str, Any]]:
    items = _sequence(value, name)
    if len(items) > MAX_CHARTS_PER_FILE:
        raise OfficeInputError(f"{name} 图表过多。", f"{name} contains too many charts.")
    parsed: list[dict[str, Any]] = []
    total_points = 0
    for index, raw in enumerate(items):
        item_name = f"{name}[{index}]"
        item = _mapping(raw, item_name)
        _strict_fields(
            item,
            {"type", "categories", "series", "title", "left_inches", "top_inches", "width_inches", "height_inches"},
            item_name,
        )
        chart_type = item.get("type")
        if chart_type not in {"bar", "line", "pie", "scatter"}:
            raise OfficeInputError(f"{item_name}.type 不支持。", f"{item_name}.type is unsupported.")
        raw_series = _sequence(item.get("series"), f"{item_name}.series")
        if not 1 <= len(raw_series) <= MAX_CHART_SERIES or (chart_type == "pie" and len(raw_series) != 1):
            raise OfficeInputError(
                f"{item_name}.series 数量无效。",
                f"{item_name}.series count is invalid.",
            )
        categories: list[str | int | float | bool] = []
        if chart_type == "scatter":
            if "categories" in item:
                raise OfficeInputError(f"{item_name} scatter 不接受 categories。", f"{item_name} scatter does not accept categories.")
        else:
            raw_categories = _sequence(item.get("categories"), f"{item_name}.categories")
            if not raw_categories:
                raise OfficeInputError(f"{item_name}.categories 不能为空。", f"{item_name}.categories cannot be empty.")
            for point, category in enumerate(raw_categories):
                scalar = _scalar(category, f"{item_name}.categories[{point}]")
                if scalar is None:
                    raise OfficeInputError(f"{item_name}.categories 不接受 null。", f"{item_name}.categories does not accept null.")
                categories.append(scalar)
        series: list[dict[str, Any]] = []
        for series_index, raw_item in enumerate(raw_series):
            series_name = f"{item_name}.series[{series_index}]"
            series_item = _mapping(raw_item, series_name)
            _strict_fields(series_item, {"name", "values", "x_values", "y_values"}, series_name)
            label = _required_text(series_item.get("name"), f"{series_name}.name")
            if chart_type == "scatter":
                if "values" in series_item:
                    raise OfficeInputError(f"{series_name} scatter 不接受 values。", f"{series_name} scatter does not accept values.")
                raw_x = _sequence(series_item.get("x_values"), f"{series_name}.x_values")
                raw_y = _sequence(series_item.get("y_values"), f"{series_name}.y_values")
                if not raw_x or len(raw_x) != len(raw_y):
                    raise OfficeInputError(
                        f"{series_name} x_values/y_values 必须非空且等长。",
                        f"{series_name} x_values/y_values must be non-empty and equal length.",
                    )
                x_values = [_finite_chart_number(point, f"{series_name}.x_values[{i}]") for i, point in enumerate(raw_x)]
                y_values = [_finite_chart_number(point, f"{series_name}.y_values[{i}]") for i, point in enumerate(raw_y)]
                series.append({"name": label, "x_values": x_values, "y_values": y_values})
                total_points += len(x_values)
            else:
                if {"x_values", "y_values"} & set(series_item):
                    raise OfficeInputError(f"{series_name} 不接受 x_values/y_values。", f"{series_name} does not accept x_values/y_values.")
                raw_values = _sequence(series_item.get("values"), f"{series_name}.values")
                if len(raw_values) != len(categories):
                    raise OfficeInputError(
                        f"{series_name}.values 必须与 categories 等长。",
                        f"{series_name}.values must match categories length.",
                    )
                values = [_finite_chart_number(point, f"{series_name}.values[{i}]") for i, point in enumerate(raw_values)]
                series.append({"name": label, "values": values})
                total_points += len(values)
        if total_points > MAX_CHART_POINTS:
            raise OfficeInputError(f"{name} 数据点过多。", f"{name} has too many data points.")
        parsed.append(
            {
                "type": chart_type,
                "categories": categories,
                "series": series,
                "title": _optional_text(item.get("title"), f"{item_name}.title"),
                "left_inches": _bounded_number(item.get("left_inches"), f"{item_name}.left_inches", minimum=0, maximum=100),
                "top_inches": _bounded_number(item.get("top_inches"), f"{item_name}.top_inches", minimum=0, maximum=100),
                "width_inches": _bounded_number(item.get("width_inches"), f"{item_name}.width_inches", minimum=0.1, maximum=100),
                "height_inches": _bounded_number(item.get("height_inches"), f"{item_name}.height_inches", minimum=0.1, maximum=100),
            }
        )
    return parsed


def _parse_bullets(value: Any, name: str) -> list[dict[str, Any]]:
    bullets = _sequence(value, name)
    if len(bullets) > MAX_BULLETS_PER_SLIDE:
        raise OfficeInputError(
            f"单张幻灯片最多支持 {MAX_BULLETS_PER_SLIDE} 个项目符号。",
            f"A slide supports at most {MAX_BULLETS_PER_SLIDE} bullets.",
        )
    result: list[dict[str, Any]] = []
    for index, raw in enumerate(bullets):
        if isinstance(raw, str):
            result.append({"text": _bounded_text(raw, f"{name}[{index}]"), "level": 0})
            continue
        item = _mapping(raw, f"{name}[{index}]")
        if _office_v2_enabled():
            _strict_fields(item, {"text", "level"}, f"{name}[{index}]")
        text = _required_text(item.get("text"), f"{name}[{index}].text")
        level = item.get("level", 0)
        if isinstance(level, bool) or not isinstance(level, int) or not 0 <= level <= 4:
            raise OfficeInputError(
                f"{name}[{index}].level 必须是 0 到 4 的整数。",
                f"{name}[{index}].level must be an integer from 0 through 4.",
            )
        result.append({"text": text, "level": level})
    return result


def _parse_replacements(value: Any) -> list[dict[str, Any]]:
    if value is None:
        return []
    replacements = _sequence(value, "replacements")
    if len(replacements) > MAX_REPLACEMENTS:
        raise OfficeInputError(
            f"单次最多执行 {MAX_REPLACEMENTS} 项文本替换。",
            f"At most {MAX_REPLACEMENTS} replacements may be made in one call.",
        )
    parsed: list[dict[str, Any]] = []
    for index, raw in enumerate(replacements):
        item = _mapping(raw, f"replacements[{index}]")
        if _office_v2_enabled():
            _strict_fields(
                item,
                {"old_text", "new_text", "replace_all"},
                f"replacements[{index}]",
            )
        old_text = _required_text(item.get("old_text"), f"replacements[{index}].old_text")
        new_text = _optional_text(item.get("new_text"), f"replacements[{index}].new_text")
        if new_text is None:
            raise OfficeInputError(
                f"replacements[{index}].new_text 不能缺失。",
                f"replacements[{index}].new_text is required.",
            )
        if old_text == new_text:
            raise OfficeInputError(
                f"replacements[{index}] 的新旧文本相同。",
                f"replacements[{index}] old_text and new_text are identical.",
            )
        replace_all = item.get("replace_all", False)
        if not isinstance(replace_all, bool):
            raise OfficeInputError(
                f"replacements[{index}].replace_all 必须是布尔值。",
                f"replacements[{index}].replace_all must be a boolean.",
            )
        parsed.append(
            {"old_text": old_text, "new_text": new_text, "replace_all": replace_all}
        )
    return parsed


def _apply_replacements(
    paragraphs: list[Any],
    replacements: Sequence[Mapping[str, Any]],
    format_name: str,
) -> int:
    total = 0
    for index, replacement in enumerate(replacements):
        old_text = replacement["old_text"]
        occurrences = sum(_paragraph_run_text(paragraph).count(old_text) for paragraph in paragraphs)
        if occurrences == 0:
            raise OfficeInputError(
                f"{format_name} 中找不到 replacements[{index}].old_text。",
                f"replacements[{index}].old_text was not found in the {format_name} file.",
            )
        if occurrences > 1 and not replacement["replace_all"]:
            raise OfficeInputError(
                f"{format_name} 中找到 {occurrences} 处匹配；请提供更唯一的文本或设置 replace_all=true。",
                (
                    f"Found {occurrences} matches in the {format_name} file; provide "
                    "more unique text or set replace_all=true."
                ),
            )
        remaining = occurrences if replacement["replace_all"] else 1
        for paragraph in paragraphs:
            if remaining <= 0:
                break
            replaced = _replace_in_runs(
                paragraph,
                old_text,
                replacement["new_text"],
                limit=remaining,
            )
            remaining -= replaced
            total += replaced
    return total


def _replace_in_runs(paragraph: Any, old: str, new: str, *, limit: int) -> int:
    runs = list(paragraph.runs)
    text = "".join(run.text or "" for run in runs)
    positions: list[int] = []
    cursor = 0
    while len(positions) < limit:
        position = text.find(old, cursor)
        if position < 0:
            break
        positions.append(position)
        cursor = position + len(old)
    if not positions:
        return 0

    for start in reversed(positions):
        end = start + len(old)
        offsets: list[tuple[int, int]] = []
        offset = 0
        for run_index, run in enumerate(runs):
            next_offset = offset + len(run.text or "")
            offsets.append((offset, next_offset))
            offset = next_offset
        start_run = next(
            index
            for index, (run_start, run_end) in enumerate(offsets)
            if run_start <= start < run_end
        )
        end_run = next(
            index
            for index, (run_start, run_end) in enumerate(offsets)
            if run_start < end <= run_end
        )
        start_offset = start - offsets[start_run][0]
        end_offset = end - offsets[end_run][0]
        if start_run == end_run:
            run_text = runs[start_run].text or ""
            runs[start_run].text = run_text[:start_offset] + new + run_text[end_offset:]
        else:
            start_text = runs[start_run].text or ""
            end_text = runs[end_run].text or ""
            runs[start_run].text = start_text[:start_offset] + new
            for run_index in range(start_run + 1, end_run):
                runs[run_index].text = ""
            runs[end_run].text = end_text[end_offset:]
        text = text[:start] + new + text[end:]
    return len(positions)


def _iter_docx_paragraphs(document: Any) -> Iterator[Any]:
    yield from document.paragraphs
    for table in document.tables:
        yield from _iter_docx_table_paragraphs(table)


def _iter_docx_table_paragraphs(table: Any) -> Iterator[Any]:
    for row in table.rows:
        for cell in row.cells:
            yield from cell.paragraphs
            for nested in cell.tables:
                yield from _iter_docx_table_paragraphs(nested)


def _iter_pptx_paragraphs(presentation: Any) -> Iterator[Any]:
    for slide in presentation.slides:
        for shape in slide.shapes:
            yield from _iter_pptx_shape_paragraphs(shape)


def _iter_pptx_shape_paragraphs(shape: Any) -> Iterator[Any]:
    if getattr(shape, "has_text_frame", False):
        yield from shape.text_frame.paragraphs
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                yield from cell.text_frame.paragraphs
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            yield from _iter_pptx_shape_paragraphs(child)


def _paragraph_run_text(paragraph: Any) -> str:
    return "".join(run.text or "" for run in paragraph.runs)


def _docx_semantic(document: Any) -> Iterable[str]:
    for paragraph in _iter_docx_paragraphs(document):
        yield _paragraph_run_text(paragraph)


def _docx_page_break_count(document: Any) -> int:
    from docx.oxml.ns import qn

    page_type = qn("w:type")
    return sum(
        1
        for element in document.element.iter(qn("w:br"))
        if element.get(page_type) == "page"
    )


def _docx_length_snapshot(value: Any) -> int | float | None:
    if value is None:
        return None
    if isinstance(value, float):
        return round(value, 6)
    return int(value)


def _docx_v2_snapshot(document: Any) -> dict[str, Any]:
    from docx.oxml.ns import qn

    sections = tuple(
        (
            int(section.start_type),
            int(section.orientation),
            int(section.page_width),
            int(section.page_height),
            int(section.top_margin),
            int(section.bottom_margin),
            int(section.left_margin),
            int(section.right_margin),
            "\n".join(paragraph.text for paragraph in section.header.paragraphs),
            "\n".join(paragraph.text for paragraph in section.footer.paragraphs),
            bool(section.header.is_linked_to_previous),
            bool(section.footer.is_linked_to_previous),
        )
        for section in document.sections
    )
    paragraphs: list[tuple[Any, ...]] = []
    for paragraph in _iter_docx_paragraphs(document):
        p_pr = paragraph._p.pPr
        num_pr = p_pr.find(qn("w:numPr")) if p_pr is not None else None
        ilvl = num_pr.find(qn("w:ilvl")) if num_pr is not None else None
        num_id = num_pr.find(qn("w:numId")) if num_pr is not None else None
        paragraph_format = paragraph.paragraph_format
        runs = tuple(
            (
                run.text,
                run.bold,
                run.italic,
                run.underline,
                str(run.font.color.rgb) if run.font.color.rgb is not None else None,
                int(run.font.size) if run.font.size is not None else None,
                run.font.name,
            )
            for run in paragraph.runs
        )
        paragraphs.append(
            (
                paragraph.style.name if paragraph.style is not None else None,
                int(paragraph.alignment) if paragraph.alignment is not None else None,
                paragraph_format.keep_with_next,
                paragraph_format.keep_together,
                paragraph_format.page_break_before,
                paragraph_format.widow_control,
                _docx_length_snapshot(paragraph_format.space_before),
                _docx_length_snapshot(paragraph_format.space_after),
                _docx_length_snapshot(paragraph_format.line_spacing),
                ilvl.get(qn("w:val")) if ilvl is not None else None,
                num_id.get(qn("w:val")) if num_id is not None else None,
                runs,
            )
        )
    table_cells: list[tuple[Any, ...]] = []
    for table in document.tables:
        cells: list[tuple[Any, ...]] = []
        for row in table.rows:
            for cell in row.cells:
                tc_pr = cell._tc.tcPr
                shading = tc_pr.find(qn("w:shd")) if tc_pr is not None else None
                borders = tc_pr.find(qn("w:tcBorders")) if tc_pr is not None else None
                grid_span = tc_pr.find(qn("w:gridSpan")) if tc_pr is not None else None
                v_merge = tc_pr.find(qn("w:vMerge")) if tc_pr is not None else None
                border_snapshot = ()
                if borders is not None:
                    border_snapshot = tuple(
                        sorted(
                            (
                                child.tag.rsplit("}", 1)[-1],
                                child.get(qn("w:val")),
                                child.get(qn("w:sz")),
                                child.get(qn("w:color")),
                            )
                            for child in borders
                        )
                    )
                cells.append(
                    (
                        cell.text,
                        grid_span.get(qn("w:val")) if grid_span is not None else None,
                        v_merge.get(qn("w:val")) if v_merge is not None else None,
                        shading.get(qn("w:fill")) if shading is not None else None,
                        border_snapshot,
                    )
                )
        table_cells.append(tuple(cells))
    drawing_metadata = tuple(
        (
            shape._inline.docPr.get("descr"),
            shape._inline.docPr.get("title"),
        )
        for shape in document.inline_shapes
    )
    return {
        "sections": sections,
        "paragraphs": tuple(paragraphs),
        "tables": tuple(table_cells),
        "drawing_metadata": drawing_metadata,
    }


def _pptx_semantic(presentation: Any) -> Iterable[str]:
    for paragraph in _iter_pptx_paragraphs(presentation):
        yield _paragraph_run_text(paragraph)


def _pptx_shape_counts(presentation: Any) -> dict[str, int]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    counts = {"pictures": 0, "tables": 0, "text_boxes": 0}

    def visit(shape: Any) -> None:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            counts["pictures"] += 1
        if getattr(shape, "has_table", False):
            counts["tables"] += 1
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            counts["text_boxes"] += 1
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                visit(child)

    for slide in presentation.slides:
        for shape in slide.shapes:
            visit(shape)
    return counts


def _pptx_v2_snapshot(presentation: Any) -> tuple[dict[str, Any], ...]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    snapshots: list[dict[str, Any]] = []
    for slide in presentation.slides:
        shapes: list[tuple[Any, ...]] = []
        for shape in slide.shapes:
            base = (
                int(shape.shape_type),
                int(shape.left),
                int(shape.top),
                int(shape.width),
                int(shape.height),
            )
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                shapes.append(base + ("chart", int(shape.chart.chart_type), len(shape.chart.series)))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shapes.append(
                    base
                    + (
                        "picture",
                        round(float(shape.crop_left), 6),
                        round(float(shape.crop_top), 6),
                        round(float(shape.crop_right), 6),
                        round(float(shape.crop_bottom), 6),
                    )
                )
            elif getattr(shape, "has_table", False):
                shapes.append(base + ("table", len(shape.table.rows), len(shape.table.columns)))
            else:
                shapes.append(base + ("shape", getattr(shape, "text", "")))
        notes = None
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
        snapshots.append(
            {
                "layout": slide.slide_layout.name,
                "shapes": tuple(shapes),
                "notes": notes,
            }
        )
    return tuple(snapshots)


def _semantic_digest(values: Iterable[str]) -> str:
    digest = hashlib.sha256()
    for value in values:
        encoded = value.encode("utf-8")
        digest.update(len(encoded).to_bytes(8, "big"))
        digest.update(encoded)
    return digest.hexdigest()


def _add_slide(
    presentation: Any,
    title: str,
    subtitle: str | None,
    bullets: Sequence[Mapping[str, Any]],
    text_boxes: Sequence[Mapping[str, Any]],
    tables: Sequence[Mapping[str, Any]],
    images: Sequence[Mapping[str, Any]],
    *,
    layout_index: Any = None,
    layout_name: Any = None,
    title_style: Mapping[str, Any] | None = None,
    shapes: Sequence[Mapping[str, Any]] = (),
    charts: Sequence[Mapping[str, Any]] = (),
    speaker_notes: str | None = None,
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.util import Inches, Pt

    if layout_index is not None and layout_name is not None:
        raise OfficeInputError(
            "PPTX 布局只能指定 layout_index 或 layout_name 之一。",
            "Specify only one of PPTX layout_index or layout_name.",
        )
    if layout_index is not None:
        if isinstance(layout_index, bool) or not isinstance(layout_index, int) or not 0 <= layout_index < len(presentation.slide_layouts):
            raise OfficeInputError(
                "PPTX layout_index 超出现有布局范围。",
                "PPTX layout_index is outside the available layout range.",
            )
        selected_layout = presentation.slide_layouts[layout_index]
    elif layout_name is not None:
        requested_name = _required_text(layout_name, "presentation.slides[].layout_name")
        matching = [layout for layout in presentation.slide_layouts if layout.name == requested_name]
        if len(matching) != 1:
            raise OfficeInputError(
                f"PPTX layout_name 未唯一匹配现有布局：{requested_name}",
                f"PPTX layout_name did not uniquely match an existing layout: {requested_name}",
            )
        selected_layout = matching[0]
    elif subtitle is not None and not bullets:
        selected_layout = presentation.slide_layouts[0]
    elif bullets:
        selected_layout = presentation.slide_layouts[1]
    else:
        default_index = 5 if len(presentation.slide_layouts) > 5 else 0
        selected_layout = presentation.slide_layouts[default_index]
    slide = presentation.slides.add_slide(selected_layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
        title_shape = slide.shapes.title
    else:
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.75))
        title_shape.text_frame.text = title
    if title_style:
        _apply_pptx_text_style(title_shape.text_frame, title_style)
    if subtitle is not None:
        for placeholder in slide.placeholders:
            if placeholder == slide.shapes.title or not getattr(
                placeholder, "has_text_frame", False
            ):
                continue
            placeholder.text = subtitle
            break
    if bullets:
        body = next(
            (
                placeholder
                for placeholder in slide.placeholders
                if placeholder != slide.shapes.title
                and getattr(placeholder, "has_text_frame", False)
            ),
            None,
        )
        if body is None:
            raise OfficeInputError(
                "PPTX 默认布局中缺少正文占位符。",
                "The default PPTX layout has no body placeholder.",
            )
        frame = body.text_frame
        frame.clear()
        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = bullet["text"]
            paragraph.level = bullet["level"]
    for text_box in text_boxes:
        shape = slide.shapes.add_textbox(
            Inches(text_box["left_inches"]),
            Inches(text_box["top_inches"]),
            Inches(text_box["width_inches"]),
            Inches(text_box["height_inches"]),
        )
        shape.text_frame.text = text_box["text"]
        if text_box["font_size"] is not None:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(text_box["font_size"])
        if text_box.get("style"):
            _apply_pptx_text_style(shape.text_frame, text_box["style"])
    for table_data in tables:
        headers = table_data["headers"]
        rows = table_data["rows"]
        column_count = max([len(headers), *(len(row) for row in rows)])
        row_count = len(rows) + (1 if headers else 0)
        shape = slide.shapes.add_table(
            row_count,
            column_count,
            Inches(table_data["left_inches"]),
            Inches(table_data["top_inches"]),
            Inches(table_data["width_inches"]),
            Inches(table_data["height_inches"]),
        )
        table = shape.table
        output_row = 0
        if headers:
            for column, value in enumerate(headers):
                table.cell(0, column).text = _cell_text(value)
            output_row = 1
        for row_index, row in enumerate(rows, output_row):
            for column, value in enumerate(row):
                table.cell(row_index, column).text = _cell_text(value)
        table_style = table_data.get("style")
        if table_style:
            for row_index, row in enumerate(table.rows):
                fill_color = (
                    table_style.get("header_fill_color")
                    if headers and row_index == 0
                    else table_style.get("body_fill_color")
                )
                for cell in row.cells:
                    if fill_color:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor.from_string(fill_color)
                    if table_style.get("font_size") is not None:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(table_style["font_size"])
    for image in images:
        width = Inches(image["width_inches"]) if image["width_inches"] else None
        height = Inches(image["height_inches"]) if image["height_inches"] else None
        picture = slide.shapes.add_picture(
            io.BytesIO(image["data"]),
            Inches(image["left_inches"]),
            Inches(image["top_inches"]),
            width=width,
            height=height,
        )
        for crop in ("crop_left", "crop_top", "crop_right", "crop_bottom"):
            if crop in image:
                setattr(picture, crop, image[crop])
        if image.get("align") in {"center", "right"} and picture.width > presentation.slide_width:
            raise OfficeInputError(
                "PPTX 对齐图片宽度超过幻灯片。",
                "An aligned PPTX picture is wider than the slide.",
            )
        if image.get("align") == "center":
            picture.left = int((presentation.slide_width - picture.width) / 2)
        elif image.get("align") == "right":
            picture.left = presentation.slide_width - picture.width
    for shape_data in shapes:
        left = Inches(shape_data["left_inches"])
        top = Inches(shape_data["top_inches"])
        width = Inches(shape_data["width_inches"])
        height = Inches(shape_data["height_inches"])
        if shape_data["type"] == "line":
            shape = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                left,
                top,
                left + width,
                top + height,
            )
        else:
            shape_type = {
                "rectangle": MSO_SHAPE.RECTANGLE,
                "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
                "ellipse": MSO_SHAPE.OVAL,
            }[shape_data["type"]]
            shape = slide.shapes.add_shape(shape_type, left, top, width, height)
            if shape_data.get("text") is not None:
                shape.text_frame.text = shape_data["text"]
                if shape_data.get("text_style"):
                    _apply_pptx_text_style(shape.text_frame, shape_data["text_style"])
            if shape_data.get("fill_color"):
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(shape_data["fill_color"])
        if shape_data.get("line_color"):
            shape.line.color.rgb = RGBColor.from_string(shape_data["line_color"])
    for chart_data in charts:
        _add_pptx_chart(slide, chart_data)
    if speaker_notes is not None:
        slide.notes_slide.notes_text_frame.text = speaker_notes


def _apply_pptx_text_style(text_frame: Any, style: Mapping[str, Any]) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(style.get("alignment"))
    for paragraph in text_frame.paragraphs:
        if alignment is not None:
            paragraph.alignment = alignment
        for run in paragraph.runs:
            if style.get("font_size") is not None:
                run.font.size = Pt(style["font_size"])
            for field in ("bold", "italic"):
                if field in style:
                    setattr(run.font, field, style[field])
            if style.get("color"):
                run.font.color.rgb = RGBColor.from_string(style["color"])
            if style.get("font"):
                run.font.name = style["font"]


def _add_pptx_chart(slide: Any, chart_data: Mapping[str, Any]) -> None:
    from pptx.chart.data import ChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    chart_type = chart_data["type"]
    if chart_type == "scatter":
        data = XyChartData()
        for series_data in chart_data["series"]:
            series = data.add_series(series_data["name"])
            for x_value, y_value in zip(
                series_data["x_values"], series_data["y_values"], strict=True
            ):
                series.add_data_point(x_value, y_value)
        native_type = XL_CHART_TYPE.XY_SCATTER
    else:
        data = ChartData()
        data.categories = chart_data["categories"]
        for series_data in chart_data["series"]:
            data.add_series(series_data["name"], tuple(series_data["values"]))
        native_type = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
        }[chart_type]
    chart = slide.shapes.add_chart(
        native_type,
        Inches(chart_data["left_inches"]),
        Inches(chart_data["top_inches"]),
        Inches(chart_data["width_inches"]),
        Inches(chart_data["height_inches"]),
        data,
    ).chart
    if chart_data.get("title") is not None:
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_data["title"]


def _validate_sheet_name(name: str) -> None:
    if len(name) > 31 or _INVALID_SHEET_TITLE.search(name) or name.startswith("'") or name.endswith("'"):
        raise OfficeInputError(
            f"无效的 XLSX 工作表名：{name}",
            f"Invalid XLSX sheet name: {name}",
        )


def _apply_xlsx_style(cell: Any, value: Any, name: str) -> None:
    from openpyxl.styles import PatternFill

    style = _mapping(value, name)
    unknown = sorted(set(style) - {"number_format", "font", "fill"})
    if unknown:
        raise OfficeInputError(
            f"{name} 包含不支持的样式项：{', '.join(unknown)}",
            f"{name} contains unsupported style fields: {', '.join(unknown)}",
        )
    if not style:
        raise OfficeInputError(f"{name} 不能为空。", f"{name} cannot be empty.")
    if "number_format" in style:
        number_format = _required_text(style["number_format"], f"{name}.number_format")
        if len(number_format) > 255:
            raise OfficeInputError(
                f"{name}.number_format 超过 255 个字符。",
                f"{name}.number_format exceeds 255 characters.",
            )
        cell.number_format = number_format
    if "font" in style:
        font_data = _mapping(style["font"], f"{name}.font")
        unknown_font = sorted(set(font_data) - {"bold", "italic", "color", "size"})
        if unknown_font:
            raise OfficeInputError(
                f"{name}.font 包含不支持的字段：{', '.join(unknown_font)}",
                f"{name}.font contains unsupported fields: {', '.join(unknown_font)}",
            )
        if not font_data:
            raise OfficeInputError(
                f"{name}.font 不能为空。", f"{name}.font cannot be empty."
            )
        font = copy.copy(cell.font)
        for field in ("bold", "italic"):
            if field in font_data:
                if not isinstance(font_data[field], bool):
                    raise OfficeInputError(
                        f"{name}.font.{field} 必须是布尔值。",
                        f"{name}.font.{field} must be a boolean.",
                    )
                setattr(font, field, font_data[field])
        if "color" in font_data:
            font.color = _validate_hex_color(font_data["color"], f"{name}.font.color")
        if "size" in font_data:
            font.size = _bounded_number(
                font_data["size"], f"{name}.font.size", minimum=1, maximum=200
            )
        cell.font = font
    if "fill" in style:
        fill_data = _mapping(style["fill"], f"{name}.fill")
        if set(fill_data) != {"color"}:
            raise OfficeInputError(
                f"{name}.fill 只支持 color。",
                f"{name}.fill supports only color.",
            )
        color = _validate_hex_color(fill_data["color"], f"{name}.fill.color")
        cell.fill = PatternFill(fill_type="solid", fgColor=color)


def _validate_hex_color(value: Any, name: str) -> str:
    if not isinstance(value, str) or not _HEX_COLOR.fullmatch(value):
        raise OfficeInputError(
            f"{name} 必须是 6 或 8 位十六进制颜色。",
            f"{name} must be a 6- or 8-digit hexadecimal color.",
        )
    return value.upper()


def _xlsx_color_snapshot(color: Any) -> dict[str, Any] | None:
    if color is None:
        return None
    return {
        "type": color.type,
        "rgb": color.rgb if color.type == "rgb" else None,
        "indexed": color.indexed if color.type == "indexed" else None,
        "theme": color.theme if color.type == "theme" else None,
        "tint": float(color.tint or 0),
    }


def _xlsx_style_snapshot(cell: Any) -> dict[str, Any]:
    return {
        "number_format": cell.number_format,
        "font": {
            "bold": bool(cell.font.bold),
            "italic": bool(cell.font.italic),
            "size": float(cell.font.size) if cell.font.size is not None else None,
            "color": _xlsx_color_snapshot(cell.font.color),
        },
        "fill": {
            "fill_type": cell.fill.fill_type,
            "fg_color": _xlsx_color_snapshot(cell.fill.fgColor),
        },
    }


def _cell_text(value: Any) -> str:
    return "" if value is None else str(value)


def _is_formula(value: Any) -> bool:
    return isinstance(value, str) and value.startswith("=")


def _flush_file(path: Path) -> None:
    # Windows implements os.fsync() with the CRT commit operation, which
    # rejects a descriptor opened read-only with EBADF.  The Office output is
    # still our private, writable temporary at this point, so open it for
    # update before making the durability barrier.
    with path.open("r+b") as handle:
        os.fsync(handle.fileno())


def _atomic_replace(source: Path, target: Path) -> None:
    """Single-filesystem install seam, kept small for fault-injection tests."""

    os.replace(source, target)


def _fsync_directory(directory: Path) -> None:
    flags = os.O_RDONLY
    if hasattr(os, "O_DIRECTORY"):
        flags |= os.O_DIRECTORY
    try:
        directory_fd = os.open(directory, flags)
    except OSError:
        return
    try:
        os.fsync(directory_fd)
    except OSError:
        pass
    finally:
        os.close(directory_fd)
