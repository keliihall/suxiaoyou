"""Safe, offline static previews for OOXML PowerPoint presentations.

The endpoint in this module deliberately does not turn Office XML into HTML.
It validates the ZIP container, parses the supported drawing primitives with
``python-pptx``, and returns a small allow-listed scene model.  The frontend
renders that model with React/SVG primitives, so hyperlinks, macros, embedded
objects, and arbitrary XML can never become executable webview content.
"""

from __future__ import annotations

import asyncio
import base64
import hashlib
import json
import logging
import math
import threading
import time
import zipfile
from concurrent.futures import ThreadPoolExecutor
from dataclasses import dataclass, field
from pathlib import Path, PurePosixPath
from typing import Any

from fastapi import APIRouter, HTTPException
from pydantic import BaseModel

from app.api.files import _resolve_requested_file_path

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/files")

# Container and rendering limits. These are intentionally independent from the
# generic binary-preview limit because an OOXML ZIP can expand substantially.
MAX_PPTX_FILE_BYTES = 50 * 1024 * 1024
MAX_PPTX_ZIP_ENTRIES = 4096
MAX_PPTX_UNCOMPRESSED_BYTES = 64 * 1024 * 1024
MAX_PPTX_ZIP_MEMBER_BYTES = 16 * 1024 * 1024
MAX_PPTX_COMPRESSION_RATIO = 250
MAX_PPTX_SLIDES = 200
MAX_PPTX_SHAPES_PER_SLIDE = 1000
MAX_PPTX_SHAPES_TOTAL = 5000
MAX_PPTX_TABLE_CELLS = 10_000
MAX_PPTX_TABLE_CELLS_PER_TABLE = 2_000
MAX_PPTX_IMAGE_BYTES = 8 * 1024 * 1024
MAX_PPTX_IMAGE_PIXELS = 12_000_000
MAX_PPTX_TOTAL_IMAGE_PIXELS = 40_000_000
MAX_PPTX_IMAGE_ASSETS = 64
MAX_PPTX_PREVIEW_IMAGE_BYTES = 8 * 1024 * 1024
MAX_PPTX_TEXT_CHARS = 2_000_000
MAX_PPTX_PARAGRAPHS = 5_000
MAX_PPTX_RUNS = 10_000
MAX_PPTX_PARAGRAPHS_SCANNED = 10_000
MAX_PPTX_RUNS_SCANNED = 20_000
MAX_PPTX_SCENE_NODES = 25_000
MAX_PPTX_SCENE_JSON_BYTES = 16 * 1024 * 1024
# Leave headroom for JSON punctuation, warning strings, and response metadata.
MAX_PPTX_SCENE_BUDGET_BYTES = 15 * 1024 * 1024
PPTX_DEADLINE_CHECK_INTERVAL = 32
PPTX_RENDER_TIMEOUT_SECONDS = 20.0
PPTX_WORKER_GRACE_SECONDS = 1.0

EMU_PER_CSS_PIXEL = 9525.0  # 914400 EMU / 96 CSS pixels per inch
POINT_TO_CSS_PIXEL = 96.0 / 72.0

SUPPORTED_IMAGE_MIME_TYPES = {
    "image/png",
    "image/jpeg",
    "image/bmp",
}

SUPPORTED_ZIP_COMPRESSION = {
    zipfile.ZIP_STORED,
    zipfile.ZIP_DEFLATED,
}

# PPTX rendering gets a dedicated worker so it cannot consume the default
# asyncio executor. The non-blocking gate prevents requests from building an
# unbounded queue behind a slow or hostile document.
_PPTX_RENDER_EXECUTOR = ThreadPoolExecutor(
    max_workers=1,
    thread_name_prefix="suxiaoyou-pptx-preview",
)
_PPTX_RENDER_GATE = threading.Lock()

THEME_COLORS = {
    "DARK_1": "#000000",
    "LIGHT_1": "#FFFFFF",
    "DARK_2": "#1F497D",
    "LIGHT_2": "#EEECE1",
    "ACCENT_1": "#4F81BD",
    "ACCENT_2": "#C0504D",
    "ACCENT_3": "#9BBB59",
    "ACCENT_4": "#8064A2",
    "ACCENT_5": "#4BACC6",
    "ACCENT_6": "#F79646",
    "HYPERLINK": "#0000FF",
    "FOLLOWED_HYPERLINK": "#800080",
    "BACKGROUND_1": "#FFFFFF",
    "TEXT_1": "#000000",
    "BACKGROUND_2": "#EEECE1",
    "TEXT_2": "#1F497D",
}


class PptxPreviewRequest(BaseModel):
    path: str
    workspace: str | None = None


class PptxPreviewFailure(Exception):
    """An expected preview failure with a stable client-facing code."""

    def __init__(self, code: str, status_code: int) -> None:
        super().__init__(code)
        self.code = code
        self.status_code = status_code


@dataclass
class PreviewContext:
    deadline: float
    warnings: set[str] = field(default_factory=lambda: {"static_preview_limitations"})
    assets: dict[str, dict[str, Any]] = field(default_factory=dict)
    image_bytes: int = 0
    image_pixels: int = 0
    text_chars: int = 0
    shape_count: int = 0
    table_cells: int = 0
    paragraph_count: int = 0
    run_count: int = 0
    paragraphs_scanned: int = 0
    runs_scanned: int = 0
    scene_nodes: int = 0
    estimated_scene_bytes: int = 4096

    def check_deadline(self) -> None:
        if time.monotonic() > self.deadline:
            raise PptxPreviewFailure("pptx_preview_timeout", 408)

    def bounded_text(self, value: str) -> str:
        remaining = MAX_PPTX_TEXT_CHARS - self.text_chars
        if remaining <= 0:
            self.warnings.add("text_truncated")
            return ""
        if len(value) > remaining:
            value = value[:remaining]
            self.warnings.add("text_truncated")
        self.text_chars += len(value)
        return value

    def reserve_bytes(self, amount: int) -> bool:
        amount = max(0, int(amount))
        if self.estimated_scene_bytes + amount > MAX_PPTX_SCENE_BUDGET_BYTES:
            self.warnings.add("scene_size_limit_exceeded")
            return False
        self.estimated_scene_bytes += amount
        return True

    def reserve_node(self, estimated_bytes: int = 192) -> bool:
        if self.scene_nodes >= MAX_PPTX_SCENE_NODES:
            self.warnings.add("scene_node_limit_exceeded")
            return False
        if not self.reserve_bytes(estimated_bytes):
            return False
        self.scene_nodes += 1
        return True

    def reserve_paragraph(self) -> bool:
        if self.paragraph_count >= MAX_PPTX_PARAGRAPHS:
            self.warnings.add("paragraph_limit_exceeded")
            return False
        if not self.reserve_node(192):
            return False
        self.paragraph_count += 1
        return True

    def reserve_run(self, value: str, family: str | None) -> bool:
        if self.run_count >= MAX_PPTX_RUNS:
            self.warnings.add("run_limit_exceeded")
            return False
        estimated_bytes = 192 + len(value.encode("utf-8"))
        if family:
            estimated_bytes += len(family.encode("utf-8"))
        if not self.reserve_node(estimated_bytes):
            return False
        self.run_count += 1
        return True


def _px(value: Any) -> float:
    try:
        return round(float(value or 0) / EMU_PER_CSS_PIXEL, 3)
    except (TypeError, ValueError):
        return 0.0


def _pt(value: Any, default: float) -> float:
    try:
        if value is None:
            return default
        # python-pptx Length values expose a ``pt`` property.
        return round(float(value.pt), 2)
    except (AttributeError, TypeError, ValueError):
        return default


def _enum_name(value: Any) -> str:
    return str(getattr(value, "name", "") or "").upper()


def _adjust_brightness(color: str, brightness: float) -> str:
    try:
        amount = max(-1.0, min(1.0, float(brightness)))
        channels = [int(color[index : index + 2], 16) for index in (1, 3, 5)]
    except (TypeError, ValueError):
        return color
    if amount >= 0:
        channels = [round(channel + (255 - channel) * amount) for channel in channels]
    else:
        channels = [round(channel * (1 + amount)) for channel in channels]
    return "#" + "".join(f"{max(0, min(255, value)):02X}" for value in channels)


def _color(color_format: Any, fallback: str) -> str:
    if color_format is None:
        return fallback
    value = fallback
    try:
        rgb = color_format.rgb
        if rgb is not None:
            value = f"#{rgb}"
        else:
            value = THEME_COLORS.get(_enum_name(color_format.theme_color), fallback)
        value = _adjust_brightness(value, color_format.brightness or 0)
    except (AttributeError, TypeError, ValueError):
        pass
    return value


def _fill_color(fill: Any, fallback: str = "transparent") -> str:
    try:
        fill_type = _enum_name(fill.type)
        if fill_type == "SOLID":
            return _color(fill.fore_color, fallback)
        if fill_type == "BACKGROUND":
            return "transparent"
    except (AttributeError, TypeError, ValueError):
        pass
    return fallback


def _line_style(shape: Any, *, is_line: bool = False) -> tuple[str, float]:
    try:
        line = shape.line
        fill_type = _enum_name(line.fill.type)
        if fill_type == "BACKGROUND":
            return "transparent", 0.0
        fallback = "#595959" if is_line else "transparent"
        stroke = _color(line.color, fallback)
        width = _px(line.width)
        if stroke != "transparent" and width <= 0:
            width = 1.0
        return stroke, width
    except (AttributeError, TypeError, ValueError):
        return ("#595959", 1.0) if is_line else ("transparent", 0.0)


def _shape_rect(shape: Any) -> dict[str, float]:
    return {
        "x": _px(getattr(shape, "left", 0)),
        "y": _px(getattr(shape, "top", 0)),
        "width": max(0.0, _px(getattr(shape, "width", 0))),
        "height": max(0.0, _px(getattr(shape, "height", 0))),
        "rotation": round(float(getattr(shape, "rotation", 0) or 0), 3),
    }


def _shape_flips(shape: Any) -> tuple[bool, bool]:
    try:
        xfrm = shape._element.xfrm
        return bool(xfrm.flipH), bool(xfrm.flipV)
    except (AttributeError, TypeError, ValueError):
        return False, False


def _default_font_size(shape: Any) -> float:
    try:
        if shape.is_placeholder:
            placeholder_type = _enum_name(shape.placeholder_format.type)
            if "TITLE" in placeholder_type:
                return 32.0
            if "SUBTITLE" in placeholder_type:
                return 20.0
    except (AttributeError, ValueError):
        pass
    return 18.0


def _serialize_font(font: Any, *, default_size: float) -> dict[str, Any]:
    name = None
    try:
        name = font.name
    except AttributeError:
        pass
    color = "#1F1F1F"
    try:
        color = _color(font.color, color)
    except AttributeError:
        pass
    return {
        "family": str(name)[:200] if name else None,
        "size": _pt(getattr(font, "size", None), default_size),
        "bold": bool(getattr(font, "bold", False) or False),
        "italic": bool(getattr(font, "italic", False) or False),
        "underline": bool(getattr(font, "underline", False) or False),
        "color": color,
    }


def _paragraph_bullet(paragraph: Any) -> bool:
    """Detect explicit bullets without resolving external layout content."""

    try:
        p_pr = paragraph._p.pPr
        if p_pr is None:
            return bool(paragraph.level)
        child_names = {child.tag.rsplit("}", 1)[-1] for child in p_pr}
        if "buNone" in child_names:
            return False
        if {"buChar", "buAutoNum"} & child_names:
            return True
    except (AttributeError, TypeError):
        pass
    return bool(getattr(paragraph, "level", 0))


def _serialize_text_frame(
    text_frame: Any,
    shape: Any,
    context: PreviewContext,
) -> dict[str, Any] | None:
    default_size = _default_font_size(shape)
    paragraphs: list[dict[str, Any]] = []
    for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
        if paragraph_index % PPTX_DEADLINE_CHECK_INTERVAL == 0:
            context.check_deadline()
        if context.paragraphs_scanned >= MAX_PPTX_PARAGRAPHS_SCANNED:
            context.warnings.add("paragraph_limit_exceeded")
            break
        context.paragraphs_scanned += 1

        runs: list[dict[str, Any]] = []
        paragraph_font = _serialize_font(paragraph.font, default_size=default_size)
        paragraph_runs = paragraph.runs
        if paragraph_runs:
            for run_index, run in enumerate(paragraph_runs):
                if run_index % PPTX_DEADLINE_CHECK_INTERVAL == 0:
                    context.check_deadline()
                if context.runs_scanned >= MAX_PPTX_RUNS_SCANNED:
                    context.warnings.add("run_limit_exceeded")
                    break
                context.runs_scanned += 1

                raw_value = str(run.text or "")
                if not raw_value:
                    continue
                value = context.bounded_text(raw_value)
                if not value:
                    break
                font = _serialize_font(
                    run.font,
                    default_size=paragraph_font["size"],
                )
                if not context.reserve_run(value, font.get("family")):
                    break
                runs.append({"text": value, **font})
        else:
            raw_value = str(paragraph.text or "")
            if raw_value:
                context.runs_scanned += 1
                value = context.bounded_text(raw_value)
                if value and context.reserve_run(value, paragraph_font.get("family")):
                    runs.append({"text": value, **paragraph_font})

        # Empty XML runs and blank paragraphs can otherwise create tens of
        # thousands of useless React nodes. Whitespace-only content is not a
        # meaningful static preview paragraph and is intentionally discarded.
        if not runs or not any(str(run["text"]).strip() for run in runs):
            continue
        if not context.reserve_paragraph():
            break

        alignment = _enum_name(paragraph.alignment)
        align = {
            "CENTER": "center",
            "RIGHT": "right",
            "JUSTIFY": "justify",
            "JUSTIFY_LOW": "justify",
            "DISTRIBUTE": "justify",
            "THAI_DISTRIBUTE": "justify",
        }.get(alignment, "left")
        paragraphs.append(
            {
                "runs": runs,
                "align": align,
                "level": max(0, min(8, int(getattr(paragraph, "level", 0) or 0))),
                "bullet": _paragraph_bullet(paragraph),
                "spaceBefore": _pt(getattr(paragraph, "space_before", None), 0),
                "spaceAfter": _pt(getattr(paragraph, "space_after", None), 0),
            }
        )

    if not paragraphs:
        return None
    if not context.reserve_node(192):
        return None

    anchor = _enum_name(getattr(text_frame, "vertical_anchor", None))
    vertical = {"MIDDLE": "middle", "BOTTOM": "bottom"}.get(anchor, "top")
    return {
        "paragraphs": paragraphs,
        "marginLeft": _px(getattr(text_frame, "margin_left", 0)),
        "marginRight": _px(getattr(text_frame, "margin_right", 0)),
        "marginTop": _px(getattr(text_frame, "margin_top", 0)),
        "marginBottom": _px(getattr(text_frame, "margin_bottom", 0)),
        "vertical": vertical,
    }


def _geometry(shape: Any, context: PreviewContext) -> str:
    shape_type = _enum_name(getattr(shape, "shape_type", None))
    if shape_type == "LINE":
        return "line"
    try:
        auto_shape = _enum_name(shape.auto_shape_type)
    except (AttributeError, ValueError):
        return "rect"

    geometry_map = {
        "RECTANGLE": "rect",
        "ROUNDED_RECTANGLE": "roundRect",
        "ROUND_1_RECTANGLE": "roundRect",
        "ROUND_2_DIAG_RECTANGLE": "roundRect",
        "ROUND_2_SAME_RECTANGLE": "roundRect",
        "SNIP_ROUND_RECTANGLE": "roundRect",
        "OVAL": "ellipse",
        "ISOSCELES_TRIANGLE": "triangle",
        "RIGHT_TRIANGLE": "rightTriangle",
        "DIAMOND": "diamond",
        "HEXAGON": "hexagon",
        "PENTAGON": "pentagon",
        "CHEVRON": "chevron",
        "PARALLELOGRAM": "parallelogram",
        "TRAPEZOID": "trapezoid",
        "LEFT_ARROW": "leftArrow",
        "RIGHT_ARROW": "rightArrow",
        "UP_ARROW": "upArrow",
        "DOWN_ARROW": "downArrow",
    }
    geometry = geometry_map.get(auto_shape)
    if geometry is None:
        context.warnings.add("shape_geometry_approximated")
        return "rect"
    return geometry


def _arrow_ends(shape: Any) -> tuple[bool, bool]:
    try:
        line = shape._element.spPr.ln
        head = line.headEnd
        tail = line.tailEnd
        return (
            bool(tail is not None and tail.get("type", "none") != "none"),
            bool(head is not None and head.get("type", "none") != "none"),
        )
    except AttributeError:
        return False, False


def _serialize_shape(shape: Any, context: PreviewContext) -> dict[str, Any]:
    geometry = _geometry(shape, context)
    is_line = geometry == "line"
    stroke, stroke_width = _line_style(shape, is_line=is_line)
    flip_h, flip_v = _shape_flips(shape)
    arrow_start, arrow_end = _arrow_ends(shape) if is_line else (False, False)
    element: dict[str, Any] = {
        "kind": "shape",
        **_shape_rect(shape),
        "geometry": geometry,
        "fill": "transparent" if is_line else _fill_color(shape.fill),
        "stroke": stroke,
        "strokeWidth": stroke_width,
        "flipH": flip_h,
        "flipV": flip_v,
        "arrowStart": arrow_start,
        "arrowEnd": arrow_end,
    }
    if bool(getattr(shape, "has_text_frame", False)):
        text = _serialize_text_frame(shape.text_frame, shape, context)
        if text is not None:
            element["text"] = text
    return element


def _image_signature_matches(content: bytes, mime_type: str) -> bool:
    if mime_type == "image/png":
        return content.startswith(b"\x89PNG\r\n\x1a\n")
    if mime_type == "image/jpeg":
        return content.startswith(b"\xff\xd8\xff")
    if mime_type == "image/bmp":
        return content.startswith(b"BM")
    return False


def _unsupported_element(shape: Any, label: str) -> dict[str, Any]:
    return {"kind": "unsupported", **_shape_rect(shape), "label": label}


def _serialize_picture(shape: Any, context: PreviewContext) -> dict[str, Any]:
    context.check_deadline()
    try:
        image = shape.image
        content = bytes(image.blob)
        mime_type = str(image.content_type or "").lower()
        width_px, height_px = image.size
    except Exception:
        context.warnings.add("unsupported_image")
        return _unsupported_element(shape, "image")

    if (
        mime_type not in SUPPORTED_IMAGE_MIME_TYPES
        or not _image_signature_matches(content, mime_type)
        or len(content) > MAX_PPTX_IMAGE_BYTES
        or width_px <= 0
        or height_px <= 0
        or width_px * height_px > MAX_PPTX_IMAGE_PIXELS
    ):
        context.warnings.add("unsupported_image")
        return _unsupported_element(shape, "image")

    digest = hashlib.sha256(content).hexdigest()
    if digest not in context.assets:
        image_pixels = int(width_px) * int(height_px)
        if len(context.assets) >= MAX_PPTX_IMAGE_ASSETS:
            context.warnings.add("asset_limit_exceeded")
            return _unsupported_element(shape, "image")
        if context.image_pixels + image_pixels > MAX_PPTX_TOTAL_IMAGE_PIXELS:
            context.warnings.add("total_image_pixels_exceeded")
            return _unsupported_element(shape, "image")
        if context.image_bytes + len(content) > MAX_PPTX_PREVIEW_IMAGE_BYTES:
            context.warnings.add("image_budget_exceeded")
            return _unsupported_element(shape, "image")
        encoded_size = 4 * ((len(content) + 2) // 3)
        if not context.reserve_node(encoded_size + 320):
            return _unsupported_element(shape, "image")
        encoded = base64.b64encode(content).decode("ascii")
        data_url = f"data:{mime_type};base64,{encoded}"
        context.assets[digest] = {
            "mimeType": mime_type,
            "dataUrl": data_url,
            "width": int(width_px),
            "height": int(height_px),
        }
        context.image_bytes += len(content)
        context.image_pixels += image_pixels

    flip_h, flip_v = _shape_flips(shape)
    return {
        "kind": "image",
        **_shape_rect(shape),
        "assetId": digest,
        "cropLeft": max(0.0, min(0.99, float(getattr(shape, "crop_left", 0) or 0))),
        "cropRight": max(0.0, min(0.99, float(getattr(shape, "crop_right", 0) or 0))),
        "cropTop": max(0.0, min(0.99, float(getattr(shape, "crop_top", 0) or 0))),
        "cropBottom": max(0.0, min(0.99, float(getattr(shape, "crop_bottom", 0) or 0))),
        "flipH": flip_h,
        "flipV": flip_v,
    }


def _serialize_table(shape: Any, context: PreviewContext) -> dict[str, Any]:
    table = shape.table
    table_cell_count = len(table.rows) * len(table.columns)
    if (
        table_cell_count > MAX_PPTX_TABLE_CELLS_PER_TABLE
        or context.table_cells + table_cell_count > MAX_PPTX_TABLE_CELLS
    ):
        context.warnings.add("table_cell_limit_exceeded")
        return _unsupported_element(shape, "table")
    context.table_cells += table_cell_count
    column_widths = [_px(column.width) for column in table.columns]
    row_heights = [_px(row.height) for row in table.rows]
    cells: list[dict[str, Any]] = []
    y = 0.0
    budget_exhausted = False
    for row_index, row in enumerate(table.rows):
        x = 0.0
        for column_index, cell in enumerate(row.cells):
            if column_index % PPTX_DEADLINE_CHECK_INTERVAL == 0:
                context.check_deadline()
            width = column_widths[column_index]
            if not cell.is_spanned:
                if not context.reserve_node(256):
                    budget_exhausted = True
                    break
                span_width = max(1, int(cell.span_width or 1))
                span_height = max(1, int(cell.span_height or 1))
                cell_data: dict[str, Any] = {
                    "x": round(x, 3),
                    "y": round(y, 3),
                    "width": round(
                        sum(column_widths[column_index : column_index + span_width]),
                        3,
                    ),
                    "height": round(
                        sum(row_heights[row_index : row_index + span_height]),
                        3,
                    ),
                    "fill": _fill_color(cell.fill, "#FFFFFF"),
                }
                text = _serialize_text_frame(cell.text_frame, shape, context)
                if text is not None:
                    cell_data["text"] = text
                cells.append(cell_data)
            x += width
        if budget_exhausted:
            break
        y += row_heights[row_index]
    return {
        "kind": "table",
        **_shape_rect(shape),
        "cells": cells,
    }


def _serialize_slide(slide: Any, index: int, context: PreviewContext) -> dict[str, Any]:
    context.check_deadline()
    background = "#FFFFFF"
    try:
        background = _fill_color(slide.background.fill, "#FFFFFF")
        if background == "transparent":
            background = "#FFFFFF"
    except (AttributeError, ValueError):
        pass

    elements: list[dict[str, Any]] = []
    for shape_index, shape in enumerate(slide.shapes):
        if shape_index >= MAX_PPTX_SHAPES_PER_SLIDE:
            context.warnings.add("shape_limit_exceeded")
            break
        if shape_index % PPTX_DEADLINE_CHECK_INTERVAL == 0:
            context.check_deadline()
        if context.shape_count >= MAX_PPTX_SHAPES_TOTAL:
            context.warnings.add("shape_limit_exceeded")
            break
        if not context.reserve_node(320):
            break
        context.shape_count += 1
        shape_type = _enum_name(getattr(shape, "shape_type", None))
        try:
            if shape_type == "PICTURE":
                elements.append(_serialize_picture(shape, context))
            elif bool(getattr(shape, "has_table", False)):
                elements.append(_serialize_table(shape, context))
            elif shape_type in {"AUTO_SHAPE", "TEXT_BOX", "PLACEHOLDER", "LINE", "FREEFORM"}:
                elements.append(_serialize_shape(shape, context))
            elif shape_type == "GROUP":
                context.warnings.add("unsupported_group")
                elements.append(_unsupported_element(shape, "group"))
            elif shape_type in {
                "CHART",
                "DIAGRAM",
                "EMBEDDED_OLE_OBJECT",
                "LINKED_OLE_OBJECT",
                "MEDIA",
                "WEB_VIDEO",
                "SCRIPT_ANCHOR",
                "FORM_CONTROL",
                "OLE_CONTROL_OBJECT",
            }:
                context.warnings.add("unsupported_embedded_content")
                elements.append(_unsupported_element(shape, shape_type.lower()))
            elif bool(getattr(shape, "has_text_frame", False)):
                elements.append(_serialize_shape(shape, context))
            else:
                context.warnings.add("unsupported_element")
                elements.append(_unsupported_element(shape, "object"))
        except PptxPreviewFailure:
            raise
        except Exception:
            logger.debug("PPTX shape %s on slide %d could not be rendered", shape_type, index, exc_info=True)
            context.warnings.add("unsupported_element")
            elements.append(_unsupported_element(shape, "object"))

    hidden = False
    try:
        hidden = slide._element.get("show") == "0"
    except AttributeError:
        pass
    return {"index": index, "background": background, "hidden": hidden, "elements": elements}


def _validate_archive(path: Path) -> set[str]:
    warnings: set[str] = set()
    try:
        with zipfile.ZipFile(path) as archive:
            infos = archive.infolist()
            if len(infos) > MAX_PPTX_ZIP_ENTRIES:
                raise PptxPreviewFailure("pptx_zip_entry_limit", 413)

            names: set[str] = set()
            total_size = 0
            slide_count = 0
            for info in infos:
                name = info.filename
                pure_name = PurePosixPath(name)
                if (
                    not name
                    or name.startswith(("/", "\\"))
                    or "\\" in name
                    or ".." in pure_name.parts
                ):
                    raise PptxPreviewFailure("pptx_invalid_archive_path", 422)
                if info.flag_bits & 0x1:
                    raise PptxPreviewFailure("pptx_encrypted_archive", 422)
                if info.compress_type not in SUPPORTED_ZIP_COMPRESSION:
                    raise PptxPreviewFailure(
                        "pptx_unsupported_zip_compression",
                        422,
                    )
                if info.file_size > MAX_PPTX_ZIP_MEMBER_BYTES:
                    raise PptxPreviewFailure("pptx_zip_member_limit", 413)
                total_size += info.file_size
                if total_size > MAX_PPTX_UNCOMPRESSED_BYTES:
                    raise PptxPreviewFailure("pptx_uncompressed_limit", 413)
                if info.file_size > 1024 * 1024:
                    ratio = info.file_size / max(1, info.compress_size)
                    if ratio > MAX_PPTX_COMPRESSION_RATIO:
                        raise PptxPreviewFailure("pptx_compression_ratio_limit", 413)

                if name in names:
                    raise PptxPreviewFailure("pptx_duplicate_zip_entry", 422)
                names.add(name)
                if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                    slide_count += 1
                lowered = name.lower()
                if (
                    "vbaproject" in lowered
                    or lowered.startswith("ppt/embeddings/")
                    or lowered.startswith("ppt/activex/")
                    or lowered.startswith("customui/")
                ):
                    warnings.add("ignored_embedded_content")

            if "[Content_Types].xml" not in names or "ppt/presentation.xml" not in names:
                raise PptxPreviewFailure("pptx_invalid_ooxml", 422)
            if slide_count > MAX_PPTX_SLIDES:
                raise PptxPreviewFailure("pptx_slide_limit", 413)

            # Relationship XML is bounded by the member limits above. Reading
            # it only detects external targets; python-pptx itself skips them.
            for info in infos:
                if info.filename.endswith(".rels"):
                    rel_bytes = archive.read(info)
                    if b'TargetMode="External"' in rel_bytes or b"TargetMode='External'" in rel_bytes:
                        warnings.add("ignored_external_links")
                        break
    except zipfile.BadZipFile as exc:
        raise PptxPreviewFailure("pptx_invalid_archive", 422) from exc
    return warnings


def _render_pptx(path: Path, archive_warnings: set[str]) -> dict[str, Any]:
    from pptx import Presentation

    deadline = time.monotonic() + PPTX_RENDER_TIMEOUT_SECONDS
    context = PreviewContext(deadline=deadline)
    context.warnings.update(archive_warnings)
    context.check_deadline()
    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        raise PptxPreviewFailure("pptx_parse_failed", 422) from exc

    slide_count = len(presentation.slides)
    if slide_count < 1:
        raise PptxPreviewFailure("pptx_empty_presentation", 422)
    if slide_count > MAX_PPTX_SLIDES:
        raise PptxPreviewFailure("pptx_slide_limit", 413)
    # Count every slide up front so later truncation cannot return 25,000
    # populated nodes plus additional uncounted empty slide containers.
    context.scene_nodes = slide_count
    if not context.reserve_bytes(slide_count * 160):
        raise PptxPreviewFailure("pptx_scene_size_limit", 413)

    width = _px(presentation.slide_width)
    height = _px(presentation.slide_height)
    if not math.isfinite(width) or not math.isfinite(height) or width <= 0 or height <= 0:
        raise PptxPreviewFailure("pptx_invalid_dimensions", 422)

    slides: list[dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides):
        context.check_deadline()
        slides.append(_serialize_slide(slide, index + 1, context))
    context.check_deadline()
    result = {
        "name": path.name,
        "path": str(path),
        "width": width,
        "height": height,
        "slideCount": slide_count,
        "slides": slides,
        "assets": context.assets,
        "warnings": sorted(context.warnings),
        "sceneNodeCount": context.scene_nodes,
    }
    encoded_size = len(
        json.dumps(
            result,
            ensure_ascii=False,
            separators=(",", ":"),
        ).encode("utf-8")
    )
    if encoded_size > MAX_PPTX_SCENE_JSON_BYTES:
        raise PptxPreviewFailure("pptx_scene_size_limit", 413)
    return result


def build_pptx_preview(path: Path) -> dict[str, Any]:
    """Validate and render one PPTX. This function runs in a worker thread."""

    warnings = _validate_archive(path)
    return _render_pptx(path, warnings)


def _release_render_gate_when_done(worker: asyncio.Future[Any]) -> None:
    """Release the gate only after the underlying executor job has ended."""

    try:
        # Retrieve a late worker exception after the HTTP request has already
        # timed out or disconnected, preventing an unhandled-future warning.
        worker.exception()
    except BaseException:
        pass
    finally:
        _PPTX_RENDER_GATE.release()


@router.post("/pptx-preview")
async def preview_pptx(body: PptxPreviewRequest) -> dict[str, Any]:
    """Return a safe scene model for a local OOXML ``.pptx`` file."""

    try:
        requested_path = _resolve_requested_file_path(body.path, body.workspace)
        path = requested_path.resolve(strict=True)
    except (OSError, RuntimeError, ValueError):
        raise HTTPException(status_code=404, detail="pptx_file_not_found")

    if not path.is_file():
        raise HTTPException(status_code=400, detail="pptx_not_a_file")
    suffix = path.suffix.lower()
    if suffix == ".ppt":
        raise HTTPException(status_code=415, detail="ppt_legacy_unsupported")
    if suffix != ".pptx":
        raise HTTPException(status_code=415, detail="pptx_file_required")

    try:
        size = path.stat().st_size
    except OSError:
        raise HTTPException(status_code=404, detail="pptx_file_not_found")
    if size > MAX_PPTX_FILE_BYTES:
        raise HTTPException(status_code=413, detail="pptx_file_size_limit")

    if not _PPTX_RENDER_GATE.acquire(blocking=False):
        raise HTTPException(status_code=429, detail="pptx_preview_busy")

    release_deferred = False
    worker: asyncio.Future[dict[str, Any]] | None = None
    try:
        loop = asyncio.get_running_loop()
        worker = loop.run_in_executor(
            _PPTX_RENDER_EXECUTOR,
            build_pptx_preview,
            path,
        )
        return await asyncio.wait_for(
            asyncio.shield(worker),
            timeout=PPTX_RENDER_TIMEOUT_SECONDS + PPTX_WORKER_GRACE_SECONDS,
        )
    except TimeoutError:
        if worker is not None:
            release_deferred = True
            worker.add_done_callback(_release_render_gate_when_done)
        raise HTTPException(status_code=408, detail="pptx_preview_timeout")
    except asyncio.CancelledError:
        if worker is not None:
            release_deferred = True
            worker.add_done_callback(_release_render_gate_when_done)
        raise
    except PptxPreviewFailure as exc:
        raise HTTPException(status_code=exc.status_code, detail=exc.code)
    except Exception:
        logger.exception("Unexpected PPTX preview failure for %s", path)
        raise HTTPException(status_code=422, detail="pptx_render_failed")
    finally:
        if not release_deferred:
            _PPTX_RENDER_GATE.release()
