"""Fail-closed import foundation for workspace-scoped user Office templates."""

from __future__ import annotations

import asyncio
import hashlib
import json
import os
import re
import shutil
import stat
import tempfile
import threading
from collections.abc import Awaitable, Callable, Iterable, Mapping
from dataclasses import dataclass
from pathlib import Path
from typing import Any, BinaryIO, Final, Literal, TypeAlias

from app.office_rendering.models import AUTHORITATIVE_QUALITY, RenderManifest
from app.office_templates.errors import (
    OfficeTemplateError,
    TemplateContractError,
    TemplateInUseError,
    TemplateNotFoundError,
)
from app.office_templates.instantiation import OfficeTemplateInstantiator
from app.office_templates.models import (
    AllowedOutputRules,
    TemplatePackageManifest,
    TemplateRecord,
)
from app.office_templates.registry import OfficeTemplateRegistry
from app.office_templates.substitution import (
    is_substitutable_part,
    placeholder_counts,
)
from app.office_templates.validation import (
    TemplateSafetyLimits,
    inspect_ooxml_package,
)
from app.office_validation.draft import OfficeDraftValidationService
from app.release_readiness import v11_capability_released
from app.utils.id import generate_ulid


USER_TEMPLATE_MAX_SOURCE_BYTES: Final = 75 * 1024 * 1024
USER_TEMPLATE_MAX_PLACEHOLDERS: Final = 256
USER_TEMPLATE_MAX_TOTAL_DECLARED_CHARS: Final = 1024 * 1024
USER_TEMPLATE_SCHEMA_VERSION: Final = 1
USER_TEMPLATE_VALIDATION_REPORT_VERSION: Final = 1
USER_TEMPLATE_RECONCILIATION_MAX_RECORDS: Final = 512
USER_TEMPLATE_RECONCILIATION_MAX_OWNERS: Final = 100_000
_COPY_CHUNK_BYTES: Final = 1024 * 1024

UserTemplateStatus: TypeAlias = Literal[
    "needs_confirmation",
    "needs_review",
]
UserTemplateFormat: TypeAlias = Literal["docx", "xlsx", "pptx"]
UserTemplateOwnerStatus: TypeAlias = Literal[
    "needs_confirmation",
    "needs_review",
    "approved",
    "tombstoned",
]
UserTemplateRegistryKey: TypeAlias = tuple[str, int]
UserTemplateRegistryOwnerRow: TypeAlias = tuple[str, int, str]

_PLACEHOLDER = re.compile(r"^[A-Za-z][A-Za-z0-9_.-]{0,63}$")
_TEMPLATE_REF = re.compile(r"^utpl-[0-9a-z]{26}$")
_FORMAT_BY_SUFFIX: Final[dict[str, UserTemplateFormat]] = {
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".pptx": "pptx",
}


class UserTemplateFeatureDisabledError(OfficeTemplateError):
    """The composed user-template Beta capability is not released."""


class UserTemplateReopenError(OfficeTemplateError):
    """A statically safe OOXML package failed an independent library reopen."""


class UserTemplateEvidenceError(OfficeTemplateError):
    """Stored source or render evidence no longer matches its approval record."""


@dataclass(frozen=True, slots=True)
class UserTemplateRegistryOwner:
    """One path-free row from a global user-template DB owner snapshot."""

    template_ref: str
    revision: int
    status: UserTemplateOwnerStatus

    def __post_init__(self) -> None:
        validate_user_template_ref(self.template_ref)
        if (
            not isinstance(self.revision, int)
            or isinstance(self.revision, bool)
            or self.revision < 1
        ):
            raise TemplateContractError("user template owner revision is invalid")
        if not isinstance(self.status, str) or self.status not in {
            "needs_confirmation",
            "needs_review",
            "approved",
            "tombstoned",
        }:
            raise TemplateContractError("user template owner status is invalid")

    @property
    def key(self) -> UserTemplateRegistryKey:
        return self.template_ref, self.revision


UserTemplateGlobalOwnerLoader: TypeAlias = Callable[
    [],
    Awaitable[
        Iterable[UserTemplateRegistryOwner | UserTemplateRegistryOwnerRow]
    ],
]


@dataclass(frozen=True, slots=True)
class UserTemplateReconciliationReport:
    """Path-free result of one conservative registry/DB reconciliation."""

    scanned_records: int
    owner_records: int
    retained_active: tuple[UserTemplateRegistryKey, ...]
    retained_referenced: tuple[UserTemplateRegistryKey, ...]
    deleted_orphans: tuple[UserTemplateRegistryKey, ...]
    deleted_tombstoned: tuple[UserTemplateRegistryKey, ...]

    @property
    def deleted_records(self) -> int:
        return len(self.deleted_orphans) + len(self.deleted_tombstoned)


@dataclass(frozen=True, slots=True)
class UserTemplatePlaceholder:
    """The Beta placeholder contract is deliberately required text only."""

    name: str
    value_type: Literal["text"]
    required: Literal[True]
    min_chars: int
    max_chars: int
    description: str = ""

    def __post_init__(self) -> None:
        if not isinstance(self.name, str) or _PLACEHOLDER.fullmatch(self.name) is None:
            raise TemplateContractError("user template placeholder name is invalid")
        if self.value_type != "text":
            raise TemplateContractError(
                "user template placeholders support only type text"
            )
        if self.required is not True:
            raise TemplateContractError(
                "user template placeholders must all be required"
            )
        if (
            not isinstance(self.min_chars, int)
            or isinstance(self.min_chars, bool)
            or not isinstance(self.max_chars, int)
            or isinstance(self.max_chars, bool)
            or not 0 <= self.min_chars <= self.max_chars <= 100_000
        ):
            raise TemplateContractError(
                "user template placeholder text bounds are invalid"
            )
        if (
            not isinstance(self.description, str)
            or len(self.description) > 256
            or any(ord(character) < 32 for character in self.description)
        ):
            raise TemplateContractError(
                "user template placeholder description is invalid"
            )

    def to_dict(self) -> dict[str, Any]:
        return {
            "name": self.name,
            "type": self.value_type,
            "required": self.required,
            "min_chars": self.min_chars,
            "max_chars": self.max_chars,
            "description": self.description,
        }


@dataclass(frozen=True, slots=True)
class UserTemplateImportCandidate:
    """Path-free evidence returned after registry import and local rendering."""

    template_ref: str
    revision: int
    display_name: str
    format: UserTemplateFormat
    source_sha256: str
    source_size_bytes: int
    manifest_sha256: str
    import_request_sha256: str
    placeholder_schema: tuple[UserTemplatePlaceholder, ...]
    placeholder_parts: tuple[str, ...]
    allowed_operations: tuple[str, ...]
    status: UserTemplateStatus
    render_manifest: RenderManifest
    render_manifest_sha256: str
    validation_report: Mapping[str, Any]

    def __post_init__(self) -> None:
        validate_user_template_ref(self.template_ref)
        if self.revision != 1:
            raise TemplateContractError("initial user template revision must be one")
        if self.status not in {"needs_confirmation", "needs_review"}:
            raise TemplateContractError("user template import status is invalid")
        if self.allowed_operations != ("instantiate_text",):
            raise TemplateContractError("user template operation set is invalid")
        if not isinstance(self.render_manifest, RenderManifest):
            raise TemplateContractError("user template render manifest is invalid")


def validate_user_template_ref(value: object) -> str:
    if not isinstance(value, str) or _TEMPLATE_REF.fullmatch(value) is None:
        raise TemplateContractError("user template reference is invalid")
    return value


def normalize_placeholder_schema(
    fields: tuple[UserTemplatePlaceholder, ...] | list[UserTemplatePlaceholder],
) -> tuple[UserTemplatePlaceholder, ...]:
    """Return one sorted, unique and bounded required-text schema."""

    try:
        normalized = tuple(fields)
    except TypeError as exc:
        raise TemplateContractError("user template placeholder schema is invalid") from exc
    if (
        not normalized
        or len(normalized) > USER_TEMPLATE_MAX_PLACEHOLDERS
        or any(not isinstance(field, UserTemplatePlaceholder) for field in normalized)
    ):
        raise TemplateContractError("user template placeholder schema is invalid")
    names = tuple(field.name for field in normalized)
    if len(names) != len(set(names)):
        raise TemplateContractError("user template placeholder names must be unique")
    if sum(field.max_chars for field in normalized) > USER_TEMPLATE_MAX_TOTAL_DECLARED_CHARS:
        raise TemplateContractError(
            "user template placeholder schema exceeds its total text budget"
        )
    return tuple(sorted(normalized, key=lambda field: field.name))


def decode_user_template_placeholder_schema(
    value: object,
) -> tuple[UserTemplatePlaceholder, ...]:
    """Decode the exact canonical schema persisted with an approved revision."""

    if not isinstance(value, list):
        raise TemplateContractError("stored user template schema is invalid")
    expected = {
        "name",
        "type",
        "required",
        "min_chars",
        "max_chars",
        "description",
    }
    fields: list[UserTemplatePlaceholder] = []
    for item in value:
        if not isinstance(item, dict) or set(item) != expected:
            raise TemplateContractError("stored user template schema is invalid")
        fields.append(
            UserTemplatePlaceholder(
                name=item["name"],
                value_type=item["type"],
                required=item["required"],
                min_chars=item["min_chars"],
                max_chars=item["max_chars"],
                description=item["description"],
            )
        )
    normalized = normalize_placeholder_schema(fields)
    if value != [field.to_dict() for field in normalized]:
        raise TemplateContractError("stored user template schema is not canonical")
    return normalized


def validate_user_template_values(
    schema: tuple[UserTemplatePlaceholder, ...],
    values: Mapping[str, object],
) -> dict[str, str]:
    """Enforce exact names plus each approved field's character bounds."""

    fields = normalize_placeholder_schema(schema)
    if not isinstance(values, Mapping) or any(
        not isinstance(name, str) for name in values
    ):
        raise TemplateContractError("user template values must be a field object")
    expected = {field.name for field in fields}
    provided = set(values)
    if provided != expected:
        raise TemplateContractError("user template value names do not match its schema")
    normalized: dict[str, str] = {}
    total = 0
    for field in fields:
        value = values[field.name]
        if not isinstance(value, str):
            raise TemplateContractError(
                "user template values support only text"
            )
        length = len(value)
        if not field.min_chars <= length <= field.max_chars:
            raise TemplateContractError(
                "user template value violates its approved text bounds"
            )
        total += length
        if total > USER_TEMPLATE_MAX_TOTAL_DECLARED_CHARS:
            raise TemplateContractError("user template values exceed their text budget")
        normalized[field.name] = value
    return normalized


class UserOfficeTemplateService:
    """Validate, render and register immutable user-supplied OOXML revisions."""

    def __init__(
        self,
        root: str | Path,
        *,
        draft_validation: OfficeDraftValidationService,
        limits: TemplateSafetyLimits | None = None,
    ) -> None:
        if not isinstance(draft_validation, OfficeDraftValidationService):
            raise TypeError("draft_validation must be OfficeDraftValidationService")
        try:
            private_root = Path(root).expanduser()
        except TypeError as exc:
            raise TemplateContractError("user template storage root is invalid") from exc
        if not private_root.is_absolute():
            raise TemplateContractError("user template storage root must be absolute")
        _ensure_private_root(private_root)
        self.root = private_root.resolve(strict=True)
        selected_limits = limits or TemplateSafetyLimits(
            max_package_bytes=USER_TEMPLATE_MAX_SOURCE_BYTES,
            max_entry_bytes=USER_TEMPLATE_MAX_SOURCE_BYTES,
        )
        if selected_limits.max_package_bytes > USER_TEMPLATE_MAX_SOURCE_BYTES:
            raise TemplateContractError(
                "user template package limit exceeds the Beta source budget"
            )
        self.registry = OfficeTemplateRegistry(
            self.root / "registry",
            limits=selected_limits,
        )
        self.import_root = self.root / "import-staging"
        _ensure_private_root(self.import_root)
        self._draft = draft_validation
        self._reconciliation_lock = asyncio.Lock()
        self._reconciliation_report: UserTemplateReconciliationReport | None = None

    @staticmethod
    def require_enabled() -> None:
        if not v11_capability_released("user_office_templates"):
            raise UserTemplateFeatureDisabledError(
                "user Office templates are not released"
            )

    async def validate_and_register(
        self,
        stream: BinaryIO,
        *,
        filename: str,
        display_name: str,
        placeholders: tuple[UserTemplatePlaceholder, ...],
    ) -> UserTemplateImportCandidate:
        """Create one immutable registry revision only after all import gates pass."""

        self.require_enabled()
        safe_name, format_name = _validate_filename(filename)
        normalized_display_name = _validate_display_name(display_name)
        schema = normalize_placeholder_schema(placeholders)
        template_ref = f"utpl-{generate_ulid().lower()}"
        revision = 1
        temporary = Path(
            tempfile.mkdtemp(prefix="import-", dir=self.import_root)
        ).resolve(strict=True)
        _assert_within(self.import_root, temporary)
        if os.name != "nt":
            os.chmod(temporary, 0o700)
        source = temporary / ("source" + Path(safe_name).suffix.casefold())
        registered = False
        try:
            source_sha256, source_size = await asyncio.to_thread(
                _copy_stream_to_private_file,
                stream,
                source,
                USER_TEMPLATE_MAX_SOURCE_BYTES,
            )
            content = await asyncio.to_thread(
                _read_private_regular_file,
                source,
                USER_TEMPLATE_MAX_SOURCE_BYTES,
            )
            if hashlib.sha256(content).hexdigest() != source_sha256:
                raise UserTemplateEvidenceError(
                    "user template source changed during import"
                )
            placeholder_names = tuple(field.name for field in schema)
            inspection = inspect_ooxml_package(
                content,
                format_name,
                expected_placeholders=placeholder_names,
                limits=self.registry.limits,
            )
            parts = _placeholder_parts(inspection.entries, format_name)
            await asyncio.to_thread(_independent_reopen, source, format_name)

            artifact = await self._draft.capture(
                boundary_root=temporary,
                source_path=source,
                expected_source_sha256=source_sha256,
            )
            render_manifest = artifact.manifest
            if (
                artifact.source_sha256 != source_sha256
                or render_manifest.source_sha256 != source_sha256
                or render_manifest.document_format != format_name
            ):
                raise UserTemplateEvidenceError(
                    "user template render evidence is not source-bound"
                )

            manifest = TemplatePackageManifest(
                template_id=template_ref,
                template_version=str(revision),
                format=format_name,
                source_sha256=source_sha256,
                license="User-provided template; rights not verified",
                provenance=f"local-user-import:{template_ref}",
                required_placeholders=placeholder_names,
                allowed_output_rules=AllowedOutputRules(
                    extensions=(f".{format_name}",),
                    max_output_bytes=USER_TEMPLATE_MAX_SOURCE_BYTES,
                    allow_overwrite=False,
                ),
            )
            await asyncio.to_thread(self.registry.import_template, manifest, source)
            registered = True

            render_digest = hashlib.sha256(
                render_manifest.canonical_bytes()
            ).hexdigest()
            request_digest = hashlib.sha256(
                _canonical_json(
                    {
                        "schema_version": USER_TEMPLATE_SCHEMA_VERSION,
                        "display_name": normalized_display_name,
                        "format": format_name,
                        "source_sha256": source_sha256,
                        "placeholder_schema": [field.to_dict() for field in schema],
                    }
                )
            ).hexdigest()
            status: UserTemplateStatus = (
                "needs_confirmation"
                if render_manifest.quality == AUTHORITATIVE_QUALITY
                else "needs_review"
            )
            report = {
                "schema_version": USER_TEMPLATE_VALIDATION_REPORT_VERSION,
                "ooxml_safety": "pass",
                "placeholder_contract": "pass",
                "independent_reopen": "pass",
                "render_completed": True,
                "render_quality": render_manifest.quality,
                "approval_eligible": status == "needs_confirmation",
            }
            return UserTemplateImportCandidate(
                template_ref=template_ref,
                revision=revision,
                display_name=normalized_display_name,
                format=format_name,
                source_sha256=source_sha256,
                source_size_bytes=source_size,
                manifest_sha256=manifest.template_sha256,
                import_request_sha256=request_digest,
                placeholder_schema=schema,
                placeholder_parts=parts,
                allowed_operations=("instantiate_text",),
                status=status,
                render_manifest=render_manifest,
                render_manifest_sha256=render_digest,
                validation_report=report,
            )
        except BaseException:
            if registered:
                await self.discard_orphan(template_ref, revision)
            raise
        finally:
            shutil.rmtree(temporary, ignore_errors=True)

    async def verify_approval_evidence(
        self,
        *,
        template_ref: str,
        revision: int,
        source_sha256: str,
        manifest_sha256: str,
        render_manifest_sha256: str,
        format_name: UserTemplateFormat,
        placeholder_schema: tuple[UserTemplatePlaceholder, ...],
        placeholder_parts: tuple[str, ...],
        render_cache_key: str,
        renderer_id: str,
        renderer_version: str,
        font_digest: str,
        render_parameters_version: str,
        render_parameters_sha256: str,
    ) -> RenderManifest:
        """Reopen source/render evidence before an approval state transition."""

        self.require_enabled()
        validate_user_template_ref(template_ref)
        if not isinstance(revision, int) or isinstance(revision, bool) or revision < 1:
            raise TemplateContractError("user template revision is invalid")
        record, _content, _parts = await asyncio.to_thread(
            self._verified_registry_contract,
            template_ref,
            revision,
            source_sha256,
            manifest_sha256,
            format_name,
            placeholder_schema,
            placeholder_parts,
        )
        artifact = await self._draft.capture(
            boundary_root=self.registry.root,
            source_path=record.content_path,
            expected_source_sha256=source_sha256,
        )
        manifest = artifact.manifest
        if (
            hashlib.sha256(manifest.canonical_bytes()).hexdigest()
            != render_manifest_sha256
            or manifest.cache_key != render_cache_key
            or manifest.renderer_id != renderer_id
            or manifest.renderer_version != renderer_version
            or manifest.font_digest != font_digest
            or manifest.parameters_version != render_parameters_version
            or manifest.parameters_sha256 != render_parameters_sha256
        ):
            raise UserTemplateEvidenceError(
                "user template renderer or render evidence changed before approval"
            )
        if manifest.quality != AUTHORITATIVE_QUALITY:
            raise UserTemplateEvidenceError(
                "approximate render evidence cannot authorize a user template"
            )
        return manifest

    async def preview_page_bytes(
        self,
        *,
        template_ref: str,
        revision: int,
        source_sha256: str,
        manifest_sha256: str,
        format_name: UserTemplateFormat,
        placeholder_schema: tuple[UserTemplatePlaceholder, ...],
        placeholder_parts: tuple[str, ...],
        render_manifest_sha256: str,
        render_cache_key: str,
        renderer_id: str,
        renderer_version: str,
        font_digest: str,
        render_parameters_version: str,
        render_parameters_sha256: str,
        render_quality: str,
        render_page_count: int,
        page_number: int,
    ) -> bytes:
        """Return one private PNG snapshot after the full evidence chain revalidates.

        Returning bytes is deliberate: Starlette opens ``FileResponse`` paths
        after the endpoint returns, which would leave a replace/symlink window
        after the cache and database checks.  The bounded no-follow read below
        binds the response body itself to the manifest digest instead.
        """

        self.require_enabled()
        validate_user_template_ref(template_ref)
        if (
            not isinstance(revision, int)
            or isinstance(revision, bool)
            or revision < 1
            or render_quality not in {"authoritative", "approximate"}
            or not isinstance(render_page_count, int)
            or isinstance(render_page_count, bool)
            or render_page_count < 1
            or not isinstance(page_number, int)
            or isinstance(page_number, bool)
            or not 1 <= page_number <= min(render_page_count, 1000)
        ):
            raise TemplateContractError("user template preview request is invalid")
        record, _content, _parts = await asyncio.to_thread(
            self._verified_registry_contract,
            template_ref,
            revision,
            source_sha256,
            manifest_sha256,
            format_name,
            placeholder_schema,
            placeholder_parts,
        )
        artifact = await self._draft.capture(
            boundary_root=self.registry.root,
            source_path=record.content_path,
            expected_source_sha256=source_sha256,
        )
        manifest = artifact.manifest
        if (
            hashlib.sha256(manifest.canonical_bytes()).hexdigest()
            != render_manifest_sha256
            or manifest.source_sha256 != source_sha256
            or manifest.document_format != format_name
            or manifest.cache_key != render_cache_key
            or manifest.renderer_id != renderer_id
            or manifest.renderer_version != renderer_version
            or manifest.font_digest != font_digest
            or manifest.parameters_version != render_parameters_version
            or manifest.parameters_sha256 != render_parameters_sha256
            or manifest.quality != render_quality
            or len(manifest.pages) != render_page_count
        ):
            raise UserTemplateEvidenceError(
                "user template preview evidence no longer matches"
            )
        page = manifest.pages[page_number - 1]
        if page.page_number != page_number or page.mime_type != "image/png":
            raise UserTemplateEvidenceError(
                "user template preview page evidence is invalid"
            )
        self._draft.revalidate(artifact)
        payload = await asyncio.to_thread(
            _read_private_regular_file,
            artifact.entry_path / page.filename,
            page.size_bytes,
        )
        if (
            len(payload) != page.size_bytes
            or hashlib.sha256(payload).hexdigest() != page.sha256
        ):
            raise UserTemplateEvidenceError(
                "user template preview page changed while being read"
            )
        self._draft.revalidate(artifact)
        return payload

    def build_precommit_plan(
        self,
        *,
        template_ref: str,
        revision: int,
        source_sha256: str,
        manifest_sha256: str,
        format_name: UserTemplateFormat,
        placeholder_schema: tuple[UserTemplatePlaceholder, ...],
        placeholder_parts: tuple[str, ...],
        renderer_id: str,
        renderer_version: str,
        font_digest: str,
        render_parameters_version: str,
        render_parameters_sha256: str,
    ) -> Any:
        """Build a server-owned golden plan from revalidated approval evidence."""

        self.require_enabled()
        record, _content, parts = self._verified_registry_contract(
            template_ref,
            revision,
            source_sha256,
            manifest_sha256,
            format_name,
            placeholder_schema,
            placeholder_parts,
        )
        from app.office_validation.draft import OfficeGoldenPolicy
        from app.office_validation.precommit import OfficeCreateValidationPlan
        from app.office_validation.visual import VisualDiffPolicy

        policy = OfficeGoldenPolicy(
            policy_id=f"user-template/{template_ref}/{revision}",
            template_id=template_ref,
            template_version=str(revision),
            template_manifest_sha256=manifest_sha256,
            baseline_sha256=source_sha256,
            renderer_id=renderer_id,
            renderer_version=renderer_version,
            font_digest=font_digest,
            parameters_version=render_parameters_version,
            parameters_sha256=render_parameters_sha256,
            allowed_changed_parts=parts,
            # Text substitution may legitimately reflow a whole page.  The
            # structural allow-list, page count/dimensions, renderer identity,
            # and blank-page guard remain mandatory and authoritative.
            visual=VisualDiffPolicy(
                max_outside_changed_ratio=0.65,
                max_total_changed_ratio=0.65,
                max_blank_fraction_increase=0.15,
            ),
        )
        return OfficeCreateValidationPlan(
            golden_root=self.registry.root,
            golden_path=record.content_path,
            golden_policy=policy,
            template_manifest=record.manifest,
        )

    def verify_registry_contract(
        self,
        *,
        template_ref: str,
        revision: int,
        source_sha256: str,
        manifest_sha256: str,
        format_name: UserTemplateFormat,
        placeholder_schema: tuple[UserTemplatePlaceholder, ...],
        placeholder_parts: tuple[str, ...],
    ) -> None:
        """Synchronously revalidate the bounded immutable source contract."""

        self.require_enabled()
        self._verified_registry_contract(
            template_ref,
            revision,
            source_sha256,
            manifest_sha256,
            format_name,
            placeholder_schema,
            placeholder_parts,
        )

    async def reconcile_registry_orphans_once(
        self,
        load_global_owners: UserTemplateGlobalOwnerLoader,
        *,
        max_records: int = USER_TEMPLATE_RECONCILIATION_MAX_RECORDS,
        max_owner_records: int = USER_TEMPLATE_RECONCILIATION_MAX_OWNERS,
    ) -> UserTemplateReconciliationReport:
        """Reconcile this private registry against one global DB owner snapshot.

        The caller must serialize this startup boundary with the full import
        publication transaction.  This service additionally serializes and
        caches its own first *successful* run.  All registry records are
        bounded and fully revalidated before the first deletion.  A failed or
        cancelled owner query therefore cannot be interpreted as an empty DB.
        """

        if not callable(load_global_owners):
            raise TemplateContractError("user template owner loader is invalid")
        _validate_reconciliation_limit(
            max_records,
            maximum=USER_TEMPLATE_RECONCILIATION_MAX_RECORDS,
            label="registry record",
        )
        _validate_reconciliation_limit(
            max_owner_records,
            maximum=USER_TEMPLATE_RECONCILIATION_MAX_OWNERS,
            label="owner record",
        )
        async with self._reconciliation_lock:
            if self._reconciliation_report is not None:
                return self._reconciliation_report
            try:
                loaded_owners = await load_global_owners()
            except asyncio.CancelledError:
                raise
            except Exception as exc:
                raise UserTemplateEvidenceError(
                    "global user template owners could not be loaded"
                ) from exc
            owners = _normalize_registry_owners(
                loaded_owners,
                maximum=max_owner_records,
            )
            worker = asyncio.create_task(
                asyncio.to_thread(
                    self._reconcile_registry_orphans,
                    owners,
                    max_records,
                )
            )
            cancellation: asyncio.CancelledError | None = None
            while True:
                try:
                    report = await asyncio.shield(worker)
                    break
                except asyncio.CancelledError as exc:
                    # The worker may already have deleted validated orphans.
                    # Keep the app lifecycle/once locks held until that bounded
                    # mutation settles, cache its result, and only then expose
                    # cancellation to the caller.
                    cancellation = cancellation or exc
                    if worker.done():
                        break
            try:
                report = worker.result()
            except BaseException as exc:
                if cancellation is not None:
                    raise cancellation from exc
                raise
            self._reconciliation_report = report
            if cancellation is not None:
                raise cancellation
            return report

    def _reconcile_registry_orphans(
        self,
        owners: Mapping[UserTemplateRegistryKey, UserTemplateOwnerStatus],
        max_records: int,
    ) -> UserTemplateReconciliationReport:
        """Run the bounded validate-then-delete phase under the registry lock."""

        # OfficeTemplateRegistry owns the mutation lock.  Holding it across
        # both phases ensures a retain/import/delete cannot change the fully
        # validated record set between classification and deletion.
        with self.registry._lock:  # noqa: SLF001 - same trust boundary
            _assert_registry_record_budget(self.registry.root, max_records)
            records = self.registry.list_templates()
            if len(records) > max_records:
                raise UserTemplateEvidenceError(
                    "user template registry exceeds its reconciliation budget"
                )

            validated: list[tuple[UserTemplateRegistryKey, TemplateRecord]] = []
            for record in records:
                validated.append(
                    (_validate_reconciliation_record(record), record)
                )

            retained_active: list[UserTemplateRegistryKey] = []
            retained_referenced: list[UserTemplateRegistryKey] = []
            orphan_candidates: list[UserTemplateRegistryKey] = []
            tombstoned_candidates: list[UserTemplateRegistryKey] = []
            for key, record in validated:
                status = owners.get(key)
                if status in {"needs_confirmation", "needs_review", "approved"}:
                    retained_active.append(key)
                    continue
                if record.reference_ids:
                    retained_referenced.append(key)
                    continue
                if status == "tombstoned":
                    tombstoned_candidates.append(key)
                else:
                    orphan_candidates.append(key)

            deleted_orphans: list[UserTemplateRegistryKey] = []
            deleted_tombstoned: list[UserTemplateRegistryKey] = []
            for candidates, destination in (
                (orphan_candidates, deleted_orphans),
                (tombstoned_candidates, deleted_tombstoned),
            ):
                for key in candidates:
                    try:
                        self.registry.delete(key[0], str(key[1]))
                    except TemplateInUseError:
                        # A defensive recheck in registry.delete wins over the
                        # earlier snapshot if an out-of-band retain ever appears.
                        retained_referenced.append(key)
                    except TemplateNotFoundError:
                        # Another trusted cleanup may already have completed this
                        # exact idempotent outcome.
                        continue
                    else:
                        destination.append(key)

            return UserTemplateReconciliationReport(
                scanned_records=len(validated),
                owner_records=len(owners),
                retained_active=tuple(sorted(retained_active)),
                retained_referenced=tuple(sorted(retained_referenced)),
                deleted_orphans=tuple(sorted(deleted_orphans)),
                deleted_tombstoned=tuple(sorted(deleted_tombstoned)),
            )

    def instantiate_approved(
        self,
        *,
        template_ref: str,
        revision: int,
        placeholder_schema: tuple[UserTemplatePlaceholder, ...],
        values: Mapping[str, object],
        staging_root: Path,
        output_path: Path,
    ) -> Any:
        """Instantiate only after the caller acquired a verified approval lease."""

        self.require_enabled()
        normalized = validate_user_template_values(placeholder_schema, values)
        return OfficeTemplateInstantiator(self.registry).instantiate(
            template_ref,
            str(revision),
            normalized,
            staging_root=staging_root,
            output_path=output_path,
        )

    def _verified_registry_contract(
        self,
        template_ref: str,
        revision: int,
        source_sha256: str,
        manifest_sha256: str,
        format_name: UserTemplateFormat,
        placeholder_schema: tuple[UserTemplatePlaceholder, ...],
        placeholder_parts: tuple[str, ...],
    ) -> tuple[Any, bytes, tuple[str, ...]]:
        schema = normalize_placeholder_schema(placeholder_schema)
        record, content = self.registry.read_source(template_ref, str(revision))
        manifest = record.manifest
        names = tuple(field.name for field in schema)
        if (
            manifest.template_id != template_ref
            or manifest.template_version != str(revision)
            or manifest.format != format_name
            or manifest.source_sha256 != source_sha256
            or manifest.template_sha256 != manifest_sha256
            or manifest.required_placeholders != names
            or manifest.license != "User-provided template; rights not verified"
            or manifest.provenance != f"local-user-import:{template_ref}"
            or manifest.allowed_output_rules.extensions != (f".{format_name}",)
            or manifest.allowed_output_rules.max_output_bytes
            != USER_TEMPLATE_MAX_SOURCE_BYTES
            or manifest.allowed_output_rules.allow_overwrite
        ):
            raise UserTemplateEvidenceError(
                "user template registry evidence changed before use"
            )
        inspection = inspect_ooxml_package(
            content,
            format_name,
            expected_placeholders=names,
            limits=self.registry.limits,
        )
        actual_parts = _placeholder_parts(inspection.entries, format_name)
        if actual_parts != tuple(placeholder_parts):
            raise UserTemplateEvidenceError(
                "user template placeholder evidence changed before use"
            )
        return record, content, actual_parts

    async def discard_orphan(self, template_ref: str, revision: int) -> None:
        """Best-effort removal for registry rows that never gained DB ownership."""

        try:
            await asyncio.to_thread(
                self.registry.delete,
                template_ref,
                str(revision),
            )
        except (TemplateNotFoundError, TemplateInUseError):
            return


_RUNTIME_LOCK = threading.RLock()
_RUNTIME_SERVICE: UserOfficeTemplateService | None = None


def set_user_office_template_service(
    service: UserOfficeTemplateService | None,
) -> None:
    """Install a trusted app-owned service; ``None`` restores fail-closed mode."""

    if service is not None and not isinstance(service, UserOfficeTemplateService):
        raise TypeError("user Office template service is invalid")
    global _RUNTIME_SERVICE
    with _RUNTIME_LOCK:
        _RUNTIME_SERVICE = service


def get_user_office_template_service() -> UserOfficeTemplateService | None:
    with _RUNTIME_LOCK:
        return _RUNTIME_SERVICE


def _validate_reconciliation_limit(
    value: object,
    *,
    maximum: int,
    label: str,
) -> int:
    if (
        not isinstance(value, int)
        or isinstance(value, bool)
        or not 1 <= value <= maximum
    ):
        raise TemplateContractError(
            f"user template reconciliation {label} limit is invalid"
        )
    return value


def _normalize_registry_owners(
    value: Iterable[UserTemplateRegistryOwner | UserTemplateRegistryOwnerRow],
    *,
    maximum: int,
) -> dict[UserTemplateRegistryKey, UserTemplateOwnerStatus]:
    try:
        iterator = iter(value)
    except TypeError as exc:
        raise TemplateContractError(
            "global user template owner snapshot is invalid"
        ) from exc
    owners: dict[UserTemplateRegistryKey, UserTemplateOwnerStatus] = {}
    for index, raw_owner in enumerate(iterator, start=1):
        if index > maximum:
            raise TemplateContractError(
                "global user template owner snapshot exceeds its budget"
            )
        if isinstance(raw_owner, UserTemplateRegistryOwner):
            owner = raw_owner
        elif isinstance(raw_owner, tuple) and len(raw_owner) == 3:
            try:
                owner = UserTemplateRegistryOwner(
                    template_ref=raw_owner[0],
                    revision=raw_owner[1],
                    status=raw_owner[2],  # type: ignore[arg-type]
                )
            except (TemplateContractError, TypeError) as exc:
                raise TemplateContractError(
                    "global user template owner snapshot is invalid"
                ) from exc
        else:
            raise TemplateContractError(
                "global user template owner snapshot is invalid"
            )
        if owner.key in owners:
            raise TemplateContractError(
                "global user template owner snapshot contains duplicate keys"
            )
        owners[owner.key] = owner.status
    return owners


def _assert_registry_record_budget(registry_root: Path, maximum: int) -> None:
    """Bound directory enumeration before OfficeTemplateRegistry materializes it."""

    records_root = registry_root / "records"
    parents_seen = 0
    records_seen = 0
    try:
        with os.scandir(records_root) as parents:
            for parent in parents:
                parents_seen += 1
                if parents_seen > maximum:
                    raise UserTemplateEvidenceError(
                        "user template registry exceeds its reconciliation budget"
                    )
                if not parent.is_dir(follow_symlinks=False):
                    # The registry's strict loader will reject this entry.  It
                    # still consumes budget so hostile junk cannot be unbounded.
                    continue
                with os.scandir(parent.path) as revisions:
                    for _revision in revisions:
                        records_seen += 1
                        if records_seen > maximum:
                            raise UserTemplateEvidenceError(
                                "user template registry exceeds its reconciliation budget"
                            )
    except UserTemplateEvidenceError:
        raise
    except OSError as exc:
        raise UserTemplateEvidenceError(
            "user template registry cannot be bounded for reconciliation"
        ) from exc


def _validate_reconciliation_record(
    record: TemplateRecord,
) -> UserTemplateRegistryKey:
    """Validate the exact user-owned manifest contract for one registry row."""

    if not isinstance(record, TemplateRecord):
        raise UserTemplateEvidenceError(
            "user template registry returned an invalid record"
        )
    manifest = record.manifest
    try:
        template_ref = validate_user_template_ref(manifest.template_id)
    except TemplateContractError as exc:
        raise UserTemplateEvidenceError(
            "user template registry contains a foreign template id"
        ) from exc
    raw_revision = manifest.template_version
    if re.fullmatch(r"[1-9][0-9]{0,9}", raw_revision) is None:
        raise UserTemplateEvidenceError(
            "user template registry revision is invalid"
        )
    revision = int(raw_revision)
    if revision > 2_147_483_647:
        raise UserTemplateEvidenceError(
            "user template registry revision exceeds its DB contract"
        )
    if (
        manifest.format not in _FORMAT_BY_SUFFIX.values()
        or not manifest.required_placeholders
        or manifest.license != "User-provided template; rights not verified"
        or manifest.provenance != f"local-user-import:{template_ref}"
        or manifest.allowed_output_rules.extensions != (f".{manifest.format}",)
        or manifest.allowed_output_rules.max_output_bytes
        != USER_TEMPLATE_MAX_SOURCE_BYTES
        or manifest.allowed_output_rules.allow_overwrite
    ):
        raise UserTemplateEvidenceError(
            "user template registry manifest is outside its owned contract"
        )
    return template_ref, revision


def _validate_filename(filename: str) -> tuple[str, UserTemplateFormat]:
    if not isinstance(filename, str):
        raise TemplateContractError("user template filename is invalid")
    normalized = filename.replace("\\", "/").rsplit("/", 1)[-1].strip()
    if (
        not normalized
        or len(normalized.encode("utf-8")) > 240
        or any(ord(character) < 32 for character in normalized)
    ):
        raise TemplateContractError("user template filename is invalid")
    suffix = Path(normalized).suffix.casefold()
    try:
        format_name = _FORMAT_BY_SUFFIX[suffix]
    except KeyError as exc:
        raise TemplateContractError(
            "user templates must be DOCX, XLSX, or PPTX"
        ) from exc
    return normalized, format_name


def _validate_display_name(value: str) -> str:
    if not isinstance(value, str):
        raise TemplateContractError("user template display name is invalid")
    normalized = " ".join(value.split()).strip()
    if (
        not normalized
        or len(normalized) > 160
        or "/" in normalized
        or "\\" in normalized
        or any(ord(character) < 32 for character in normalized)
    ):
        raise TemplateContractError("user template display name is invalid")
    return normalized


def _copy_stream_to_private_file(
    stream: BinaryIO,
    destination: Path,
    maximum: int,
) -> tuple[str, int]:
    try:
        stream.seek(0)
    except (AttributeError, OSError) as exc:
        raise TemplateContractError("user template upload cannot be read") from exc
    flags = os.O_WRONLY | os.O_CREAT | os.O_EXCL | getattr(os, "O_BINARY", 0)
    descriptor: int | None = None
    digest = hashlib.sha256()
    total = 0
    try:
        descriptor = os.open(destination, flags, 0o600)
        while True:
            chunk = stream.read(_COPY_CHUNK_BYTES)
            if not chunk:
                break
            if not isinstance(chunk, bytes):
                raise TemplateContractError("user template upload returned invalid bytes")
            total += len(chunk)
            if total > maximum:
                raise TemplateContractError("user template exceeds the source byte budget")
            digest.update(chunk)
            view = memoryview(chunk)
            while view:
                written = os.write(descriptor, view)
                if written <= 0:
                    raise OSError("short write")
                view = view[written:]
        if total < 1:
            raise TemplateContractError("user template upload is empty")
        os.fsync(descriptor)
        os.close(descriptor)
        descriptor = None
        if os.name != "nt":
            os.chmod(destination, 0o600, follow_symlinks=False)
        return digest.hexdigest(), total
    except TemplateContractError:
        raise
    except OSError as exc:
        raise UserTemplateEvidenceError(
            "user template upload could not be staged privately"
        ) from exc
    finally:
        if descriptor is not None:
            os.close(descriptor)
        if total < 1 or total > maximum:
            try:
                destination.unlink()
            except FileNotFoundError:
                pass


def _read_private_regular_file(path: Path, maximum: int) -> bytes:
    if path.is_symlink():
        raise UserTemplateEvidenceError("user template staging file is redirected")
    flags = os.O_RDONLY | getattr(os, "O_BINARY", 0)
    if hasattr(os, "O_NOFOLLOW"):
        flags |= os.O_NOFOLLOW
    try:
        descriptor = os.open(path, flags)
    except OSError as exc:
        raise UserTemplateEvidenceError("user template staging file is unavailable") from exc
    try:
        before = os.fstat(descriptor)
        if not stat.S_ISREG(before.st_mode) or not 1 <= before.st_size <= maximum:
            raise UserTemplateEvidenceError("user template staging file is invalid")
        chunks: list[bytes] = []
        total = 0
        while True:
            chunk = os.read(descriptor, min(_COPY_CHUNK_BYTES, maximum + 1 - total))
            if not chunk:
                break
            total += len(chunk)
            if total > maximum:
                raise UserTemplateEvidenceError(
                    "user template staging file exceeds its budget"
                )
            chunks.append(chunk)
        after = os.fstat(descriptor)
        if (
            before.st_dev != after.st_dev
            or before.st_ino != after.st_ino
            or before.st_size != after.st_size
            or before.st_mtime_ns != after.st_mtime_ns
            or total != after.st_size
        ):
            raise UserTemplateEvidenceError(
                "user template staging file changed while reading"
            )
        return b"".join(chunks)
    except OSError as exc:
        raise UserTemplateEvidenceError("user template staging file cannot be read") from exc
    finally:
        os.close(descriptor)


def _placeholder_parts(
    entries: Mapping[str, bytes],
    format_name: UserTemplateFormat,
) -> tuple[str, ...]:
    parts: list[str] = []
    for part_name, payload in entries.items():
        if not is_substitutable_part(part_name, format_name):
            continue
        if placeholder_counts(part_name, payload, format_name):
            parts.append(part_name)
    if not parts:
        raise TemplateContractError(
            "user template contains no supported placeholder locations"
        )
    return tuple(sorted(parts))


def _independent_reopen(path: Path, format_name: UserTemplateFormat) -> None:
    try:
        if format_name == "docx":
            from docx import Document

            Document(str(path))
        elif format_name == "xlsx":
            from openpyxl import load_workbook

            workbook = load_workbook(
                str(path),
                read_only=False,
                data_only=False,
                keep_vba=False,
                keep_links=False,
            )
            workbook.close()
        else:
            from pptx import Presentation

            Presentation(str(path))
    except Exception as exc:
        raise UserTemplateReopenError(
            "user template failed independent reopen validation"
        ) from exc


def _canonical_json(value: object) -> bytes:
    return json.dumps(
        value,
        ensure_ascii=False,
        allow_nan=False,
        separators=(",", ":"),
        sort_keys=True,
    ).encode("utf-8")


def _ensure_private_root(path: Path) -> None:
    if path.is_symlink():
        raise UserTemplateEvidenceError("user template storage root is redirected")
    try:
        path.mkdir(mode=0o700, parents=True, exist_ok=True)
    except OSError as exc:
        raise UserTemplateEvidenceError(
            "user template storage root is unavailable"
        ) from exc
    if path.is_symlink() or not path.is_dir():
        raise UserTemplateEvidenceError("user template storage root is invalid")
    if os.name != "nt":
        try:
            os.chmod(path, 0o700, follow_symlinks=False)
        except OSError as exc:
            raise UserTemplateEvidenceError(
                "user template storage root cannot be made private"
            ) from exc


def _assert_within(boundary: Path, candidate: Path) -> None:
    try:
        candidate.resolve(strict=True).relative_to(boundary.resolve(strict=True))
    except (OSError, RuntimeError, ValueError) as exc:
        raise UserTemplateEvidenceError(
            "user template staging directory escaped its private root"
        ) from exc


__all__ = [
    "USER_TEMPLATE_MAX_SOURCE_BYTES",
    "USER_TEMPLATE_RECONCILIATION_MAX_OWNERS",
    "USER_TEMPLATE_RECONCILIATION_MAX_RECORDS",
    "USER_TEMPLATE_SCHEMA_VERSION",
    "UserOfficeTemplateService",
    "UserTemplateEvidenceError",
    "UserTemplateFeatureDisabledError",
    "UserTemplateGlobalOwnerLoader",
    "UserTemplateImportCandidate",
    "UserTemplateReconciliationReport",
    "UserTemplateRegistryOwner",
    "UserTemplateRegistryOwnerRow",
    "UserTemplatePlaceholder",
    "UserTemplateReopenError",
    "decode_user_template_placeholder_schema",
    "get_user_office_template_service",
    "normalize_placeholder_schema",
    "set_user_office_template_service",
    "validate_user_template_values",
    "validate_user_template_ref",
]
