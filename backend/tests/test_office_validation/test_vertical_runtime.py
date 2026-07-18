"""Vertical proof for signed Office authoring, versioning, and rewind.

The renderer is intentionally a deterministic authoritative test double: this
test proves production orchestration and evidence binding, not the fidelity of
the native renderer deployment.
"""

from __future__ import annotations

import hashlib
from pathlib import Path

import pytest
from pptx import Presentation
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession, async_sessionmaker

from app import release_features
from app.models.checkpoint_change import CheckpointChange
from app.models.message import Message
from app.models.session import Session
from app.models.session_checkpoint import SessionCheckpoint
from app.office_rendering import OfficeRenderCache, RendererDescriptor
from app.office_rendering.models import canonical_json_bytes
from app.office_templates.policies import FirstPartyOfficePrecommitPolicyResolver
from app.office_validation import (
    DeterministicOfficePrecommitCoordinator,
    OfficeDraftValidationService,
)
from app.runtime.checkpoint_runtime import (
    TurnCheckpointBinding,
    admit_turn_checkpoint,
    finish_turn_checkpoint,
    record_tool_checkpoint_effects,
)
from app.runtime.rewind import RewindService
from app.schemas.agent import AgentInfo
from app.storage.file_versions import FileVersionStore
from app.streaming.manager import GenerationJob, StreamManager
from app.tool import workspace_transaction as transaction_module
from app.tool.builtin.office import OfficeTool
from app.tool.context import ToolContext
from tests.test_office_rendering.helpers import FakeProvider


pytestmark = [
    pytest.mark.asyncio,
    pytest.mark.skipif(
        transaction_module.guarded_file_mutation_unavailable_reason() is not None,
        reason="guarded mutation primitive unavailable",
    ),
]


_PARAMETERS = {"dpi": 144}
_PARAMETERS_VERSION = "vertical-precommit-v1"
_PROJECT_NAME = "v1.1 vertical Office proof"
_INITIAL_TITLE = _PROJECT_NAME
_EDITED_TITLE = "v1.1 vertical Office proof (rewind me)"


async def _admit(
    session_factory: async_sessionmaker[AsyncSession],
    workspace: Path,
    *,
    stream_id: str,
    call_id: str,
    coordinator: DeterministicOfficePrecommitCoordinator,
) -> tuple[GenerationJob, TurnCheckpointBinding, ToolContext]:
    anchor_id = f"anchor-{stream_id}"
    async with session_factory() as db:
        async with db.begin():
            db.add(
                Message(
                    id=anchor_id,
                    session_id="office-vertical-session",
                    data={"role": "user"},
                )
            )
    job = GenerationJob(
        stream_id,
        "office-vertical-session",
        invocation_source="desktop",
        invocation_source_id="desktop",
    )
    binding = await admit_turn_checkpoint(
        session_factory,
        job=job,
        workspace=str(workspace),
        request_message_id=anchor_id,
        todo_snapshot=[],
    )
    assert binding is not None
    context = ToolContext(
        session_id=job.session_id,
        message_id=f"message-{stream_id}",
        agent=AgentInfo(name="test", description="", mode="primary"),
        call_id=call_id,
        workspace=str(workspace),
        root_turn_id=binding.root_turn_id,
        turn_run_id=binding.turn_run_id,
        checkpoint_id=binding.checkpoint_id,
        workspace_instance_id=binding.workspace_instance_id,
    )
    context._app_state = {  # type: ignore[attr-defined]
        "office_precommit_coordinator": coordinator,
    }
    return job, binding, context


async def _persist_and_finish(
    session_factory: async_sessionmaker[AsyncSession],
    *,
    job: GenerationJob,
    binding: TurnCheckpointBinding,
    call_id: str,
    metadata: dict[str, object],
    expected_changes: int = 1,
) -> None:
    recorded = await record_tool_checkpoint_effects(
        session_factory,
        job=job,
        binding=binding,
        tool_id="office",
        call_id=call_id,
        metadata=metadata,
    )
    assert recorded == expected_changes
    await finish_turn_checkpoint(
        session_factory,
        job=job,
        binding=binding,
        status="completed",
        response_message_id=None,
    )


def _presentation_text(path: Path) -> str:
    presentation = Presentation(path)
    return "\n".join(
        getattr(shape, "text", "")
        for slide in presentation.slides
        for shape in slide.shapes
    )


@pytest.mark.parametrize(
    ("filename", "payload"),
    (
        (
            "ordinary.docx",
            {"document": {"title": "Ordinary authoritative document"}},
        ),
        (
            "ordinary.xlsx",
            {
                "workbook": {
                    "sheets": [
                        {"name": "Data", "rows": [["Metric", "Value"], ["A", 1]]}
                    ]
                }
            },
        ),
        (
            "ordinary.pptx",
            {
                "presentation": {
                    "slides": [{"title": "Ordinary authoritative deck"}]
                }
            },
        ),
    ),
)
async def test_ordinary_create_uses_production_authoritative_policy_for_all_formats(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    filename: str,
    payload: dict[str, object],
) -> None:
    for gate in (
        "V11_CHECKPOINTS_RELEASED",
        "V11_REWIND_RELEASED",
        "V11_VALIDATION_AGENT_RELEASED",
        "V11_OFFICE_V2_RELEASED",
    ):
        monkeypatch.setattr(release_features, gate, True)
    private = (tmp_path / "private").resolve()
    workspace = (tmp_path / "workspace").resolve()
    workspace.mkdir()
    monkeypatch.setenv("SUXIAOYOU_PRIVATE_DATA_DIR", str(private))
    descriptor = RendererDescriptor(
        renderer_id="fake-attested-ordinary-renderer",
        renderer_version="1",
        font_digest="f" * 64,
        quality="authoritative",
    )
    provider = FakeProvider(descriptor)
    service = OfficeDraftValidationService(
        cache=OfficeRenderCache((private / "render-cache").resolve()),
        provider=provider,
        parameters_version=_PARAMETERS_VERSION,
        parameters=_PARAMETERS,
    )
    resolver = FirstPartyOfficePrecommitPolicyResolver(
        registry_root=(private / "policy-registry").resolve(),
        renderer=descriptor,
        parameters_version=_PARAMETERS_VERSION,
        parameters_sha256=hashlib.sha256(
            canonical_json_bytes(_PARAMETERS)
        ).hexdigest(),
    )
    coordinator = DeterministicOfficePrecommitCoordinator(
        service=service,
        policies=resolver,
    )
    context = ToolContext(
        session_id="ordinary-session",
        message_id=f"ordinary-{filename}",
        agent=AgentInfo(name="test", description="", mode="primary"),
        call_id=f"ordinary-{filename}-call",
        workspace=str(workspace),
        root_turn_id="ordinary-root-turn",
        turn_run_id=f"ordinary-{filename}-run",
        checkpoint_id=f"ordinary-{filename}-checkpoint",
        workspace_instance_id="ordinary-workspace",
    )
    context._app_state = {  # type: ignore[attr-defined]
        "office_precommit_coordinator": coordinator,
    }

    result = await OfficeTool().execute(
        {
            "file_path": filename,
            "operation": "create",
            **payload,
        },
        context,
    )

    assert result.success, result.error
    assert result.metadata["office_visual_validation"] == "authoritative"
    assert result.metadata["office_validation_checkpoint_id"] == (
        context.checkpoint_id
    )
    assert (workspace / "suxiaoyou_written" / filename).is_file()
    assert provider.calls == 1


async def test_signed_template_edit_version_and_rewind_share_one_evidence_chain(
    session_factory: async_sessionmaker[AsyncSession],
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    for gate in (
        "V11_CHECKPOINTS_RELEASED",
        "V11_REWIND_RELEASED",
        "V11_VALIDATION_AGENT_RELEASED",
        "V11_OFFICE_V2_RELEASED",
    ):
        monkeypatch.setattr(release_features, gate, True)
    private = (tmp_path / "private").resolve()
    monkeypatch.setenv("SUXIAOYOU_PRIVATE_DATA_DIR", str(private))
    workspace = (tmp_path / "workspace").resolve()
    workspace.mkdir()
    async with session_factory() as db:
        async with db.begin():
            db.add(
                Session(
                    id="office-vertical-session",
                    directory=str(workspace),
                    title="Office vertical proof",
                    version="1.1.0",
                )
            )

    descriptor = RendererDescriptor(
        renderer_id="fake-attested-vertical-renderer",
        renderer_version="1",
        font_digest="f" * 64,
        quality="authoritative",
    )
    provider = FakeProvider(descriptor)
    cache = OfficeRenderCache((private / "render-cache").resolve())
    service = OfficeDraftValidationService(
        cache=cache,
        provider=provider,
        parameters_version=_PARAMETERS_VERSION,
        parameters=_PARAMETERS,
    )
    resolver = FirstPartyOfficePrecommitPolicyResolver(
        registry_root=(private / "first-party-policy-registry").resolve(),
        renderer=descriptor,
        parameters_version=_PARAMETERS_VERSION,
        parameters_sha256=hashlib.sha256(
            canonical_json_bytes(_PARAMETERS)
        ).hexdigest(),
    )
    coordinator = DeterministicOfficePrecommitCoordinator(
        service=service,
        policies=resolver,
    )
    tool = OfficeTool()
    target = workspace / "suxiaoyou_written" / "vertical-status.pptx"

    create_job, create_binding, create_context = await _admit(
        session_factory,
        workspace,
        stream_id="office-create",
        call_id="office-create-call",
        coordinator=coordinator,
    )
    created = await tool.execute(
        {
            "file_path": "vertical-status.pptx",
            "operation": "create",
            "first_party_template": {
                "template_id": "status-update",
                "template_version": "1.0.0",
                "values": {
                    "next_step": "Exercise the rewind boundary",
                    "owner": "Office runtime",
                    "period": "2026 Q3",
                    "project_name": _PROJECT_NAME,
                    "status": "Vertical validation",
                    "summary": (
                        "The signed template, authoritative precommit contract, "
                        "file version, and rewind ledger share one test chain."
                    ),
                },
            },
        },
        create_context,
    )
    assert created.success, created.error
    assert created.metadata["office_visual_validation"] == "authoritative"
    assert created.metadata["office_validation_checkpoint_id"] == (
        create_binding.checkpoint_id
    )
    assert created.metadata["first_party_template"] is True
    assert created.metadata["template_id"] == "status-update"
    assert "_office_validation_report" in created.metadata
    create_digest = hashlib.sha256(target.read_bytes()).hexdigest()
    create_bytes = target.read_bytes()
    await _persist_and_finish(
        session_factory,
        job=create_job,
        binding=create_binding,
        call_id=create_context.call_id,
        metadata=created.metadata,
        # The first write also creates the conventional output directory.
        expected_changes=2,
    )
    assert "_office_validation_report" not in created.metadata

    edit_job, edit_binding, edit_context = await _admit(
        session_factory,
        workspace,
        stream_id="office-edit",
        call_id="office-edit-call",
        coordinator=coordinator,
    )
    edited = await tool.execute(
        {
            "file_path": "vertical-status.pptx",
            "operation": "edit",
            "replacements": [
                {
                    "old_text": _INITIAL_TITLE,
                    "new_text": _EDITED_TITLE,
                }
            ],
        },
        edit_context,
    )
    assert edited.success, edited.error
    assert edited.metadata["office_visual_validation"] == "authoritative"
    assert edited.metadata["office_validation_checkpoint_id"] == (
        edit_binding.checkpoint_id
    )
    previous_version_id = edited.metadata["previous_version_id"]
    assert isinstance(previous_version_id, str)
    assert edited.metadata["previous_sha256"] == create_digest
    assert hashlib.sha256(target.read_bytes()).hexdigest() != create_digest
    assert _EDITED_TITLE in _presentation_text(target)
    await _persist_and_finish(
        session_factory,
        job=edit_job,
        binding=edit_binding,
        call_id=edit_context.call_id,
        metadata=edited.metadata,
    )

    async with session_factory() as db:
        changes = list(
            (
                await db.execute(
                    select(CheckpointChange).order_by(CheckpointChange.time_created)
                )
            ).scalars()
        )
    file_changes = [
        change
        for change in changes
        if change.relative_path == "suxiaoyou_written/vertical-status.pptx"
    ]
    assert [(change.operation, change.relative_path) for change in file_changes] == [
        ("created", "suxiaoyou_written/vertical-status.pptx"),
        ("modified", "suxiaoyou_written/vertical-status.pptx"),
    ]
    assert any(
        change.operation == "created"
        and change.node_kind == "directory"
        and change.relative_path == "suxiaoyou_written"
        for change in changes
    )
    assert file_changes[0].after_sha256 == create_digest
    assert file_changes[0].details["office_validation"]["candidate_sha256"] == (
        create_digest
    )
    assert file_changes[1].before_version_id == previous_version_id
    assert file_changes[1].details["office_validation"]["checkpoint_id"] == (
        edit_binding.checkpoint_id
    )
    previous = FileVersionStore(workspace).get_version(previous_version_id)
    assert previous.sha256 == create_digest
    assert previous.relative_path == "suxiaoyou_written/vertical-status.pptx"

    rewind = RewindService(session_factory, stream_manager=StreamManager())
    preview = await rewind.preview(
        session_id="office-vertical-session",
        workspace_instance_id=edit_binding.workspace_instance_id,
        checkpoint_id=edit_binding.checkpoint_id,
    )
    assert preview.can_execute
    assert preview.affected_checkpoint_ids == (edit_binding.checkpoint_id,)
    assert [(path.relative_path, path.action) for path in preview.paths] == [
        ("suxiaoyou_written/vertical-status.pptx", "restore_file")
    ]
    result = await rewind.execute(
        session_id="office-vertical-session",
        workspace_instance_id=edit_binding.workspace_instance_id,
        checkpoint_id=edit_binding.checkpoint_id,
    )
    assert not result.already_rewound
    assert target.read_bytes() == create_bytes
    assert _INITIAL_TITLE in _presentation_text(target)
    assert _EDITED_TITLE not in _presentation_text(target)
    async with session_factory() as db:
        create_checkpoint = await db.get(
            SessionCheckpoint,
            create_binding.checkpoint_id,
        )
        edit_checkpoint = await db.get(
            SessionCheckpoint,
            edit_binding.checkpoint_id,
        )
    assert create_checkpoint is not None and create_checkpoint.state == "finalized"
    assert edit_checkpoint is not None and edit_checkpoint.state == "rewound"
    restored = edit_checkpoint.details["rewind_result"]["restored_paths"]
    assert restored == [
        {
            "relative_path": "suxiaoyou_written/vertical-status.pptx",
            "exists": True,
            "node_kind": "file",
            "sha256": create_digest,
            "mode": restored[0]["mode"],
            "size": len(create_bytes),
        }
    ]
    # The golden and candidate create renders, followed by the edit candidate,
    # all flowed through one provider/cache identity.  The edit baseline reused
    # the exact create candidate render rather than minting new evidence.
    assert provider.calls == 3
