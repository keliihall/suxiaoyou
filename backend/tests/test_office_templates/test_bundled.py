from __future__ import annotations

import base64
import hashlib
import json
import shutil
import zipfile
from importlib import resources
from pathlib import Path
from typing import Any, Callable

import pytest
from cryptography.hazmat.primitives import serialization
from cryptography.hazmat.primitives.asymmetric.ed25519 import Ed25519PrivateKey
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from app import release_features
from app.office_templates import (
    BundledOfficeTemplateCatalog,
    BundledOfficeTemplateService,
    TemplateContractError,
    TemplateFeatureDisabledError,
    TemplateIntegrityError,
)
from tests.test_office_templates.helpers import rewrite_zip


KEY_ID = "test-office-template-key"


def _canonical_json(value: object) -> bytes:
    return (
        json.dumps(
            value,
            ensure_ascii=False,
            allow_nan=False,
            separators=(",", ":"),
            sort_keys=True,
        ).encode("utf-8")
        + b"\n"
    )


def _packaged_asset_root() -> Path:
    root = resources.files("app.office_templates").joinpath("assets")
    assert isinstance(root, Path)
    return root.resolve(strict=True)


def _copy_assets(tmp_path: Path) -> Path:
    destination = tmp_path / "assets"
    shutil.copytree(_packaged_asset_root(), destination)
    return destination.resolve(strict=True)


def _resign_catalog(
    asset_root: Path,
    mutate: Callable[[dict[str, Any]], None],
) -> bytes:
    catalog_path = asset_root / "catalog.json"
    catalog = json.loads(catalog_path.read_text(encoding="utf-8"))
    mutate(catalog)
    catalog_bytes = _canonical_json(catalog)
    private = Ed25519PrivateKey.generate()
    public = private.public_key().public_bytes(
        encoding=serialization.Encoding.Raw,
        format=serialization.PublicFormat.Raw,
    )
    catalog_path.write_bytes(catalog_bytes)
    signature = {
        "schema_version": 1,
        "algorithm": "Ed25519",
        "key_id": KEY_ID,
        "catalog_sha256": hashlib.sha256(catalog_bytes).hexdigest(),
        "signature": base64.b64encode(private.sign(catalog_bytes)).decode("ascii"),
    }
    (asset_root / "catalog.sig.json").write_bytes(_canonical_json(signature))
    return public


def _docx_text(document: Document) -> str:
    values = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                values.append(cell.text)
    for section in document.sections:
        values.extend(paragraph.text for paragraph in section.header.paragraphs)
        values.extend(paragraph.text for paragraph in section.footer.paragraphs)
    return "\n".join(values)


def _pptx_text(presentation: Presentation) -> str:
    values: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                values.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values.extend(cell.text for cell in row.cells)
    return "\n".join(values)


def test_packaged_signed_catalog_contains_three_real_deterministic_assets() -> None:
    catalog = BundledOfficeTemplateCatalog()

    descriptors = catalog.list_templates()

    assert [item.immutable_key for item in descriptors] == [
        ("business-brief", "1.0.0"),
        ("project-tracker", "1.0.0"),
        ("status-update", "1.0.0"),
    ]
    assert {item.manifest.format for item in descriptors} == {
        "docx",
        "xlsx",
        "pptx",
    }
    assert all(item.allowed_operations == ("instantiate_text",) for item in descriptors)
    root = _packaged_asset_root()
    assert (root / "catalog.json").is_file()
    assert (root / "catalog.sig.json").is_file()
    for descriptor in descriptors:
        selected, content = catalog.read_template(*descriptor.immutable_key)
        assert selected == descriptor
        assert hashlib.sha256(content).hexdigest() == descriptor.manifest.source_sha256
        assert root.joinpath(*descriptor.asset_path.split("/")).is_file()
        with zipfile.ZipFile(root / descriptor.asset_path) as archive:
            infos = archive.infolist()
            assert infos
            assert [info.filename for info in infos] == sorted(
                info.filename for info in infos
            )
            assert all(info.date_time == (1980, 1, 1, 0, 0, 0) for info in infos)
            assert all(not info.is_dir() for info in infos)


def test_service_gate_is_released_by_default_and_remains_dynamic(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    service = BundledOfficeTemplateService(tmp_path / "registry")
    assert release_features.V11_OFFICE_V2_RELEASED is True
    assert len(service.list_templates()) == 3

    monkeypatch.setattr(release_features, "V11_OFFICE_V2_RELEASED", False)
    staging = tmp_path / "staging"
    staging.mkdir()
    output = staging / "closed.docx"
    with pytest.raises(TemplateFeatureDisabledError, match="not released"):
        service.instantiate(
            "business-brief",
            "1.0.0",
            {},
            staging_root=staging,
            output_path=output,
        )
    assert not output.exists()


def test_all_first_party_templates_instantiate_as_copies_and_reopen_independently(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(release_features, "V11_OFFICE_V2_RELEASED", True)
    asset_root = _packaged_asset_root()
    before = {
        path.name: hashlib.sha256(path.read_bytes()).hexdigest()
        for path in (asset_root / "templates").iterdir()
    }
    service = BundledOfficeTemplateService(tmp_path / "registry")
    staging = tmp_path / "staging"
    staging.mkdir()

    docx_output = staging / "brief.docx"
    docx_result = service.instantiate(
        "business-brief",
        "1.0.0",
        {
            "classification": "内部资料",
            "next_step": "完成高保真预览验收",
            "owner": "交付团队",
            "recipient": "产品委员会",
            "report_date": "2026-07-17",
            "summary": "苏小有 v1.1 已进入 Office 视觉闭环验证阶段。",
            "title": "v1.1 Office 交付简报",
        },
        staging_root=staging,
        output_path=docx_output,
    )
    reopened_docx = Document(docx_output)
    docx_text = _docx_text(reopened_docx)
    assert "v1.1 Office 交付简报" in docx_text
    assert "苏小有 v1.1 已进入 Office 视觉闭环验证阶段。" in docx_text
    assert "内部资料" in docx_text

    xlsx_output = staging / "tracker.xlsx"
    xlsx_result = service.instantiate(
        "project-tracker",
        "1.0.0",
        {
            "owner": "Office 小组",
            "project_name": "苏小有 v1.1",
            "report_date": "2026-07-17",
        },
        staging_root=staging,
        output_path=xlsx_output,
    )
    reopened_xlsx = load_workbook(xlsx_output, data_only=False)
    sheet = reopened_xlsx["项目跟踪"]
    assert sheet["A1"].value == "苏小有 v1.1 项目跟踪"
    assert sheet["H3"].value == "=AVERAGE(F5:F7)"
    assert sheet.freeze_panes == "A5"
    assert len(sheet._charts) == 1
    assert len(sheet.data_validations.dataValidation) == 1
    assert sum(
        len(rules) for rules in sheet.conditional_formatting._cf_rules.values()
    ) == 2

    pptx_output = staging / "update.pptx"
    pptx_result = service.instantiate(
        "status-update",
        "1.0.0",
        {
            "next_step": "完成五平台抽检",
            "owner": "Office 小组",
            "period": "2026 Q3",
            "project_name": "苏小有 v1.1",
            "status": "视觉验证中",
            "summary": "高保真预览、签名模板和 rewind 已纳入同一版本链路。",
        },
        staging_root=staging,
        output_path=pptx_output,
    )
    reopened_pptx = Presentation(pptx_output)
    pptx_text = _pptx_text(reopened_pptx)
    assert len(reopened_pptx.slides) == 1
    assert "苏小有 v1.1" in pptx_text
    assert "视觉验证中" in pptx_text
    assert "完成五平台抽检" in pptx_text
    assert "{{" not in docx_text + pptx_text

    assert docx_result.output_path == docx_output.resolve(strict=True)
    assert xlsx_result.output_path == xlsx_output.resolve(strict=True)
    assert pptx_result.output_path == pptx_output.resolve(strict=True)
    assert all(
        result.source_sha256 != result.output_sha256
        for result in (docx_result, xlsx_result, pptx_result)
    )
    for descriptor in service.list_templates():
        registry_record = service.registry.read(*descriptor.immutable_key)
        assert registry_record.content_path != asset_root / descriptor.asset_path
        assert hashlib.sha256(registry_record.content_path.read_bytes()).hexdigest() == (
            descriptor.manifest.source_sha256
        )
    after = {
        path.name: hashlib.sha256(path.read_bytes()).hexdigest()
        for path in (asset_root / "templates").iterdir()
    }
    assert after == before


def test_placeholder_schema_rejects_type_length_and_missing_values(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(release_features, "V11_OFFICE_V2_RELEASED", True)
    service = BundledOfficeTemplateService(tmp_path / "registry")
    staging = tmp_path / "staging"
    staging.mkdir()
    base = {
        "owner": "Office 小组",
        "project_name": "苏小有 v1.1",
        "report_date": "2026-07-17",
    }

    with pytest.raises(TemplateContractError, match="type text"):
        service.instantiate(
            "project-tracker",
            "1.0.0",
            {**base, "owner": 7},
            staging_root=staging,
            output_path=staging / "type.xlsx",
        )
    with pytest.raises(TemplateContractError, match="length contract"):
        service.instantiate(
            "project-tracker",
            "1.0.0",
            {**base, "project_name": "x" * 121},
            staging_root=staging,
            output_path=staging / "length.xlsx",
        )
    with pytest.raises(TemplateContractError, match="missing required"):
        service.instantiate(
            "project-tracker",
            "1.0.0",
            {key: value for key, value in base.items() if key != "owner"},
            staging_root=staging,
            output_path=staging / "missing.xlsx",
        )
    assert not tuple(staging.iterdir())


def test_catalog_asset_signature_missing_and_digest_tampering_fail_closed(
    tmp_path: Path,
) -> None:
    catalog_tamper = _copy_assets(tmp_path / "catalog-tamper")
    catalog_path = catalog_tamper / "catalog.json"
    catalog_path.write_bytes(catalog_path.read_bytes().replace(b"1.0.0", b"1.0.1", 1))
    with pytest.raises(TemplateIntegrityError, match="digest|trusted"):
        BundledOfficeTemplateCatalog(catalog_tamper).list_templates()

    asset_tamper = _copy_assets(tmp_path / "asset-tamper")
    asset_path = asset_tamper / "templates" / "business-brief.docx"
    asset_path.write_bytes(asset_path.read_bytes() + b"tampered")
    with pytest.raises(TemplateIntegrityError, match="digest"):
        BundledOfficeTemplateCatalog(asset_tamper).list_templates()

    missing_asset = _copy_assets(tmp_path / "missing")
    (missing_asset / "templates" / "project-tracker.xlsx").unlink()
    with pytest.raises(TemplateIntegrityError, match="missing"):
        BundledOfficeTemplateCatalog(missing_asset).list_templates()

    missing_signature = _copy_assets(tmp_path / "signature")
    (missing_signature / "catalog.sig.json").unlink()
    with pytest.raises(TemplateIntegrityError, match="missing"):
        BundledOfficeTemplateCatalog(missing_signature).list_templates()


@pytest.mark.parametrize(
    "mutate",
    [
        lambda catalog: catalog.__setitem__("schema_version", True),
        lambda catalog: catalog["templates"][0]["placeholders"][0].__setitem__(
            "max_chars", "40"
        ),
        lambda catalog: catalog["templates"][0].__setitem__("format", "xlsx"),
        lambda catalog: catalog["templates"][0]["allowed_operations"].append(
            "run_macro"
        ),
    ],
    ids=["boolean-schema", "placeholder-type", "format-mismatch", "operation-expansion"],
)
def test_validly_signed_catalog_contract_errors_still_fail_closed(
    tmp_path: Path,
    mutate: Callable[[dict[str, Any]], None],
) -> None:
    root = _copy_assets(tmp_path)
    public = _resign_catalog(root, mutate)

    with pytest.raises(TemplateIntegrityError, match="schema|contract"):
        BundledOfficeTemplateCatalog(
            root,
            trusted_public_keys={KEY_ID: public},
        ).list_templates()


def test_validly_signed_dangerous_ooxml_is_rejected_after_digest_rebinding(
    tmp_path: Path,
) -> None:
    root = _copy_assets(tmp_path)
    asset = root / "templates" / "business-brief.docx"
    malicious = rewrite_zip(
        asset.read_bytes(),
        additions={"word/vbaProject.bin": b"macro"},
    )
    asset.write_bytes(malicious)

    def mutate(catalog: dict[str, Any]) -> None:
        catalog["templates"][0]["source_sha256"] = hashlib.sha256(
            malicious
        ).hexdigest()

    public = _resign_catalog(root, mutate)

    with pytest.raises(TemplateIntegrityError, match="safety validation"):
        BundledOfficeTemplateCatalog(
            root,
            trusted_public_keys={KEY_ID: public},
        ).list_templates()


def test_unknown_signing_key_is_not_accepted_even_for_a_valid_signature(
    tmp_path: Path,
) -> None:
    root = _copy_assets(tmp_path)
    _resign_catalog(root, lambda _catalog: None)

    with pytest.raises(TemplateIntegrityError, match="not trusted"):
        BundledOfficeTemplateCatalog(root).list_templates()
