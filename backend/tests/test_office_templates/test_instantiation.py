from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pytest
from docx import Document
from lxml import etree
from openpyxl import load_workbook
from pptx import Presentation

from app.office_templates import (
    OfficeTemplateInstantiator,
    OfficeTemplateRegistry,
    TemplateContractError,
    TemplateInstantiationError,
    inspect_ooxml_package,
)
from tests.test_office_templates.helpers import (
    manifest_for,
    write_source,
    zip_entries,
)


DOCX_PLACEHOLDERS = ("body", "client", "footer", "header", "table")
DOCX_VALUES = {
    "body": "这是高保真正文",
    "client": "苏小有客户",
    "footer": "内部资料",
    "header": "2026 年度报告",
    "table": "华东区",
}


def _setup(
    tmp_path: Path,
    content: bytes,
    format_name: str,
    placeholders: tuple[str, ...],
    *,
    max_output_bytes: int = 10 * 1024 * 1024,
) -> tuple[OfficeTemplateRegistry, OfficeTemplateInstantiator]:
    registry = OfficeTemplateRegistry(tmp_path / "registry")
    manifest = manifest_for(
        content,
        format_name,
        placeholders,
        max_output_bytes=max_output_bytes,
    )
    source = write_source(tmp_path / "source", f"fixture.{format_name}", content)
    registry.import_template(manifest, source)
    return registry, OfficeTemplateInstantiator(registry)


def _xml_text(payload: bytes, namespace: str, tag: str) -> str:
    root = etree.fromstring(payload)
    return "".join(root.xpath(f".//x:{tag}/text()", namespaces={"x": namespace}))


def test_docx_cjk_body_table_header_footer_and_run_styles_are_preserved(
    tmp_path: Path,
    docx_template_bytes: bytes,
) -> None:
    _registry, instantiator = _setup(
        tmp_path,
        docx_template_bytes,
        "docx",
        DOCX_PLACEHOLDERS,
    )
    staging = tmp_path / "staging"
    staging.mkdir()
    output = staging / "report.docx"

    result = instantiator.instantiate(
        "quarterly-report",
        "1.0.0",
        DOCX_VALUES,
        staging_root=staging,
        output_path=output,
    )

    assert result.output_path == output.resolve(strict=True)
    assert result.source_sha256 != result.output_sha256
    assert len(result.template_sha256) == len(result.output_sha256) == 64
    assert {(change.part_name, change.placeholder, change.occurrences) for change in result.changes} == {
        ("word/document.xml", "body", 1),
        ("word/document.xml", "client", 1),
        ("word/document.xml", "table", 1),
        ("word/header1.xml", "header", 1),
        ("word/footer1.xml", "footer", 1),
    }
    output_bytes = output.read_bytes()
    inspect_ooxml_package(
        output_bytes,
        "docx",
        expected_placeholders=(),
    )
    reopened = Document(BytesIO(output_bytes))
    assert "苏小有客户" in "\n".join(p.text for p in reopened.paragraphs)
    assert "这是高保真正文" in "\n".join(p.text for p in reopened.paragraphs)
    assert reopened.tables[0].cell(0, 0).text == "表格 华东区"
    section = reopened.sections[0]
    assert section.header.paragraphs[0].text == "页眉 2026 年度报告"
    assert section.footer.paragraphs[0].text == "页脚 内部资料"

    entries = zip_entries(output_bytes)
    word_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    root = etree.fromstring(entries["word/document.xml"])
    runs = root.xpath("(.//w:p)[1]/w:r", namespaces={"w": word_ns})
    assert len(runs) == 3
    replacement_run = runs[1]
    emptied_run = runs[2]
    assert replacement_run.xpath("boolean(w:rPr/w:b)", namespaces={"w": word_ns})
    assert replacement_run.xpath(
        "string(w:rPr/w:color/@w:val)",
        namespaces={"w": word_ns},
    ) == "112233"
    assert replacement_run.xpath("string(w:t)", namespaces={"w": word_ns}) == "苏小有客户"
    assert emptied_run.xpath("boolean(w:rPr/w:i)", namespaces={"w": word_ns})
    assert emptied_run.xpath("string(w:t)", namespaces={"w": word_ns}) == ""


def test_xlsx_formula_chart_cell_format_and_unmodified_parts_survive(
    tmp_path: Path,
    xlsx_template_bytes: bytes,
) -> None:
    _registry, instantiator = _setup(
        tmp_path,
        xlsx_template_bytes,
        "xlsx",
        ("company",),
    )
    staging = tmp_path / "staging"
    staging.mkdir()
    output = staging / "book.xlsx"

    result = instantiator.instantiate(
        "quarterly-report",
        "1.0.0",
        {"company": "苏小有科技"},
        staging_root=staging,
        output_path=output,
    )

    workbook = load_workbook(BytesIO(output.read_bytes()), data_only=False)
    sheet = workbook["数据"]
    assert sheet["A1"].value == "公司 苏小有科技"
    assert sheet["A1"].font.bold is True
    assert sheet["A1"].font.color is not None
    assert sheet["C2"].value == "=SUM(B2:B4)"
    assert sheet["C2"].number_format == "0.00"
    assert len(sheet._charts) == 1
    assert [(change.part_name, change.placeholder, change.occurrences) for change in result.changes] == [
        ("xl/worksheets/sheet1.xml", "company", 1)
    ]
    original_entries = zip_entries(xlsx_template_bytes)
    output_entries = zip_entries(output.read_bytes())
    chart_parts = [name for name in original_entries if name.startswith("xl/charts/")]
    assert chart_parts
    for part_name in chart_parts:
        assert output_entries[part_name] == original_entries[part_name]
    original_sheet = etree.fromstring(original_entries["xl/worksheets/sheet1.xml"])
    output_sheet = etree.fromstring(output_entries["xl/worksheets/sheet1.xml"])
    spreadsheet_ns = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
    assert original_sheet.xpath(
        "string(.//s:c[@r='C2']/s:f)", namespaces={"s": spreadsheet_ns}
    ) == output_sheet.xpath(
        "string(.//s:c[@r='C2']/s:f)", namespaces={"s": spreadsheet_ns}
    )


def test_pptx_shape_table_cjk_and_complex_run_format_survive(
    tmp_path: Path,
    pptx_template_bytes: bytes,
) -> None:
    _registry, instantiator = _setup(
        tmp_path,
        pptx_template_bytes,
        "pptx",
        ("table", "topic"),
    )
    staging = tmp_path / "staging"
    staging.mkdir()
    output = staging / "deck.pptx"

    result = instantiator.instantiate(
        "quarterly-report",
        "1.0.0",
        {"topic": "下一代办公", "table": "高保真预览"},
        staging_root=staging,
        output_path=output,
    )

    presentation = Presentation(BytesIO(output.read_bytes()))
    slide = presentation.slides[0]
    text_shape = next(shape for shape in slide.shapes if getattr(shape, "has_text_frame", False))
    paragraph = text_shape.text_frame.paragraphs[0]
    assert paragraph.text == "下一代办公"
    assert paragraph.runs[0].font.bold is True
    assert paragraph.runs[0].font.size == 28 * 12700
    assert paragraph.runs[0].font.color.rgb == (0x22, 0x44, 0x66)
    assert paragraph.runs[1].font.italic is True
    table_shape = next(shape for shape in slide.shapes if getattr(shape, "has_table", False))
    assert table_shape.table.cell(0, 0).text == "表格 高保真预览"
    assert {(item.placeholder, item.occurrences) for item in result.changes} == {
        ("topic", 1),
        ("table", 1),
    }
    original_entries = zip_entries(pptx_template_bytes)
    output_entries = zip_entries(output.read_bytes())
    for part_name, payload in original_entries.items():
        if part_name != "ppt/slides/slide1.xml":
            assert output_entries[part_name] == payload


def test_instantiation_is_deterministic_for_same_source_and_values(
    tmp_path: Path,
    docx_template_bytes: bytes,
) -> None:
    _registry, instantiator = _setup(
        tmp_path,
        docx_template_bytes,
        "docx",
        DOCX_PLACEHOLDERS,
    )
    first_staging = tmp_path / "first"
    second_staging = tmp_path / "second"
    first_staging.mkdir()
    second_staging.mkdir()
    first_path = first_staging / "output.docx"
    second_path = second_staging / "output.docx"

    first = instantiator.instantiate(
        "quarterly-report",
        "1.0.0",
        DOCX_VALUES,
        staging_root=first_staging,
        output_path=first_path,
    )
    second = instantiator.instantiate(
        "quarterly-report",
        "1.0.0",
        DOCX_VALUES,
        staging_root=second_staging,
        output_path=second_path,
    )

    assert first_path.read_bytes() == second_path.read_bytes()
    assert first.output_sha256 == second.output_sha256
    assert first.source_sha256 == second.source_sha256
    assert first.template_sha256 == second.template_sha256
    assert first.changes == second.changes


def test_value_and_staging_contracts_reject_unsafe_or_ambiguous_requests(
    tmp_path: Path,
    docx_template_bytes: bytes,
) -> None:
    _registry, instantiator = _setup(
        tmp_path,
        docx_template_bytes,
        "docx",
        DOCX_PLACEHOLDERS,
    )
    staging = tmp_path / "staging"
    staging.mkdir()

    with pytest.raises(TemplateContractError, match="missing required"):
        instantiator.instantiate(
            "quarterly-report",
            "1.0.0",
            {name: value for name, value in DOCX_VALUES.items() if name != "body"},
            staging_root=staging,
            output_path=staging / "missing.docx",
        )
    with pytest.raises(TemplateContractError, match="unknown placeholders"):
        instantiator.instantiate(
            "quarterly-report",
            "1.0.0",
            {**DOCX_VALUES, "execute": "{{ dangerous }}"},
            staging_root=staging,
            output_path=staging / "unknown.docx",
        )
    with pytest.raises(TemplateContractError, match="beneath"):
        instantiator.instantiate(
            "quarterly-report",
            "1.0.0",
            DOCX_VALUES,
            staging_root=staging,
            output_path=tmp_path / "escaped.docx",
        )
    with pytest.raises(TemplateContractError, match="extension"):
        instantiator.instantiate(
            "quarterly-report",
            "1.0.0",
            DOCX_VALUES,
            staging_root=staging,
            output_path=staging / "wrong.xlsx",
        )
    existing = staging / "existing.docx"
    existing.write_bytes(b"do not overwrite")
    with pytest.raises(TemplateContractError, match="already exists"):
        instantiator.instantiate(
            "quarterly-report",
            "1.0.0",
            DOCX_VALUES,
            staging_root=staging,
            output_path=existing,
        )
    assert existing.read_bytes() == b"do not overwrite"

    target_parent = tmp_path / "outside-parent"
    target_parent.mkdir()
    redirected = staging / "redirected"
    try:
        redirected.symlink_to(target_parent, target_is_directory=True)
    except (OSError, NotImplementedError):
        return
    with pytest.raises(TemplateContractError, match="non-symlink"):
        instantiator.instantiate(
            "quarterly-report",
            "1.0.0",
            DOCX_VALUES,
            staging_root=staging,
            output_path=redirected / "linked.docx",
        )


def test_output_budget_failure_leaves_no_partial_artifact(
    tmp_path: Path,
    docx_template_bytes: bytes,
) -> None:
    _registry, instantiator = _setup(
        tmp_path,
        docx_template_bytes,
        "docx",
        DOCX_PLACEHOLDERS,
        max_output_bytes=100,
    )
    staging = tmp_path / "staging"
    staging.mkdir()
    output = staging / "too-large.docx"
    with pytest.raises(TemplateInstantiationError, match="byte budget"):
        instantiator.instantiate(
            "quarterly-report",
            "1.0.0",
            DOCX_VALUES,
            staging_root=staging,
            output_path=output,
        )
    assert not output.exists()
