from __future__ import annotations

import hashlib
import zipfile
from copy import copy
from collections.abc import Mapping
from io import BytesIO
from pathlib import Path

from docx import Document
from docx.shared import RGBColor as DocxRGBColor
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.util import Inches, Pt

from app.office_templates import AllowedOutputRules, TemplatePackageManifest


def make_docx_template() -> bytes:
    document = Document()
    paragraph = document.add_paragraph("致：")
    first = paragraph.add_run("{{cli")
    first.bold = True
    first.font.color.rgb = DocxRGBColor(0x11, 0x22, 0x33)
    second = paragraph.add_run("ent}}")
    second.italic = True
    document.add_paragraph("正文：{{body}}")
    table = document.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "表格 {{table}}"
    section = document.sections[0]
    section.header.paragraphs[0].text = "页眉 {{header}}"
    section.footer.paragraphs[0].text = "页脚 {{footer}}"
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def make_xlsx_template() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "数据"
    sheet["A1"] = "公司 {{company}}"
    cell_font = copy(sheet["A1"].font)
    cell_font.bold = True
    cell_font.color = "123456"
    sheet["A1"].font = cell_font
    sheet["B1"] = "销量"
    sheet["B2"] = 3
    sheet["B3"] = 5
    sheet["B4"] = 7
    sheet["C2"] = "=SUM(B2:B4)"
    sheet["C2"].number_format = "0.00"
    chart = BarChart()
    chart.title = "销量图"
    chart.add_data(
        Reference(sheet, min_col=2, min_row=1, max_row=4),
        titles_from_data=True,
    )
    sheet.add_chart(chart, "E2")
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def make_pptx_template() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(1))
    paragraph = shape.text_frame.paragraphs[0]
    first = paragraph.add_run()
    first.text = "{{to"
    first.font.bold = True
    first.font.size = Pt(28)
    first.font.color.rgb = PptxRGBColor(0x22, 0x44, 0x66)
    second = paragraph.add_run()
    second.text = "pic}}"
    second.font.italic = True
    table_shape = slide.shapes.add_table(
        1,
        1,
        Inches(1),
        Inches(2.5),
        Inches(6),
        Inches(1),
    )
    table_shape.table.cell(0, 0).text = "表格 {{table}}"
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def manifest_for(
    content: bytes,
    format_name: str,
    placeholders: tuple[str, ...],
    *,
    template_id: str = "quarterly-report",
    version: str = "1.0.0",
    provenance: str = "unit-test fixture generated with the OOXML format library",
    max_output_bytes: int = 10 * 1024 * 1024,
) -> TemplatePackageManifest:
    return TemplatePackageManifest(
        template_id=template_id,
        template_version=version,
        format=format_name,  # type: ignore[arg-type]
        source_sha256=hashlib.sha256(content).hexdigest(),
        license="CC0-1.0",
        provenance=provenance,
        required_placeholders=tuple(sorted(placeholders)),
        allowed_output_rules=AllowedOutputRules(
            extensions=(f".{format_name}",),
            max_output_bytes=max_output_bytes,
        ),
    )


def write_source(directory: Path, filename: str, content: bytes) -> Path:
    directory.mkdir(parents=True, exist_ok=True)
    path = directory / filename
    path.write_bytes(content)
    return path.resolve(strict=True)


def zip_entries(content: bytes) -> dict[str, bytes]:
    with zipfile.ZipFile(BytesIO(content), "r") as archive:
        return {info.filename: archive.read(info) for info in archive.infolist()}


def rewrite_zip(
    content: bytes,
    *,
    replacements: Mapping[str, bytes] | None = None,
    additions: Mapping[str, bytes] | None = None,
) -> bytes:
    replacements = replacements or {}
    additions = additions or {}
    output = BytesIO()
    with zipfile.ZipFile(BytesIO(content), "r") as source:
        with zipfile.ZipFile(output, "w", allowZip64=False) as target:
            for info in source.infolist():
                payload = replacements.get(info.filename, source.read(info))
                target.writestr(info, payload)
            for name, payload in additions.items():
                info = zipfile.ZipInfo(name, date_time=(1980, 1, 1, 0, 0, 0))
                info.compress_type = zipfile.ZIP_DEFLATED
                target.writestr(info, payload, compresslevel=9)
    return output.getvalue()
