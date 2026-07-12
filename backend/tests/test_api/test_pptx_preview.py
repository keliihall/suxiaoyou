"""Security and rendering tests for the offline PPTX preview endpoint."""

from __future__ import annotations

import json
import inspect
import asyncio
import threading
import time
import warnings
import zipfile
from io import BytesIO
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from app.api import pptx as pptx_api
from app.api.pptx import PreviewContext, _serialize_text_frame

pytestmark = pytest.mark.asyncio


def _png_bytes(color: str = "red") -> bytes:
    stream = BytesIO()
    Image.new("RGB", (32, 18), color).save(stream, format="PNG")
    return stream.getvalue()


def _write_sample_presentation(path: Path, *, slides: int = 2) -> None:
    presentation = Presentation()
    for slide_index in range(slides):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.5),
            Inches(0.5),
            Inches(5),
            Inches(1.25),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(40, 90, 170)
        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = f"Slide {slide_index + 1} / 中文"
        run.font.size = Pt(26)
        # An external hyperlink is intentionally present in the OOXML. The
        # preview must preserve only its text and never expose/fetch the URL.
        run.hyperlink.address = "https://preview-must-not-fetch.invalid/"

        slide.shapes.add_picture(
            BytesIO(_png_bytes()),
            Inches(6),
            Inches(0.5),
            width=Inches(2),
        )
        table = slide.shapes.add_table(
            2,
            2,
            Inches(0.5),
            Inches(2.25),
            Inches(6),
            Inches(2),
        ).table
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "苏小有"
        table.cell(1, 1).text = "42"
    presentation.save(path)


class TestPptxStaticPreview:
    async def test_renders_text_images_tables_and_basic_shapes_without_external_urls(
        self,
        app_client,
        tmp_path,
    ):
        target = tmp_path / "sample deck.pptx"
        _write_sample_presentation(target)

        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )

        assert response.status_code == 200
        body = response.json()
        assert body["name"] == target.name
        assert body["slideCount"] == 2
        assert len(body["slides"]) == 2
        assert body["width"] > body["height"] > 0
        assert {item["kind"] for item in body["slides"][0]["elements"]} >= {
            "shape",
            "image",
            "table",
        }
        assert "ignored_external_links" in body["warnings"]
        assert "static_preview_limitations" in body["warnings"]

        [asset] = body["assets"].values()
        assert asset["mimeType"] == "image/png"
        assert asset["dataUrl"].startswith("data:image/png;base64,")
        assert "preview-must-not-fetch.invalid" not in json.dumps(body)
        assert "Slide 1 / \\u4e2d\\u6587" in json.dumps(body)

    async def test_relative_path_resolves_against_workspace(self, app_client, tmp_path):
        target = tmp_path / "relative.pptx"
        _write_sample_presentation(target, slides=1)

        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": target.name, "workspace": str(tmp_path)},
        )

        assert response.status_code == 200
        assert response.json()["path"] == str(target.resolve())

    async def test_legacy_ppt_is_rejected_before_ooxml_parser(
        self,
        app_client,
        tmp_path,
        monkeypatch,
    ):
        target = tmp_path / "legacy.ppt"
        target.write_bytes(b"not-ooxml")
        called = False

        def fail_if_called(_path):
            nonlocal called
            called = True
            raise AssertionError("legacy PPT reached OOXML parser")

        monkeypatch.setattr(pptx_api, "build_pptx_preview", fail_if_called)
        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )

        assert response.status_code == 415
        assert response.json()["detail"] == "ppt_legacy_unsupported"
        assert called is False

    async def test_corrupt_pptx_has_stable_error_code(self, app_client, tmp_path):
        target = tmp_path / "corrupt.pptx"
        target.write_bytes(b"not a ZIP")

        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )

        assert response.status_code == 422
        assert response.json()["detail"] == "pptx_invalid_archive"

    async def test_unauthenticated_preview_is_rejected(self, app_client, tmp_path):
        target = tmp_path / "private.pptx"
        _write_sample_presentation(target, slides=1)

        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
            headers={"Authorization": ""},
        )

        assert response.status_code == 401


class TestPptxPreviewLimits:
    async def test_image_count_and_pixel_budgets_are_hard_limits(
        self,
        app_client,
        tmp_path,
        monkeypatch,
    ):
        assert pptx_api.MAX_PPTX_IMAGE_PIXELS == 12_000_000
        assert pptx_api.MAX_PPTX_TOTAL_IMAGE_PIXELS == 40_000_000
        assert pptx_api.MAX_PPTX_IMAGE_ASSETS == 64
        target = tmp_path / "image-limits.pptx"
        _write_sample_presentation(target, slides=1)
        monkeypatch.setattr(pptx_api, "MAX_PPTX_IMAGE_ASSETS", 0)

        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )

        assert response.status_code == 200
        body = response.json()
        assert body["assets"] == {}
        assert "asset_limit_exceeded" in body["warnings"]

    async def test_twenty_thousand_empty_runs_do_not_create_scene_nodes(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(0, 0, Inches(2), Inches(2))
        paragraph = shape.text_frame.paragraphs[0]
        for _ in range(20_000):
            paragraph.add_run().text = ""

        context = PreviewContext(deadline=time.monotonic() + 5)
        rendered = _serialize_text_frame(shape.text_frame, shape, context)

        assert rendered is None
        assert context.run_count == 0
        assert context.scene_nodes == 0
        assert context.runs_scanned <= pptx_api.MAX_PPTX_RUNS_SCANNED

    async def test_thousands_of_empty_paragraphs_are_discarded(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(0, 0, Inches(2), Inches(2))
        for _ in range(6_000):
            shape.text_frame.add_paragraph()

        context = PreviewContext(deadline=time.monotonic() + 5)
        rendered = _serialize_text_frame(shape.text_frame, shape, context)

        assert rendered is None
        assert context.paragraph_count == 0
        assert context.scene_nodes == 0
        assert context.paragraphs_scanned <= pptx_api.MAX_PPTX_PARAGRAPHS_SCANNED

    async def test_shape_iteration_is_bounded_without_materializing_all_shapes(
        self,
        app_client,
        tmp_path,
        monkeypatch,
    ):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        for index in range(12):
            slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(index / 10),
                Inches(index / 10),
                Inches(1),
                Inches(0.5),
            )
        target = tmp_path / "many-shapes.pptx"
        presentation.save(target)
        monkeypatch.setattr(pptx_api, "MAX_PPTX_SHAPES_PER_SLIDE", 3)

        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )

        assert response.status_code == 200
        body = response.json()
        assert len(body["slides"][0]["elements"]) == 3
        assert "shape_limit_exceeded" in body["warnings"]
        assert "list(slide.shapes)" not in inspect.getsource(pptx_api._serialize_slide)

    async def test_scene_node_budget_truncates_before_frontend_dom_growth(
        self,
        app_client,
        tmp_path,
        monkeypatch,
    ):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        for index in range(20):
            shape = slide.shapes.add_textbox(
                Inches(0.1),
                Inches(index / 10),
                Inches(2),
                Inches(0.25),
            )
            shape.text = f"node {index}"
        target = tmp_path / "scene-nodes.pptx"
        presentation.save(target)
        monkeypatch.setattr(pptx_api, "MAX_PPTX_SCENE_NODES", 12)

        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )

        assert response.status_code == 200
        body = response.json()
        assert body["sceneNodeCount"] <= 12
        assert "scene_node_limit_exceeded" in body["warnings"]

    async def test_final_json_size_is_a_hard_limit(
        self,
        app_client,
        tmp_path,
        monkeypatch,
    ):
        target = tmp_path / "response-budget.pptx"
        _write_sample_presentation(target, slides=1)
        monkeypatch.setattr(pptx_api, "MAX_PPTX_SCENE_JSON_BYTES", 128)

        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )

        assert response.status_code == 413
        assert response.json()["detail"] == "pptx_scene_size_limit"

    async def test_file_size_limit_is_checked_before_zip_open(
        self,
        app_client,
        tmp_path,
        monkeypatch,
    ):
        target = tmp_path / "large.pptx"
        target.write_bytes(b"12")
        monkeypatch.setattr(pptx_api, "MAX_PPTX_FILE_BYTES", 1)

        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )

        assert response.status_code == 413
        assert response.json()["detail"] == "pptx_file_size_limit"

    async def test_slide_limit_is_checked_before_presentation_parse(
        self,
        app_client,
        tmp_path,
        monkeypatch,
    ):
        target = tmp_path / "too-many-slides.pptx"
        _write_sample_presentation(target, slides=2)
        monkeypatch.setattr(pptx_api, "MAX_PPTX_SLIDES", 1)

        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )

        assert response.status_code == 413
        assert response.json()["detail"] == "pptx_slide_limit"

    async def test_archive_member_path_traversal_is_rejected(self, app_client, tmp_path):
        target = tmp_path / "traversal.pptx"
        with zipfile.ZipFile(target, "w") as archive:
            archive.writestr("[Content_Types].xml", "<Types />")
            archive.writestr("ppt/presentation.xml", "<presentation />")
            archive.writestr("../outside.bin", b"payload")

        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )

        assert response.status_code == 422
        assert response.json()["detail"] == "pptx_invalid_archive_path"

    async def test_duplicate_archive_members_are_rejected(self, app_client, tmp_path):
        target = tmp_path / "duplicate.pptx"
        with warnings.catch_warnings():
            warnings.simplefilter("ignore", UserWarning)
            with zipfile.ZipFile(target, "w") as archive:
                archive.writestr("[Content_Types].xml", "<Types />")
                archive.writestr("ppt/presentation.xml", "<presentation />")
                archive.writestr("ppt/presentation.xml", "<different />")

        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )

        assert response.status_code == 422
        assert response.json()["detail"] == "pptx_duplicate_zip_entry"

    async def test_zip_entry_count_is_bounded(self, app_client, tmp_path, monkeypatch):
        target = tmp_path / "entries.pptx"
        _write_sample_presentation(target, slides=1)
        monkeypatch.setattr(pptx_api, "MAX_PPTX_ZIP_ENTRIES", 2)

        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )

        assert response.status_code == 413
        assert response.json()["detail"] == "pptx_zip_entry_limit"

    async def test_zip_compression_is_limited_to_store_or_deflate(
        self,
        app_client,
        tmp_path,
    ):
        target = tmp_path / "bzip2.pptx"
        with zipfile.ZipFile(target, "w", compression=zipfile.ZIP_BZIP2) as archive:
            archive.writestr("[Content_Types].xml", "<Types />")
            archive.writestr("ppt/presentation.xml", "<presentation />")

        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )

        assert response.status_code == 422
        assert response.json()["detail"] == "pptx_unsupported_zip_compression"

    async def test_uncompressed_and_member_limits_are_independent(
        self,
        app_client,
        tmp_path,
        monkeypatch,
    ):
        target = tmp_path / "archive-limits.pptx"
        _write_sample_presentation(target, slides=1)
        monkeypatch.setattr(pptx_api, "MAX_PPTX_ZIP_MEMBER_BYTES", 1)

        member_response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )
        assert member_response.status_code == 413
        assert member_response.json()["detail"] == "pptx_zip_member_limit"

        monkeypatch.setattr(pptx_api, "MAX_PPTX_ZIP_MEMBER_BYTES", 32 * 1024 * 1024)
        monkeypatch.setattr(pptx_api, "MAX_PPTX_UNCOMPRESSED_BYTES", 1)
        total_response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )
        assert total_response.status_code == 413
        assert total_response.json()["detail"] == "pptx_uncompressed_limit"

    async def test_table_cell_budget_degrades_to_placeholder(
        self,
        app_client,
        tmp_path,
        monkeypatch,
    ):
        target = tmp_path / "table-limit.pptx"
        _write_sample_presentation(target, slides=1)
        monkeypatch.setattr(pptx_api, "MAX_PPTX_TABLE_CELLS_PER_TABLE", 1)

        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )

        assert response.status_code == 200
        body = response.json()
        assert "table_cell_limit_exceeded" in body["warnings"]
        assert any(
            element["kind"] == "unsupported" and element["label"] == "table"
            for element in body["slides"][0]["elements"]
        )

    async def test_animated_image_format_is_never_returned_as_data_url(
        self,
        app_client,
        tmp_path,
    ):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        gif = BytesIO()
        frames = [Image.new("RGB", (8, 8), "red"), Image.new("RGB", (8, 8), "blue")]
        frames[0].save(gif, format="GIF", save_all=True, append_images=frames[1:], duration=20, loop=0)
        gif.seek(0)
        slide.shapes.add_picture(gif, Inches(1), Inches(1))
        target = tmp_path / "animated.pptx"
        presentation.save(target)

        response = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )

        assert response.status_code == 200
        body = response.json()
        assert body["assets"] == {}
        assert "unsupported_image" in body["warnings"]
        assert "data:image/gif" not in json.dumps(body)

    async def test_timed_out_worker_keeps_gate_until_real_completion(
        self,
        app_client,
        tmp_path,
        monkeypatch,
    ):
        target = tmp_path / "slow.pptx"
        _write_sample_presentation(target, slides=1)
        started = threading.Event()
        release = threading.Event()
        state_lock = threading.Lock()
        active = 0
        peak = 0
        calls = 0

        def slow_preview(_path):
            nonlocal active, peak, calls
            with state_lock:
                calls += 1
                call_number = calls
                active += 1
                peak = max(peak, active)
            try:
                if call_number == 1:
                    started.set()
                    assert release.wait(timeout=2)
                return {"worker": call_number}
            finally:
                with state_lock:
                    active -= 1

        monkeypatch.setattr(pptx_api, "build_pptx_preview", slow_preview)
        monkeypatch.setattr(pptx_api, "PPTX_RENDER_TIMEOUT_SECONDS", 0.001)
        monkeypatch.setattr(pptx_api, "PPTX_WORKER_GRACE_SECONDS", 0)

        timed_out = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )
        assert started.is_set()
        assert timed_out.status_code == 408
        assert timed_out.json()["detail"] == "pptx_preview_timeout"

        busy = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )
        assert busy.status_code == 429
        assert busy.json()["detail"] == "pptx_preview_busy"
        assert peak == 1

        release.set()
        for _ in range(100):
            if not pptx_api._PPTX_RENDER_GATE.locked():
                break
            await asyncio.sleep(0.01)
        assert not pptx_api._PPTX_RENDER_GATE.locked()

        monkeypatch.setattr(pptx_api, "PPTX_RENDER_TIMEOUT_SECONDS", 1)
        recovered = await app_client.post(
            "/api/files/pptx-preview",
            json={"path": str(target)},
        )
        assert recovered.status_code == 200
        assert recovered.json() == {"worker": 2}
        assert peak == 1
