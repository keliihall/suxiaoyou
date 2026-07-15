"""Tests for the restricted cross-platform Office tool."""

from __future__ import annotations

import hashlib
import zipfile
from pathlib import Path

import pytest
from docx import Document
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.agent.agent import AgentRegistry
from app.agent.permission import GLOBAL_DEFAULTS, evaluate, merge_rulesets
from app.schemas.agent import AgentInfo
from app.tool import workspace_transaction as workspace_transaction_module
from app.tool.builtin import office as office_module
from app.tool.builtin.office import OfficeInputError, OfficeTool
from app.tool.context import ToolContext


def _context(workspace: Path, *, language: str = "zh") -> ToolContext:
    return ToolContext(
        session_id="office-session",
        message_id="office-message",
        agent=AgentInfo(name="test", description="", mode="primary"),
        call_id="office-call",
        language=language,  # type: ignore[arg-type]
        workspace=str(workspace),
    )


@pytest.fixture
def workspace(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    project = tmp_path / "workspace"
    project.mkdir()
    private = tmp_path / "private"
    private.mkdir()
    monkeypatch.setenv("SUXIAOYOU_PRIVATE_DATA_DIR", str(private))
    return project


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _make_png(path: Path) -> None:
    Image.new("RGB", (16, 12), color=(20, 80, 160)).save(path, format="PNG")


def test_flush_file_opens_office_output_for_update(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    output = tmp_path / "office-output.tmp"
    output.write_bytes(b"office")
    opened_modes: list[str] = []
    real_open = Path.open

    def tracked_open(path: Path, mode: str = "r", *args, **kwargs):
        opened_modes.append(mode)
        return real_open(path, mode, *args, **kwargs)

    monkeypatch.setattr(Path, "open", tracked_open)
    office_module._flush_file(output)

    assert opened_modes == ["r+b"]


def _rewrite_archive_without(path: Path, removed_name: str) -> None:
    temporary = path.with_name(f".{path.name}.rewrite")
    with zipfile.ZipFile(path) as source, zipfile.ZipFile(temporary, "w") as destination:
        for info in source.infolist():
            if info.filename != removed_name:
                destination.writestr(info, source.read(info.filename))
    temporary.replace(path)


async def _create_docx(tool: OfficeTool, workspace: Path, path: str = "report.docx"):
    return await tool.execute(
        {
            "file_path": path,
            "operation": "create",
            "document": {
                "title": "Quarterly report",
                "paragraphs": [
                    {"text": "Summary", "style": "heading1"},
                    {"text": "Revenue increased.", "style": "normal"},
                ],
                "tables": [
                    {"headers": ["Metric", "Value"], "rows": [["Revenue", 42]]}
                ],
            },
        },
        _context(workspace),
    )


class TestOfficeContract:
    def test_builtin_agents_keep_office_behind_write_approval(self):
        agents = AgentRegistry()
        build_rules = merge_rulesets(GLOBAL_DEFAULTS, agents.get("build").permissions)
        plan_rules = merge_rulesets(GLOBAL_DEFAULTS, agents.get("plan").permissions)
        general_rules = merge_rulesets(GLOBAL_DEFAULTS, agents.get("general").permissions)

        assert evaluate("office", "report.docx", build_rules) == "ask"
        assert evaluate("office", "report.docx", general_rules) == "ask"
        assert evaluate("office", "report.docx", plan_rules) == "deny"

    def test_schema_is_declarative_and_documents_safety_constraints(self):
        tool = OfficeTool()
        schema = tool.parameters_schema()

        assert tool.id == "office"
        assert schema["additionalProperties"] is False
        assert set(schema["properties"]) == {
            "file_path",
            "operation",
            "overwrite",
            "document",
            "workbook",
            "presentation",
            "replacements",
        }
        assert "command" not in schema["properties"]
        assert "Python or shell" in tool.description
        assert "never recalculated" in tool.description
        assert "external templates" in tool.description

    @pytest.mark.asyncio
    async def test_requires_workspace_and_localizes_error(self):
        result = await OfficeTool().execute(
            {
                "file_path": "report.docx",
                "operation": "create",
                "document": {"paragraphs": [{"text": "hello"}]},
            },
            ToolContext(
                session_id="s",
                message_id="m",
                agent=AgentInfo(name="test", description="", mode="primary"),
                call_id="c",
                language="en",
            ),
        )

        assert not result.success
        assert result.error == "The Office tool requires a selected workspace."

    @pytest.mark.asyncio
    async def test_rejects_workspace_escape(self, workspace: Path, tmp_path: Path):
        outside = tmp_path / "outside.docx"
        result = await OfficeTool().execute(
            {
                "file_path": str(outside),
                "operation": "create",
                "document": {"paragraphs": [{"text": "no"}]},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "inside the current workspace" in (result.error or "")
        assert not outside.exists()

    @pytest.mark.asyncio
    async def test_rejects_symlink_escape(self, workspace: Path, tmp_path: Path):
        outside = tmp_path / "outside"
        outside.mkdir()
        link = workspace / "linked"
        try:
            link.symlink_to(outside, target_is_directory=True)
        except OSError:
            pytest.skip("symlinks are unavailable on this platform")

        result = await OfficeTool().execute(
            {
                "file_path": str(link / "escaped.docx"),
                "operation": "create",
                "document": {"paragraphs": [{"text": "no"}]},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "inside the current workspace" in (result.error or "")
        assert not (outside / "escaped.docx").exists()

    @pytest.mark.asyncio
    async def test_rejects_image_outside_workspace(self, workspace: Path, tmp_path: Path):
        outside_image = tmp_path / "outside.png"
        _make_png(outside_image)
        result = await OfficeTool().execute(
            {
                "file_path": "blocked.docx",
                "operation": "create",
                "document": {"images": [{"path": str(outside_image)}]},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "inside the current workspace" in (result.error or "")
        assert not (workspace / "suxiaoyou_written" / "blocked.docx").exists()

    @pytest.mark.asyncio
    @pytest.mark.parametrize("extension", ["docm", "xlsm", "pptm", "dotx", "xls"])
    async def test_rejects_macro_template_and_legacy_extensions(
        self,
        workspace: Path,
        extension: str,
    ):
        result = await OfficeTool().execute(
            {
                "file_path": f"unsafe.{extension}",
                "operation": "create",
                "document": {},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "Macro-enabled, template, and legacy" in (result.error or "")

    @pytest.mark.asyncio
    async def test_rejects_external_template_parameter(self, workspace: Path):
        result = await OfficeTool().execute(
            {
                "file_path": "report.docx",
                "operation": "create",
                "template_path": "https://example.com/template.dotx",
                "document": {"paragraphs": [{"text": "unsafe"}]},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert result.error == "External templates are not accepted by the Office tool."

    @pytest.mark.asyncio
    async def test_requires_payload_matching_extension(self, workspace: Path):
        result = await OfficeTool().execute(
            {
                "file_path": "wrong.docx",
                "operation": "create",
                "workbook": {"sheets": []},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert ".docx requires exactly one document payload" in (result.error or "")


class TestDocx:
    @pytest.mark.asyncio
    async def test_create_docx_with_paragraphs_table_reopen_and_atomic_metadata(
        self,
        workspace: Path,
    ):
        result = await _create_docx(OfficeTool(), workspace)
        target = workspace / "suxiaoyou_written" / "report.docx"

        assert result.success, result.error
        assert target.is_file()
        assert result.metadata["reopened_and_validated"] is True
        assert result.metadata["atomic_install"] is True
        assert result.metadata["workspace_transaction"] is True
        assert result.metadata["atomic_file_install"] is True
        assert result.metadata["paragraphs_added"] == 3
        assert result.metadata["tables_added"] == 1
        reopened = Document(target)
        assert [paragraph.text for paragraph in reopened.paragraphs] == [
            "Quarterly report",
            "Summary",
            "Revenue increased.",
        ]
        assert reopened.tables[0].cell(1, 1).text == "42"
        with zipfile.ZipFile(target) as archive:
            assert not any(
                "customxml" in {part.casefold() for part in Path(name).parts}
                for name in archive.namelist()
            )

    @pytest.mark.asyncio
    async def test_create_docx_with_page_break_and_workspace_local_image(
        self,
        workspace: Path,
    ):
        image_path = workspace / "chart.png"
        _make_png(image_path)
        result = await OfficeTool().execute(
            {
                "file_path": "illustrated.docx",
                "operation": "create",
                "document": {
                    "paragraphs": [
                        {"text": "Page one", "page_break_after": True},
                        {"text": "Page two"},
                    ],
                    "images": [
                        {"path": "chart.png", "width_inches": 2, "caption": "Chart"}
                    ],
                },
            },
            _context(workspace),
        )

        assert result.success, result.error
        assert result.metadata["images_added"] == 1
        assert result.metadata["page_breaks_added"] == 1
        reopened = Document(result.metadata["file_path"])
        assert len(reopened.inline_shapes) == 1
        assert reopened.paragraphs[-1].text == "Chart"
        assert "w:type=\"page\"" in reopened.element.xml

    @pytest.mark.asyncio
    async def test_create_refuses_implicit_overwrite(self, workspace: Path):
        tool = OfficeTool()
        first = await _create_docx(tool, workspace)
        target = Path(first.metadata["file_path"])
        before = target.read_bytes()

        second = await _create_docx(tool, workspace)

        assert not second.success
        assert "overwrite=true" in (second.error or "")
        assert target.read_bytes() == before

    @pytest.mark.asyncio
    async def test_edit_replaces_text_across_runs_appends_and_versions_original(
        self,
        workspace: Path,
    ):
        target = workspace / "styled.docx"
        document = Document()
        office_module._drop_default_docx_custom_xml(document)
        paragraph = document.add_paragraph()
        paragraph.add_run("Hello ").bold = True
        paragraph.add_run("world").italic = True
        document.save(target)
        original_hash = _sha256(target)

        result = await OfficeTool().execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "document": {"paragraphs": [{"text": "Tail", "style": "normal"}]},
                "replacements": [{"old_text": "lo wo", "new_text": "p, Wo"}],
            },
            _context(workspace),
        )

        assert result.success, result.error
        assert result.metadata["replacements"] == 1
        assert result.metadata["previous_sha256"] == original_hash
        reopened = Document(target)
        assert reopened.paragraphs[0].text == "Help, World"
        assert reopened.paragraphs[0].runs[0].bold is True
        assert reopened.paragraphs[-1].text == "Tail"

    @pytest.mark.asyncio
    async def test_edit_requires_unique_match_unless_replace_all(self, workspace: Path):
        target = workspace / "duplicates.docx"
        document = Document()
        office_module._drop_default_docx_custom_xml(document)
        document.add_paragraph("same")
        document.add_paragraph("same")
        document.save(target)
        before = target.read_bytes()

        rejected = await OfficeTool().execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "document": {},
                "replacements": [{"old_text": "same", "new_text": "new"}],
            },
            _context(workspace, language="en"),
        )

        assert not rejected.success
        assert "Found 2 matches" in (rejected.error or "")
        assert target.read_bytes() == before

        accepted = await OfficeTool().execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "document": {},
                "replacements": [
                    {"old_text": "same", "new_text": "new", "replace_all": True}
                ],
            },
            _context(workspace),
        )
        assert accepted.success, accepted.error
        assert [paragraph.text for paragraph in Document(target).paragraphs] == ["new", "new"]


class TestXlsx:
    @pytest.mark.asyncio
    async def test_create_xlsx_stores_formula_without_recalculation(self, workspace: Path):
        result = await OfficeTool().execute(
            {
                "file_path": "metrics.xlsx",
                "operation": "create",
                "workbook": {
                    "sheets": [
                        {
                            "name": "Data",
                            "rows": [["Amount", 2], ["Double", "=B1*2"]],
                        }
                    ],
                    "cells": [{"sheet": "Data", "cell": "C1", "value": True}],
                },
            },
            _context(workspace, language="en"),
        )
        target = workspace / "suxiaoyou_written" / "metrics.xlsx"

        assert result.success, result.error
        assert result.metadata["formulas_written"] == 1
        assert result.metadata["formulas_recalculated"] is False
        assert "does not recalculate" in result.output
        formula_view = load_workbook(target, data_only=False)
        assert formula_view["Data"]["B2"].value == "=B1*2"
        formula_view.close()
        cached_view = load_workbook(target, data_only=True)
        assert cached_view["Data"]["B2"].value is None
        cached_view.close()

    @pytest.mark.asyncio
    async def test_edit_xlsx_appends_creates_sheet_updates_cells_and_versions(
        self,
        workspace: Path,
    ):
        target = workspace / "book.xlsx"
        created = await OfficeTool().execute(
            {
                "file_path": str(target),
                "operation": "create",
                "workbook": {"sheets": [{"name": "Data", "rows": [["old", 1]]}]},
            },
            _context(workspace),
        )
        assert created.success, created.error
        original_hash = _sha256(target)

        edited = await OfficeTool().execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "workbook": {
                    "sheets": [
                        {"name": "Data", "action": "append", "rows": [["new", 2]]},
                        {"name": "Notes", "action": "create", "rows": [["note"]]},
                    ],
                    "cells": [{"sheet": "Data", "cell": "B1", "value": 9}],
                },
            },
            _context(workspace),
        )

        assert edited.success, edited.error
        assert edited.metadata["previous_sha256"] == original_hash
        assert edited.metadata["sheets_created"] == 1
        workbook = load_workbook(target, data_only=False)
        assert workbook.sheetnames == ["Data", "Notes"]
        assert workbook["Data"]["B1"].value == 9
        assert workbook["Data"]["A2"].value == "new"
        assert workbook["Notes"]["A1"].value == "note"
        workbook.close()

    @pytest.mark.asyncio
    async def test_edit_xlsx_applies_basic_style_and_deletes_named_sheet(
        self,
        workspace: Path,
    ):
        tool = OfficeTool()
        created = await tool.execute(
            {
                "file_path": "styled.xlsx",
                "operation": "create",
                "workbook": {
                    "sheets": [
                        {"name": "Data", "rows": [[0.25]]},
                        {"name": "Remove", "rows": [["old"]]},
                    ]
                },
            },
            _context(workspace),
        )
        assert created.success, created.error

        edited = await tool.execute(
            {
                "file_path": created.metadata["file_path"],
                "operation": "edit",
                "workbook": {
                    "delete_sheets": ["Remove"],
                    "cells": [
                        {
                            "sheet": "Data",
                            "cell": "A1",
                            "style": {
                                "number_format": "0.0%",
                                "font": {
                                    "bold": True,
                                    "italic": True,
                                    "color": "FF0000",
                                    "size": 14,
                                },
                                "fill": {"color": "00FF00"},
                            },
                        }
                    ],
                },
            },
            _context(workspace),
        )

        assert edited.success, edited.error
        assert edited.metadata["sheets_deleted"] == 1
        assert edited.metadata["styles_applied"] == 1
        workbook = load_workbook(edited.metadata["file_path"], data_only=False)
        assert workbook.sheetnames == ["Data"]
        cell = workbook["Data"]["A1"]
        assert cell.value == 0.25
        assert cell.number_format == "0.0%"
        assert cell.font.bold is True
        assert cell.font.italic is True
        assert cell.font.sz == 14
        assert cell.fill.fill_type == "solid"
        workbook.close()

    @pytest.mark.asyncio
    async def test_rejects_external_formula_reference(self, workspace: Path):
        result = await OfficeTool().execute(
            {
                "file_path": "external.xlsx",
                "operation": "create",
                "workbook": {
                    "sheets": [
                        {"name": "Data", "rows": [["='[other.xlsx]Sheet1'!A1"]]}
                    ]
                },
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "external workbook or network reference" in (result.error or "")

    @pytest.mark.asyncio
    @pytest.mark.parametrize("coordinate", ["A0", "XFE1", "not-a-cell"])
    async def test_rejects_invalid_or_out_of_range_cell_coordinate(
        self,
        workspace: Path,
        coordinate: str,
    ):
        result = await OfficeTool().execute(
            {
                "file_path": "bad-cell.xlsx",
                "operation": "create",
                "workbook": {
                    "sheets": [{"name": "Data", "rows": [[1]]}],
                    "cells": [{"sheet": "Data", "cell": coordinate, "value": 2}],
                },
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "cell coordinate" in (result.error or "").lower()

    @pytest.mark.asyncio
    @pytest.mark.parametrize("sheet_name", ["bad/name", "x" * 32, "'quoted'"])
    async def test_rejects_invalid_sheet_name(self, workspace: Path, sheet_name: str):
        result = await OfficeTool().execute(
            {
                "file_path": "bad.xlsx",
                "operation": "create",
                "workbook": {"sheets": [{"name": sheet_name, "rows": [[1]]}]},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "Invalid XLSX sheet name" in (result.error or "")


class TestPptx:
    @pytest.mark.asyncio
    async def test_create_and_edit_pptx_with_reopen_validation_and_version(
        self,
        workspace: Path,
    ):
        tool = OfficeTool()
        created = await tool.execute(
            {
                "file_path": "deck.pptx",
                "operation": "create",
                "presentation": {
                    "slides": [
                        {
                            "title": "Roadmap",
                            "bullets": ["Phase one", {"text": "Detail", "level": 1}],
                        }
                    ]
                },
            },
            _context(workspace),
        )
        target = Path(created.metadata["file_path"])
        original_hash = _sha256(target)

        edited = await tool.execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "presentation": {
                    "slides": [{"title": "Next", "subtitle": "Validated"}]
                },
                "replacements": [{"old_text": "Phase one", "new_text": "Phase 1"}],
            },
            _context(workspace),
        )

        assert edited.success, edited.error
        assert edited.metadata["previous_sha256"] == original_hash
        assert edited.metadata["slides_added"] == 1
        presentation = Presentation(target)
        assert len(presentation.slides) == 2
        assert presentation.slides[0].shapes.title.text == "Roadmap"
        first_slide_text = "\n".join(
            shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")
        )
        assert "Phase 1" in first_slide_text
        assert presentation.slides[1].shapes.title.text == "Next"

    @pytest.mark.asyncio
    async def test_rejects_ambiguous_subtitle_and_bullets(self, workspace: Path):
        result = await OfficeTool().execute(
            {
                "file_path": "ambiguous.pptx",
                "operation": "create",
                "presentation": {
                    "slides": [
                        {"title": "Title", "subtitle": "Sub", "bullets": ["Item"]}
                    ]
                },
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "cannot provide subtitle and bullets" in (result.error or "")

    @pytest.mark.asyncio
    async def test_create_pptx_with_text_box_table_and_workspace_local_image(
        self,
        workspace: Path,
    ):
        image_path = workspace / "logo.png"
        _make_png(image_path)
        result = await OfficeTool().execute(
            {
                "file_path": "rich-deck.pptx",
                "operation": "create",
                "presentation": {
                    "slides": [
                        {
                            "title": "Evidence",
                            "text_boxes": [
                                {
                                    "text": "Callout",
                                    "left_inches": 1,
                                    "top_inches": 1.5,
                                    "width_inches": 3,
                                    "height_inches": 0.8,
                                    "font_size": 18,
                                }
                            ],
                            "tables": [
                                {
                                    "left_inches": 1,
                                    "top_inches": 2.5,
                                    "width_inches": 5,
                                    "height_inches": 1.5,
                                    "headers": ["Metric", "Value"],
                                    "rows": [["A", 1]],
                                }
                            ],
                            "images": [
                                {
                                    "path": "logo.png",
                                    "left_inches": 7,
                                    "top_inches": 1,
                                    "width_inches": 1,
                                }
                            ],
                        }
                    ]
                },
            },
            _context(workspace),
        )

        assert result.success, result.error
        assert result.metadata["text_boxes_added"] == 1
        assert result.metadata["tables_added"] == 1
        assert result.metadata["images_added"] == 1
        presentation = Presentation(result.metadata["file_path"])
        shapes = list(presentation.slides[0].shapes)
        assert sum(shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX for shape in shapes) == 1
        assert sum(getattr(shape, "has_table", False) for shape in shapes) == 1
        assert sum(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in shapes) == 1
        table_shape = next(shape for shape in shapes if getattr(shape, "has_table", False))
        assert table_shape.table.cell(1, 1).text == "1"


class TestOOXMLAndAtomicSafety:
    @pytest.mark.asyncio
    async def test_rejects_disguised_macro_part_without_modifying_file(
        self,
        workspace: Path,
    ):
        created = await _create_docx(OfficeTool(), workspace, "macro.docx")
        target = Path(created.metadata["file_path"])
        with zipfile.ZipFile(target, "a") as archive:
            archive.writestr("word/vbaProject.bin", b"not-a-real-macro")
        before = target.read_bytes()

        result = await OfficeTool().execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "document": {"paragraphs": [{"text": "blocked"}]},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "macros were detected" in (result.error or "")
        assert target.read_bytes() == before

    @pytest.mark.asyncio
    async def test_rejects_external_template_relationship(self, workspace: Path):
        created = await _create_docx(OfficeTool(), workspace, "template.docx")
        target = Path(created.metadata["file_path"])
        relationship_xml = b"""<?xml version='1.0' encoding='UTF-8'?>
<Relationships xmlns='http://schemas.openxmlformats.org/package/2006/relationships'>
  <Relationship Id='rId1'
    Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/attachedTemplate'
    Target='https://example.com/template.dotx' TargetMode='External'/>
</Relationships>"""
        with zipfile.ZipFile(target, "a") as archive:
            archive.writestr("custom/_rels/external.rels", relationship_xml)
        before = target.read_bytes()

        result = await OfficeTool().execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "document": {"paragraphs": [{"text": "blocked"}]},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "external template" in (result.error or "")
        assert target.read_bytes() == before

    @pytest.mark.asyncio
    async def test_rejects_embedded_package_relationship_without_modifying_file(
        self,
        workspace: Path,
    ):
        created = await _create_docx(OfficeTool(), workspace, "package.docx")
        target = Path(created.metadata["file_path"])
        relationship_xml = b"""<?xml version='1.0' encoding='UTF-8'?>
<Relationships xmlns='http://schemas.openxmlformats.org/package/2006/relationships'>
  <Relationship Id='rId1'
    Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/package'
    Target='../embeddings/embedded.xlsx'/>
</Relationships>"""
        with zipfile.ZipFile(target, "a") as archive:
            archive.writestr("custom/_rels/package.rels", relationship_xml)
        before = target.read_bytes()

        result = await OfficeTool().execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "document": {"paragraphs": [{"text": "blocked"}]},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "embedded object" in (result.error or "")
        assert target.read_bytes() == before

    @pytest.mark.asyncio
    @pytest.mark.parametrize(
        ("extension", "payload_name", "create_payload", "edit_payload", "part_name"),
        [
            (
                "docx",
                "document",
                {"paragraphs": [{"text": "original"}]},
                {"paragraphs": [{"text": "blocked"}]},
                "word/embeddings/oleObject1.bin",
            ),
            (
                "xlsx",
                "workbook",
                {"sheets": [{"name": "Sheet1", "rows": [["original"]]}]},
                {"cells": [{"sheet": "Sheet1", "cell": "A2", "value": "blocked"}]},
                "xl/embeddings/embeddedWorkbook.xlsx",
            ),
            (
                "pptx",
                "presentation",
                {"slides": [{"title": "original"}]},
                {"slides": [{"title": "blocked"}]},
                "ppt/embeddings/oleObject1.bin",
            ),
        ],
    )
    async def test_rejects_embedded_parts_for_each_ooxml_format(
        self,
        workspace: Path,
        extension: str,
        payload_name: str,
        create_payload: dict,
        edit_payload: dict,
        part_name: str,
    ):
        tool = OfficeTool()
        created = await tool.execute(
            {
                "file_path": f"embedded.{extension}",
                "operation": "create",
                payload_name: create_payload,
            },
            _context(workspace),
        )
        assert created.success, created.error
        target = Path(created.metadata["file_path"])
        with zipfile.ZipFile(target, "a") as archive:
            archive.writestr(part_name, b"unsupported-embedded-content")
        before = target.read_bytes()

        result = await tool.execute(
            {
                "file_path": str(target),
                "operation": "edit",
                payload_name: edit_payload,
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "cannot safely preserve" in (result.error or "")
        assert target.read_bytes() == before

    @pytest.mark.asyncio
    @pytest.mark.parametrize(
        "part_name",
        [
            "xl/activeX/activeX1.bin",
            "xl/controls/control1.xml",
            "xl/ctrlProps/ctrlProp1.xml",
        ],
    )
    async def test_rejects_embedded_control_parts_without_modifying_file(
        self,
        workspace: Path,
        part_name: str,
    ):
        tool = OfficeTool()
        created = await tool.execute(
            {
                "file_path": "controls.xlsx",
                "operation": "create",
                "workbook": {"sheets": [{"name": "Sheet1", "rows": [[1]]}]},
            },
            _context(workspace),
        )
        assert created.success, created.error
        target = Path(created.metadata["file_path"])
        with zipfile.ZipFile(target, "a") as archive:
            archive.writestr(part_name, b"unsupported-control-content")
        before = target.read_bytes()

        result = await tool.execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "workbook": {
                    "cells": [{"sheet": "Sheet1", "cell": "A2", "value": 2}]
                },
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "cannot safely preserve" in (result.error or "")
        assert target.read_bytes() == before

    @pytest.mark.asyncio
    async def test_rejects_custom_xml_xlsx_before_openpyxl_can_silently_drop_it(
        self,
        workspace: Path,
    ):
        tool = OfficeTool()
        created = await tool.execute(
            {
                "file_path": "custom-xml.xlsx",
                "operation": "create",
                "workbook": {"sheets": [{"name": "Sheet1", "rows": [["original"]]}]},
            },
            _context(workspace),
        )
        assert created.success, created.error
        target = Path(created.metadata["file_path"])

        # Synthetic form of the custom XML store found in real Office files.
        # A direct openpyxl round trip demonstrates the regression: these parts
        # disappear even though the workbook still opens successfully.
        with zipfile.ZipFile(target, "a") as archive:
            archive.writestr("customXml/item1.xml", b"<customer><id>42</id></customer>")
            archive.writestr(
                "customXml/itemProps1.xml",
                b"""<?xml version='1.0' encoding='UTF-8'?>
<ds:datastoreItem ds:itemID='{01234567-89AB-CDEF-0123-456789ABCDEF}'
 xmlns:ds='http://schemas.openxmlformats.org/officeDocument/2006/customXml'>
 <ds:schemaRefs/>
</ds:datastoreItem>""",
            )
            archive.writestr(
                "customXml/_rels/item1.xml.rels",
                b"""<?xml version='1.0' encoding='UTF-8'?>
<Relationships xmlns='http://schemas.openxmlformats.org/package/2006/relationships'>
 <Relationship Id='rId1'
  Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/customXmlProps'
  Target='itemProps1.xml'/>
</Relationships>""",
            )
        before = target.read_bytes()

        probe = workspace / "openpyxl-roundtrip.xlsx"
        workbook = load_workbook(target, data_only=False)
        try:
            workbook.save(probe)
        finally:
            workbook.close()
        with zipfile.ZipFile(probe) as archive:
            assert "customXml/item1.xml" not in archive.namelist()

        result = await tool.execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "workbook": {
                    "cells": [{"sheet": "Sheet1", "cell": "A2", "value": "blocked"}]
                },
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "customXml data" in (result.error or "")
        assert target.read_bytes() == before

    @pytest.mark.asyncio
    async def test_rejects_unknown_relationship_type_before_edit(
        self,
        workspace: Path,
    ):
        tool = OfficeTool()
        created = await tool.execute(
            {
                "file_path": "unknown-relationship.xlsx",
                "operation": "create",
                "workbook": {"sheets": [{"name": "Sheet1", "rows": [[1]]}]},
            },
            _context(workspace),
        )
        assert created.success, created.error
        target = Path(created.metadata["file_path"])
        with zipfile.ZipFile(target, "a") as archive:
            archive.writestr(
                "xl/worksheets/_rels/sheet1.xml.rels",
                b"""<?xml version='1.0' encoding='UTF-8'?>
<Relationships xmlns='http://schemas.openxmlformats.org/package/2006/relationships'>
 <Relationship Id='rIdFuture' Type='https://example.invalid/ooxml/futureFeature'
  Target='../workbook.xml'/>
</Relationships>""",
            )
        before = target.read_bytes()

        result = await tool.execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "workbook": {"cells": [{"sheet": "Sheet1", "cell": "A2", "value": 2}]},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "unsupported relationship type" in (result.error or "")
        assert target.read_bytes() == before

    @pytest.mark.asyncio
    async def test_rejects_unknown_package_part_before_edit(
        self,
        workspace: Path,
    ):
        tool = OfficeTool()
        created = await tool.execute(
            {
                "file_path": "unknown-part.xlsx",
                "operation": "create",
                "workbook": {"sheets": [{"name": "Sheet1", "rows": [[1]]}]},
            },
            _context(workspace),
        )
        assert created.success, created.error
        target = Path(created.metadata["file_path"])
        with zipfile.ZipFile(target, "a") as archive:
            archive.writestr("xl/futureFeature.xml", b"<future/>")
        before = target.read_bytes()

        result = await tool.execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "workbook": {"cells": [{"sheet": "Sheet1", "cell": "A2", "value": 2}]},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "package part" in (result.error or "")
        assert target.read_bytes() == before

    @pytest.mark.asyncio
    async def test_edit_fails_if_a_supported_source_part_disappears(
        self,
        workspace: Path,
        monkeypatch: pytest.MonkeyPatch,
    ):
        tool = OfficeTool()
        created = await tool.execute(
            {
                "file_path": "dropped-part.xlsx",
                "operation": "create",
                "workbook": {"sheets": [{"name": "Sheet1", "rows": [[1]]}]},
            },
            _context(workspace),
        )
        assert created.success, created.error
        target = Path(created.metadata["file_path"])
        before = target.read_bytes()
        original_writer = office_module._write_xlsx

        def dropping_writer(*args, **kwargs):
            summary, expected = original_writer(*args, **kwargs)
            _rewrite_archive_without(args[0], "docProps/app.xml")
            return summary, expected

        monkeypatch.setattr(office_module, "_write_xlsx", dropping_writer)
        result = await tool.execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "workbook": {"cells": [{"sheet": "Sheet1", "cell": "A2", "value": 2}]},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "would remove existing package parts" in (result.error or "")
        assert target.read_bytes() == before

    @pytest.mark.asyncio
    async def test_ancestor_symlink_swap_cannot_redirect_commit_outside_workspace(
        self,
        workspace: Path,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ):
        reports = workspace / "reports"
        reports.mkdir()
        outside = tmp_path / "outside"
        outside.mkdir()
        probe = workspace / "symlink-probe"
        try:
            probe.symlink_to(outside, target_is_directory=True)
            probe.unlink()
        except OSError:
            pytest.skip("symlinks are unavailable on this platform")

        tool = OfficeTool()
        target = reports / "race.docx"
        created = await tool.execute(
            {
                "file_path": str(target),
                "operation": "create",
                "document": {"paragraphs": [{"text": "baseline"}]},
            },
            _context(workspace),
        )
        assert created.success, created.error
        baseline = target.read_bytes()
        outside_target = outside / target.name
        outside_target.write_bytes(b"outside-sentinel")
        moved_reports = workspace / "reports-concurrent"
        original_runner = office_module._run_office_operation

        def swap_ancestor_after_staging(*args, **kwargs):
            summary = original_runner(*args, **kwargs)
            reports.rename(moved_reports)
            reports.symlink_to(outside, target_is_directory=True)
            return summary

        monkeypatch.setattr(
            office_module,
            "_run_office_operation",
            swap_ancestor_after_staging,
        )
        result = await tool.execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "document": {"paragraphs": [{"text": "must-not-commit"}]},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert outside_target.read_bytes() == b"outside-sentinel"
        assert (moved_reports / target.name).read_bytes() == baseline

    @pytest.mark.asyncio
    async def test_concurrent_edit_wins_and_staged_office_content_is_not_merged(
        self,
        workspace: Path,
        monkeypatch: pytest.MonkeyPatch,
    ):
        tool = OfficeTool()
        created = await _create_docx(tool, workspace, "concurrent.docx")
        assert created.success, created.error
        target = Path(created.metadata["file_path"])
        original_runner = office_module._run_office_operation
        concurrent_bytes: dict[str, bytes] = {}

        def write_concurrently_after_staging(*args, **kwargs):
            summary = original_runner(*args, **kwargs)
            concurrent = Document()
            concurrent.add_paragraph("concurrent writer content")
            concurrent.save(target)
            concurrent_bytes["value"] = target.read_bytes()
            return summary

        monkeypatch.setattr(
            office_module,
            "_run_office_operation",
            write_concurrently_after_staging,
        )
        result = await tool.execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "document": {"paragraphs": [{"text": "staged Office content"}]},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "changed outside" in (result.error or "")
        assert target.read_bytes() == concurrent_bytes["value"]
        assert [paragraph.text for paragraph in Document(target).paragraphs] == [
            "concurrent writer content"
        ]

    @pytest.mark.asyncio
    async def test_missing_guarded_rename_fails_closed_without_unsafe_fallback(
        self,
        workspace: Path,
        monkeypatch: pytest.MonkeyPatch,
    ):
        tool = OfficeTool()
        created = await _create_docx(tool, workspace, "guarded-rename.docx")
        assert created.success, created.error
        target = Path(created.metadata["file_path"])
        before = target.read_bytes()

        def guarded_rename_unavailable(*_args, **_kwargs):
            raise workspace_transaction_module.WorkspaceMutationError(
                "Atomic guarded workspace replacement is unavailable on this platform"
            )

        monkeypatch.setattr(
            workspace_transaction_module,
            "_renameat_with_flags",
            guarded_rename_unavailable,
        )
        result = await tool.execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "document": {"paragraphs": [{"text": "must-not-fallback"}]},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert "guarded workspace replacement is unavailable" in (result.error or "")
        assert target.read_bytes() == before

    @pytest.mark.asyncio
    async def test_reopen_failure_preserves_existing_destination_and_cleans_temp(
        self,
        workspace: Path,
        monkeypatch: pytest.MonkeyPatch,
    ):
        tool = OfficeTool()
        created = await _create_docx(tool, workspace, "validation.docx")
        target = Path(created.metadata["file_path"])
        before = target.read_bytes()

        def fail_validation(*_args, **_kwargs):
            raise OfficeInputError("validation failed", "validation failed")

        monkeypatch.setattr(office_module, "_reopen_and_verify", fail_validation)
        result = await tool.execute(
            {
                "file_path": str(target),
                "operation": "create",
                "overwrite": True,
                "document": {"paragraphs": [{"text": "replacement"}]},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert target.read_bytes() == before
        assert not list(target.parent.glob(f".{target.name}.*.docx"))

    @pytest.mark.asyncio
    async def test_atomic_replace_failure_preserves_original_and_reports_error(
        self,
        workspace: Path,
        monkeypatch: pytest.MonkeyPatch,
    ):
        tool = OfficeTool()
        created = await _create_docx(tool, workspace, "atomic.docx")
        target = Path(created.metadata["file_path"])
        before = target.read_bytes()

        def fail_replace(_source: Path, _target: Path):
            raise OSError("simulated replace failure")

        monkeypatch.setattr(office_module, "_atomic_replace", fail_replace)
        result = await tool.execute(
            {
                "file_path": str(target),
                "operation": "edit",
                "document": {"paragraphs": [{"text": "not installed"}]},
            },
            _context(workspace, language="en"),
        )

        assert not result.success
        assert result.error == "Office file write failed; the original file was not changed."
        assert target.read_bytes() == before
        assert not list(target.parent.glob(f".{target.name}.*.docx"))
