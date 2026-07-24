"""Focused release-gated tests for the Office v1.1 frozen subset."""

from __future__ import annotations

import asyncio
import hashlib
import io
import stat
import threading
import zipfile
from pathlib import Path

import pytest
from docx import Document
from docx.enum.section import WD_ORIENT
from docx.oxml.ns import qn
from openpyxl import Workbook, load_workbook
from PIL import Image
from pptx import Presentation

from app.schemas.agent import AgentInfo
from app.office_validation import (
    OfficeDraftSeal,
    OfficeDraftValidationResult,
    OfficeValidationReport,
    ValidationCheck,
)
from app.office_validation.precommit import OfficePrecommitRequest
from app.storage.workspace_identity import ensure_workspace_identity
from app.tool.builtin import office as office_module
from app.tool.builtin.office import OfficeTool
from app.tool.context import ToolContext
from app.tool.workspace_transaction import WorkspaceOfficePrecommitView


class _TrustedTestPrecommitSession:
    """Test-only seal issuer; deterministic draft tests cover real rendering."""

    def __init__(
        self,
        request: OfficePrecommitRequest,
        view: WorkspaceOfficePrecommitView,
    ) -> None:
        self._request = request
        self._view = view
        self._result: OfficeDraftValidationResult | None = None
        self._state = "begun"

    async def validate_candidate(self) -> OfficeDraftValidationResult:
        assert self._state == "begun"
        root_info = self._view.staged_root.lstat()
        source_info = self._view.staged_target.lstat()
        payload = self._view.staged_target.read_bytes()
        source_sha256 = hashlib.sha256(payload).hexdigest()
        baseline_sha256 = (
            self._view.baseline.sha256
            if self._view.baseline is not None
            and self._view.baseline.sha256 is not None
            else source_sha256
        )
        check = ValidationCheck(
            code="test_authoritative",
            outcome="pass",
            message="Trusted test coordinator approved the staged candidate.",
        )
        report = OfficeValidationReport(
            document_format=self._request.document_format,
            baseline_sha256=baseline_sha256,
            candidate_sha256=source_sha256,
            renderer_id="trusted-test-renderer",
            renderer_version="1",
            font_digest="f" * 64,
            verdict="pass",
            checks=(check,),
            checkpoint_id=self._request.checkpoint_id,
            root_turn_id=self._request.root_turn_id,
        )
        seal = OfficeDraftSeal(
            relative_path=self._view.relative_path,
            source_sha256=source_sha256,
            source_mode=stat.S_IMODE(source_info.st_mode),
            source_size=source_info.st_size,
            root_identity=(root_info.st_dev, root_info.st_ino),
            source_identity=(source_info.st_dev, source_info.st_ino),
            validation_generation=self._view.validation_generation,
            renderer_id=report.renderer_id,
            renderer_version=report.renderer_version,
            font_digest=report.font_digest,
            parameters_version="trusted-test-v1",
            parameters_sha256="a" * 64,
            quality="authoritative",
            cache_key=hashlib.sha256(b"trusted-test-cache" + payload).hexdigest(),
        )
        self._result = OfficeDraftValidationResult(report=report, candidate=seal)
        self._state = "validated"
        return self._result

    def consume_commit_seal(
        self,
        result: OfficeDraftValidationResult,
    ) -> OfficeDraftSeal:
        assert self._state == "validated"
        assert result is self._result
        seal = result.commit_seal
        assert seal is not None
        self._state = "committing"
        return seal

    def mark_committed(self, result: OfficeDraftValidationResult) -> None:
        assert self._state == "committing"
        assert result is self._result
        self._state = "committed"

    def abort(self) -> None:
        if self._state != "committed":
            self._state = "aborted"


class _TrustedTestPrecommitCoordinator:
    async def begin(
        self,
        *,
        request: OfficePrecommitRequest,
        view: WorkspaceOfficePrecommitView,
    ) -> _TrustedTestPrecommitSession:
        return _TrustedTestPrecommitSession(request, view)


_TRUSTED_TEST_PRECOMMIT_COORDINATOR = _TrustedTestPrecommitCoordinator()


def _context(workspace: Path) -> ToolContext:
    identity = ensure_workspace_identity(workspace)
    context = ToolContext(
        session_id="office-v2-session",
        message_id="office-v2-message",
        agent=AgentInfo(name="test", description="", mode="primary"),
        call_id="office-v2-call",
        language="en",
        workspace=str(workspace),
        root_turn_id="office-v2-root-turn",
        turn_run_id="office-v2-turn-run",
        checkpoint_id="office-v2-checkpoint",
        workspace_instance_id="office-v2-workspace-instance",
        workspace_identity_token=identity.durable_token,
    )
    context._app_state = {  # type: ignore[attr-defined]
        "office_precommit_coordinator": _TRUSTED_TEST_PRECOMMIT_COORDINATOR,
    }
    return context


@pytest.fixture
def workspace(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    private = tmp_path / "private"
    private.mkdir()
    monkeypatch.setenv("SUXIAOYOU_PRIVATE_DATA_DIR", str(private))
    return workspace


@pytest.fixture
def office_v2(monkeypatch: pytest.MonkeyPatch) -> None:
    for gate in (
        "V11_CHECKPOINTS_RELEASED",
        "V11_REWIND_RELEASED",
        "V11_VALIDATION_AGENT_RELEASED",
        "V11_OFFICE_V2_RELEASED",
    ):
        monkeypatch.setattr(office_module.release_features, gate, True)


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _output(workspace: Path, name: str) -> Path:
    return workspace / "suxiaoyou_written" / name


def test_edit_mutation_intent_is_derived_from_writer_summary() -> None:
    pptx = office_module._office_edit_mutation_intent(
        {"operation": "edit", "replacements": []},
        {"slides_added": 2},
        "pptx",
    )
    xlsx = office_module._office_edit_mutation_intent(
        {"operation": "edit"},
        {
            "sheets_created": 0,
            "sheets_deleted": 0,
            "rows_appended": 0,
            "cells_written": 1,
            "charts_added": 0,
        },
        "xlsx",
    )
    docx = office_module._office_edit_mutation_intent(
        {"operation": "edit"},
        {
            "paragraphs_added": 1,
            "page_breaks_added": 0,
            "tables_added": 0,
            "images_added": 0,
        },
        "docx",
    )

    assert pptx.required_page_delta == 2
    assert pptx.expected_logical_unit_delta == 2
    assert pptx.max_total_changed_ratio == 0.0
    assert xlsx.expected_logical_unit_delta == 0
    assert xlsx.max_removed_pages <= 1
    assert xlsx.max_total_changed_ratio == 0.30
    assert docx.max_removed_pages == 0
    assert docx.max_total_changed_ratio == 0.45


def _png(path: Path) -> None:
    Image.new("RGB", (120, 80), color=(30, 100, 180)).save(path, "PNG")


def _metadata_contains_key(value: object, forbidden: set[str]) -> bool:
    if isinstance(value, dict):
        return any(
            key in forbidden or _metadata_contains_key(item, forbidden)
            for key, item in value.items()
        )
    if isinstance(value, (list, tuple)):
        return any(_metadata_contains_key(item, forbidden) for item in value)
    return False


def _basic_embedded_workbook() -> bytes:
    stream = io.BytesIO()
    workbook = Workbook()
    workbook.active["A1"] = 1
    workbook.save(stream)
    workbook.close()
    return stream.getvalue()


def _rewrite_zip_bytes(
    payload: bytes,
    *,
    replacements: dict[str, bytes] | None = None,
    additions: dict[str, bytes] | None = None,
) -> bytes:
    updates = replacements or {}
    output = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(payload)) as source, zipfile.ZipFile(
        output,
        "w",
    ) as destination:
        for info in source.infolist():
            destination.writestr(info, updates.get(info.filename, source.read(info)))
        for name, data in (additions or {}).items():
            destination.writestr(name, data)
    return output.getvalue()


def _embedded_outer(workbooks: dict[str, bytes]) -> io.BytesIO:
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w") as archive:
        for name, payload in workbooks.items():
            archive.writestr(name, payload)
    output.seek(0)
    return output


@pytest.mark.asyncio
async def test_gate_false_keeps_schema_hidden_and_rejects_v2_without_mutation(
    workspace: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(
        office_module.release_features,
        "V11_OFFICE_V2_RELEASED",
        False,
    )
    tool = OfficeTool()
    assert "charts" not in tool.parameters_schema()["properties"]["workbook"]["properties"]
    created = await tool.execute(
        {
            "file_path": "gate.xlsx",
            "operation": "create",
            "workbook": {"sheets": [{"name": "Data", "rows": [[1, 2]]}]},
        },
        _context(workspace),
    )
    assert created.success
    path = _output(workspace, "gate.xlsx")
    before = _sha256(path)

    rejected = await tool.execute(
        {
            "file_path": "gate.xlsx",
            "operation": "edit",
            "workbook": {
                "merged_cells": [{"sheet": "Data", "range": "A1:B1"}]
            },
        },
        _context(workspace),
    )

    assert not rejected.success
    assert "not released" in (rejected.error or "")
    assert _sha256(path) == before

    template_rejected = await tool.execute(
        {
            "file_path": "closed-template.docx",
            "operation": "create",
            "first_party_template": {
                "template_id": "business-brief",
                "template_version": "1.0.0",
                "values": {},
            },
        },
        _context(workspace),
    )
    assert not template_rejected.success
    assert "not released" in (template_rejected.error or "")
    assert not _output(workspace, "closed-template.docx").exists()


def test_office_gate_alone_does_not_expose_authoring_schema(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    for dependency_gate in (
        "V11_CHECKPOINTS_RELEASED",
        "V11_REWIND_RELEASED",
        "V11_VALIDATION_AGENT_RELEASED",
    ):
        monkeypatch.setattr(
            office_module.release_features,
            dependency_gate,
            False,
        )
    monkeypatch.setattr(
        office_module.release_features,
        "V11_OFFICE_V2_RELEASED",
        True,
    )

    schema = OfficeTool().parameters_schema()

    assert "first_party_template" not in schema["properties"]
    assert "charts" not in schema["properties"]["workbook"]["properties"]


@pytest.mark.asyncio
async def test_gate_on_requires_checkpoint_and_authoritative_coordinator(
    workspace: Path,
    office_v2: None,
) -> None:
    tool = OfficeTool()
    without_coordinator = _context(workspace)
    without_coordinator._app_state = {}  # type: ignore[attr-defined]
    unavailable = await tool.execute(
        {
            "file_path": "no-validator.docx",
            "operation": "create",
            "document": {"paragraphs": [{"text": "blocked"}]},
        },
        without_coordinator,
    )
    assert not unavailable.success
    assert "authoritative Office v1.1 precommit validator" in (
        unavailable.error or ""
    )
    assert not _output(workspace, "no-validator.docx").exists()

    without_checkpoint = _context(workspace)
    without_checkpoint.checkpoint_id = None
    missing_identity = await tool.execute(
        {
            "file_path": "no-checkpoint.docx",
            "operation": "create",
            "document": {"paragraphs": [{"text": "blocked"}]},
        },
        without_checkpoint,
    )
    assert not missing_identity.success
    assert "trusted checkpoint identity" in (missing_identity.error or "")
    assert not _output(workspace, "no-checkpoint.docx").exists()


@pytest.mark.asyncio
async def test_gate_on_uses_sealed_commit_and_hides_validation_capability(
    workspace: Path,
    office_v2: None,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    def forbidden_plain_commit(_self: object) -> object:
        raise AssertionError("plain transaction.commit() bypassed the seal")

    monkeypatch.setattr(
        office_module.WorkspaceMutationTransaction,
        "commit",
        forbidden_plain_commit,
    )
    result = await OfficeTool().execute(
        {
            "file_path": "sealed.docx",
            "operation": "create",
            "document": {"paragraphs": [{"text": "sealed"}]},
        },
        _context(workspace),
    )

    assert result.success, result.error
    assert result.metadata is not None
    assert result.metadata["office_visual_validation"] == "authoritative"
    assert result.metadata["office_validation_checkpoint_id"] == (
        "office-v2-checkpoint"
    )
    assert not _metadata_contains_key(
        result.metadata,
        {
            "cache_key",
            "commit_seal",
            "entry_path",
            "golden_path",
            "seal",
            "staged_path",
            "staging_path",
        },
    )
    assert _output(workspace, "sealed.docx").is_file()


@pytest.mark.asyncio
async def test_cancellation_waits_for_private_writer_before_abort(
    workspace: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    entered = threading.Event()
    release = threading.Event()

    def blocking_writer(
        target: Path,
        _args: object,
        _ctx: object,
        _staged_workspace: Path,
    ) -> dict[str, object]:
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_bytes(b"private-only")
        entered.set()
        assert release.wait(5)
        return {"format": "docx"}

    monkeypatch.setattr(office_module, "_run_office_operation", blocking_writer)
    task = asyncio.create_task(
        OfficeTool().execute(
            {
                "file_path": "cancelled-writer.docx",
                "operation": "create",
                "document": {"paragraphs": [{"text": "cancel"}]},
            },
            _context(workspace),
        )
    )
    assert await asyncio.to_thread(entered.wait, 2)
    task.cancel()
    await asyncio.sleep(0)
    assert not task.done()
    release.set()
    with pytest.raises(asyncio.CancelledError):
        await task
    assert not _output(workspace, "cancelled-writer.docx").exists()


@pytest.mark.asyncio
async def test_cancellation_during_sealed_commit_waits_for_known_outcome(
    workspace: Path,
    office_v2: None,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    entered = threading.Event()
    release = threading.Event()
    original = (
        office_module.WorkspaceMutationTransaction.commit_with_precommit_office_seal
    )

    def blocking_commit(self: object, seal: object) -> object:
        entered.set()
        assert release.wait(5)
        return original(self, seal)  # type: ignore[arg-type]

    monkeypatch.setattr(
        office_module.WorkspaceMutationTransaction,
        "commit_with_precommit_office_seal",
        blocking_commit,
    )
    task = asyncio.create_task(
        OfficeTool().execute(
            {
                "file_path": "cancelled-commit.docx",
                "operation": "create",
                "document": {"paragraphs": [{"text": "commit"}]},
            },
            _context(workspace),
        )
    )
    assert await asyncio.to_thread(entered.wait, 2)
    task.cancel()
    await asyncio.sleep(0)
    assert not task.done()
    release.set()
    with pytest.raises(asyncio.CancelledError):
        await task
    assert _output(workspace, "cancelled-commit.docx").is_file()


@pytest.mark.parametrize(
    ("template_id", "filename", "values"),
    [
        (
            "business-brief",
            "brief.docx",
            {
                "classification": "Internal",
                "next_step": "Complete native renderer evidence",
                "owner": "Office team",
                "recipient": "Release review",
                "report_date": "2026-07-17",
                "summary": "The signed template is installed through private staging.",
                "title": "v1.1 delivery brief",
            },
        ),
        (
            "project-tracker",
            "tracker.xlsx",
            {
                "owner": "Office team",
                "project_name": "Suxiaoyou v1.1",
                "report_date": "2026-07-17",
            },
        ),
        (
            "status-update",
            "status.pptx",
            {
                "next_step": "Complete five-target validation",
                "owner": "Office team",
                "period": "2026 Q3",
                "project_name": "Suxiaoyou v1.1",
                "status": "Validation in progress",
                "summary": "Preview, templates, versions, and rewind share one boundary.",
            },
        ),
    ],
)
@pytest.mark.asyncio
async def test_signed_first_party_templates_use_the_office_atomic_commit(
    workspace: Path,
    office_v2: None,
    template_id: str,
    filename: str,
    values: dict[str, str],
) -> None:
    tool = OfficeTool()
    schemas = tool.parameters_schema()["properties"]["first_party_template"]["oneOf"]
    assert {item["properties"]["template_id"]["const"] for item in schemas} == {
        "business-brief",
        "project-tracker",
        "status-update",
    }

    result = await tool.execute(
        {
            "file_path": filename,
            "operation": "create",
            "first_party_template": {
                "template_id": template_id,
                "template_version": "1.0.0",
                "values": values,
            },
        },
        _context(workspace),
    )

    assert result.success, result.error
    assert result.metadata is not None
    assert result.metadata["first_party_template"] is True
    assert result.metadata["template_id"] == template_id
    assert result.metadata["template_version"] == "1.0.0"
    assert result.metadata["atomic_install"] is True
    assert result.metadata["reopened_and_validated"] is True
    assert len(result.metadata["template_source_sha256"]) == 64
    assert len(result.metadata["template_manifest_sha256"]) == 64
    assert len(result.metadata["template_output_sha256"]) == 64
    path = _output(workspace, filename)
    assert result.metadata["template_output_sha256"] == _sha256(path)
    if filename.endswith(".docx"):
        assert "v1.1 delivery brief" in "\n".join(
            paragraph.text for paragraph in Document(path).paragraphs
        )
    elif filename.endswith(".xlsx"):
        workbook = load_workbook(path, data_only=False)
        try:
            assert workbook["项目跟踪"]["A1"].value == "Suxiaoyou v1.1 项目跟踪"
            assert len(workbook["项目跟踪"]._charts) == 1
        finally:
            workbook.close()
    else:
        presentation = Presentation(path)
        assert len(presentation.slides) == 1
        assert any(
            "Suxiaoyou v1.1" in getattr(shape, "text", "")
            for shape in presentation.slides[0].shapes
        )


@pytest.mark.asyncio
async def test_first_party_template_mismatch_and_overwrite_fail_atomically(
    workspace: Path,
    office_v2: None,
) -> None:
    payload = {
        "template_id": "business-brief",
        "template_version": "1.0.0",
        "values": {
            "classification": "Internal",
            "next_step": "Review",
            "owner": "Office team",
            "recipient": "Release review",
            "report_date": "2026-07-17",
            "summary": "Summary",
            "title": "Brief",
        },
    }
    tool = OfficeTool()
    mismatch = await tool.execute(
        {
            "file_path": "wrong.xlsx",
            "operation": "create",
            "first_party_template": payload,
        },
        _context(workspace),
    )
    assert not mismatch.success
    assert "does not match" in (mismatch.error or "")
    assert not _output(workspace, "wrong.xlsx").exists()

    created = await tool.execute(
        {
            "file_path": "preserve.docx",
            "operation": "create",
            "first_party_template": payload,
        },
        _context(workspace),
    )
    assert created.success, created.error
    path = _output(workspace, "preserve.docx")
    before = _sha256(path)
    rejected = await tool.execute(
        {
            "file_path": "preserve.docx",
            "operation": "create",
            "overwrite": True,
            "first_party_template": payload,
        },
        _context(workspace),
    )
    assert not rejected.success
    assert "cannot overwrite" in (rejected.error or "")
    assert _sha256(path) == before


@pytest.mark.asyncio
async def test_xlsx_v2_complex_formats_and_all_chart_types_reopen(
    workspace: Path,
    office_v2: None,
) -> None:
    schema = OfficeTool().parameters_schema()
    assert "charts" in schema["properties"]["workbook"]["properties"]
    result = await OfficeTool().execute(
        {
            "file_path": "analysis.xlsx",
            "operation": "create",
            "workbook": {
                "sheets": [
                    {
                        "name": "Data",
                        "rows": [
                            ["Category", "North", "South", "X", "Y1", "Y2"],
                            ["Q1", 10, 11, 1, 2, 4],
                            ["Q2", 12, 13, 2, 5, 7],
                            ["Q3", 14, 15, 3, 8, 9],
                        ],
                    }
                ],
                "merged_cells": [{"sheet": "Data", "range": "H1:I1"}],
                "row_heights": [{"sheet": "Data", "row": 1, "height": 28}],
                "column_widths": [{"sheet": "Data", "column": "A", "width": 20}],
                "freeze_panes": [{"sheet": "Data", "cell": "B2"}],
                "auto_filters": [{"sheet": "Data", "range": "A1:F4"}],
                "conditional_formats": [
                    {
                        "sheet": "Data",
                        "range": "B2:C4",
                        "operator": "greaterThan",
                        "formula": ["12"],
                        "fill_color": "FFF2CC",
                    }
                ],
                "data_validations": [
                    {
                        "sheet": "Data",
                        "range": "G2:G4",
                        "type": "list",
                        "formula1": '"Open,Closed"',
                        "allow_blank": True,
                    }
                ],
                "named_ranges": [
                    {"name": "QuarterValues", "sheet": "Data", "range": "B2:C4"}
                ],
                "charts": [
                    {
                        "sheet": "Data",
                        "type": "bar",
                        "data_range": "B1:C4",
                        "categories_range": "A2:A4",
                        "anchor": "H2",
                    },
                    {
                        "sheet": "Data",
                        "type": "line",
                        "data_range": "B1:C4",
                        "categories_range": "A2:A4",
                        "anchor": "H18",
                    },
                    {
                        "sheet": "Data",
                        "type": "pie",
                        "data_range": "B1:B4",
                        "categories_range": "A2:A4",
                        "anchor": "N2",
                    },
                    {
                        "sheet": "Data",
                        "type": "scatter",
                        "x_range": "D2:D4",
                        "y_ranges": ["E2:E4", "F2:F4"],
                        "anchor": "N18",
                    },
                ],
            },
        },
        _context(workspace),
    )

    assert result.success, result.error
    assert result.metadata["charts_added"] == 4
    workbook = load_workbook(_output(workspace, "analysis.xlsx"), data_only=False)
    try:
        sheet = workbook["Data"]
        assert "H1:I1" in {str(value) for value in sheet.merged_cells.ranges}
        assert sheet.row_dimensions[1].height == 28
        assert sheet.column_dimensions["A"].width == 20
        assert sheet.freeze_panes == "B2"
        assert sheet.auto_filter.ref == "A1:F4"
        assert len(sheet.conditional_formatting) == 1
        assert len(sheet.data_validations.dataValidation) == 1
        assert len(sheet._charts) == 4
        assert workbook.defined_names["QuarterValues"].attr_text == "'Data'!$B$2:$C$4"
    finally:
        workbook.close()


@pytest.mark.asyncio
async def test_xlsx_edit_proves_unrelated_sheet_digest_and_invalid_input_is_atomic(
    workspace: Path,
    office_v2: None,
) -> None:
    tool = OfficeTool()
    created = await tool.execute(
        {
            "file_path": "preserve.xlsx",
            "operation": "create",
            "workbook": {
                "sheets": [
                    {"name": "Target", "rows": [[1]]},
                    {"name": "Untouched", "rows": [[2]]},
                ]
            },
        },
        _context(workspace),
    )
    assert created.success
    edited = await tool.execute(
        {
            "file_path": "preserve.xlsx",
            "operation": "edit",
            "workbook": {
                "cells": [{"sheet": "Target", "cell": "A1", "value": 3}]
            },
        },
        _context(workspace),
    )
    assert edited.success, edited.error
    assert edited.metadata["package_parts_verified_unchanged"] == [
        "xl/worksheets/sheet2.xml"
    ]

    path = _output(workspace, "preserve.xlsx")
    before = _sha256(path)
    rejected = await tool.execute(
        {
            "file_path": "preserve.xlsx",
            "operation": "edit",
            "workbook": {
                "merged_cells": [
                    {"sheet": "Target", "range": "A1:B2"},
                    {"sheet": "Target", "range": "B2:C3"},
                ]
            },
        },
        _context(workspace),
    )
    assert not rejected.success
    assert "overlaps" in (rejected.error or "")
    assert _sha256(path) == before


@pytest.mark.asyncio
async def test_pptx_v2_layout_styles_crop_notes_and_all_chart_types(
    workspace: Path,
    office_v2: None,
) -> None:
    image_path = workspace / "photo.png"
    _png(image_path)
    chart_geometry = {
        "left_inches": 0.5,
        "top_inches": 4.5,
        "width_inches": 2.5,
        "height_inches": 2,
    }
    charts = []
    for index, chart_type in enumerate(("bar", "line", "pie")):
        charts.append(
            {
                "type": chart_type,
                "categories": ["A", "B"],
                "series": [{"name": "Series", "values": [1, 2]}],
                **chart_geometry,
                "left_inches": 0.5 + index * 3,
            }
        )
    charts.append(
        {
            "type": "scatter",
            "series": [{"name": "XY", "x_values": [1, 2], "y_values": [3, 4]}],
            **chart_geometry,
            "left_inches": 9.5,
        }
    )
    result = await OfficeTool().execute(
        {
            "file_path": "deck.pptx",
            "operation": "create",
            "presentation": {
                "slides": [
                    {
                        "title": "Executive view",
                        "layout_name": "Title Only",
                        "title_style": {"font_size": 30, "bold": True, "color": "123456"},
                        "text_boxes": [
                            {
                                "text": "Styled body",
                                "left_inches": 0.5,
                                "top_inches": 1,
                                "width_inches": 3,
                                "height_inches": 1,
                                "style": {"italic": True, "alignment": "center"},
                            }
                        ],
                        "shapes": [
                            {
                                "type": "rounded_rectangle",
                                "left_inches": 4,
                                "top_inches": 1,
                                "width_inches": 2,
                                "height_inches": 1,
                                "text": "Decision",
                                "fill_color": "D9EAF7",
                                "line_color": "123456",
                            }
                        ],
                        "tables": [
                            {
                                "left_inches": 0.5,
                                "top_inches": 2.2,
                                "width_inches": 4,
                                "height_inches": 1.5,
                                "headers": ["Metric", "Value"],
                                "rows": [["ARR", 42]],
                                "style": {"header_fill_color": "123456", "font_size": 11},
                            }
                        ],
                        "images": [
                            {
                                "path": "photo.png",
                                "left_inches": 6,
                                "top_inches": 1,
                                "width_inches": 2,
                                "height_inches": 2,
                                "crop_left": 0.1,
                                "crop_right": 0.1,
                                "align": "right",
                            }
                        ],
                        "charts": charts,
                        "speaker_notes": "Present the evidence, then the decision.",
                    }
                ]
            },
        },
        _context(workspace),
    )

    assert result.success, result.error
    assert result.metadata["charts_added"] == 4
    presentation = Presentation(_output(workspace, "deck.pptx"))
    slide = presentation.slides[0]
    assert slide.slide_layout.name == "Title Only"
    assert slide.notes_slide.notes_text_frame.text == "Present the evidence, then the decision."
    assert sum(int(getattr(shape, "has_chart", False)) for shape in slide.shapes) == 4
    picture = next(shape for shape in slide.shapes if shape.shape_type == 13)
    assert round(picture.crop_left, 3) == 0.1
    assert round(picture.crop_right, 3) == 0.1


@pytest.mark.asyncio
async def test_docx_v2_sections_runs_lists_table_merge_and_static_chart_metadata(
    workspace: Path,
    office_v2: None,
) -> None:
    chart = workspace / "chart.png"
    _png(chart)
    result = await OfficeTool().execute(
        {
            "file_path": "report.docx",
            "operation": "create",
            "document": {
                "paragraphs": [
                    {
                        "runs": [
                            {"text": "Revenue ", "bold": True, "color": "C00000", "size": 14},
                            {"text": "grew", "italic": True, "underline": True},
                        ],
                        "format": {"alignment": "center", "keep_with_next": True},
                    },
                    {"text": "Parent", "list": {"level": 0, "ordered": True}},
                    {
                        "text": "Child",
                        "list": {"level": 1, "ordered": True},
                        "format": {"page_break_before": True},
                    },
                ],
                "tables": [
                    {
                        "headers": ["Metric", "Value"],
                        "rows": [["ARR", 42], ["NRR", 110]],
                        "merges": [{"start": [1, 0], "end": [2, 0]}],
                        "format": {
                            "border_color": "123456",
                            "border_size": 8,
                            "header_fill_color": "D9EAF7",
                            "body_fill_color": "F2F2F2",
                        },
                    }
                ],
                "charts": [
                    {
                        "path": "chart.png",
                        "width_inches": 3,
                        "alt_text": "Quarterly revenue chart",
                        "source": "Finance warehouse, FY26Q2",
                    }
                ],
                "sections": [
                    {
                        "action": "configure",
                        "index": 0,
                        "orientation": "landscape",
                        "paper_size": "a4",
                        "margins": {"top_inches": 0.5, "left_inches": 0.6},
                        "header": "Confidential",
                        "footer": "Board packet",
                    }
                ],
            },
        },
        _context(workspace),
    )

    assert result.success, result.error
    document = Document(_output(workspace, "report.docx"))
    assert document.sections[0].orientation == WD_ORIENT.LANDSCAPE
    assert document.sections[0].header.paragraphs[0].text == "Confidential"
    assert document.sections[0].footer.paragraphs[0].text == "Board packet"
    assert document.paragraphs[0].runs[0].bold is True
    assert str(document.paragraphs[0].runs[0].font.color.rgb) == "C00000"
    for paragraph, level in zip(document.paragraphs[1:3], ("0", "1"), strict=True):
        num_pr = paragraph._p.pPr.find(qn("w:numPr"))
        assert num_pr.find(qn("w:ilvl")).get(qn("w:val")) == level
    assert document.tables[0].cell(1, 0)._tc is document.tables[0].cell(2, 0)._tc
    doc_pr = document.inline_shapes[0]._inline.docPr
    assert doc_pr.get("descr") == "Quarterly revenue chart"
    assert doc_pr.get("title") == "Finance warehouse, FY26Q2"


@pytest.mark.asyncio
async def test_unknown_chart_extension_is_rejected_before_source_changes(
    workspace: Path,
    office_v2: None,
) -> None:
    tool = OfficeTool()
    created = await tool.execute(
        {
            "file_path": "unsafe-chart.xlsx",
            "operation": "create",
            "workbook": {
                "sheets": [{"name": "Data", "rows": [["Y"], [1], [2]]}],
                "charts": [
                    {
                        "sheet": "Data",
                        "type": "line",
                        "data_range": "A1:A3",
                        "anchor": "C2",
                    }
                ],
            },
        },
        _context(workspace),
    )
    assert created.success
    path = _output(workspace, "unsafe-chart.xlsx")
    rewritten = path.with_suffix(".rewrite")
    with zipfile.ZipFile(path) as source, zipfile.ZipFile(rewritten, "w") as destination:
        for info in source.infolist():
            data = source.read(info)
            if info.filename == "xl/charts/chart1.xml":
                data = data.replace(b"</chartSpace>", b"<extLst/></chartSpace>")
            destination.writestr(info, data)
    rewritten.replace(path)
    before = _sha256(path)

    rejected = await tool.execute(
        {
            "file_path": "unsafe-chart.xlsx",
            "operation": "edit",
            "workbook": {"cells": [{"sheet": "Data", "cell": "A2", "value": 3}]},
        },
        _context(workspace),
    )

    assert not rejected.success
    assert "unknown extension" in (rejected.error or "")
    assert _sha256(path) == before


@pytest.mark.asyncio
async def test_unknown_standard_chart_element_and_lossy_round_trip_fail_closed(
    workspace: Path,
    office_v2: None,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    tool = OfficeTool()
    created = await tool.execute(
        {
            "file_path": "future-chart.xlsx",
            "operation": "create",
            "workbook": {
                "sheets": [
                    {"name": "Data", "rows": [["Y"], [1], [2]]},
                    {"name": "Other", "rows": [[10]]},
                ],
                "charts": [
                    {
                        "sheet": "Data",
                        "type": "line",
                        "data_range": "A1:A3",
                        "anchor": "C2",
                    }
                ],
            },
        },
        _context(workspace),
    )
    assert created.success, created.error
    path = _output(workspace, "future-chart.xlsx")
    rewritten = path.with_suffix(".rewrite")
    marker = b"<futureFeature><value>preserve-me</value></futureFeature>"
    with zipfile.ZipFile(path) as source, zipfile.ZipFile(rewritten, "w") as destination:
        for info in source.infolist():
            data = source.read(info)
            if info.filename == "xl/charts/chart1.xml":
                data = data.replace(b"</chartSpace>", marker + b"</chartSpace>")
            destination.writestr(info, data)
    rewritten.replace(path)
    before = _sha256(path)
    request = {
        "file_path": "future-chart.xlsx",
        "operation": "edit",
        "workbook": {"cells": [{"sheet": "Other", "cell": "A1", "value": 11}]},
    }

    rejected = await tool.execute(request, _context(workspace))

    assert not rejected.success
    assert "unsupported XML element" in (rejected.error or "")
    assert _sha256(path) == before

    # The digest invariant is a second line of defence if a future parser
    # accidentally widens the chart-element allowlist.
    monkeypatch.setattr(office_module, "_inspect_chart_part", lambda *_args: None)
    digest_rejected = await tool.execute(request, _context(workspace))
    assert not digest_rejected.success
    assert "declared untouched" in (digest_rejected.error or "")
    assert _sha256(path) == before


def test_embedded_workbook_rejects_excel4_macro_surfaces(office_v2: None) -> None:
    base = _basic_embedded_workbook()
    with zipfile.ZipFile(io.BytesIO(base)) as archive:
        workbook_relationships = archive.read("xl/_rels/workbook.xml.rels")
    macro_relationship = b"""<Relationship Id="rIdXlm"
 Type="http://schemas.microsoft.com/office/2006/relationships/xlMacrosheet"
 Target="worksheets/sheet1.xml"/>"""
    relationship_variant = _rewrite_zip_bytes(
        base,
        replacements={
            "xl/_rels/workbook.xml.rels": workbook_relationships.replace(
                b"</Relationships>",
                macro_relationship + b"</Relationships>",
            )
        },
    )
    macro_formula = b"""<worksheet
 xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
 <sheetData><row r="1"><c r="A1"><f>EXEC(&quot;review-probe&quot;)</f></c></row></sheetData>
</worksheet>"""
    formula_variant = _rewrite_zip_bytes(
        base,
        replacements={"xl/worksheets/sheet1.xml": macro_formula},
    )
    nested_macro_formula = b"""<worksheet
 xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
 <sheetData><row r="1"><c r="A1"><f>SUM(_xlfn._xlws.CALL(&quot;review-probe&quot;))</f></c></row></sheetData>
</worksheet>"""
    nested_formula_variant = _rewrite_zip_bytes(
        base,
        replacements={"xl/worksheets/sheet1.xml": nested_macro_formula},
    )
    dde_formula = b"""<worksheet
 xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
 <sheetData><row r="1"><c r="A1"><f>cmd|' /C review-probe'!A0</f></c></row></sheetData>
</worksheet>"""
    dde_formula_variant = _rewrite_zip_bytes(
        base,
        replacements={"xl/worksheets/sheet1.xml": dde_formula},
    )
    macrosheet_variant = _rewrite_zip_bytes(
        base,
        additions={"xl/macrosheets/sheet1.xml": macro_formula},
    )
    embedded_name = "ppt/embeddings/Microsoft_Excel_Sheet1.xlsx"

    for payload, error_fragment in (
        (macrosheet_variant, "unsafe part"),
        (relationship_variant, "macro"),
        (formula_variant, "Excel 4.0 macro formula"),
        (nested_formula_variant, "Excel 4.0 macro formula"),
        (dde_formula_variant, "external"),
    ):
        outer = _embedded_outer({embedded_name: payload})
        with zipfile.ZipFile(outer) as archive:
            with pytest.raises(office_module.OfficeInputError) as error:
                office_module._inspect_embedded_chart_workbook(
                    archive,
                    embedded_name,
                )
        assert error_fragment in error.value.en


@pytest.mark.asyncio
async def test_embedded_xlm_macrosheet_rejects_pptx_edit_atomically(
    workspace: Path,
    office_v2: None,
) -> None:
    tool = OfficeTool()
    created = await tool.execute(
        {
            "file_path": "macro-chart.pptx",
            "operation": "create",
            "presentation": {
                "slides": [
                    {
                        "title": "Chart",
                        "charts": [
                            {
                                "type": "bar",
                                "categories": ["A"],
                                "series": [{"name": "Series", "values": [1]}],
                                "left_inches": 1,
                                "top_inches": 1,
                                "width_inches": 2,
                                "height_inches": 2,
                            }
                        ],
                    }
                ]
            },
        },
        _context(workspace),
    )
    assert created.success, created.error
    path = _output(workspace, "macro-chart.pptx")
    rewritten = path.with_suffix(".rewrite")
    macro_formula = b"""<worksheet
 xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
 <sheetData><row r="1"><c r="A1"><f>EXEC(&quot;review-probe&quot;)</f></c></row></sheetData>
</worksheet>"""
    with zipfile.ZipFile(path) as source, zipfile.ZipFile(rewritten, "w") as destination:
        for info in source.infolist():
            data = source.read(info)
            if info.filename.startswith("ppt/embeddings/") and info.filename.endswith(
                ".xlsx"
            ):
                data = _rewrite_zip_bytes(
                    data,
                    additions={"xl/macrosheets/sheet1.xml": macro_formula},
                )
            destination.writestr(info, data)
    rewritten.replace(path)
    before = _sha256(path)

    rejected = await tool.execute(
        {
            "file_path": "macro-chart.pptx",
            "operation": "edit",
            "presentation": {"slides": [{"title": "Blocked"}]},
        },
        _context(workspace),
    )

    assert not rejected.success
    assert "unsafe part" in (rejected.error or "")
    assert _sha256(path) == before


def test_embedded_workbook_uses_shared_size_and_compression_budgets(
    office_v2: None,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    base = _basic_embedded_workbook()
    with zipfile.ZipFile(io.BytesIO(base)) as archive:
        expanded_size = sum(info.file_size for info in archive.infolist())
    first = "ppt/embeddings/Microsoft_Excel_Sheet1.xlsx"
    second = "ppt/embeddings/Microsoft_Excel_Sheet2.xlsx"
    outer = _embedded_outer({first: base, second: base})
    monkeypatch.setattr(
        office_module,
        "MAX_EMBEDDED_WORKBOOK_TOTAL_BYTES",
        expanded_size * 2 - 1,
    )
    budget = [0]
    with zipfile.ZipFile(outer) as archive:
        office_module._inspect_embedded_chart_workbook(
            archive,
            first,
            aggregate_uncompressed=budget,
        )
        with pytest.raises(office_module.OfficeInputError, match="aggregate"):
            office_module._inspect_embedded_chart_workbook(
                archive,
                second,
                aggregate_uncompressed=budget,
            )

    monkeypatch.setattr(
        office_module,
        "MAX_EMBEDDED_WORKBOOK_TOTAL_BYTES",
        50 * 1024 * 1024,
    )
    highly_compressible_sheet = (
        b'<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
        + b" " * (1024 * 1024)
        + b"</worksheet>"
    )
    compressed_bomb = _rewrite_zip_bytes(
        base,
        replacements={"xl/worksheets/sheet1.xml": highly_compressible_sheet},
    )
    outer = _embedded_outer({first: compressed_bomb})
    with zipfile.ZipFile(outer) as archive:
        with pytest.raises(office_module.OfficeInputError, match="compression ratio"):
            office_module._inspect_embedded_chart_workbook(archive, first)


@pytest.mark.asyncio
async def test_xlsx_merge_area_and_aggregate_budgets_are_atomic(
    workspace: Path,
    office_v2: None,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    tool = OfficeTool()
    default_per_range = office_module.MAX_MERGED_CELLS_PER_RANGE
    default_total = office_module.MAX_TOTAL_MERGED_CELLS
    created = await tool.execute(
        {
            "file_path": "merge-budget.xlsx",
            "operation": "create",
            "workbook": {"sheets": [{"name": "Data", "rows": [[1]]}]},
        },
        _context(workspace),
    )
    assert created.success, created.error
    path = _output(workspace, "merge-budget.xlsx")
    before = _sha256(path)

    oversized = await tool.execute(
        {
            "file_path": "merge-budget.xlsx",
            "operation": "edit",
            "workbook": {
                "merged_cells": [{"sheet": "Data", "range": "A1:XFD13"}]
            },
        },
        _context(workspace),
    )
    assert not oversized.success
    assert "too many merged cells" in (oversized.error or "")
    assert _sha256(path) == before

    monkeypatch.setattr(office_module, "MAX_MERGED_CELLS_PER_RANGE", 4)
    monkeypatch.setattr(office_module, "MAX_TOTAL_MERGED_CELLS", 5)
    aggregate = await tool.execute(
        {
            "file_path": "merge-budget.xlsx",
            "operation": "edit",
            "workbook": {
                "merged_cells": [
                    {"sheet": "Data", "range": "A1:C1"},
                    {"sheet": "Data", "range": "A2:C2"},
                ]
            },
        },
        _context(workspace),
    )
    assert not aggregate.success
    assert "aggregate merged-cell budget" in (aggregate.error or "")
    assert _sha256(path) == before

    # Existing ranges and requested ranges share the same whole-file budget.
    monkeypatch.setattr(office_module, "MAX_MERGED_CELLS_PER_RANGE", 5)
    monkeypatch.setattr(office_module, "MAX_TOTAL_MERGED_CELLS", 5)
    existing = await tool.execute(
        {
            "file_path": "existing-merge-budget.xlsx",
            "operation": "create",
            "workbook": {
                "sheets": [{"name": "Data", "rows": [[1]]}],
                "merged_cells": [{"sheet": "Data", "range": "A1:B2"}],
            },
        },
        _context(workspace),
    )
    assert existing.success, existing.error
    existing_path = _output(workspace, "existing-merge-budget.xlsx")
    existing_before = _sha256(existing_path)
    existing_plus_new = await tool.execute(
        {
            "file_path": "existing-merge-budget.xlsx",
            "operation": "edit",
            "workbook": {
                "merged_cells": [{"sheet": "Data", "range": "C1:D1"}]
            },
        },
        _context(workspace),
    )
    assert not existing_plus_new.success
    assert "aggregate merged-cell budget" in (existing_plus_new.error or "")
    assert _sha256(existing_path) == existing_before

    # An existing hostile merge must be rejected by the ZIP preflight before
    # openpyxl allocates cells while loading the workbook.
    monkeypatch.setattr(
        office_module,
        "MAX_MERGED_CELLS_PER_RANGE",
        default_per_range,
    )
    monkeypatch.setattr(office_module, "MAX_TOTAL_MERGED_CELLS", default_total)
    rewritten = path.with_suffix(".rewrite")
    with zipfile.ZipFile(path) as source, zipfile.ZipFile(rewritten, "w") as destination:
        for info in source.infolist():
            data = source.read(info)
            if info.filename == "xl/worksheets/sheet1.xml":
                data = data.replace(
                    b"</worksheet>",
                    (
                        b'<mergeCells count="1"><mergeCell ref="A1:XFD13"/>'
                        b"</mergeCells></worksheet>"
                    ),
                )
            destination.writestr(info, data)
    rewritten.replace(path)
    hostile_before = _sha256(path)
    hostile = await tool.execute(
        {
            "file_path": "merge-budget.xlsx",
            "operation": "edit",
            "workbook": {"cells": [{"sheet": "Data", "cell": "A1", "value": 2}]},
        },
        _context(workspace),
    )
    assert not hostile.success
    assert "merged range contains too many cells" in (hostile.error or "")
    assert _sha256(path) == hostile_before


@pytest.mark.asyncio
async def test_pptx_chart_limit_is_global_and_checked_before_generation(
    workspace: Path,
    office_v2: None,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    chart = {
        "type": "bar",
        "categories": ["A"],
        "series": [{"name": "Series", "values": [1]}],
        "left_inches": 1,
        "top_inches": 1,
        "width_inches": 2,
        "height_inches": 2,
    }
    generated = False

    def unexpected_generation(*_args: object) -> None:
        nonlocal generated
        generated = True

    monkeypatch.setattr(office_module, "_add_pptx_chart", unexpected_generation)
    result = await OfficeTool().execute(
        {
            "file_path": "too-many-charts.pptx",
            "operation": "create",
            "presentation": {
                "slides": [
                    {
                        "title": "First",
                        "charts": [dict(chart) for _ in range(26)],
                    },
                    {
                        "title": "Second",
                        "charts": [dict(chart) for _ in range(25)],
                    },
                ]
            },
        },
        _context(workspace),
    )

    assert not result.success
    assert "per PPTX file" in (result.error or "")
    assert generated is False
    assert not _output(workspace, "too-many-charts.pptx").exists()


@pytest.mark.asyncio
async def test_existing_and_new_charts_share_whole_file_budget(
    workspace: Path,
    office_v2: None,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    tool = OfficeTool()
    chart = {
        "type": "bar",
        "categories": ["A"],
        "series": [{"name": "Series", "values": [1]}],
        "left_inches": 1,
        "top_inches": 1,
        "width_inches": 2,
        "height_inches": 2,
    }
    xlsx_created = await tool.execute(
        {
            "file_path": "existing-chart.xlsx",
            "operation": "create",
            "workbook": {
                "sheets": [{"name": "Data", "rows": [["Y"], [1]]}],
                "charts": [
                    {
                        "sheet": "Data",
                        "type": "bar",
                        "data_range": "A1:A2",
                        "anchor": "C2",
                    }
                ],
            },
        },
        _context(workspace),
    )
    assert xlsx_created.success, xlsx_created.error
    pptx_created = await tool.execute(
        {
            "file_path": "existing-chart.pptx",
            "operation": "create",
            "presentation": {"slides": [{"title": "Chart", "charts": [chart]}]},
        },
        _context(workspace),
    )
    assert pptx_created.success, pptx_created.error

    monkeypatch.setattr(office_module, "MAX_CHARTS_PER_FILE", 1)
    xlsx_path = _output(workspace, "existing-chart.xlsx")
    pptx_path = _output(workspace, "existing-chart.pptx")
    xlsx_before = _sha256(xlsx_path)
    pptx_before = _sha256(pptx_path)

    xlsx_result = await tool.execute(
        {
            "file_path": "existing-chart.xlsx",
            "operation": "edit",
            "workbook": {
                "charts": [
                    {
                        "sheet": "Data",
                        "type": "line",
                        "data_range": "A1:A2",
                        "anchor": "C18",
                    }
                ]
            },
        },
        _context(workspace),
    )
    pptx_result = await tool.execute(
        {
            "file_path": "existing-chart.pptx",
            "operation": "edit",
            "presentation": {"slides": [{"title": "Second", "charts": [chart]}]},
        },
        _context(workspace),
    )

    assert not xlsx_result.success
    assert "after this operation" in (xlsx_result.error or "")
    assert not pptx_result.success
    assert "after this operation" in (pptx_result.error or "")
    assert _sha256(xlsx_path) == xlsx_before
    assert _sha256(pptx_path) == pptx_before
